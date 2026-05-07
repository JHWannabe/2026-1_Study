# -*- coding: utf-8 -*-
"""
6_generate_ppt.py
────────────────────────────────────────────────────────────────
0508 분석 결과 그래프들을 배치하여 PPT를 생성한다.

출력: results/0508_research_report.pptx

실행: python 6_generate_ppt.py
────────────────────────────────────────────────────────────────
"""

import sys, io, re
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
else:
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import warnings
warnings.filterwarnings('ignore')

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── 경로 ──────────────────────────────────────────────────────
SCRIPT_DIR  = Path(__file__).parent
RESULT_ROOT = SCRIPT_DIR.parent.parent / "results"
BASE_DIR    = RESULT_ROOT / "0508"
FIG_DIR     = BASE_DIR / "figures"
OUT_PPT     = RESULT_ROOT / "0508_research_report.pptx"

METHODS = [
    ("clinic_only",          "임상 변수 (Age, Sex, BMI)"),
    ("aec_only",             "AEC 원시 피처 256개"),
    ("aec_feature_only",     "AEC 추출 피처 전체"),
    ("aec_feature_selected", "AEC 추출 피처 선택"),
    ("clinic_aec_feature",   "임상 변수 + AEC 피처 선택"),
]
METHODS_RESNET1D = [
    ("clinic_only",        "임상 변수"),
    ("aec_only",           "AEC 원시"),
    ("aec_feature_only",   "AEC 피처(fold 선택)"),
    ("clinic_aec_feature", "임상+AEC(fold 선택)"),
]
TARGETS = ["TAMA", "SMI"]

# ── 색상 ──────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x35, 0x5E)
BLUE    = RGBColor(0x2E, 0x75, 0xB6)
SKYBLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE  = RGBColor(0xFF, 0x7F, 0x27)
GREEN   = RGBColor(0x00, 0x8B, 0x45)

W = Inches(13.33)
H = Inches(7.5)


# ── MD 파싱 ──────────────────────────────────────────────────
def _parse_resnet1d(path: Path) -> dict:
    if not path.exists():
        return {}
    text = path.read_text(encoding='utf-8')
    result = {}
    m = re.search(
        r'\|\s*\*\*Mean\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|',
        text)
    if m:
        result['mae_mean'], result['r2_mean'] = float(m.group(1)), float(m.group(2))
    m = re.search(
        r'\|\s*\*\*Std\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|',
        text)
    if m:
        result['mae_std'], result['r2_std'] = float(m.group(1)), float(m.group(2))
    m = re.search(r'Test R² = \*\*([-0-9.]+)\*\*', text)
    if m:
        result['test_r2'] = float(m.group(1))
    m = re.search(r'Test MAE = \*\*([-0-9.]+)\*\*', text)
    if m:
        result['test_mae'] = float(m.group(1))
    return result


def _parse_linear(path: Path) -> dict:
    if not path.exists():
        return {}
    text = path.read_text(encoding='utf-8')
    result = {}
    m = re.search(
        r'\|\s*\*\*Mean\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|',
        text)
    if m:
        result['mse_mean'], result['rmse_mean'], result['r2_mean'] = \
            float(m.group(1)), float(m.group(2)), float(m.group(3))
    m = re.search(
        r'\|\s*\*\*Std\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|',
        text)
    if m:
        result['mse_std'], result['rmse_std'], result['r2_std'] = \
            float(m.group(1)), float(m.group(2)), float(m.group(3))
    m = re.search(
        r'Test Set 성능.*?\n\|[^\n]+\n\|---[^\n]+\n\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*\s*\|\s*\*\*([-0-9.]+)\*\*',
        text, re.DOTALL)
    if m:
        result['test_mse'], result['test_rmse'], result['test_r2'] = \
            float(m.group(1)), float(m.group(2)), float(m.group(3))
    m = re.search(r'\*\*Input features:\*\* (.+)', text)
    if m:
        result['features'] = m.group(1).strip()
    return result


def _parse_logistic(path: Path) -> dict:
    if not path.exists():
        return {}
    text = path.read_text(encoding='utf-8')
    result = {}
    m = re.search(
        r'\|\s*\*\*Mean\*\*\s*\|\s*\*\*([0-9.]+)\*\*\s*\|\s*\*\*([0-9.]+)\*\*\s*\|',
        text)
    if m:
        result['acc_mean'], result['auc_mean'] = float(m.group(1)), float(m.group(2))
    m = re.search(
        r'\|\s*\*\*Std\*\*\s*\|\s*\*\*([0-9.]+)\*\*\s*\|\s*\*\*([0-9.]+)\*\*\s*\|',
        text)
    if m:
        result['acc_std'], result['auc_std'] = float(m.group(1)), float(m.group(2))
    m = re.search(
        r'Test Set 성능.*?\n\|[^\n]+\n\|---[^\n]+\n\|\s*\*\*([0-9.]+)\*\*\s*\|\s*\*\*([0-9.]+)\*\*',
        text, re.DOTALL)
    if m:
        result['test_acc'], result['test_auc'] = float(m.group(1)), float(m.group(2))
    thresholds = {}
    for sex in ['Male', 'Female']:
        t = re.search(rf'\|\s*{sex}\s*\|\s*([0-9.]+)\s*\|', text)
        if t:
            thresholds[sex] = float(t.group(1))
    if thresholds:
        result['thresholds'] = thresholds
    return result


def collect_data() -> dict:
    data = {}
    for method_key, _ in METHODS:
        data[method_key] = {}
        for target in TARGETS:
            data[method_key][target] = {
                'linear':   _parse_linear(BASE_DIR / method_key / target / "Linear_Report.md"),
                'logistic': _parse_logistic(BASE_DIR / method_key / target / "Logistic_Report.md"),
                'resnet1d': _parse_resnet1d(BASE_DIR / method_key / target / "ResNet1D_Report.md"),
            }
    return data


# ── PPT 헬퍼 ──────────────────────────────────────────────────
def new_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_header(slide, title, subtitle=''):
    add_rect(slide, 0, 0, W, Inches(1.05), NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.07), Inches(12.6), Inches(0.62),
             size=20, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.68), Inches(12.6), Inches(0.33),
                 size=11, color=SKYBLUE)


def add_img(slide, path: Path, left, top, width, height=None):
    if path.exists():
        if height:
            slide.shapes.add_picture(str(path), left, top, width, height)
        else:
            slide.shapes.add_picture(str(path), left, top, width)
    else:
        h = height or Inches(3)
        add_rect(slide, left, top, width, h, GRAY)
        add_text(slide, f'(그래프 없음)\n{path.name}',
                 left + Inches(0.1), top + Inches(0.2),
                 width - Inches(0.2), h - Inches(0.4),
                 size=9, color=RGBColor(0x88, 0x88, 0x88))


def add_table(slide, headers, rows, left, top, col_widths,
              row_h=Inches(0.38), font_size=9):
    if not rows:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, left, top,
        sum(col_widths), row_h * n_rows
    ).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    def _cell(cell, text, bold=False, bg=None, txt_color=DARK):
        cell.text = str(text)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        runs = cell.text_frame.paragraphs[0].runs
        if runs:
            runs[0].font.size = Pt(font_size)
            runs[0].font.bold = bold
            runs[0].font.color.rgb = txt_color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    for ci, h in enumerate(headers):
        _cell(tbl.cell(0, ci), h, bold=True, bg=NAVY, txt_color=WHITE)
    for ri, row in enumerate(rows):
        bg = SKYBLUE if ri % 2 == 0 else WHITE
        for ci in range(n_cols):
            val = row[ci] if ci < len(row) else ''
            _cell(tbl.cell(ri + 1, ci), val, bg=bg)


def _fmt(v, d=4):
    try:
        return f'{float(v):.{d}f}'
    except (TypeError, ValueError):
        return 'N/A'


# ── 공통: 5방법 이미지 1행 슬라이드 ──────────────────────────
def _slide_5images(prs, title, subtitle, img_fname_fn, label_key='label'):
    """5개 방법의 이미지를 한 행으로 가득 채운 슬라이드."""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, title, subtitle)

    img_w, img_h = Inches(2.5), Inches(6.0)
    for col, (method_key, method_label) in enumerate(METHODS):
        lx = Inches(0.25) + col * (img_w + Inches(0.06))
        add_text(slide, method_label,
                 lx, Inches(1.1), img_w, Inches(0.28),
                 size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_img(slide, img_fname_fn(method_key),
                lx, Inches(1.4), img_w, img_h)
    return slide


def _slide_4images(prs, title, subtitle, img_fname_fn):
    """4개 방법(METHODS_RESNET1D)의 이미지를 한 행으로 가득 채운 슬라이드."""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, title, subtitle)

    img_w, img_h = Inches(3.1), Inches(6.0)
    for col, (method_key, method_label) in enumerate(METHODS_RESNET1D):
        lx = Inches(0.25) + col * (img_w + Inches(0.09))
        add_text(slide, method_label,
                 lx, Inches(1.1), img_w, Inches(0.28),
                 size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_img(slide, img_fname_fn(method_key),
                lx, Inches(1.4), img_w, img_h)
    return slide


# ── 슬라이드 빌더 ─────────────────────────────────────────────
def slide_title(prs):
    slide = new_blank_slide(prs)
    add_rect(slide, 0, 0, W, H, NAVY)
    add_text(slide, 'TAMA / SMI 예측 회귀분석 연구 보고서',
             Inches(1), Inches(1.8), Inches(11.3), Inches(1.4),
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, '0508 버전',
             Inches(1), Inches(3.3), Inches(11.3), Inches(0.7),
             size=20, color=SKYBLUE, align=PP_ALIGN.CENTER)
    add_text(slide,
             '분석 방법: clinic_only / aec_only / aec_feature_only / aec_feature_selected / clinic_aec_feature\n'
             '타겟: TAMA, SMI  |  알고리즘: Linear / Logistic Regression (5-Fold CV) + ResNet1D\n'
             '데이터: 강남 (aec_feature_filtered 시트)',
             Inches(1), Inches(4.2), Inches(11.3), Inches(1.8),
             size=13, color=SKYBLUE, align=PP_ALIGN.CENTER)
    print('  [슬라이드 1] 표지')


def slide_overview(prs):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, '분석 방법 구성',
               'clinic_only / aec_only / aec_feature_only / aec_feature_selected / clinic_aec_feature')

    headers = ['방법', '피처 구성', '피처 수']
    rows = [
        ['clinic_only',          'PatientAge, PatientSex, BMI', '3'],
        ['aec_only',             'aec1 ~ aec256 (원시 AEC 시계열 통계)', '256'],
        ['aec_feature_only',     'aec_feature_filtered 시트 전체 사용', '~65'],
        ['aec_feature_selected', 'AEC 피처 fold별 선택 (|r|<0.8, VIF<10)', '가변'],
        ['clinic_aec_feature',   'Age/Sex/BMI (고정) + AEC 피처 fold별 선택', '가변'],
    ]
    col_widths = [Inches(3.0), Inches(7.0), Inches(2.0)]
    add_table(slide, headers, rows, Inches(0.35), Inches(1.3), col_widths,
              row_h=Inches(0.6), font_size=12)

    add_text(slide,
             '평가 방법: Train 80% 에서 5-Fold CV 수행 → Test 20% 에서 최종 성능 측정',
             Inches(0.35), Inches(4.5), Inches(12.6), Inches(0.45),
             size=12, color=BLUE)

    box_data = [
        ('임상 변수 기준선', 'Age, Sex, BMI만 사용\n→ AEC 추가 효과의 기준점', NAVY),
        ('AEC 원시 피처', '256개 aec 시계열\n→ 원시 특성의 예측력 평가', BLUE),
        ('AEC 피처 전체', '추출된 통계/주파수/웨이블릿\n→ 전체 사용 (~65개)', GREEN),
        ('AEC 피처 선택', 'corr+VIF fold별 선택\n→ 전체 vs 선택 성능 비교', ORANGE),
        ('임상+AEC 선택', 'Age/Sex/BMI + AEC 선택\n→ 임상정보 추가 효과 확인', RGBColor(0x70, 0x30, 0xA0)),
    ]
    bw, gap, ty = Inches(2.4), Inches(0.12), Inches(5.1)
    for i, (ttl, desc, color) in enumerate(box_data):
        lx = Inches(0.35) + i * (bw + gap)
        add_rect(slide, lx, ty, bw, Inches(1.6), color)
        add_text(slide, ttl, lx + Inches(0.08), ty + Inches(0.1), bw - Inches(0.16), Inches(0.44),
                 size=10, bold=True, color=WHITE)
        add_text(slide, desc, lx + Inches(0.08), ty + Inches(0.55), bw - Inches(0.16), Inches(0.95),
                 size=9, color=WHITE, wrap=True)
    print('  [슬라이드 2] 분석 개요')


# ── Linear Regression ────────────────────────────────────────
def slide_linear_table(prs, target: str, data: dict):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide,
               f'[{target}] 선형 회귀 — 성능 지표',
               '5-Fold CV (Train 80%) 평균±표준편차  |  Test Set (20%) 최종 성능')

    headers = ['방법', 'CV R² (mean±std)', 'CV RMSE (mean±std)', 'Test R²', 'Test RMSE', '피처']
    rows = []
    for method_key, method_label in METHODS:
        d = data[method_key][target]['linear']
        if not d:
            rows.append([method_label, 'N/A', 'N/A', 'N/A', 'N/A', 'N/A'])
            continue
        r2        = f"{_fmt(d.get('r2_mean'))} ± {_fmt(d.get('r2_std'))}"
        rmse      = f"{_fmt(d.get('rmse_mean'), 3)} ± {_fmt(d.get('rmse_std'), 3)}"
        test_r2   = _fmt(d.get('test_r2'))
        test_rmse = _fmt(d.get('test_rmse'), 3)
        feat      = d.get('features', '—')[:45]
        rows.append([method_label, r2, rmse, test_r2, test_rmse, feat])

    col_widths = [Inches(2.6), Inches(2.4), Inches(2.4), Inches(1.1), Inches(1.1), Inches(2.6)]
    add_table(slide, headers, rows, Inches(0.35), Inches(1.3), col_widths,
              row_h=Inches(0.65), font_size=11)
    print(f'  [슬라이드] {target} 선형 회귀 — 표')


def slide_linear_actual_vs_pred(prs, target: str):
    _slide_5images(
        prs,
        f'[{target}] 선형 회귀 — Actual vs Predicted',
        '각 방법별 예측값 vs 실제값 산점도 (Test Set 20%)',
        lambda mk: BASE_DIR / mk / target / 'Linear_Actual_vs_Predicted.png',
    )
    print(f'  [슬라이드] {target} 선형 회귀 — Actual vs Predicted')


def slide_linear_coefficients(prs, target: str):
    _slide_5images(
        prs,
        f'[{target}] 선형 회귀 — 회귀 계수',
        '각 방법별 회귀 계수(Coefficient) 시각화',
        lambda mk: BASE_DIR / mk / target / 'Linear_Coefficients.png',
    )
    print(f'  [슬라이드] {target} 선형 회귀 — 회귀 계수')


# ── Logistic Regression ──────────────────────────────────────
def slide_logistic_table(prs, target: str, data: dict):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide,
               f'[{target}] 로지스틱 회귀 — 성능 지표',
               '5-Fold CV (Train 80%)  |  Test Set (20%)  |  이진화: 성별 내 하위 25% = Low')

    headers = ['방법', 'CV AUC (mean±std)', 'CV Acc (mean±std)',
               'Test AUC', 'Test Acc', 'Male Thr', 'Female Thr']
    rows = []
    for method_key, method_label in METHODS:
        d = data[method_key][target]['logistic']
        if not d:
            rows.append([method_label, 'N/A', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A'])
            continue
        auc      = f"{_fmt(d.get('auc_mean'))} ± {_fmt(d.get('auc_std'))}"
        acc      = f"{_fmt(d.get('acc_mean'))} ± {_fmt(d.get('acc_std'))}"
        test_auc = _fmt(d.get('test_auc'))
        test_acc = _fmt(d.get('test_acc'))
        thr      = d.get('thresholds', {})
        rows.append([method_label, auc, acc, test_auc, test_acc,
                     _fmt(thr.get('Male'), 2), _fmt(thr.get('Female'), 2)])

    col_widths = [Inches(2.4), Inches(2.2), Inches(2.2), Inches(1.1), Inches(1.1), Inches(1.2), Inches(1.2)]
    add_table(slide, headers, rows, Inches(0.35), Inches(1.3), col_widths,
              row_h=Inches(0.65), font_size=11)
    print(f'  [슬라이드] {target} 로지스틱 회귀 — 표')


def slide_logistic_roc(prs, target: str):
    _slide_5images(
        prs,
        f'[{target}] 로지스틱 회귀 — ROC Curve',
        '각 방법별 ROC Curve (AUC-ROC 비교, Test Set 20%)',
        lambda mk: BASE_DIR / mk / target / 'Logistic_ROC_Curve.png',
    )
    print(f'  [슬라이드] {target} 로지스틱 회귀 — ROC Curve')


def slide_logistic_confusion(prs, target: str):
    _slide_5images(
        prs,
        f'[{target}] 로지스틱 회귀 — Confusion Matrix',
        '각 방법별 Confusion Matrix (Test Set 20%)',
        lambda mk: BASE_DIR / mk / target / 'Logistic_Confusion_Matrix.png',
    )
    print(f'  [슬라이드] {target} 로지스틱 회귀 — Confusion Matrix')


# ── ResNet1D ─────────────────────────────────────────────────
def slide_resnet1d_table(prs, target: str, data: dict):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide,
               f'[{target}] ResNet1D (1D CNN) — 성능 지표',
               '5-Fold CV (Train 80%) 평균±표준편차  |  Test Set (20%) 최종 성능')

    headers = ['방법', 'CV R² (mean±std)', 'CV MAE (mean±std)', 'Test R²', 'Test MAE']
    rows = []
    for method_key, method_label in METHODS_RESNET1D:
        d = data[method_key][target]['resnet1d']
        if not d:
            rows.append([method_label, 'N/A', 'N/A', 'N/A', 'N/A'])
            continue
        r2       = f"{_fmt(d.get('r2_mean'))} ± {_fmt(d.get('r2_std'))}"
        mae      = f"{_fmt(d.get('mae_mean'), 3)} ± {_fmt(d.get('mae_std'), 3)}"
        test_r2  = _fmt(d.get('test_r2'))
        test_mae = _fmt(d.get('test_mae'), 3)
        rows.append([method_label, r2, mae, test_r2, test_mae])

    col_widths = [Inches(3.2), Inches(3.0), Inches(3.0), Inches(1.6), Inches(1.6)]
    add_table(slide, headers, rows, Inches(0.35), Inches(1.3), col_widths,
              row_h=Inches(0.7), font_size=12)
    print(f'  [슬라이드] {target} ResNet1D — 표')


def slide_resnet1d_plots(prs, target: str):
    _slide_4images(
        prs,
        f'[{target}] ResNet1D (1D CNN) — 예측 결과',
        '각 방법별 Test Set 예측 결과 (Actual vs Predicted)',
        lambda mk: BASE_DIR / mk / target / 'ResNet1D_Results.png',
    )
    print(f'  [슬라이드] {target} ResNet1D — 예측 결과')


# ── 분포 & 산점도 ─────────────────────────────────────────────
def slide_distribution(prs, target: str):
    _slide_5images(
        prs,
        f'[{target}] 성별 분포',
        '방법별 성별(Sex) 분포 비교',
        lambda mk: BASE_DIR / mk / target / 'Distribution_by_Sex.png',
    )
    print(f'  [슬라이드] {target} 성별 분포')


def slide_scatter(prs, target: str):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide,
               f'[{target}] 산점도 — BMI / Age vs {target}',
               '방법별 BMI, PatientAge vs 타겟 산점도 (성별 구분)')

    img_w, img_h = Inches(2.4), Inches(2.82)
    row_data = [
        (f'BMI_vs_{target}_by_Sex.png',        f'BMI vs {target}'),
        (f'PatientAge_vs_{target}_by_Sex.png',  f'Age vs {target}'),
    ]
    ty_label = Inches(1.12)
    label_h  = Inches(0.28)
    row_gap  = Inches(0.15)

    # 열 방법 레이블 (한 번만)
    for col, (_, method_label) in enumerate(METHODS):
        lx = Inches(0.85) + col * (img_w + Inches(0.05))
        add_text(slide, method_label, lx, ty_label, img_w, label_h,
                 size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    for ri, (fname, row_label) in enumerate(row_data):
        ty_img = ty_label + label_h + ri * (img_h + row_gap)
        add_text(slide, row_label,
                 Inches(0.05), ty_img + Inches(0.9), Inches(0.76), Inches(1.0),
                 size=8, bold=True, color=NAVY)
        for col, (method_key, _) in enumerate(METHODS):
            lx = Inches(0.85) + col * (img_w + Inches(0.05))
            add_img(slide, BASE_DIR / method_key / target / fname,
                    lx, ty_img, img_w, img_h)

    print(f'  [슬라이드] {target} 산점도')


# ── 비교 그래프 ───────────────────────────────────────────────
def slide_comparison_linear(prs):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, '방법별 성능 비교 — 선형 회귀',
               '5-Fold CV R² / RMSE  ×  5 방법  ×  TAMA / SMI')
    add_img(slide, FIG_DIR / '01_linear_comparison.png',
            Inches(0.35), Inches(1.15), Inches(12.6), Inches(6.1))
    print('  [슬라이드] 비교 그래프 — Linear')


def slide_comparison_logistic(prs):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, '방법별 성능 비교 — 로지스틱 회귀',
               '5-Fold CV AUC / Accuracy  ×  5 방법  ×  TAMA / SMI')
    add_img(slide, FIG_DIR / '02_logistic_comparison.png',
            Inches(0.35), Inches(1.15), Inches(12.6), Inches(6.1))
    print('  [슬라이드] 비교 그래프 — Logistic')


def slide_comparison_resnet1d(prs):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, '방법별 성능 비교 — ResNet1D',
               '5-Fold CV R² / MAE  ×  4 방법  ×  TAMA / SMI')
    add_img(slide, FIG_DIR / '04_resnet1d_comparison.png',
            Inches(0.35), Inches(1.15), Inches(12.6), Inches(6.1))
    print('  [슬라이드] 비교 그래프 — ResNet1D')


def slide_overview_plot(prs):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, '전체 성능 지표 개요',
               '6개 지표 × 5 방법 × 2 타겟 한 눈에 보기')
    add_img(slide, FIG_DIR / '03_overview.png',
            Inches(0.35), Inches(1.15), Inches(12.6), Inches(6.1))
    print('  [슬라이드] 전체 개요')


def slide_conclusion(prs, data: dict):
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header(slide, '결론 및 핵심 성과 (0508)', '')

    achievements = []
    for target in TARGETS:
        best_r2_m, best_r2 = None, -999
        best_auc_m, best_auc = None, -999
        for mk, _ in METHODS:
            r2  = data[mk][target]['linear'].get('r2_mean', -999)
            auc = data[mk][target]['logistic'].get('auc_mean', -999)
            if r2 > best_r2:
                best_r2, best_r2_m = r2, mk
            if auc > best_auc:
                best_auc, best_auc_m = auc, mk
        achievements.append(f'{target} 선형 CV R²  최고: {best_r2_m} ({_fmt(best_r2)})')
        achievements.append(f'{target} AUC-ROC CV 최고: {best_auc_m} ({_fmt(best_auc)})')

    y_pos = Inches(1.2)
    for ach in achievements:
        add_rect(slide, Inches(0.35), y_pos, Inches(12.6), Inches(0.72), SKYBLUE)
        add_text(slide, ach,
                 Inches(0.55), y_pos + Inches(0.12), Inches(12.2), Inches(0.5),
                 size=13, color=DARK)
        y_pos += Inches(0.85)

    add_text(slide, '한계 및 향후 과제',
             Inches(0.35), y_pos + Inches(0.1), Inches(12.6), Inches(0.4),
             size=13, bold=True, color=NAVY)
    limits = (
        '• 단일 기관(강남) 데이터 — 외부 검증(신촌) 필요\n'
        '• aec_feature_only 피처 선택의 통계적 근거 보강 필요\n'
        '• 1D CNN / LSTM 등 딥러닝 모델과의 성능 비교 향후 과제'
    )
    add_text(slide, limits,
             Inches(0.35), y_pos + Inches(0.55), Inches(12.6), Inches(1.2),
             size=11, color=DARK)

    add_text(slide, '자동 생성: 6_generate_ppt.py (0508)',
             Inches(0.35), Inches(7.2), Inches(12.6), Inches(0.25),
             size=8, italic=True, color=RGBColor(0xBA, 0xB0, 0xAC))
    print('  [슬라이드] 결론')


# ── 메인 ─────────────────────────────────────────────────────
def main():
    print('=' * 55)
    print('  6_generate_ppt.py — PPT 생성')
    print('=' * 55)

    print('\n[1] 데이터 수집 중...')
    data = collect_data()

    print('\n[2] 슬라이드 빌드 중...')
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)    # 1
    slide_overview(prs) # 2

    for target in TARGETS:
        # 선형 회귀 (3장)
        slide_linear_table(prs, target, data)
        slide_linear_actual_vs_pred(prs, target)
        slide_linear_coefficients(prs, target)
        # 로지스틱 회귀 (3장)
        slide_logistic_table(prs, target, data)
        slide_logistic_roc(prs, target)
        slide_logistic_confusion(prs, target)
        # ResNet1D (2장)
        slide_resnet1d_table(prs, target, data)
        slide_resnet1d_plots(prs, target)
        # 분포 & 산점도 (2장)
        slide_distribution(prs, target)
        slide_scatter(prs, target)

    # 비교 그래프 (3장) + 개요 (1장) + 결론 (1장)
    slide_comparison_linear(prs)
    slide_comparison_logistic(prs)
    slide_comparison_resnet1d(prs)
    slide_overview_plot(prs)
    slide_conclusion(prs, data)

    print(f'\n[3] PPT 저장 중...')
    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f'  저장 완료: {OUT_PPT}')
    print(f'  슬라이드 수: {len(prs.slides)}')
    print('=' * 55)


if __name__ == '__main__':
    main()
