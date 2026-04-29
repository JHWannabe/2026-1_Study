"""
0424 vs 0430 연구 설계 변경사항 비교 PPT 생성 스크립트
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 경로 ─────────────────────────────────────────────────────────────────────
SCRIPT_DIR  = Path(__file__).parent
RESULT_ROOT = SCRIPT_DIR.parent.parent / "results"
FS_DIR      = RESULT_ROOT / "feature_selection"
REG_DIR     = RESULT_ROOT / "regression"

OUTPUT_PPT  = RESULT_ROOT / "0424_vs_0430_comparison.pptx"

# ── 색상 팔레트 ───────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x1A, 0x35, 0x5E)
BLUE     = RGBColor(0x2E, 0x75, 0xB6)
SKYBLUE  = RGBColor(0xBD, 0xD7, 0xEE)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1A, 0x1A, 0x2E)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
ORANGE   = RGBColor(0xED, 0x7D, 0x31)
GREEN    = RGBColor(0x00, 0x8B, 0x45)
RED      = RGBColor(0xC0, 0x00, 0x00)
GOLD     = RGBColor(0xFF, 0xC0, 0x00)

# ── 슬라이드 크기 (16:9 widescreen) ─────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# 유틸리티
# ─────────────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, italic=False,
             color: RGBColor = DARK, align=PP_ALIGN.LEFT,
             wrap=True, valign=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if valign:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = valign
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txb


def add_image_safe(slide, img_path, l, t, w, h=None):
    if not os.path.exists(img_path):
        add_rect(slide, l, t, w, h or 2, LGRAY)
        add_text(slide, f"[이미지 없음]\n{os.path.basename(img_path)}",
                 l + 0.1, t + 0.1, w - 0.2, (h or 2) - 0.2,
                 font_size=10, color=RED)
        return
    pic_args = dict(image_file=str(img_path), left=Inches(l), top=Inches(t), width=Inches(w))
    if h:
        pic_args["height"] = Inches(h)
    slide.shapes.add_picture(**pic_args)


def slide_header(slide, title, subtitle=None):
    """상단 네이비 헤더 바 + 제목 텍스트"""
    add_rect(slide, 0, 0, 13.33, 1.05, NAVY)
    add_text(slide, title, 0.3, 0.1, 12.5, 0.75,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.82, 12.5, 0.35,
                 font_size=13, color=SKYBLUE, align=PP_ALIGN.LEFT)


def add_table(slide, headers, rows, l, t, w, h,
              header_fill=NAVY, header_color=WHITE,
              row_fills=None, font_size=12, header_font_size=12):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    col_w = int(Inches(w) / n_cols)
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    def cell_fill(cell, color: RGBColor):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def cell_text(cell, text, bold=False, color=DARK, size=12, align=PP_ALIGN.CENTER):
        cell.text = text
        tf = cell.text_frame
        tf.paragraphs[0].alignment = align
        for p in tf.paragraphs:
            p.alignment = align
            for r in p.runs:
                r.font.size  = Pt(size)
                r.font.bold  = bold
                r.font.color.rgb = color

    for ci, h_txt in enumerate(headers):
        c = tbl.cell(0, ci)
        cell_fill(c, header_fill)
        cell_text(c, h_txt, bold=True, color=header_color,
                  size=header_font_size, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        bg = (row_fills[ri] if row_fills and ri < len(row_fills) else
              (LGRAY if ri % 2 == 0 else WHITE))
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            cell_fill(c, bg)
            cell_text(c, str(val), size=font_size,
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    return tbl


def badge(slide, text, l, t, fill: RGBColor = BLUE):
    w = max(1.2, len(text) * 0.18 + 0.4)
    add_rect(slide, l, t, w, 0.35, fill)
    add_text(slide, text, l + 0.05, t + 0.02, w - 0.1, 0.31,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return w


def change_block(slide, num, title, l, t, w=6.2):
    """변경사항 번호 + 제목 블록"""
    badge(slide, f"변경 {num}", l, t, ORANGE)
    add_text(slide, title, l + 1.55, t, w - 1.55, 0.4,
             font_size=15, bold=True, color=NAVY)


def before_after_label(slide, l, t):
    badge(slide, "0424 이전", l, t, RED)
    badge(slide, "0430 이후", l + 3.3, t, GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1. Title
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # 중앙 흰 카드
    add_rect(sl, 1.5, 1.8, 10.33, 4.0, WHITE)

    add_text(sl, "연구 설계 변경사항 비교", 1.7, 2.1, 9.9, 0.9,
             font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, "0424  →  0430", 1.7, 3.0, 9.9, 0.7,
             font_size=30, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(sl, "TAMA 예측 회귀분석 연구  |  AEC 신호 피처 기반 근감소증 예측", 1.7, 3.75, 9.9, 0.5,
             font_size=16, color=DARK, align=PP_ALIGN.CENTER)

    bullets = [
        "① 피처 선택: 수동  →  자동 파이프라인 (4단계)",
        "② 임상 기준선: Age+Sex  →  Age+Sex+BMI",
        "③ Case 구조: 3단계  →  5단계 (AEC_prev vs AEC_new)",
        "④ 다병원 자동 순회 + 교차 병원 비교",
    ]
    for i, b in enumerate(bullets):
        add_text(sl, b, 2.1, 4.35 + i * 0.32, 9.0, 0.3,
                 font_size=13, color=DARK, align=PP_ALIGN.LEFT)

    add_text(sl, "2026-04-30", 11.8, 7.1, 1.5, 0.3,
             font_size=11, color=SKYBLUE, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2. 5가지 변경사항 한눈에 보기
# ─────────────────────────────────────────────────────────────────────────────

def slide_overview(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "변경사항 한눈에 보기",
                 "0424 → 0430  |  총 4가지 설계 변경")

    headers = ["#", "항목", "0424 이전", "0430 이후", "핵심 목적"]
    rows = [
        ["①", "피처 선택",
         "수동 선택\n(상관계수+VIF → 연구자 결정)\n['mean','CV','skewness','slope_abs_mean']",
         "자동 파이프라인\n(4단계 필터링 + 앙상블 투표 + CV R²)",
         "다중공선성·과적합 위험 제거,\n객관적 피처 선택"],
        ["②", "임상 기준선",
         "Age + Sex",
         "Age + Sex + BMI",
         "BMI 교란변수 보정\n→ AEC 순수 기여도 분리"],
        ["③", "Case 구조",
         "Case 1(임상) / 2(+AEC) / 3(+AEC+스캐너)\n단일 AEC 세트",
         "Case 1~5\nAEC_prev vs AEC_new 교차 비교",
         "'어떤 AEC 피처 세트가 더 유용한가'\n직접 정량 비교"],
        ["④", "다병원 분석",
         "강남 단독\n(config.SITE 수동 변경)",
         "강남·신촌 자동 순회 +\n교차 병원 비교(외부 검증)",
         "재현성(generalizability) 확인"],
    ]

    fills = [LGRAY, WHITE, LGRAY, WHITE]
    add_table(sl, headers, rows,
              l=0.15, t=1.15, w=13.0, h=6.15,
              header_fill=NAVY, header_font_size=13, font_size=11,
              row_fills=fills)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3. ① 피처 선택: 수동 → 자동 파이프라인
# ─────────────────────────────────────────────────────────────────────────────

def slide_feature_selection(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "① 피처 선택: 수동 → 자동 파이프라인")

    # 왼쪽: 0424
    add_rect(sl, 0.2, 1.15, 5.5, 5.85, RGBColor(0xFF, 0xF0, 0xF0))
    badge(sl, "0424 수동 선택", 0.3, 1.2, RED)
    lines_old = [
        "• 전체 AEC 피처(60개+) vs TAMA",
        "  Pearson / Spearman 상관계수 + VIF 계산",
        "",
        "• feature_selection_report.xlsx 검토 후",
        "  연구자가 수동으로 결정",
        "",
        "• 단일 병원(강남)만 분석",
        "",
        "• 선택 결과 (config.py에 하드코딩):",
        "  ['mean', 'CV', 'skewness', 'slope_abs_mean']",
        "",
        "▲ 한계: 60개+ 피처에서 단순 상관계수만으로",
        "   선택 시 과적합·다중공선성 위험",
    ]
    for i, ln in enumerate(lines_old):
        add_text(sl, ln, 0.3, 1.65 + i * 0.33, 5.2, 0.33,
                 font_size=11, color=DARK if "▲" not in ln else RED)

    # 가운데 화살표
    add_text(sl, "→", 5.85, 3.8, 0.6, 0.5,
             font_size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # 오른쪽: 0430 파이프라인
    add_rect(sl, 6.55, 1.15, 6.6, 5.85, RGBColor(0xF0, 0xFF, 0xF0))
    badge(sl, "0430 자동 파이프라인", 6.65, 1.2, GREEN)

    steps = [
        ("Step 1", "Near-zero variance 제거\n(표준화 후 threshold=0.01)"),
        ("Step 2", "Pearson |r| ≥ 0.95 중복 피처 제거\n(VIF는 참고용 보고만)"),
        ("Step 3", "단변량 필터: MI > 0  OR  Spearman p < 0.05\n(상한 없음 — 후보 최대 보존)"),
        ("Step 4", "모델 기반 앙상블 투표\n  LASSO + RFECV(Ridge) + RF Permutation\n→ pool 구성 후 CV R² 최적화\n  (pool ≤ 20: 완전탐색 | >20: SFS)"),
        ("Final", "VIF > 10 반복 제거 (FIXED: 'mean' 강제 보존)\n→ 강남 / 신촌 / 병합 각각 독립 실행"),
    ]
    colors = [BLUE, BLUE, BLUE, ORANGE, GREEN]
    y_pos = 1.65
    for (tag, desc), clr in zip(steps, colors):
        badge(sl, tag, 6.65, y_pos, clr)
        add_text(sl, desc, 8.25, y_pos - 0.03, 4.7, 0.55,
                 font_size=11, color=DARK)
        y_pos += 0.72

    add_text(sl, "★ 고정 피처: 'mean'은 임상 해석을 위해 모든 단계에서 강제 보존",
             6.65, 6.65, 6.3, 0.35,
             font_size=11, bold=True, color=NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4. ① 피처 선택 결과 이미지
# ─────────────────────────────────────────────────────────────────────────────

def slide_fs_results(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "① 피처 선택 결과: AEC_prev vs AEC_new (CV R²)",
                 "3개 데이터셋(강남 / 신촌 / 병합)에서 이전 수동 선택과 파이프라인 결과 비교")

    # 좌: cross-dataset r2 비교
    add_text(sl, "Previous vs Pipeline — 3개 데이터셋 CV R² 비교",
             0.2, 1.1, 6.5, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, FS_DIR / "cross_dataset_comparison_r2.png",
                   0.2, 1.5, 6.5, 3.8)

    # 우: feature heatmap
    add_text(sl, "데이터셋별 최종 선택 피처 히트맵",
             7.0, 1.1, 6.1, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, FS_DIR / "cross_dataset_feature_heatmap.png",
                   7.0, 1.5, 6.1, 3.8)

    # 하단 설명
    add_rect(sl, 0.2, 5.5, 12.9, 1.75, RGBColor(0xEE, 0xF5, 0xFF))
    add_text(sl, "해석 포인트",
             0.35, 5.55, 3.0, 0.35, font_size=12, bold=True, color=NAVY)
    notes = [
        "• 파이프라인 CV R²가 이전 수동 선택 대비 동등 또는 개선",
        "• 히트맵: 1=선택됨 / 0=미선택. 데이터셋 간 공통 피처가 재현성 지표",
        "• 강남·신촌 모두에서 선택된 피처 → 병합 데이터에서도 안정적으로 기여",
    ]
    for i, n in enumerate(notes):
        add_text(sl, n, 0.35, 5.95 + i * 0.3, 12.5, 0.3, font_size=11, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5. ② BMI 공변량 추가
# ─────────────────────────────────────────────────────────────────────────────

def slide_bmi(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "② 임상 기준선: BMI 공변량 추가")

    # 비교 블록
    add_rect(sl, 0.2, 1.2, 5.9, 5.8, RGBColor(0xFF, 0xF0, 0xF0))
    badge(sl, "0424", 0.35, 1.28, RED)
    add_text(sl, "기준선 = PatientAge + PatientSex",
             0.35, 1.7, 5.5, 0.45, font_size=15, bold=True, color=DARK)

    lines_old = [
        "Case 1: [Age, Sex]",
        "Case 2: [Age, Sex, AEC 피처]",
        "Case 3: [Age, Sex, AEC 피처, Scanner]",
        "",
        "→ BMI 정보 미포함",
        "→ 체지방량·근육량과 직접 연관된",
        "   주요 교란변수 미보정",
        "",
        "▲ 문제: AEC 피처의 독립 예측력이",
        "   BMI에 의해 과대추정될 가능성",
    ]
    for i, ln in enumerate(lines_old):
        add_text(sl, ln, 0.4, 2.2 + i * 0.35, 5.5, 0.35,
                 font_size=12,
                 color=RED if "▲" in ln else DARK)

    add_text(sl, "→", 6.25, 3.8, 0.7, 0.5,
             font_size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_rect(sl, 7.1, 1.2, 6.0, 5.8, RGBColor(0xF0, 0xFF, 0xF0))
    badge(sl, "0430", 7.25, 1.28, GREEN)
    add_text(sl, "기준선 = PatientAge + PatientSex + BMI",
             7.25, 1.7, 5.7, 0.45, font_size=15, bold=True, color=DARK)

    lines_new = [
        "Case 1: [Age, Sex, BMI]",
        "Case 2: [Age, Sex, BMI, AEC_prev]",
        "Case 3: [Age, Sex, BMI, AEC_new]",
        "Case 4: [Age, Sex, BMI, AEC_prev, Scanner]",
        "Case 5: [Age, Sex, BMI, AEC_new, Scanner]",
        "",
        "→ BMI = 체지방량·제지방량과 직접 연관",
        "   근감소증의 핵심 교란변수",
        "",
        "★ 효과: AEC 피처의 순수 기여도를",
        "   BMI 보정 후 분리하여 정량화",
    ]
    for i, ln in enumerate(lines_new):
        add_text(sl, ln, 7.25, 2.2 + i * 0.35, 5.6, 0.35,
                 font_size=12,
                 color=GREEN if "★" in ln else DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6. ② BMI 기여도 분석 결과
# ─────────────────────────────────────────────────────────────────────────────

def slide_bmi_results(prs):
    """BMI 유무에 따른 Case 1/2/4 성능 비교 결과 슬라이드"""
    import subprocess, tempfile
    from pathlib import Path as _Path

    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "② BMI 기여도 분석 결과 (Case 1 / 2 / 4)",
                 "Clinical baseline에 BMI 추가 시 R² / AUC 변화  |  강남·신촌 비교")

    # ── 이미지: 강남 좌 / 신촌 우 ────────────────────────────────────────────
    add_text(sl, "강남 — R² & AUC 비교 (BMI 없음 vs 있음)",
             0.2, 1.1, 6.4, 0.32, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "bmi_comparison_r2_auc.png",
                   0.2, 1.45, 6.4, 2.85)

    add_text(sl, "신촌 — R² & AUC 비교 (BMI 없음 vs 있음)",
             6.85, 1.1, 6.3, 0.32, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "sinchon" / "bmi_comparison_r2_auc.png",
                   6.85, 1.45, 6.3, 2.85)

    # ── Delta 표: Excel에서 로드 ──────────────────────────────────────────────
    def _load_bmi_delta():
        bmi_path = REG_DIR / "bmi_comparison_summary.xlsx"
        if not bmi_path.exists():
            return None
        tmp = _Path(tempfile.gettempdir()) / "ppt_bmi_delta.xlsx"
        subprocess.run(
            ["powershell", "-Command",
             f'Copy-Item -Path "{bmi_path}" -Destination "{tmp}" -Force'],
            capture_output=True,
        )
        if not tmp.exists():
            return None
        import pandas as pd
        return pd.read_excel(str(tmp), sheet_name="delta_bmi")

    delta_df = _load_bmi_delta()

    add_text(sl, "BMI 추가 효과 (Delta = with BMI − without BMI)",
             0.2, 4.42, 12.9, 0.32, font_size=12, bold=True, color=NAVY)

    if delta_df is not None and len(delta_df) > 0:
        # 컬럼 이름 유연하게 처리
        cols = delta_df.columns.tolist()
        hospital_col = cols[0]
        case_col = cols[1] if len(cols) > 1 else None

        headers_d = ["병원", "Case",
                     "ΔR² (선형)", "ΔRMSE (선형)",
                     "ΔAUC (로지스틱)", "ΔAcc (로지스틱)"]
        rows_d = []
        fills_d = []
        for _, row in delta_df.iterrows():
            hosp = str(row.get(hospital_col, ""))
            case = str(row.get(case_col, "")) if case_col else ""
            dr2  = float(row.get("Delta_Lin_R2",   row.get("delta_lin_r2",   0)))
            drmse = float(row.get("Delta_Lin_RMSE", row.get("delta_lin_rmse", 0)))
            dauc = float(row.get("Delta_Log_AUC",  row.get("delta_log_auc",  0)))
            dacc = float(row.get("Delta_Log_Acc",  row.get("delta_log_acc",  0)))
            rows_d.append([
                hosp, case,
                f"{dr2:+.4f}", f"{drmse:+.4f}",
                f"{dauc:+.4f}", f"{dacc:+.4f}",
            ])
            if dr2 >= 0.05:
                fills_d.append(RGBColor(0xD6, 0xF5, 0xD6))   # 진한 녹: 유효
            elif dr2 >= 0.01:
                fills_d.append(RGBColor(0xEB, 0xFF, 0xEB))   # 연녹: 소폭
            else:
                fills_d.append(RGBColor(0xF8, 0xF8, 0xF8))   # 회색: 미미

        add_table(sl, headers_d, rows_d,
                  l=0.2, t=4.78, w=12.9, h=2.55,
                  header_fill=NAVY, font_size=11, header_font_size=11,
                  row_fills=fills_d)
    else:
        # 이미지 fallback: delta effect bars
        add_image_safe(sl, REG_DIR / "gangnam" / "bmi_delta_effect.png",
                       0.2, 4.78, 6.4, 2.55)
        add_image_safe(sl, REG_DIR / "sinchon" / "bmi_delta_effect.png",
                       6.85, 4.78, 6.3, 2.55)

    # 범례 설명 (표 아래 footer)
    add_rect(sl, 0.2, 7.25, 12.9, 0.22, RGBColor(0xEE, 0xF5, 0xFF))
    add_text(sl,
             "진녹(ΔR²≥0.05): BMI 보정 유효   연녹(ΔR²≥0.01): 소폭 향상   회색(ΔR²<0.01): 미미 | "
             "Case1=임상기준선, Case2=+AEC_prev, Case4=+AEC_prev+Scanner",
             0.3, 7.27, 12.7, 0.2, font_size=9, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7. ③ Case 구조 확장 (3 → 5 Cases)
# ─────────────────────────────────────────────────────────────────────────────

def slide_case_structure(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "③ Case 구조 확장: AEC_prev vs AEC_new 이중 비교",
                 "0424: 3-Case (단일 AEC)  →  0430: 5-Case (AEC_prev × AEC_new 교차 비교)")

    # 0424 케이스 표
    add_text(sl, "▼ 0424 Case 구조 (3단계, 단일 AEC 세트)",
             0.2, 1.15, 6.4, 0.35, font_size=12, bold=True, color=RED)
    old_headers = ["Case", "포함 변수", "목적"]
    old_rows = [
        ["Case 1", "Age, Sex", "임상 기준선"],
        ["Case 2", "Age, Sex + AEC 4개", "AEC 기여도 확인"],
        ["Case 3", "Age, Sex + AEC + Scanner", "스캐너 보정 효과"],
    ]
    add_table(sl, old_headers, old_rows,
              l=0.2, t=1.55, w=6.4, h=1.55,
              header_fill=RED, font_size=11, header_font_size=11)

    add_text(sl, "▼ 0430 Case 구조 (5단계, AEC_prev vs AEC_new 교차)",
             0.2, 3.25, 12.9, 0.35, font_size=12, bold=True, color=GREEN)
    new_headers = ["Case", "포함 변수", "AEC 세트", "목적"]
    new_rows = [
        ["Case 1", "Age, Sex, BMI", "—", "BMI 포함 임상 기준선"],
        ["Case 2", "Case 1 + AEC_prev", "수동 4개\n(mean,CV,skewness,slope_abs_mean)", "이전 AEC 기여도"],
        ["Case 3", "Case 1 + AEC_new", "자동 파이프라인 선택", "새 AEC 기여도"],
        ["Case 4", "Case 2 + Scanner", "AEC_prev + 스캐너", "이전 AEC + 스캐너"],
        ["Case 5", "Case 3 + Scanner", "AEC_new + 스캐너", "새 AEC + 스캐너 (최대 모델)"],
    ]
    fills_new = [LGRAY, RGBColor(0xFF, 0xEB, 0xEB), RGBColor(0xEB, 0xFF, 0xEB),
                 RGBColor(0xFF, 0xEB, 0xEB), RGBColor(0xEB, 0xFF, 0xEB)]
    add_table(sl, new_headers, new_rows,
              l=0.2, t=3.65, w=12.9, h=3.45,
              header_fill=GREEN, font_size=11, header_font_size=11,
              row_fills=fills_new)

    add_text(sl, "Case 2 vs Case 3, Case 4 vs Case 5 → AEC_prev vs AEC_new 직접 성능 비교 가능",
             0.2, 7.1, 12.9, 0.3,
             font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7. ③ Case 비교 결과 이미지
# ─────────────────────────────────────────────────────────────────────────────

def slide_case_results(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "③ Case 비교 결과 (강남 – 전체)",
                 "선형 회귀: R²/RMSE/MAE  |  로지스틱 회귀: AUC/Accuracy")

    # 선형 회귀 메트릭 비교
    add_text(sl, "선형 회귀 — Case별 메트릭 비교",
             0.2, 1.1, 6.5, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "02_linear_metrics_comparison.png",
                   0.2, 1.5, 6.5, 2.7)

    # 로지스틱 ROC
    add_text(sl, "로지스틱 회귀 — Case별 메트릭 비교",
             7.0, 1.1, 6.1, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "05_logistic_metrics_comparison.png",
                   7.0, 1.5, 6.1, 2.7)

    # 케이스 전체 개요
    add_text(sl, "Case 비교 전체 개요",
             0.2, 4.35, 6.5, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "08_case_comparison_overview.png",
                   0.2, 4.75, 6.5, 2.6)

    # ROC 곡선
    add_text(sl, "로지스틱 ROC 곡선 (Case 1~5)",
             7.0, 4.35, 6.1, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "04_logistic_roc.png",
                   7.0, 4.75, 6.1, 2.6)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8. ④ 다병원 순회 + 교차 비교
# ─────────────────────────────────────────────────────────────────────────────

def slide_multi_hospital(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "④ 다병원 자동 순회 + 교차 병원 비교",
                 "0424: 강남 단독  →  0430: 강남·신촌 자동 순회 + 외부 검증")

    # 좌: 설명 블록
    add_rect(sl, 0.2, 1.2, 5.0, 5.8, RGBColor(0xF5, 0xF5, 0xFF))
    badge(sl, "0424", 0.35, 1.28, RED)
    old_items = [
        "• config.py의 SITE 변수를",
        "  연구자가 수동 변경",
        "• 강남만 분석 (신촌 별도 실행 필요)",
        "• 재현성 비교 구조 없음",
    ]
    for i, ln in enumerate(old_items):
        add_text(sl, ln, 0.4, 1.7 + i * 0.33, 4.7, 0.33, font_size=12, color=DARK)

    add_rect(sl, 0.2, 3.2, 5.0, 3.8, RGBColor(0xF0, 0xFF, 0xF0))
    badge(sl, "0430", 0.35, 3.28, GREEN)
    new_items = [
        "• data/ 폴더 자동 스캔",
        "  강남·신촌 모두 순회",
        "• 강남 학습 → 신촌 외부 검증",
        "  (Cross-hospital validation)",
        "• 피처 선택도 3개 데이터셋 독립 실행:",
        "  강남 / 신촌 / 병합(공통 피처)",
        "• 재현성 히트맵으로 공통 피처 확인",
    ]
    for i, ln in enumerate(new_items):
        add_text(sl, ln, 0.4, 3.7 + i * 0.33, 4.7, 0.33, font_size=12, color=DARK)

    # 우: 교차 병원 비교 이미지
    add_text(sl, "교차 병원 비교 (강남 → 신촌 외부 검증)",
             5.45, 1.15, 7.7, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "09_cross_hospital_comparison.png",
                   5.45, 1.55, 7.7, 3.0)

    add_text(sl, "외부 검증 상세 결과",
             5.45, 4.65, 7.7, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "10_external_validation.png",
                   5.45, 5.05, 7.7, 2.3)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10. 선형 회귀 결과 (강남 전체)
# ─────────────────────────────────────────────────────────────────────────────

def slide_linear_results(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "선형 회귀 결과 (강남 – 전체)",
                 "Actual vs Predicted  |  계수 (Coeff)")

    add_text(sl, "Actual vs Predicted (5-fold CV)",
             0.2, 1.1, 6.5, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "01_linear_actual_vs_pred.png",
                   0.2, 1.5, 6.5, 5.7)

    add_text(sl, "회귀 계수 (Standardized Coefficients)",
             7.0, 1.1, 6.1, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "03_linear_coefficients.png",
                   7.0, 1.5, 6.1, 5.7)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 11. 로지스틱 회귀 결과 (강남 전체)
# ─────────────────────────────────────────────────────────────────────────────

def slide_logistic_results(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "로지스틱 회귀 결과 (강남 – 전체)",
                 "ROC  |  Confusion Matrix  |  계수")

    add_text(sl, "ROC Curves (Case 1~5)",
             0.2, 1.1, 4.2, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "04_logistic_roc.png",
                   0.2, 1.5, 4.2, 3.0)

    add_text(sl, "Confusion Matrix (Best Case)",
             4.7, 1.1, 4.2, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "06_logistic_confusion.png",
                   4.7, 1.5, 4.2, 3.0)

    add_text(sl, "로지스틱 계수 (OR ± 95% CI)",
             9.1, 1.1, 4.0, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, REG_DIR / "gangnam" / "all" / "07_logistic_coefficients.png",
                   9.1, 1.5, 4.0, 3.0)

    # 하단: feature selection 비교 이미지 (gangnam)
    add_text(sl, "피처 선택 — CV R² 후보 세트 비교 (강남)",
             0.2, 4.6, 6.5, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, FS_DIR / "gangnam" / "04_cv_r2_comparison.png",
                   0.2, 5.0, 6.5, 2.3)

    add_text(sl, "앙상블 투표 + Spearman 상관 (강남 최종 피처)",
             7.0, 4.6, 6.1, 0.35, font_size=12, bold=True, color=NAVY)
    add_image_safe(sl, FS_DIR / "gangnam" / "05_final_features_summary.png",
                   7.0, 5.0, 6.1, 2.3)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 12. 핵심 수치 요약 (보고서 기반)
# ─────────────────────────────────────────────────────────────────────────────

def slide_key_metrics(prs):
    """generate_report.py 결과 수치를 요약 슬라이드로 정리"""
    import subprocess, tempfile
    from pathlib import Path as _Path

    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    slide_header(sl, "핵심 수치 요약 (강남 전체, 5-Fold CV)",
                 "generate_report.py 결과 기반  |  AEC_prev vs AEC_new 정량 비교")

    # 결과 로드
    def _load(path, sheet=0):
        tmp = _Path(tempfile.gettempdir()) / f"ppt_tmp_{_Path(path).name}"
        subprocess.run(
            ["powershell", "-Command",
             f'Copy-Item -Path "{path}" -Destination "{tmp}" -Force'],
            capture_output=True,
        )
        if not tmp.exists():
            return None
        import pandas as pd
        return pd.read_excel(str(tmp), sheet_name=sheet)

    reg_gn = _load(REG_DIR / "gangnam" / "all" / "regression_results.xlsx", "summary")
    fs_cross = _load(FS_DIR / "cross_dataset_comparison.xlsx", "summary")

    # 왼쪽: 회귀 결과 표
    add_text(sl, "회귀 분석 결과 (강남, 전체, n=1,365)",
             0.2, 1.1, 8.5, 0.35, font_size=13, bold=True, color=NAVY)

    if reg_gn is not None:
        headers = ["Case", "Lin R²", "Lin RMSE", "Log AUC"]
        rows_reg = []
        for _, row in reg_gn.iterrows():
            rows_reg.append([
                str(row.get("Case","?")),
                f"{float(row.get('Lin_R2',0)):.4f}",
                f"{float(row.get('Lin_RMSE',0)):.2f}",
                f"{float(row.get('Log_AUC',0)):.4f}",
            ])
        # 강조: best R2 행 (Case 2)
        best_r2_idx = max(range(len(rows_reg)), key=lambda i: float(rows_reg[i][1]))
        best_auc_idx = max(range(len(rows_reg)), key=lambda i: float(rows_reg[i][3]))
        row_fills_reg = []
        for i in range(len(rows_reg)):
            if i == best_r2_idx and i == best_auc_idx:
                row_fills_reg.append(RGBColor(0xEB, 0xFF, 0xEB))
            elif i == best_r2_idx:
                row_fills_reg.append(RGBColor(0xEB, 0xF5, 0xFF))
            elif i == best_auc_idx:
                row_fills_reg.append(RGBColor(0xFF, 0xF5, 0xEB))
            else:
                row_fills_reg.append(RGBColor(0xF8, 0xF8, 0xF8))

        add_table(sl, headers, rows_reg,
                  l=0.2, t=1.5, w=8.5, h=2.4,
                  header_fill=NAVY, font_size=11, header_font_size=11,
                  row_fills=row_fills_reg)

        # AEC 기여도 하이라이트
        try:
            c1 = reg_gn[reg_gn["Case"].str.contains("Clinical")].iloc[0]
            c2 = reg_gn[reg_gn["Case"].str.contains("AEC_prev") & ~reg_gn["Case"].str.contains("Scanner")].iloc[0]
            c3 = reg_gn[reg_gn["Case"].str.contains("AEC_new") & ~reg_gn["Case"].str.contains("Scanner")].iloc[0]

            add_rect(sl, 0.2, 4.0, 8.5, 0.85, RGBColor(0xEE, 0xF5, 0xFF))
            add_text(sl, "AEC_prev 기여도 (Case 2 - Case 1):",
                     0.35, 4.05, 4.0, 0.3, font_size=11, bold=True, color=BLUE)
            add_text(sl,
                     f"Lin R² {float(c2['Lin_R2'])-float(c1['Lin_R2']):+.4f}  |  "
                     f"AUC {float(c2['Log_AUC'])-float(c1['Log_AUC']):+.4f}",
                     0.35, 4.35, 4.0, 0.3, font_size=12, bold=True, color=BLUE)

            add_text(sl, "AEC_new 기여도 (Case 3 - Case 1):",
                     4.6, 4.05, 4.0, 0.3, font_size=11, bold=True, color=GREEN)
            add_text(sl,
                     f"Lin R² {float(c3['Lin_R2'])-float(c1['Lin_R2']):+.4f}  |  "
                     f"AUC {float(c3['Log_AUC'])-float(c1['Log_AUC']):+.4f}",
                     4.6, 4.35, 4.0, 0.3, font_size=12, bold=True, color=GREEN)
        except Exception:
            pass

    # 오른쪽: 피처 선택 요약
    add_text(sl, "피처 선택 CV R² 비교 (AEC_prev vs Pipeline)",
             9.0, 1.1, 4.1, 0.35, font_size=13, bold=True, color=NAVY)

    if fs_cross is not None:
        headers_fs = ["데이터셋", "선택 피처 수", "Pipeline R²", "Prev R²", "Delta"]
        rows_fs = []
        for _, row in fs_cross.iterrows():
            delta = float(row.get("delta_r2", 0))
            rows_fs.append([
                str(row.get("dataset","?"))[:4],
                str(row.get("n_feats","?")),
                f"{float(row.get('pipeline_r2',0)):.4f}",
                f"{float(row.get('prev_r2',0)):.4f}",
                f"{delta:+.4f}",
            ])
        fills_fs = [RGBColor(0xEB, 0xFF, 0xEB) if float(r[4]) > 0
                    else RGBColor(0xFF, 0xEB, 0xEB) for r in rows_fs]
        add_table(sl, headers_fs, rows_fs,
                  l=9.0, t=1.5, w=4.1, h=1.8,
                  header_fill=NAVY, font_size=10, header_font_size=10,
                  row_fills=fills_fs)

    # 오른쪽 하단: 공통 피처 3개 강조
    add_rect(sl, 9.0, 3.45, 4.1, 1.45, RGBColor(0xF0, 0xFF, 0xF0))
    add_text(sl, "3개 데이터셋 공통 선택 피처",
             9.1, 3.5, 3.9, 0.3, font_size=11, bold=True, color=GREEN)
    common_feats = ["mean", "spectral_energy"]
    for i, feat in enumerate(common_feats):
        badge(sl, feat, 9.1, 3.85 + i * 0.4, GREEN)

    # 하단 note
    add_rect(sl, 0.2, 5.0, 12.9, 2.3, RGBColor(0xF5, 0xF5, 0xFF))
    add_text(sl, "연구 보고서: results/research_report_0430.md",
             0.35, 5.05, 12.5, 0.3, font_size=11, bold=True, color=NAVY)
    notes = [
        "• 강남 선형 회귀: Case 2(+AEC_prev)가 최고 R²=0.6780 — AEC_prev의 기여 (Δ+0.0160)",
        "• 강남 로지스틱 AUC: Case 3(+AEC_new)=0.8352, Case 2(+AEC_prev)=0.8346 — 두 세트 동등",
        "• 피처 선택: 강남 9개, 신촌 13개, 병합 11개 — 'mean', 'spectral_energy' 3개 데이터셋 공통",
        "• 전체 보고서(회귀+피처+교차병원+BMI기여도)는 research_report_0430.md에서 확인",
    ]
    for i, n in enumerate(notes):
        add_text(sl, n, 0.35, 5.4 + i * 0.35, 12.5, 0.35, font_size=11, color=DARK)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 13. 결론 & 기대 효과
# ─────────────────────────────────────────────────────────────────────────────

def slide_conclusion(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, 0, 0, 13.33, 1.05, RGBColor(0x10, 0x25, 0x48))
    add_text(sl, "결론 및 기대 효과", 0.3, 0.15, 12.5, 0.75,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    conclusions = [
        ("①", "피처 선택 객관화",
         "60개+ AEC 피처에서 과적합·다중공선성 없이\n최적 세트를 데이터 기반으로 자동 결정",
         GREEN),
        ("②", "BMI 보정으로 인과 추론 강화",
         "체지방·근육량 교란을 통제하여\nAEC 신호의 독립적 예측력을 분리",
         ORANGE),
        ("③", "AEC_prev vs AEC_new 직접 비교",
         "5-Case 구조로 수동 vs 자동 피처 세트의\n회귀 성능을 정량적으로 비교",
         BLUE),
        ("④", "재현성 & 외부 검증",
         "강남·신촌 교차 검증으로\n피처 선택과 모델의 일반화 가능성 확인",
         RGBColor(0x00, 0xB0, 0xD0)),
    ]

    for i, (num, title, desc, clr) in enumerate(conclusions):
        col = i % 2
        row = i // 2
        x = 1.5 + col * 5.2
        y = 1.5 + row * 2.8
        add_rect(sl, x, y, 4.8, 2.5, RGBColor(0x22, 0x44, 0x70))
        add_rect(sl, x, y, 4.8, 0.5, clr)
        add_text(sl, f"{num} {title}", x + 0.1, y + 0.05, 4.6, 0.4,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, desc, x + 0.1, y + 0.6, 4.6, 1.8,
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "0430 설계는 기존 0424 대비 방법론적 엄밀성을 전반적으로 향상,\n"
                 "BMI 보정·자동 피처 선택·다병원 검증을 통해 연구 결론의 신뢰성을 높임.",
             0.4, 7.0, 12.5, 0.45,
             font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_title(prs)
    print("  [1/12] Title")

    slide_overview(prs)
    print("  [2/12] Overview")

    slide_feature_selection(prs)
    print("  [3/12] Feature Selection 방법")

    slide_fs_results(prs)
    print("  [4/12] Feature Selection 결과")

    slide_bmi(prs)
    print("  [5/14] BMI 추가")

    slide_bmi_results(prs)
    print("  [6/14] BMI 기여도 분석 결과")

    slide_case_structure(prs)
    print("  [7/14] Case 구조")

    slide_case_results(prs)
    print("  [8/14] Case 비교 결과")

    slide_multi_hospital(prs)
    print("  [9/13] 다병원 분석")

    slide_linear_results(prs)
    print("  [10/13] 선형 회귀 결과")

    slide_logistic_results(prs)
    print("  [11/13] 로지스틱 결과")

    slide_key_metrics(prs)
    print("  [12/13] 핵심 수치 요약")

    slide_conclusion(prs)
    print("  [13/13] 결론")

    prs.save(str(OUTPUT_PPT))
    print(f"\nPPT 저장 완료: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
