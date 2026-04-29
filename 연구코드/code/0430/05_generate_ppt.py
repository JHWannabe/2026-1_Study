# -*- coding: utf-8 -*-
"""
generate_ppt.py  (0430 버전)
────────────────────────────────────────────────────────────────
results/research_report_0430.md 를 파싱하고,
results/regression/ 및 results/feature_selection/ 의 그래프들을
배치하여 PPT를 생성한다.

출력: results/0430_research_report.pptx

실행: python generate_ppt.py
────────────────────────────────────────────────────────────────
"""

import sys, io, re, os, subprocess, tempfile
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
else:
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import warnings
warnings.filterwarnings('ignore')

import numpy as np
import pandas as pd
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── 경로 설정 ──────────────────────────────────────────────────
SCRIPT_DIR  = Path(__file__).parent
RESULT_ROOT = SCRIPT_DIR.parent.parent / "results"
FS_DIR      = RESULT_ROOT / "feature_selection"
REG_DIR     = RESULT_ROOT / "regression"
FIG_DIR     = RESULT_ROOT / "figures_0430"
REPORT_MD   = RESULT_ROOT / "research_report_0430.md"
OUT_PPT     = RESULT_ROOT / "0430_research_report.pptx"

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x35, 0x5E)
BLUE    = RGBColor(0x2E, 0x75, 0xB6)
SKYBLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE  = RGBColor(0xFF, 0x7F, 0x27)
GREEN   = RGBColor(0x00, 0x8B, 0x45)
RED     = RGBColor(0xC0, 0x00, 0x00)

W = Inches(13.33)   # 와이드 슬라이드 폭
H = Inches(7.5)     # 높이


# ═══════════════════════════════════════════════════════════════
# 1. 보고서 파싱
# ═══════════════════════════════════════════════════════════════

def parse_report_0430() -> dict:
    """research_report_0430.md를 파싱하여 슬라이드에 필요한 데이터를 dict로 반환."""
    if not REPORT_MD.exists():
        print(f"  [경고] 보고서 파일 없음: {REPORT_MD}")
        return {}

    with open(str(REPORT_MD), encoding='utf-8') as fh:
        lines = fh.readlines()

    D = {}

    def strip_bold(s):
        return re.sub(r'\*\*([^*]+)\*\*', r'\1', s).strip()

    def flt(s):
        cleaned = re.sub(r'[^\d.\-]', '', strip_bold(str(s)))
        try:
            return float(cleaned)
        except ValueError:
            return 0.0

    def find_section(header):
        for i, ln in enumerate(lines):
            if header in ln:
                return i
        return -1

    def table_rows(start, max_rows=25):
        if start < 0:
            return []
        rows, seen_sep, count = [], False, 0
        for ln in lines[start + 1: start + 1 + max_rows + 20]:
            stripped = ln.strip()
            if not stripped.startswith('|'):
                if seen_sep:
                    break
                continue
            cols = [c.strip() for c in stripped.split('|')[1:-1]]
            if cols and all(re.fullmatch(r'[-: ]+', c) for c in cols if c):
                seen_sep = True
                continue
            if seen_sep:
                rows.append(cols)
                count += 1
                if count >= max_rows:
                    break
        return rows

    # ── 피처 선택 요약 ────────────────────────────────────────
    sec = find_section('### 2.1 데이터셋별 최종 선택 피처 요약')
    D['fs_rows'] = table_rows(sec, max_rows=5)

    sec = find_section('### 2.2 최종 선택 AEC 피처')
    D['feat_rows'] = table_rows(sec, max_rows=5)

    sec = find_section('### 2.3 데이터셋별 피처 선택 일치도')
    D['matrix_rows'] = table_rows(sec, max_rows=20)

    # ── 회귀 결과 (강남) ──────────────────────────────────────
    sec_gn_lin = find_section('#### 선형 회귀 (5-Fold CV)')
    D['gn_lin_rows'] = table_rows(sec_gn_lin, max_rows=6)

    sec_gn_log = find_section('#### 로지스틱 회귀 (5-Fold CV)')
    D['gn_log_rows'] = table_rows(sec_gn_log, max_rows=6)

    # ── 교차 병원 비교 ────────────────────────────────────────
    sec_ch_lin = find_section('### 4.1 선형 회귀 R² 비교')
    D['ch_lin_rows'] = table_rows(sec_ch_lin, max_rows=6)

    sec_ch_log = find_section('### 4.2 로지스틱 AUC 비교')
    D['ch_log_rows'] = table_rows(sec_ch_log, max_rows=6)

    # ── 결론 핵심 성과 ────────────────────────────────────────
    sec_concl = find_section('### 8.3 0430 핵심 성과')
    D['core_achievements'] = []
    for ln in lines[sec_concl+1: sec_concl+10]:
        stripped = ln.strip()
        if stripped.startswith(('1.', '2.', '3.', '4.')):
            D['core_achievements'].append(strip_bold(stripped))

    return D


# ═══════════════════════════════════════════════════════════════
# 2. PPT 공통 헬퍼
# ═══════════════════════════════════════════════════════════════

def new_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # 완전 빈 레이아웃
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor = WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_colored_rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _safe_img(slide, img_path: Path,
              left, top, width, height=None):
    """이미지가 있으면 삽입, 없으면 플레이스홀더 텍스트박스."""
    if img_path.exists():
        if height:
            slide.shapes.add_picture(str(img_path), left, top, width, height)
        else:
            slide.shapes.add_picture(str(img_path), left, top, width)
    else:
        ph_h = height if height else Inches(3)
        add_colored_rect(slide, left, top, width, ph_h, GRAY)
        add_textbox(slide, f"(그래프 없음)\n{img_path.name}",
                    left + Inches(0.1), top + Inches(0.2),
                    width - Inches(0.2), ph_h - Inches(0.4),
                    font_size=10, color=RGBColor(0x88, 0x88, 0x88))


def add_header_bar(slide, title: str, subtitle: str = ''):
    add_colored_rect(slide, 0, 0, W, Inches(1.1), NAVY)
    add_textbox(slide, title,
                Inches(0.35), Inches(0.08), Inches(12.5), Inches(0.65),
                font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.35),
                    font_size=12, color=SKYBLUE)


def add_md_table(slide, headers: list, rows: list,
                 left, top, col_widths: list, row_height=Inches(0.38),
                 font_size=9, header_color=NAVY, alt_color=SKYBLUE):
    """마크다운 파싱 결과(headers, rows)를 PPT 표로 변환."""
    if not rows:
        return
    n_rows = len(rows) + 1   # header 포함
    n_cols = len(headers)
    total_w = sum(col_widths)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top,
                                  total_w, row_height * n_rows).table

    # 열 너비 설정
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    def _set_cell(cell, text, bold=False, bg_color=None, txt_color=DARK):
        cell.text = str(text)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = cell.text_frame.paragraphs[0].runs
        if run:
            run[0].font.size = Pt(font_size)
            run[0].font.bold = bold
            run[0].font.color.rgb = txt_color
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    # 헤더 행
    for ci, h in enumerate(headers):
        _set_cell(tbl.cell(0, ci), h, bold=True,
                  bg_color=header_color, txt_color=WHITE)

    # 데이터 행
    for ri, row in enumerate(rows):
        bg = alt_color if ri % 2 == 0 else WHITE
        for ci in range(n_cols):
            val = row[ci] if ci < len(row) else ''
            _set_cell(tbl.cell(ri + 1, ci), val, bg_color=bg)


def _load_excel(path: Path, sheet_name=0) -> pd.DataFrame:
    tmp = Path(tempfile.gettempdir()) / f"ppt0430_{path.name}"
    subprocess.run(
        ["powershell", "-Command",
         f'Copy-Item -Path "{path}" -Destination "{tmp}" -Force'],
        capture_output=True,
    )
    if not tmp.exists():
        return pd.DataFrame()
    try:
        return pd.read_excel(str(tmp), sheet_name=sheet_name)
    except Exception:
        return pd.DataFrame()


def load_reg(hosp_key: str, sex_key: str) -> pd.DataFrame:
    p = REG_DIR / hosp_key / sex_key / "regression_results.xlsx"
    return _load_excel(p, sheet_name="summary") if p.exists() else pd.DataFrame()


# ═══════════════════════════════════════════════════════════════
# 3. 슬라이드 빌더
# ═══════════════════════════════════════════════════════════════

def slide_title(prs: Presentation):
    """슬라이드 1: 표지"""
    slide = new_blank_slide(prs)
    add_colored_rect(slide, 0, 0, W, H, NAVY)

    add_textbox(slide, "TAMA 예측 회귀분석 연구 보고서",
                Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "0430 버전 — 설계 변경 적용",
                Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
                font_size=22, color=SKYBLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "분석 도구: Python (statsmodels, scikit-learn)\n"
                "병원: 강남 / 신촌  |  성별 그룹: 전체 / 여성(F) / 남성(M)\n"
                "AEC 세트: AEC_prev(수동 4개) vs AEC_new(파이프라인 자동 선택)",
                Inches(1), Inches(4.5), Inches(11.3), Inches(1.8),
                font_size=14, color=SKYBLUE, align=PP_ALIGN.CENTER)
    print("  [슬라이드 1] 표지")


def slide_design_changes(prs: Presentation):
    """슬라이드 2: 0424 → 0430 설계 변경사항"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "0424 → 0430 설계 변경사항",
                   "6가지 핵심 변경으로 다병원·성별 층화·자동 피처 선택 체계 도입")

    headers = ['#', '항목', '0424 이전', '0430 이후']
    rows = [
        ['①', '피처 선택', '수동 (상관계수+VIF → 연구자 결정)', '자동 파이프라인 (4단계 필터 + 앙상블 투표 + CV R²)'],
        ['②', '임상 기준선', 'PatientAge + PatientSex', 'PatientAge + PatientSex + BMI (교란변수 통제)'],
        ['③', 'Case 구조', 'Case 1~3 (단일 AEC 세트)', 'Case 1~5 (AEC_prev vs AEC_new 교차 비교)'],
        ['④', '성별 층화', '성별 = 공변량(더미)만', '전체 / 여성 / 남성 독립 모델 (3개 서브그룹)'],
        ['⑤', '다병원 분석', '강남 단독 (SITE 수동 변경)', '강남·신촌 자동 순회 + 교차 병원 비교'],
        ['⑥', '이진화 기준', '성별 특이적 P25 (남/여 별도)', '분석 그룹 내 하위 25% 동적 산출'],
    ]
    col_widths = [Inches(0.5), Inches(1.6), Inches(4.0), Inches(6.3)]
    add_md_table(slide, headers, rows,
                 Inches(0.35), Inches(1.25), col_widths, row_height=Inches(0.42),
                 font_size=10)

    # Case 구조 추가 설명
    add_textbox(slide,
                "Case 1 (Clinical: Age+Sex+BMI) | Case 2 (+AEC_prev) | Case 3 (+AEC_new) | "
                "Case 4 (+AEC_prev+Scanner) | Case 5 (+AEC_new+Scanner)",
                Inches(0.35), Inches(6.85), Inches(12.6), Inches(0.5),
                font_size=9, color=BLUE, italic=True)
    print("  [슬라이드 2] 설계 변경사항")


def slide_feature_pipeline(prs: Presentation, D: dict):
    """슬라이드 3: 피처 선택 파이프라인 4단계"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "피처 선택 파이프라인 (0430 신설)",
                   "60개+ AEC 피처 → 4단계 자동 필터링 → 앙상블 투표 + CV R² 최적화")

    step_data = [
        ("Step 1", "Near-zero variance 제거", "표준화 후 variance < 0.01", NAVY),
        ("Step 2", "Pearson 상관 중복 제거", "|r| ≥ 0.95인 후순위 피처 제거", BLUE),
        ("Step 3", "단변량 필터 (OR 결합)", "MI > 0 OR Spearman p < 0.05", ORANGE),
        ("Step 4", "앙상블 투표 + 완전/SFS 탐색", "LASSO + RFECV + RF Permutation → 5-fold CV R² 최대화", GREEN),
        ("Final",  "VIF Pruning", "VIF > 10 반복 제거 (mean 고정 보호)", RED),
    ]

    box_w = Inches(2.4)
    box_h = Inches(1.2)
    gap   = Inches(0.15)
    top_y = Inches(1.3)

    for i, (label, title, desc, color) in enumerate(step_data):
        lx = Inches(0.35) + i * (box_w + gap)
        add_colored_rect(slide, lx, top_y, box_w, box_h, color)
        add_textbox(slide, label, lx, top_y + Inches(0.05), box_w, Inches(0.35),
                    font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, lx, top_y + Inches(0.38), box_w, Inches(0.45),
                    font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

        if i < 4:
            arrow_x = lx + box_w + Inches(0.02)
            add_textbox(slide, "→", arrow_x, top_y + Inches(0.4), gap, Inches(0.4),
                        font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 단계별 설명 텍스트
    for i, (label, title, desc, color) in enumerate(step_data):
        lx = Inches(0.35) + i * (box_w + gap)
        add_textbox(slide, desc, lx, top_y + box_h + Inches(0.1), box_w, Inches(0.55),
                    font_size=8, color=DARK, wrap=True)

    # CV R² 비교 그래프
    _safe_img(slide, FIG_DIR / "01_aec_pipeline_cv_r2.png",
              Inches(0.35), Inches(3.15), Inches(8.0), Inches(4.0))

    # 피처 선택 요약 표
    fs_rows = D.get('fs_rows', [])
    if fs_rows:
        headers = ['데이터셋', 'N', '피처수', 'Pipeline R²', 'Prev R²', 'ΔR²']
        col_widths = [Inches(1.4), Inches(0.6), Inches(0.7), Inches(1.4),
                      Inches(1.0), Inches(0.9)]
        add_md_table(slide, headers, fs_rows,
                     Inches(8.6), Inches(3.15), col_widths, row_height=Inches(0.42),
                     font_size=9)

    print("  [슬라이드 3] 피처 선택 파이프라인")


def slide_feature_matrix(prs: Presentation):
    """슬라이드 4: 최종 피처 선택 결과"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "최종 선택 AEC 피처",
                   "강남 / 신촌 / 병합 데이터셋에서 파이프라인이 선택한 피처 비교")

    _safe_img(slide, FIG_DIR / "02_aec_feature_matrix.png",
              Inches(0.35), Inches(1.2), Inches(6.5), Inches(5.9))

    # 피처 목록 (cross_dataset_comparison.xlsx에서 로드)
    fs_path = FS_DIR / "cross_dataset_comparison.xlsx"
    feat_df = _load_excel(fs_path, sheet_name="summary")
    if not feat_df.empty:
        y_pos = Inches(1.3)
        for _, row in feat_df.iterrows():
            dataset = str(row.get('dataset', '?'))
            feats   = str(row.get('features', ''))
            add_textbox(slide, f"▶ {dataset}",
                        Inches(7.1), y_pos, Inches(5.8), Inches(0.35),
                        font_size=10, bold=True, color=NAVY)
            add_textbox(slide, feats,
                        Inches(7.1), y_pos + Inches(0.32), Inches(5.8), Inches(0.65),
                        font_size=9, color=DARK, wrap=True)
            y_pos += Inches(1.05)

    add_textbox(slide,
                "mean 피처는 AEC 신호의 평균 진폭 수준 = 가장 해석 가능한 피처 → 모든 최종 세트에 강제 포함",
                Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.35),
                font_size=9, italic=True, color=BLUE)
    print("  [슬라이드 4] 최종 피처 선택 결과")


def slide_regression_hospital(prs: Presentation, D: dict,
                               hosp_key: str, hosp_label: str):
    """슬라이드 5-6: 병원별 회귀 결과 (선형 + 로지스틱)"""

    # ── 5a: 선형 회귀 ──
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide,
                   f"[{hosp_label}] 선형 회귀 결과 (5-Fold CV)",
                   "Case 1~5 R²·MAE·RMSE 비교 | 임상 기준선 → AEC_prev/new → 스캐너 추가")

    df = load_reg(hosp_key, 'all')
    if not df.empty:
        headers = ['Case', 'N feat', 'R² (mean±std)', 'MAE (cm²)', 'RMSE (cm²)']
        rows = []
        for _, row in df.iterrows():
            r2_std = f"±{row['Lin_R2_std']:.4f}" if 'Lin_R2_std' in row else ''
            rows.append([str(row.get('Case', '?')),
                         str(row.get('N_features', '?')),
                         f"{row.get('Lin_R2', 0):.4f} {r2_std}",
                         f"{row.get('Lin_MAE', 0):.2f}",
                         f"{row.get('Lin_RMSE', 0):.2f}"])
        col_widths = [Inches(3.2), Inches(0.8), Inches(2.6), Inches(1.3), Inches(1.3)]
        add_md_table(slide, headers, rows,
                     Inches(0.35), Inches(1.25), col_widths, row_height=Inches(0.45),
                     font_size=9)

    _safe_img(slide,
              REG_DIR / hosp_key / "all" / "01_linear_actual_vs_pred.png",
              Inches(0.35), Inches(4.1), Inches(7.5), Inches(3.1))
    _safe_img(slide,
              REG_DIR / hosp_key / "all" / "02_linear_metrics_comparison.png",
              Inches(7.0), Inches(4.1), Inches(6.0), Inches(3.1))
    print(f"  [슬라이드] {hosp_label} 선형 회귀")

    # ── 5b: 로지스틱 회귀 ──
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide,
                   f"[{hosp_label}] 로지스틱 회귀 결과 (5-Fold CV)",
                   "Case 1~5 AUC·Sensitivity·Specificity 비교 | 그룹 내 하위 25% = Low TAMA")

    if not df.empty:
        headers = ['Case', 'AUC (mean±std)', 'Accuracy', 'Sensitivity', 'Specificity']
        rows = []
        for _, row in df.iterrows():
            auc_std = f"±{row['Log_AUC_std']:.4f}" if 'Log_AUC_std' in row else ''
            rows.append([str(row.get('Case', '?')),
                         f"{row.get('Log_AUC', 0):.4f} {auc_std}",
                         f"{row.get('Log_Acc', 0):.4f}",
                         f"{row.get('Log_Sens', 0):.4f}",
                         f"{row.get('Log_Spec', 0):.4f}"])
        col_widths = [Inches(3.2), Inches(2.6), Inches(1.5), Inches(1.5), Inches(1.5)]
        add_md_table(slide, headers, rows,
                     Inches(0.35), Inches(1.25), col_widths, row_height=Inches(0.45),
                     font_size=9)

    _safe_img(slide,
              REG_DIR / hosp_key / "all" / "04_logistic_roc.png",
              Inches(0.35), Inches(4.1), Inches(6.0), Inches(3.1))
    _safe_img(slide,
              REG_DIR / hosp_key / "all" / "05_logistic_metrics_comparison.png",
              Inches(6.5), Inches(4.1), Inches(6.5), Inches(3.1))
    print(f"  [슬라이드] {hosp_label} 로지스틱 회귀")


def slide_diagnostics(prs: Presentation, hosp_key: str, hosp_label: str):
    """슬라이드: 진단 플롯 (잔차 진단, Forest Plot)"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide,
                   f"[{hosp_label}] 모델 진단 — 잔차 & Forest Plot",
                   "선형 잔차 4-Panel + 유의 계수 Forest Plot | 대표 모델: Case 4 (임상+AEC_prev+스캐너)")

    _safe_img(slide,
              REG_DIR / hosp_key / "05_linear_residuals.png",
              Inches(0.35), Inches(1.2), Inches(7.2), Inches(3.5))
    _safe_img(slide,
              REG_DIR / hosp_key / "06_linear_forest.png",
              Inches(7.7), Inches(1.2), Inches(5.3), Inches(3.5))
    _safe_img(slide,
              REG_DIR / hosp_key / "09_logistic_calibration.png",
              Inches(0.35), Inches(4.9), Inches(5.5), Inches(2.4))
    _safe_img(slide,
              REG_DIR / hosp_key / "10_logistic_confusion.png",
              Inches(6.1), Inches(4.9), Inches(6.9), Inches(2.4))
    print(f"  [슬라이드] {hosp_label} 진단 플롯")


def slide_cross_hospital(prs: Presentation):
    """슬라이드: 교차 병원 비교"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "교차 병원 비교 (강남 vs 신촌)",
                   "Case 1~5 × 두 병원의 선형 R²·AUC 재현성 비교")

    _safe_img(slide, FIG_DIR / "03_cross_hospital_linear.png",
              Inches(0.35), Inches(1.2), Inches(12.6), Inches(2.9))
    _safe_img(slide, FIG_DIR / "04_cross_hospital_logistic.png",
              Inches(0.35), Inches(4.3), Inches(12.6), Inches(2.9))
    print("  [슬라이드] 교차 병원 비교")


def slide_sex_stratification(prs: Presentation):
    """슬라이드: 성별 층화 분석"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "성별 층화 분석 (전체 / 여성(F) / 남성(M))",
                   "각 성별 그룹 독립 모델 — 이질성 탐색 | 강남 기준 5-Fold CV")

    _safe_img(slide, FIG_DIR / "05_sex_strat_linear.png",
              Inches(0.35), Inches(1.2), Inches(12.6), Inches(2.9))
    _safe_img(slide, FIG_DIR / "06_sex_strat_logistic.png",
              Inches(0.35), Inches(4.3), Inches(12.6), Inches(2.9))
    print("  [슬라이드] 성별 층화 분석")


def slide_aec_comparison(prs: Presentation):
    """슬라이드: AEC_prev vs AEC_new"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "AEC_prev vs AEC_new — 피처 선택 전략 성능 비교",
                   "AEC_prev: 수동 4개 (mean, CV, skewness, slope_abs_mean) | AEC_new: 파이프라인 자동 선택")

    _safe_img(slide, FIG_DIR / "07_aec_prev_vs_new.png",
              Inches(0.35), Inches(1.2), Inches(12.6), Inches(3.7))

    add_textbox(slide,
                "Case 2 vs Case 3: AEC_prev vs AEC_new (스캐너 없음)\n"
                "Case 4 vs Case 5: AEC_prev vs AEC_new (스캐너 포함)\n"
                "→ AEC_new가 일관되게 우수하면 파이프라인 선택의 유효성 입증",
                Inches(0.35), Inches(5.2), Inches(12.6), Inches(1.0),
                font_size=12, color=DARK)

    _safe_img(slide, FS_DIR / "cross_dataset_comparison_r2.png",
              Inches(0.35), Inches(5.9), Inches(6.5), Inches(1.4))
    _safe_img(slide, FS_DIR / "cross_dataset_feature_heatmap.png",
              Inches(7.0), Inches(5.9), Inches(6.0), Inches(1.4))
    print("  [슬라이드] AEC_prev vs AEC_new")


def slide_bmi_analysis(prs: Presentation):
    """슬라이드: BMI 기여도 분석"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "BMI 기여도 분석 (0430 신설)",
                   "Case 1/2/4 × no BMI vs +BMI — BMI가 TAMA 예측에 기여하는 독립적 효과 정량화")

    _safe_img(slide, FIG_DIR / "08_bmi_contribution_gangnam.png",
              Inches(0.35), Inches(1.2), Inches(12.6), Inches(2.9))
    _safe_img(slide, FIG_DIR / "09_bmi_contribution_sinchon.png",
              Inches(0.35), Inches(4.3), Inches(12.6), Inches(2.9))

    add_textbox(slide,
                "BMI 추가 효과가 AEC 투입 후 감소(감쇠)하면 → AEC와 BMI가 일부 공통 정보를 공유함을 시사",
                Inches(0.35), Inches(7.15), Inches(12.6), Inches(0.28),
                font_size=9, italic=True, color=BLUE)
    print("  [슬라이드] BMI 기여도")


def slide_case_progression(prs: Presentation):
    """슬라이드: Case 1→5 성능 추이"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "Case 1 → 5 예측 성능 지표 추이",
                   "임상 기준선 → AEC 추가 → 스캐너 추가에 따른 R²·AUC·Sensitivity·RMSE 변화 정량화")

    _safe_img(slide, FIG_DIR / "10_case_progression_gangnam.png",
              Inches(0.35), Inches(1.2), Inches(6.2), Inches(5.9))
    _safe_img(slide, FIG_DIR / "11_case_progression_sinchon.png",
              Inches(6.75), Inches(1.2), Inches(6.2), Inches(5.9))
    print("  [슬라이드] Case 1→5 추이")


def slide_eda(prs: Presentation, hosp_key: str, hosp_label: str):
    """슬라이드: EDA — CT 스캐너·kVp·상관행렬"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, f"[{hosp_label}] CT 장비 & 피처 상관 분석 (EDA)",
                   "스캐너 모델 분포 / kVp 분포 / Clinical+AEC_prev 상관행렬")

    _safe_img(slide,
              REG_DIR / hosp_key / "16_scanner_distribution.png",
              Inches(0.35), Inches(1.2), Inches(5.5), Inches(2.8))
    _safe_img(slide,
              REG_DIR / hosp_key / "17_kvp_distribution.png",
              Inches(6.0), Inches(1.2), Inches(7.0), Inches(2.8))
    _safe_img(slide,
              REG_DIR / hosp_key / "18_correlation_matrix.png",
              Inches(0.35), Inches(4.2), Inches(12.6), Inches(3.0))
    print(f"  [슬라이드] {hosp_label} EDA")


def slide_n_threshold(prs: Presentation):
    """슬라이드: N수 & 임계값 요약"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "병원 × 성별 그룹별 분석 대상 N수 & Low-TAMA 임계값",
                   "이진화 임계값 = 각 그룹 내 하위 25% (그룹-내 P25)")

    _safe_img(slide, FIG_DIR / "12_sex_n_threshold.png",
              Inches(0.35), Inches(1.2), Inches(12.6), Inches(5.9))
    print("  [슬라이드] N수 & 임계값")


def slide_aec_prev_vs_new_detail(prs: Presentation):
    """슬라이드: Case 별 AIC/BIC & Case 추이"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "Case 1~5 모델 적합도 — AIC/BIC & 계수 Forest Plot",
                   "AIC/BIC 낮을수록 우수 | 계수 크기로 기여 변수 파악")

    _safe_img(slide,
              REG_DIR / "gangnam" / "all" / "14_case_aic_bar.png",
              Inches(0.35), Inches(1.2), Inches(6.2), Inches(2.8))
    _safe_img(slide,
              REG_DIR / "gangnam" / "all" / "08_case_comparison_overview.png",
              Inches(6.75), Inches(1.2), Inches(6.2), Inches(2.8))
    _safe_img(slide,
              REG_DIR / "gangnam" / "all" / "03_linear_coefficients.png",
              Inches(0.35), Inches(4.2), Inches(12.6), Inches(3.0))
    print("  [슬라이드] AIC/BIC & 계수")


def slide_conclusion(prs: Presentation, D: dict):
    """슬라이드: 결론"""
    slide = new_blank_slide(prs)
    fill_bg(slide)
    add_header_bar(slide, "결론 및 핵심 성과 (0430)", "")

    achievements = D.get('core_achievements', [
        "1. BMI 보정: Case 1 기준선에서 R² 향상 — BMI가 강력한 TAMA 예측 변수임을 확인",
        "2. 자동 피처 선택: 60개+ AEC 피처에서 과적합·다중공선성 없이 객관적 세트 도출",
        "3. 성별 층화: 전체/여성/남성 독립 모델로 이질성 탐색 — 그룹별 예측 패턴 확인",
        "4. 다병원 검증: 강남·신촌 교차 검증으로 피처 선택과 모델의 재현성 확인",
    ])

    y_pos = Inches(1.3)
    for ach in achievements:
        add_colored_rect(slide, Inches(0.35), y_pos, Inches(12.6), Inches(0.75), SKYBLUE)
        add_textbox(slide, ach,
                    Inches(0.55), y_pos + Inches(0.1), Inches(12.2), Inches(0.55),
                    font_size=13, color=DARK, bold=False)
        y_pos += Inches(0.88)

    add_textbox(slide, "한계 및 향후 과제",
                Inches(0.35), y_pos + Inches(0.1), Inches(12.6), Inches(0.4),
                font_size=14, bold=True, color=NAVY)
    limits = [
        "• AEC_new vs AEC_prev 성능 차이가 미미한 경우 — 더 많은 환자 데이터 필요",
        "• 성별 층화 시 소그룹 표본 크기에 따른 불안정성 주의 (특히 남성/여성 단독)",
        "• 단면 연구 설계 — 인과 추론을 위한 전향적 코호트 연구 권장",
        "• Raw AEC 시계열(200포인트)을 1D CNN / LSTM으로 직접 학습 시 추가 성능 향상 가능성",
    ]
    add_textbox(slide, "\n".join(limits),
                Inches(0.35), y_pos + Inches(0.55), Inches(12.6), Inches(1.3),
                font_size=11, color=DARK)
    add_textbox(slide, "자동 생성: generate_ppt.py (0430)",
                Inches(0.35), Inches(7.2), Inches(12.6), Inches(0.25),
                font_size=8, color=RGBColor(0xBA, 0xB0, 0xAC), italic=True)
    print("  [슬라이드] 결론")


# ═══════════════════════════════════════════════════════════════
# 4. 메인
# ═══════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("  generate_ppt.py (0430) — PPT 생성")
    print("=" * 60)

    print("\n[1] 보고서 파싱 중...")
    D = parse_report_0430()

    print("\n[2] 슬라이드 빌드 중...")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)                                         # 1
    slide_design_changes(prs)                                # 2
    slide_feature_pipeline(prs, D)                          # 3
    slide_feature_matrix(prs)                               # 4

    # 강남 회귀 결과
    slide_regression_hospital(prs, D, 'gangnam', '강남')   # 5a, 5b
    slide_diagnostics(prs, 'gangnam', '강남')              # 6

    # 신촌 회귀 결과
    slide_regression_hospital(prs, D, 'sinchon', '신촌')  # 7a, 7b
    slide_diagnostics(prs, 'sinchon', '신촌')             # 8

    slide_cross_hospital(prs)                               # 9
    slide_sex_stratification(prs)                           # 10
    slide_aec_comparison(prs)                               # 11
    slide_bmi_analysis(prs)                                 # 12
    slide_case_progression(prs)                             # 13
    slide_aec_prev_vs_new_detail(prs)                       # 14
    slide_n_threshold(prs)                                  # 15
    slide_eda(prs, 'gangnam', '강남')                      # 16
    slide_eda(prs, 'sinchon', '신촌')                      # 17
    slide_conclusion(prs, D)                                # 18

    print(f"\n[3] PPT 저장 중...")
    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    n_slides = len(prs.slides)
    print(f"  저장 완료: {OUT_PPT}")
    print(f"  슬라이드 수: {n_slides}")
    print("=" * 60)


if __name__ == '__main__':
    main()
