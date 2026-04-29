# -*- coding: utf-8 -*-
"""
07_generate_sharing_ppt.py  (0430 공유용 발표 자료)
────────────────────────────────────────────────
260430_Linear_Logistic_Regression.pptx 스타일로
연구 결과 공유 자료를 생성한다.

출력: results/0430_sharing_report.pptx
실행: python 07_generate_sharing_ppt.py
────────────────────────────────────────────────
"""
import sys, io, os
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
else:
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import warnings
warnings.filterwarnings('ignore')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── 경로 설정 ──────────────────────────────────────────────────
SCRIPT_DIR   = Path(__file__).parent
RESULT_ROOT  = SCRIPT_DIR.parent.parent / "results"
STUDY_ROOT   = SCRIPT_DIR.parent.parent.parent
REF_PPT_PATH = STUDY_ROOT / "연구자료" / "260430_Linear_Logistic_Regression.pptx"
FS_DIR       = RESULT_ROOT / "feature_selection"
REG_DIR      = RESULT_ROOT / "regression"
FIG_DIR      = RESULT_ROOT / "figures_0430"
OUT_PPT      = RESULT_ROOT / "0430_sharing_report.pptx"

# ── 색상 팔레트 (참조 PPT 기반) ──────────────────────────────
NAVY    = RGBColor(0x1A, 0x35, 0x5E)   # #1A355E
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # #1A1A2E
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE  = RGBColor(0xFF, 0x7F, 0x27)   # #FF7F27
BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # #2E75B6
SKYBLUE = RGBColor(0xD6, 0xE4, 0xF0)
GRAY    = RGBColor(0xAA, 0xAA, 0xAA)

FONT_NAME = "KoPub돋움체 Medium"
FONT_BOLD = "KoPub돋움체 Bold"

W = Inches(13.33)
H = Inches(7.50)

SLIDE_NUM_X = Inches(12.47)
SLIDE_NUM_Y = Inches(7.13)
TITLE_X     = Inches(0.56)
TITLE_Y     = Inches(0.40)
TITLE_W     = Inches(12.22)
TITLE_H     = Inches(0.58)
CONTENT_Y   = Inches(0.98)
CONTENT_H   = Inches(5.54)


# ═══════════════════════════════════════════════════════════════
# 1. 공통 헬퍼
# ═══════════════════════════════════════════════════════════════

def new_slide(prs: Presentation):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def fill_bg(slide, color: RGBColor = WHITE):
    bg = slide.background
    f  = bg.fill
    f.solid()
    f.fore_color.rgb = color


def txb(slide, text, left, top, width, height,
        size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, font=None):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf    = shape.text_frame
    tf.word_wrap = wrap
    para  = tf.paragraphs[0]
    para.alignment = align
    run   = para.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font or FONT_NAME
    return shape


def rect(slide, left, top, width, height, fill_color: RGBColor):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.fill.background()
    return sh


def img(slide, path: Path, left, top, width, height=None):
    if path.exists():
        if height:
            slide.shapes.add_picture(str(path), left, top, width, height)
        else:
            slide.shapes.add_picture(str(path), left, top, width)
    else:
        ph_h = height or Inches(3)
        rect(slide, left, top, width, ph_h, RGBColor(0x33, 0x33, 0x33))
        txb(slide, f"[{path.name}]",
            left + Inches(0.1), top + Inches(0.1),
            width - Inches(0.2), ph_h - Inches(0.2),
            size=9, color=GRAY, align=PP_ALIGN.CENTER)


def slide_num(slide, n: int, color: RGBColor = DARK):
    txb(slide, str(n), SLIDE_NUM_X, SLIDE_NUM_Y, Inches(0.61), Inches(0.35),
        size=12, color=color, align=PP_ALIGN.RIGHT)


def slide_title_bar(slide, title: str, n: int, bg_dark=False):
    """참조 PPT 형식: 제목 Placeholder 위치의 텍스트박스 + 슬라이드 번호"""
    txt_color = WHITE if bg_dark else DARK
    txb(slide, title, TITLE_X, TITLE_Y, TITLE_W, TITLE_H,
        size=20, bold=True, color=txt_color, font=FONT_BOLD)
    slide_num(slide, n, color=txt_color)


def note_line(slide, text: str, y_pos=None, color=ORANGE, size=11):
    y = y_pos or Inches(6.82)
    txb(slide, text, Inches(0.56), y, Inches(12.22), Inches(0.36),
        size=size, color=color, italic=False)


def key_finding(slide, text: str, y_pos=None, color=BLUE, size=13):
    y = y_pos or Inches(6.3)
    txb(slide, text, Inches(0.56), y, Inches(12.22), Inches(0.40),
        size=size, bold=True, color=color)


def table(slide, headers, rows, left, top, col_widths,
          row_h=Inches(0.38), font_size=10,
          header_bg=NAVY, alt_bg=SKYBLUE, text_dark=True):
    if not rows:
        return
    nc = len(headers)
    nr = len(rows) + 1
    tw = sum(col_widths)
    tbl = slide.shapes.add_table(nr, nc, left, top, tw, row_h * nr).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    def _cell(row_i, col_i, text, bold=False, bg=None, fc=DARK):
        c = tbl.cell(row_i, col_i)
        c.text = str(text)
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        runs = c.text_frame.paragraphs[0].runs
        if runs:
            runs[0].font.size  = Pt(font_size)
            runs[0].font.bold  = bold
            runs[0].font.color.rgb = fc
            runs[0].font.name  = FONT_NAME
        if bg:
            c.fill.solid()
            c.fill.fore_color.rgb = bg

    for ci, h in enumerate(headers):
        _cell(0, ci, h, bold=True, bg=header_bg, fc=WHITE)
    for ri, row in enumerate(rows):
        bg = alt_bg if ri % 2 == 0 else WHITE
        for ci in range(nc):
            val = row[ci] if ci < len(row) else ''
            _cell(ri + 1, ci, val, bg=bg, fc=DARK if text_dark else WHITE)


# ═══════════════════════════════════════════════════════════════
# 2. 슬라이드 빌더
# ═══════════════════════════════════════════════════════════════

def s01_title(prs):
    """슬라이드 1: 표지"""
    slide = new_slide(prs)
    fill_bg(slide, WHITE)

    # 상단 네이비 바
    rect(slide, 0, 0, W, Inches(1.8), NAVY)
    txb(slide, "Dept. of Medical Device Engineering & Management",
        Inches(0.56), Inches(0.25), Inches(10), Inches(0.5),
        size=13, color=SKYBLUE, font=FONT_NAME)
    txb(slide, "Gangnam Severance Hospital, Yonsei University",
        Inches(0.56), Inches(0.7), Inches(10), Inches(0.45),
        size=13, color=SKYBLUE, font=FONT_NAME)
    txb(slide, "2026. 04. 30",
        Inches(0.56), Inches(1.2), Inches(4), Inches(0.45),
        size=13, color=WHITE, font=FONT_NAME)

    # 메인 타이틀
    txb(slide, "TAMA 예측 회귀분석 결과 보고",
        Inches(0.56), Inches(2.3), Inches(12.22), Inches(0.8),
        size=32, bold=True, color=NAVY, align=PP_ALIGN.LEFT, font=FONT_BOLD)

    # 구분선
    rect(slide, Inches(0.56), Inches(3.2), Inches(8), Inches(0.04), ORANGE)

    # 부제
    txb(slide, "0430 버전  |  자동 피처 선택 + BMI 보정 + 성별 층화 + 다병원 검증",
        Inches(0.56), Inches(3.4), Inches(12.22), Inches(0.5),
        size=16, color=DARK, font=FONT_NAME)

    # 저자/소속
    txb(slide, "JangHoon Chun",
        Inches(0.56), Inches(4.3), Inches(6), Inches(0.45),
        size=18, bold=True, color=NAVY, font=FONT_BOLD)
    txb(slide,
        "분석 도구: Python (statsmodels, scikit-learn, scipy)\n"
        "병원: 강남 (N=1,365) / 신촌 (N=1,269)\n"
        "AEC 세트: AEC_prev (수동 4개) vs AEC_new (파이프라인 자동 선택)",
        Inches(0.56), Inches(4.9), Inches(12.22), Inches(1.5),
        size=14, color=RGBColor(0x44, 0x44, 0x44), font=FONT_NAME)
    print("  [S01] 표지")


def s02_last_meeting(prs, n):
    """슬라이드 2: 지난 미팅 이후 변경사항"""
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "1.  지난 미팅 이후 변경사항 (0424 → 0430)", n)

    bullet_y = CONTENT_Y
    items = [
        ("①  피처 선택",
         "수동 (상관계수+VIF → 연구자 결정)",
         "자동 파이프라인 (4단계 필터 + 앙상블 투표 + CV R² 최적화)"),
        ("②  임상 기준선",
         "PatientAge + PatientSex",
         "PatientAge + PatientSex + BMI  (교란변수 보정)"),
        ("③  Case 구조",
         "Case 1~3 (단일 AEC 세트)",
         "Case 1~5 (AEC_prev vs AEC_new 교차 비교 + Scanner)"),
        ("④  성별 층화",
         "성별 = 공변량(더미)만",
         "전체 / 여성(F) / 남성(M) 독립 모델 (3개 서브그룹)"),
        ("⑤  다병원 분석",
         "강남 단독",
         "강남·신촌 자동 순회 + 교차 병원 비교"),
        ("⑥  이진화 기준",
         "성별 특이적 P25 (남/여 별도)",
         "분석 그룹 내 하위 25% 동적 산출"),
    ]

    row_h = Inches(0.78)
    for i, (label, before, after) in enumerate(items):
        y = bullet_y + i * row_h
        bg_color = SKYBLUE if i % 2 == 0 else WHITE
        rect(slide, Inches(0.35), y, Inches(12.6), row_h - Inches(0.04), bg_color)
        txb(slide, label, Inches(0.45), y + Inches(0.08),
            Inches(2.0), Inches(0.55), size=12, bold=True, color=NAVY, font=FONT_BOLD)
        txb(slide, f"이전: {before}", Inches(2.5), y + Inches(0.05),
            Inches(4.5), Inches(0.35), size=11, color=RGBColor(0x88, 0x44, 0x44), font=FONT_NAME)
        txb(slide, f"이후: {after}", Inches(2.5), y + Inches(0.38),
            Inches(10.0), Inches(0.35), size=11, color=RGBColor(0x00, 0x55, 0x00), font=FONT_NAME)

    print(f"  [S{n:02d}] 변경사항")


def s03_research_summary(prs, n):
    """슬라이드 3: Research Summary"""
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "2.  Research Summary", n)

    # Case 구조 표
    headers = ["Case", "포함 변수", "AEC 세트", "목적"]
    rows = [
        ["Case 1  (Clinical)", "Age, Sex, BMI", "—", "기준선"],
        ["Case 2  (+AEC_prev)", "Case 1 + AEC_prev", "mean, CV, skewness, slope_abs_mean", "수동 AEC 기여도"],
        ["Case 3  (+AEC_new)", "Case 1 + AEC_new", "파이프라인 자동 선택", "자동 AEC 기여도"],
        ["Case 4  (+Scanner)", "Case 2 + Scanner", "AEC_prev + CT Model + kVp", "장비 효과 (prev)"],
        ["Case 5  (+Scanner)", "Case 3 + Scanner", "AEC_new + CT Model + kVp", "장비 효과 (new)"],
    ]
    col_w = [Inches(2.4), Inches(2.8), Inches(4.4), Inches(2.7)]
    table(slide, headers, rows,
          Inches(0.35), CONTENT_Y, col_w, row_h=Inches(0.45), font_size=10)

    # 데이터셋 개요
    txb(slide, "데이터셋 개요",
        Inches(0.35), Inches(4.0), Inches(5), Inches(0.4),
        size=13, bold=True, color=NAVY, font=FONT_BOLD)
    d_headers = ["데이터셋", "N", "여성", "남성", "P25 (cm²)"]
    d_rows = [
        ["강남", "1,365", "841", "524", "100.0"],
        ["신촌", "1,269", "631", "638", "103.0"],
    ]
    d_col_w = [Inches(2.2), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.5)]
    table(slide, d_headers, d_rows,
          Inches(0.35), Inches(4.5), d_col_w, row_h=Inches(0.42), font_size=10)

    # 비교 포인트
    txb(slide, "Case 2 vs Case 3: AEC_prev vs AEC_new (스캐너 없음)  |  "
               "Case 4 vs Case 5: AEC_prev vs AEC_new (스캐너 포함)",
        Inches(5.5), Inches(4.05), Inches(7.5), Inches(0.4),
        size=11, bold=True, color=BLUE, font=FONT_NAME)
    txb(slide,
        "분석 전략: 5-Fold CV (선형 → KFold, 로지스틱 → StratifiedKFold)\n"
        "이진화 기준: 그룹 내 하위 25%(P25) = Low TAMA = 1",
        Inches(5.5), Inches(4.55), Inches(7.5), Inches(0.9),
        size=11, color=DARK, font=FONT_NAME)

    note_line(slide,
              "* 성별 층화 시 Sex 변수 제거 (Female/Male 독립 모델에서 공선성 방지)")
    print(f"  [S{n:02d}] Research Summary")


def s04_section(prs, n, title, bg_color=ORANGE):
    """섹션 구분 슬라이드"""
    slide = new_slide(prs)
    fill_bg(slide, bg_color)
    txb(slide, title,
        Inches(0.56), Inches(3.0), Inches(12.22), Inches(1.0),
        size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font=FONT_BOLD)
    slide_num(slide, n, color=WHITE)
    print(f"  [S{n:02d}] 섹션 — {title}")


def s05_data_gangnam(prs, n):
    """슬라이드 5: 강남 데이터"""
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  Data — 강남", n)

    # 텍스트 설명
    lines = [
        "총 환자수 N = 1,365 도출 경위",
        "  Clinic Data: 3,672건 → 중복 PatientID 제거 → 2,025명",
        "  TAMA outlier(0, 1) 및 BMI 결측치 제거 → 1건 제거",
        "  Axial Data: 1,642건 → Fixed Tube Current Value만 유지 → 1,365명",
        "",
        "성별 분포: 여성(F) = 841명  |  남성(M) = 524명",
        "Low-TAMA 임계값 (그룹 내 P25): 전체 100.0 cm²  |  F: 95.0 cm²  |  M: 132.0 cm²",
    ]
    txb(slide, "\n".join(lines),
        Inches(0.56), CONTENT_Y, Inches(5.8), Inches(2.8),
        size=12, color=DARK, font=FONT_NAME)

    img(slide, REG_DIR / "gangnam" / "03_tama_distribution.png",
        Inches(0.35), Inches(3.6), Inches(6.2), Inches(3.6))
    img(slide, REG_DIR / "gangnam" / "19_tama_sex_distribution.png",
        Inches(6.7), Inches(1.05), Inches(6.3), Inches(6.15))
    print(f"  [S{n:02d}] 강남 데이터")


def s06_scanner_gangnam(prs, n):
    """슬라이드 6: 강남 CT 스캐너 & kVp"""
    slide = new_slide(prs)
    fill_bg(slide, BLACK)
    slide_title_bar(slide, "3.  Detailed Progress  |  CT 스캐너 & kVp — 강남", n, bg_dark=True)

    img(slide, REG_DIR / "gangnam" / "16_scanner_distribution.png",
        Inches(0.35), Inches(1.05), Inches(6.4), Inches(6.2))
    img(slide, REG_DIR / "gangnam" / "17_kvp_distribution.png",
        Inches(6.9), Inches(1.05), Inches(6.1), Inches(6.2))
    print(f"  [S{n:02d}] 강남 스캐너")


def s07_data_sinchon(prs, n):
    """슬라이드 7: 신촌 데이터"""
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  Data — 신촌", n)

    lines = [
        "총 환자수 N = 1,269 도출 경위",
        "  Clinic Data: 3,162건 → 중복 PatientID 제거 → 3,162명",
        "  TAMA 및 BMI 결측치 제거 → 2,254명",
        "  Axial Data: 2,254건 → Fixed/No Tube Current Value 제거 → 1,269명",
        "",
        "성별 분포: 여성(F) = 631명  |  남성(M) = 638명",
        "Low-TAMA 임계값 (그룹 내 P25): 전체 103.0 cm²  |  F: 95.0 cm²  |  M: 131.0 cm²",
    ]
    txb(slide, "\n".join(lines),
        Inches(0.56), CONTENT_Y, Inches(5.8), Inches(2.8),
        size=12, color=DARK, font=FONT_NAME)

    img(slide, REG_DIR / "sinchon" / "03_tama_distribution.png",
        Inches(0.35), Inches(3.6), Inches(6.2), Inches(3.6))
    img(slide, REG_DIR / "sinchon" / "19_tama_sex_distribution.png",
        Inches(6.7), Inches(1.05), Inches(6.3), Inches(6.15))
    print(f"  [S{n:02d}] 신촌 데이터")


def s08_scanner_sinchon(prs, n):
    """슬라이드 8: 신촌 CT 스캐너 & kVp"""
    slide = new_slide(prs)
    fill_bg(slide, BLACK)
    slide_title_bar(slide, "3.  Detailed Progress  |  CT 스캐너 & kVp — 신촌", n, bg_dark=True)

    img(slide, REG_DIR / "sinchon" / "16_scanner_distribution.png",
        Inches(0.35), Inches(1.05), Inches(6.4), Inches(6.2))
    img(slide, REG_DIR / "sinchon" / "17_kvp_distribution.png",
        Inches(6.9), Inches(1.05), Inches(6.1), Inches(6.2))
    print(f"  [S{n:02d}] 신촌 스캐너")


def s10_feature_pipeline(prs, n):
    """슬라이드 10: AEC 피처 선택 파이프라인"""
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    slide_title_bar(slide, "3.  Detailed Progress  |  AEC 피처 선택 파이프라인 (0430 신설)", n, bg_dark=True)

    steps = [
        ("Step 1", "Near-zero\nVariance 제거", "표준화 후\nvariance < 0.01", RGBColor(0x1A, 0x55, 0x8E)),
        ("Step 2", "Pearson 상관\n중복 제거", "|r| ≥ 0.95\n후순위 제거", BLUE),
        ("Step 3", "단변량 필터\n(OR 결합)", "MI > 0\nOR Spearman p<0.05", ORANGE),
        ("Step 4", "앙상블 투표\n+ 완전/SFS 탐색", "LASSO + RFECV\n+ RF Permutation", RGBColor(0x00, 0x8B, 0x45)),
        ("Final", "VIF Pruning\n(mean 고정)", "VIF > 10\n반복 제거", RGBColor(0xC0, 0x00, 0x00)),
    ]

    box_w  = Inches(2.3)
    box_h  = Inches(1.5)
    gap    = Inches(0.15)
    top_y  = Inches(1.2)
    start_x = Inches(0.35)

    for i, (label, title, desc, color) in enumerate(steps):
        lx = start_x + i * (box_w + gap)
        rect(slide, lx, top_y, box_w, box_h, color)
        txb(slide, label, lx, top_y + Inches(0.08), box_w, Inches(0.35),
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BOLD)
        txb(slide, title, lx, top_y + Inches(0.4), box_w, Inches(0.55),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BOLD)
        if i < 4:
            txb(slide, "→", lx + box_w + Inches(0.02), top_y + Inches(0.5), gap, Inches(0.4),
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (label, title, desc, color) in enumerate(steps):
        lx = start_x + i * (box_w + gap)
        txb(slide, desc, lx, top_y + box_h + Inches(0.08), box_w, Inches(0.65),
            size=9, color=SKYBLUE, wrap=True, font=FONT_NAME)

    # 피처 수 변화 표 (강남)
    txb(slide, "파이프라인 단계별 피처 수 변화 (강남 기준)",
        Inches(0.35), Inches(3.15), Inches(12.6), Inches(0.38),
        size=12, bold=True, color=WHITE, font=FONT_BOLD)

    fs_headers = ["단계", "제거 수", "잔여"]
    fs_rows = [
        ["Step 1 - Near-zero var", "0", "65개 전체"],
        ["Step 2 - High correlation", "20", "45개"],
        ["Step 3 - Union pass", "제거 2", "43개 후보"],
        ["Step 4 - Best search", "0", "22개 풀"],
        ["VIF pruning (VIF>10)", "13", "최종 9개"],
    ]
    table(slide, fs_headers, fs_rows,
          Inches(0.35), Inches(3.6), [Inches(5.0), Inches(1.5), Inches(2.5)],
          row_h=Inches(0.4), font_size=10, header_bg=ORANGE, alt_bg=DARK)

    # 최종 선택 피처 요약
    txb(slide, "최종 선택 피처 수",
        Inches(9.5), Inches(3.15), Inches(3.4), Inches(0.38),
        size=12, bold=True, color=WHITE, font=FONT_BOLD)
    fs2_headers = ["데이터셋", "피처수", "Pipeline R²", "Prev R²"]
    fs2_rows = [
        ["강남", "9", "0.1818", "0.1768"],
        ["신촌", "13", "0.0775", "0.0305"],
        ["병합", "11", "0.0881", "0.0965"],
    ]
    table(slide, fs2_headers, fs2_rows,
          Inches(9.5), Inches(3.6), [Inches(1.3), Inches(0.9), Inches(1.3), Inches(0.9)],
          row_h=Inches(0.42), font_size=10, header_bg=ORANGE, alt_bg=DARK)

    note_line(slide, "* mean 피처는 AEC 신호 평균 진폭 = 해석 가능성 최우선 → 모든 최종 세트에 강제 포함",
              color=ORANGE)
    print(f"  [S{n:02d}] 피처 파이프라인")


def s11_feature_gangnam(prs, n):
    """슬라이드 11: 강남 피처 선택 결과"""
    slide = new_slide(prs)
    fill_bg(slide, BLACK)
    slide_title_bar(slide, "3.  Detailed Progress  |  AEC Feature Extracting — 강남", n, bg_dark=True)

    img(slide, FS_DIR / "gangnam" / "01_correlation_heatmap.png",
        Inches(0.35), Inches(1.05), Inches(6.2), Inches(6.2))
    img(slide, FS_DIR / "gangnam" / "05_final_features_summary.png",
        Inches(6.7), Inches(1.05), Inches(6.3), Inches(6.2))

    note_line(slide,
              "최종 선택 (9개): IQR, band2_energy, dominant_freq, mean, slope_max, "
              "spectral_energy, wavelet_cD1_energy, wavelet_cD2_energy, wavelet_energy_ratio_D1",
              color=ORANGE)
    print(f"  [S{n:02d}] 강남 피처 선택")


def s12_feature_sinchon(prs, n):
    """슬라이드 12: 신촌 피처 선택 결과"""
    slide = new_slide(prs)
    fill_bg(slide, BLACK)
    slide_title_bar(slide, "3.  Detailed Progress  |  AEC Feature Extracting — 신촌", n, bg_dark=True)

    img(slide, FS_DIR / "sinchon" / "01_correlation_heatmap.png",
        Inches(0.35), Inches(1.05), Inches(6.2), Inches(6.2))
    img(slide, FS_DIR / "sinchon" / "05_final_features_summary.png",
        Inches(6.7), Inches(1.05), Inches(6.3), Inches(6.2))

    note_line(slide,
              "최종 선택 (13개): IQR, mean, peak_count, peak_max/mean_width, skewness, "
              "slope_mean/min, spectral_centroid/energy, wavelet_cD2_std, cD3_energy, zero_crossing_rate",
              color=ORANGE)
    print(f"  [S{n:02d}] 신촌 피처 선택")


def s13_feature_comparison(prs, n):
    """슬라이드 13: 데이터셋간 피처 비교"""
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    slide_title_bar(slide, "3.  Detailed Progress  |  최종 선택 AEC 피처 비교 (강남 / 신촌 / 병합)", n, bg_dark=True)

    img(slide, FS_DIR / "cross_dataset_feature_heatmap.png",
        Inches(0.35), Inches(1.05), Inches(6.8), Inches(6.2))

    # 피처 일치 요약
    txb(slide, "공통 선택 피처",
        Inches(7.4), Inches(1.1), Inches(5.5), Inches(0.4),
        size=13, bold=True, color=ORANGE, font=FONT_BOLD)
    txb(slide,
        "• mean  (모든 데이터셋 공통, 강제 포함)\n"
        "• spectral_energy  (강남 + 신촌 + 병합)\n"
        "• IQR  (강남 + 신촌)\n"
        "• zero_crossing_rate  (신촌 + 병합)\n"
        "• slope_mean  (신촌 + 병합)",
        Inches(7.4), Inches(1.6), Inches(5.5), Inches(2.0),
        size=12, color=WHITE, font=FONT_NAME)

    txb(slide, "데이터셋별 R² (5-Fold CV)",
        Inches(7.4), Inches(3.9), Inches(5.5), Inches(0.4),
        size=13, bold=True, color=ORANGE, font=FONT_BOLD)
    r2_headers = ["데이터셋", "AEC_new R²", "AEC_prev R²", "ΔR²"]
    r2_rows = [
        ["강남", "0.1818", "0.1768", "+0.0050 ↑"],
        ["신촌", "0.0775", "0.0305", "+0.0469 ↑"],
        ["병합", "0.0881", "0.0965", "-0.0084 ↓"],
    ]
    table(slide, r2_headers, r2_rows,
          Inches(7.4), Inches(4.4), [Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.3)],
          row_h=Inches(0.42), font_size=10, header_bg=ORANGE, alt_bg=DARK)

    note_line(slide,
              "⚠  병합 데이터셋은 강남·신촌 공통 피처만 허용 → 선택 범위 제한으로 AEC_new R² 소폭 저하",
              color=ORANGE)
    print(f"  [S{n:02d}] 피처 비교")


def s15_univariate(prs, n):
    """슬라이드 15: 단변량 분석"""
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  단변량 분석 (AEC_new 최종 피처 vs TAMA)", n)

    img(slide, FS_DIR / "gangnam" / "07_final_scatter_vs_tama.png",
        Inches(0.35), Inches(1.05), Inches(6.4), Inches(6.2))
    img(slide, FS_DIR / "sinchon" / "07_final_scatter_vs_tama.png",
        Inches(6.85), Inches(1.05), Inches(6.1), Inches(6.2))

    note_line(slide,
              "* p < 0.05 유의  |  mean: 양의 상관 (진폭 클수록 TAMA 큼) "
              "  |  spectral_energy, wavelet_energy: 주파수 에너지 관련 피처",
              color=ORANGE)
    print(f"  [S{n:02d}] 단변량")


def _reg_table_linear(slide, hosp_key: str, group: str,
                       left, top, col_w):
    """선형 회귀 결과 표 (하드코딩 데이터)"""
    data = {
        "gangnam": {
            "all": {
                "N": "1,365", "P25": "100.0",
                "rows": [
                    ["Case 1  Clinical", "3", "0.662 ± 0.037", "12.83", "17.12"],
                    ["Case 2  +AEC_prev", "7", "0.678 ± 0.036", "12.53", "16.71"],
                    ["Case 3  +AEC_new", "14", "0.669 ± 0.038", "12.59", "16.93"],
                    ["Case 4  +AEC_prev +Scanner", "46", "0.668 ± 0.032", "12.77", "16.96"],
                    ["Case 5  +AEC_new +Scanner", "53", "0.662 ± 0.035", "12.74", "17.11"],
                ],
            },
        },
        "sinchon": {
            "all": {
                "N": "1,269", "P25": "103.0",
                "rows": [
                    ["Case 1  Clinical", "3", "0.640 ± 0.063", "13.82", "18.42"],
                    ["Case 2  +AEC_prev", "7", "0.641 ± 0.061", "13.83", "18.37"],
                    ["Case 3  +AEC_new", "14", "0.637 ± 0.067", "13.81", "18.48"],
                    ["Case 4  +AEC_prev +Scanner", "59", "0.634 ± 0.059", "14.00", "18.56"],
                    ["Case 5  +AEC_new +Scanner", "66", "0.638 ± 0.054", "13.94", "18.48"],
                ],
            },
        },
    }
    d = data.get(hosp_key, {}).get(group, {})
    headers = ["Case", "N feat", "R² (mean±std)", "MAE (cm²)", "RMSE (cm²)"]
    rows = d.get("rows", [])
    table(slide, headers, rows, left, top, col_w, row_h=Inches(0.42), font_size=10)


def _reg_table_logistic(slide, hosp_key: str, group: str,
                         left, top, col_w):
    data = {
        "gangnam": {
            "all": {
                "rows": [
                    ["Case 1  Clinical", "0.8293 ± 0.0208", "0.7868", "0.3414", "0.9243"],
                    ["Case 2  +AEC_prev", "0.8346 ± 0.0247", "0.7846", "0.3447", "0.9204"],
                    ["Case 3  +AEC_new", "0.8352 ± 0.0229", "0.7853", "0.3571", "0.9175"],
                    ["Case 4  +AEC_prev +Scanner", "0.8306 ± 0.0209", "0.7883", "0.3758", "0.9156"],
                    ["Case 5  +AEC_new +Scanner", "0.8320 ± 0.0182", "0.7941", "0.3913", "0.9185"],
                ],
            },
        },
        "sinchon": {
            "all": {
                "rows": [
                    ["Case 1  Clinical", "0.8577 ± 0.0333", "0.7999", "0.5256", "0.8901"],
                    ["Case 2  +AEC_prev", "0.8587 ± 0.0347", "0.7983", "0.5351", "0.8848"],
                    ["Case 3  +AEC_new", "0.8533 ± 0.0398", "0.8046", "0.5413", "0.8911"],
                    ["Case 4  +AEC_prev +Scanner", "0.8579 ± 0.0328", "0.7943", "0.5382", "0.8785"],
                    ["Case 5  +AEC_new +Scanner", "0.8531 ± 0.0357", "0.7943", "0.5414", "0.8775"],
                ],
            },
        },
    }
    d = data.get(hosp_key, {}).get(group, {})
    headers = ["Case", "AUC (mean±std)", "Accuracy", "Sensitivity", "Specificity"]
    rows = d.get("rows", [])
    table(slide, headers, rows, left, top, col_w, row_h=Inches(0.42), font_size=10)


def s17_linear_gangnam(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  Linear Regression — 강남 (전체, N=1,365)", n)

    col_w = [Inches(3.5), Inches(0.8), Inches(2.5), Inches(1.3), Inches(1.3)]
    _reg_table_linear(slide, "gangnam", "all",
                       Inches(0.35), CONTENT_Y, col_w)

    img(slide, REG_DIR / "gangnam" / "all" / "01_linear_actual_vs_pred.png",
        Inches(0.35), Inches(3.75), Inches(6.5), Inches(3.5))
    img(slide, REG_DIR / "gangnam" / "all" / "02_linear_metrics_comparison.png",
        Inches(7.0), Inches(3.75), Inches(6.0), Inches(3.5))

    key_finding(slide,
                "AEC_prev 추가(Case1→2)로 R² +0.016 향상  |  "
                "AEC_new vs AEC_prev (Case3 vs Case2): Δ = -0.009  |  "
                "임상 변수(BMI 포함) 기반 R²=0.662가 기준선",
                y_pos=Inches(3.45))
    note_line(slide,
              "* p < 0.05 유의  |  표준화 계수: 1 SD 변화 시 TAMA 변화량 (cm²)")
    print(f"  [S{n:02d}] 강남 선형 회귀")


def s18_logistic_gangnam(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    slide_title_bar(slide, "3.  Detailed Progress  |  Logistic Regression — 강남 (전체, P25=100 cm²)", n, bg_dark=True)

    col_w = [Inches(3.5), Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.5)]
    _reg_table_logistic(slide, "gangnam", "all",
                         Inches(0.35), CONTENT_Y, col_w)

    img(slide, REG_DIR / "gangnam" / "all" / "04_logistic_roc.png",
        Inches(0.35), Inches(3.75), Inches(6.5), Inches(3.5))
    img(slide, REG_DIR / "gangnam" / "all" / "05_logistic_metrics_comparison.png",
        Inches(7.0), Inches(3.75), Inches(6.0), Inches(3.5))

    key_finding(slide,
                "AEC_new 추가(Case3) AUC = 0.8352  |  "
                "Case 1 기준선 AUC = 0.8293  |  ΔAUC(new-prev) = +0.0006 (미미)",
                y_pos=Inches(3.45), color=WHITE)
    note_line(slide,
              "* Low TAMA = 그룹 내 하위 25% 기준  |  Youden Index 최적 임계값에서 혼동행렬 산출",
              color=ORANGE)
    print(f"  [S{n:02d}] 강남 로지스틱")


def s19_diag_gangnam(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, BLACK)
    slide_title_bar(slide, "3.  Detailed Progress  |  모델 진단 — 강남", n, bg_dark=True)

    img(slide, REG_DIR / "gangnam" / "05_linear_residuals.png",
        Inches(0.35), Inches(1.1), Inches(7.2), Inches(3.4))
    img(slide, REG_DIR / "gangnam" / "06_linear_forest.png",
        Inches(7.7), Inches(1.1), Inches(5.3), Inches(3.4))
    img(slide, REG_DIR / "gangnam" / "09_logistic_calibration.png",
        Inches(0.35), Inches(4.7), Inches(5.5), Inches(2.6))
    img(slide, REG_DIR / "gangnam" / "10_logistic_confusion.png",
        Inches(6.1), Inches(4.7), Inches(6.9), Inches(2.6))

    note_line(slide,
              "잔차 진단 4-Panel (좌상단) + 유의 계수 Forest Plot (우상단) | "
              "Calibration Plot (좌하단) + Confusion Matrix Youden (우하단)",
              color=ORANGE)
    print(f"  [S{n:02d}] 강남 진단")


def s20_linear_sinchon(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  Linear Regression — 신촌 (전체, N=1,269)", n)

    col_w = [Inches(3.5), Inches(0.8), Inches(2.5), Inches(1.3), Inches(1.3)]
    _reg_table_linear(slide, "sinchon", "all",
                       Inches(0.35), CONTENT_Y, col_w)

    img(slide, REG_DIR / "sinchon" / "all" / "01_linear_actual_vs_pred.png",
        Inches(0.35), Inches(3.75), Inches(6.5), Inches(3.5))
    img(slide, REG_DIR / "sinchon" / "all" / "02_linear_metrics_comparison.png",
        Inches(7.0), Inches(3.75), Inches(6.0), Inches(3.5))

    key_finding(slide,
                "Case 1 기준선 R²=0.640 (강남 0.662와 유사)  |  "
                "AEC 기여도 미미 (Δ ≤ 0.002)  |  스캐너 추가 시 오히려 소폭 저하",
                y_pos=Inches(3.45))
    note_line(slide,
              "* 신촌 데이터에서 AEC_new R² 소폭 저하 → 병원별 AEC 신호 특성 차이 시사")
    print(f"  [S{n:02d}] 신촌 선형 회귀")


def s21_logistic_sinchon(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    slide_title_bar(slide, "3.  Detailed Progress  |  Logistic Regression — 신촌 (전체, P25=103 cm²)", n, bg_dark=True)

    col_w = [Inches(3.5), Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.5)]
    _reg_table_logistic(slide, "sinchon", "all",
                         Inches(0.35), CONTENT_Y, col_w)

    img(slide, REG_DIR / "sinchon" / "all" / "04_logistic_roc.png",
        Inches(0.35), Inches(3.75), Inches(6.5), Inches(3.5))
    img(slide, REG_DIR / "sinchon" / "all" / "05_logistic_metrics_comparison.png",
        Inches(7.0), Inches(3.75), Inches(6.0), Inches(3.5))

    key_finding(slide,
                "Case 1 AUC = 0.8577 (강남보다 높음)  |  "
                "AEC 추가 효과 미미 (ΔAUC ≤ 0.001)  |  AEC_new Case3 AUC = 0.8533",
                y_pos=Inches(3.45), color=WHITE)
    note_line(slide,
              "* Sensitivity 강남(0.34) vs 신촌(0.53): 신촌 임계값 기준 Low TAMA 탐지력 우수",
              color=ORANGE)
    print(f"  [S{n:02d}] 신촌 로지스틱")


def s22_diag_sinchon(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, BLACK)
    slide_title_bar(slide, "3.  Detailed Progress  |  모델 진단 — 신촌", n, bg_dark=True)

    img(slide, REG_DIR / "sinchon" / "05_linear_residuals.png",
        Inches(0.35), Inches(1.1), Inches(7.2), Inches(3.4))
    img(slide, REG_DIR / "sinchon" / "06_linear_forest.png",
        Inches(7.7), Inches(1.1), Inches(5.3), Inches(3.4))
    img(slide, REG_DIR / "sinchon" / "09_logistic_calibration.png",
        Inches(0.35), Inches(4.7), Inches(5.5), Inches(2.6))
    img(slide, REG_DIR / "sinchon" / "10_logistic_confusion.png",
        Inches(6.1), Inches(4.7), Inches(6.9), Inches(2.6))

    note_line(slide,
              "잔차 진단 4-Panel + Forest Plot | Calibration + Confusion Matrix",
              color=ORANGE)
    print(f"  [S{n:02d}] 신촌 진단")


def s23_cross_hospital(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  교차 병원 비교 (강남 vs 신촌)", n)

    # R² 비교 표
    txb(slide, "선형 회귀 R² (전체 그룹)",
        Inches(0.35), CONTENT_Y, Inches(6.0), Inches(0.38),
        size=12, bold=True, color=NAVY, font=FONT_BOLD)
    lin_headers = ["Case", "강남", "신촌"]
    lin_rows = [
        ["Case 1  Clinical", "0.6620", "0.6398"],
        ["Case 2  +AEC_prev", "0.6780", "0.6415"],
        ["Case 3  +AEC_new", "0.6694", "0.6371"],
        ["Case 4  +AEC_prev +Scanner", "0.6682", "0.6340"],
        ["Case 5  +AEC_new +Scanner", "0.6620", "0.6376"],
    ]
    table(slide, lin_headers, lin_rows,
          Inches(0.35), Inches(1.42), [Inches(4.0), Inches(1.3), Inches(1.3)],
          row_h=Inches(0.42), font_size=10)

    # AUC 비교 표
    txb(slide, "로지스틱 회귀 AUC (전체 그룹)",
        Inches(7.0), CONTENT_Y, Inches(6.0), Inches(0.38),
        size=12, bold=True, color=NAVY, font=FONT_BOLD)
    log_rows = [
        ["Case 1  Clinical", "0.8293", "0.8577"],
        ["Case 2  +AEC_prev", "0.8346", "0.8587"],
        ["Case 3  +AEC_new", "0.8352", "0.8533"],
        ["Case 4  +AEC_prev +Scanner", "0.8306", "0.8579"],
        ["Case 5  +AEC_new +Scanner", "0.8320", "0.8531"],
    ]
    table(slide, lin_headers, log_rows,
          Inches(7.0), Inches(1.42), [Inches(4.0), Inches(1.3), Inches(1.3)],
          row_h=Inches(0.42), font_size=10)

    img(slide, REG_DIR / "09_cross_hospital_comparison.png",
        Inches(0.35), Inches(3.85), Inches(12.6), Inches(3.4))

    key_finding(slide,
                "강남 R² ≈ 신촌 R² (≈0.63~0.68)  |  신촌 AUC가 강남보다 소폭 높음 (0.857 vs 0.829)  |  "
                "두 병원 모두 AEC 기여도는 미미",
                y_pos=Inches(3.6))
    print(f"  [S{n:02d}] 교차 병원 비교")


def s24_sex_strat(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    slide_title_bar(slide, "3.  Detailed Progress  |  성별 층화 분석 (전체 / 여성(F) / 남성(M))", n, bg_dark=True)

    # 성별 층화 R² 표 (강남)
    txb(slide, "선형 R² — 강남",
        Inches(0.35), CONTENT_Y, Inches(6.5), Inches(0.35),
        size=11, bold=True, color=ORANGE, font=FONT_BOLD)
    s_headers = ["Case", "전체", "여성(F)", "남성(M)"]
    s_lin_rows = [
        ["Case 1  Clinical", "0.662", "0.241", "0.324"],
        ["Case 2  +AEC_prev", "0.678", "0.266", "0.350"],
        ["Case 3  +AEC_new", "0.669", "0.262", "0.317"],
        ["Case 4  +AEC_prev +Scanner", "0.668", "0.264", "0.347"],
        ["Case 5  +AEC_new +Scanner", "0.662", "0.243", "0.314"],
    ]
    table(slide, s_headers, s_lin_rows,
          Inches(0.35), Inches(1.38), [Inches(4.0), Inches(0.9), Inches(1.1), Inches(1.1)],
          row_h=Inches(0.40), font_size=9, header_bg=ORANGE, alt_bg=DARK)

    txb(slide, "AUC — 강남",
        Inches(7.5), CONTENT_Y, Inches(5.5), Inches(0.35),
        size=11, bold=True, color=ORANGE, font=FONT_BOLD)
    s_log_rows = [
        ["Case 1  Clinical", "0.829", "0.706", "0.792"],
        ["Case 2  +AEC_prev", "0.835", "0.716", "0.798"],
        ["Case 3  +AEC_new", "0.835", "0.710", "0.787"],
        ["Case 4  +AEC_prev +Scanner", "0.831", "0.729", "0.795"],
        ["Case 5  +AEC_new +Scanner", "0.832", "0.729", "0.795"],
    ]
    table(slide, s_headers, s_log_rows,
          Inches(7.5), Inches(1.38), [Inches(4.0), Inches(0.9), Inches(1.1), Inches(1.1)],
          row_h=Inches(0.40), font_size=9, header_bg=ORANGE, alt_bg=DARK)

    img(slide, FIG_DIR / "05_sex_strat_linear.png",
        Inches(0.35), Inches(3.85), Inches(12.6), Inches(3.4))

    note_line(slide,
              "* 전체 그룹이 R²·AUC 모두 가장 높음 (N이 가장 크기 때문)  |  "
              "성별 단독 모델에서 여성 R² < 남성 R²",
              color=ORANGE)
    print(f"  [S{n:02d}] 성별 층화")


def s25_bmi(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "3.  Detailed Progress  |  BMI 기여도 분석 (0430 신설)", n)

    # BMI 효과 표
    txb(slide, "BMI 추가 효과 Δ (+BMI − no BMI)  —  강남",
        Inches(0.35), CONTENT_Y, Inches(7.5), Inches(0.38),
        size=12, bold=True, color=NAVY, font=FONT_BOLD)
    b_headers = ["Case", "ΔR²", "ΔRMSE (cm²)", "ΔAUC", "해석"]
    b_rows = [
        ["Case 1  (Clinical)", "+0.1190", "-2.79", "+0.0649", "BMI 보정 유효"],
        ["Case 2  (+AEC_prev)", "+0.0480", "-1.20", "+0.0168", "AEC 투입 후 효과 감쇠"],
        ["Case 4  (+Scanner)", "+0.0472", "-1.17", "+0.0152", "Scanner 추가 시에도 유지"],
    ]
    table(slide, b_headers, b_rows,
          Inches(0.35), Inches(1.42), [Inches(2.8), Inches(1.2), Inches(1.5), Inches(1.2), Inches(2.6)],
          row_h=Inches(0.45), font_size=10)

    # 신촌
    txb(slide, "BMI 추가 효과 Δ (+BMI − no BMI)  —  신촌",
        Inches(0.35), Inches(3.4), Inches(7.5), Inches(0.38),
        size=12, bold=True, color=NAVY, font=FONT_BOLD)
    b_rows2 = [
        ["Case 1  (Clinical)", "+0.1182", "-2.82", "+0.0560", "BMI 보정 유효"],
        ["Case 2  (+AEC_prev)", "+0.0971", "-2.36", "+0.0381", "AEC 투입 후 효과 감쇠"],
        ["Case 4  (+Scanner)", "+0.0990", "-2.36", "+0.0237", "Scanner 추가 시에도 유지"],
    ]
    table(slide, b_headers, b_rows2,
          Inches(0.35), Inches(3.85), [Inches(2.8), Inches(1.2), Inches(1.5), Inches(1.2), Inches(2.6)],
          row_h=Inches(0.45), font_size=10)

    img(slide, REG_DIR / "gangnam" / "bmi_delta_effect.png",
        Inches(7.8), Inches(1.05), Inches(5.3), Inches(6.2))

    key_finding(slide,
                "Case 1에서 BMI 추가로 R² +0.119 향상 (강남·신촌 동일)  |  "
                "AEC 투입 후 BMI 효과 감쇠 → AEC와 BMI가 일부 공통 정보 공유",
                y_pos=Inches(6.65))
    print(f"  [S{n:02d}] BMI 기여도")


def s26_aec_comparison(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    slide_title_bar(slide, "3.  Detailed Progress  |  AEC_prev vs AEC_new — 피처 선택 전략 비교", n, bg_dark=True)

    img(slide, FIG_DIR / "07_aec_prev_vs_new.png",
        Inches(0.35), Inches(1.1), Inches(12.6), Inches(4.0))

    txb(slide,
        "• Case 2 vs Case 3: AEC_prev vs AEC_new (스캐너 없음)\n"
        "• Case 4 vs Case 5: AEC_prev vs AEC_new (스캐너 포함)\n"
        "• 강남: AEC_prev ≥ AEC_new (Δ = −0.009 for R², −0.0006 for AUC)\n"
        "• 신촌: AEC_prev ≈ AEC_new (선형 Δ = −0.004, AUC Δ = −0.005)",
        Inches(0.35), Inches(5.3), Inches(12.6), Inches(1.1),
        size=12, color=WHITE, font=FONT_NAME)

    img(slide, FS_DIR / "cross_dataset_comparison_r2.png",
        Inches(0.35), Inches(6.2), Inches(6.5), Inches(1.1))

    note_line(slide,
              "⚠  AEC_new가 AEC_prev 대비 우위를 보이지 않음 → 더 많은 환자 데이터 확보 후 재검토 필요",
              color=ORANGE)
    print(f"  [S{n:02d}] AEC 비교")


def s28_achievements(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "4.  결론  |  0430 핵심 성과", n)

    ach = [
        ("1. BMI 보정",
         "Case 1 기준선에서 R² +0.12 향상 — BMI가 TAMA의 강력한 독립 예측 변수임을 두 병원 모두에서 확인"),
        ("2. 자동 피처 선택",
         "60개+ AEC 피처에서 4단계 파이프라인으로 과적합·다중공선성 없는 객관적 세트 도출 (강남 9개, 신촌 13개)"),
        ("3. 성별 층화",
         "전체/여성/남성 독립 모델로 이질성 탐색 — 성별 그룹 내 예측력은 전체 모델 대비 낮음 (소표본 영향)"),
        ("4. 다병원 검증",
         "강남·신촌 교차 검증으로 임상 변수 기반 기준선(R²≈0.64~0.66, AUC≈0.83~0.86) 재현성 확인"),
    ]

    y = CONTENT_Y
    for label, content in ach:
        rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.88), SKYBLUE)
        txb(slide, label, Inches(0.5), y + Inches(0.05),
            Inches(2.8), Inches(0.38), size=13, bold=True, color=NAVY, font=FONT_BOLD)
        txb(slide, content, Inches(3.4), y + Inches(0.05),
            Inches(9.2), Inches(0.72), size=12, color=DARK, wrap=True, font=FONT_NAME)
        y += Inches(0.98)

    print(f"  [S{n:02d}] 핵심 성과")


def s29_next_steps(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, WHITE)
    slide_title_bar(slide, "4.  결론  |  한계 및 향후 과제 (Next Plan)", n)

    txb(slide, "한계",
        Inches(0.35), CONTENT_Y, Inches(6.0), Inches(0.4),
        size=14, bold=True, color=NAVY, font=FONT_BOLD)
    limits = [
        "• AEC_new vs AEC_prev 성능 차이가 미미 — 더 많은 환자 데이터 필요",
        "• 성별 층화 시 소그룹(여성 단독, 남성 단독) 표본 크기에 따른 불안정성",
        "• 단면 연구 설계 — 인과 추론을 위한 전향적 코호트 연구 권장",
    ]
    txb(slide, "\n".join(limits),
        Inches(0.35), Inches(1.45), Inches(12.3), Inches(1.3),
        size=12, color=DARK, font=FONT_NAME)

    txb(slide, "향후 과제",
        Inches(0.35), Inches(2.9), Inches(6.0), Inches(0.4),
        size=14, bold=True, color=NAVY, font=FONT_BOLD)
    next_items = [
        ("1.", "ML 앙상블 모델 비교", "Random Forest, XGBoost 등으로 선형 회귀 대비 성능 비교"),
        ("2.", "SMI(Skeletal Muscle Index) 적용", "TAMA / Height² → 키 보정 지표로 재분석"),
        ("3.", "Raw AEC 시계열 딥러닝", "200포인트 신호를 1D CNN / LSTM으로 직접 학습 → 피처 엔지니어링 없는 end-to-end 방식"),
        ("4.", "외부 검증 확장", "강남↔신촌 외 추가 병원 데이터로 일반화 검증"),
    ]
    y_pos = Inches(3.35)
    for num, title, desc in next_items:
        rect(slide, Inches(0.35), y_pos, Inches(12.6), Inches(0.74), SKYBLUE)
        txb(slide, f"{num}  {title}", Inches(0.5), y_pos + Inches(0.05),
            Inches(4.0), Inches(0.35), size=12, bold=True, color=NAVY, font=FONT_BOLD)
        txb(slide, desc, Inches(4.6), y_pos + Inches(0.05),
            Inches(8.0), Inches(0.58), size=11, color=DARK, font=FONT_NAME)
        y_pos += Inches(0.84)

    note_line(slide,
              "* 자동 생성: 07_generate_sharing_ppt.py (0430)",
              color=RGBColor(0xBB, 0xBB, 0xBB))
    print(f"  [S{n:02d}] Next Plan")


def s30_end(prs, n):
    slide = new_slide(prs)
    fill_bg(slide, NAVY)
    txb(slide, "Thank you",
        Inches(0.56), Inches(2.8), Inches(12.22), Inches(1.0),
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BOLD)
    txb(slide, "JangHoon Chun  |  jhjun0328@naver.com\n"
               "Dept. of Medical Device Engineering & Management\n"
               "Gangnam Severance Hospital, Yonsei University",
        Inches(0.56), Inches(4.2), Inches(12.22), Inches(1.0),
        size=15, color=SKYBLUE, align=PP_ALIGN.CENTER, font=FONT_NAME)
    slide_num(slide, n, color=SKYBLUE)
    print(f"  [S{n:02d}] 마지막")


# ═══════════════════════════════════════════════════════════════
# 3. 메인
# ═══════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("  07_generate_sharing_ppt.py (0430) — 공유용 발표 자료")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("\n[슬라이드 생성 중...]")

    s01_title(prs)                              # 1
    s02_last_meeting(prs, 2)                    # 2
    s03_research_summary(prs, 3)               # 3

    s04_section(prs, 4, "데이터 정리", ORANGE)  # 4

    s05_data_gangnam(prs, 5)                    # 5
    s06_scanner_gangnam(prs, 6)                 # 6
    s07_data_sinchon(prs, 7)                    # 7
    s08_scanner_sinchon(prs, 8)                 # 8

    s04_section(prs, 9, "AEC 피처 선택", NAVY)  # 9

    s10_feature_pipeline(prs, 10)              # 10
    s11_feature_gangnam(prs, 11)               # 11
    s12_feature_sinchon(prs, 12)               # 12
    s13_feature_comparison(prs, 13)            # 13

    s04_section(prs, 14, "Univariate Analysis", ORANGE)  # 14

    s15_univariate(prs, 15)                    # 15

    s04_section(prs, 16, "Multivariable Analysis", ORANGE)  # 16

    s17_linear_gangnam(prs, 17)                # 17
    s18_logistic_gangnam(prs, 18)              # 18
    s19_diag_gangnam(prs, 19)                  # 19
    s20_linear_sinchon(prs, 20)                # 20
    s21_logistic_sinchon(prs, 21)              # 21
    s22_diag_sinchon(prs, 22)                  # 22
    s23_cross_hospital(prs, 23)                # 23
    s24_sex_strat(prs, 24)                     # 24
    s25_bmi(prs, 25)                           # 25
    s26_aec_comparison(prs, 26)               # 26

    s04_section(prs, 27, "결론 & 핵심 성과", ORANGE)  # 27

    s28_achievements(prs, 28)                  # 28
    s29_next_steps(prs, 29)                    # 29
    s30_end(prs, 30)                           # 30

    print(f"\n[PPT 저장 중...]")
    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPT))
    print(f"  저장 완료: {OUT_PPT}")
    print(f"  슬라이드 수: {len(prs.slides)}")
    print("=" * 60)


if __name__ == '__main__':
    main()
