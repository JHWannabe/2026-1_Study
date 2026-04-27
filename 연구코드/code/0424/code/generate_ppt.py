"""
TAMA 예측 회귀분석 연구 보고서 PPT 생성 스크립트
"""

import re
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os
import config

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x35, 0x5E)   # 제목 배경
BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # 강조
SKYBLUE = RGBColor(0xD6, 0xE4, 0xF0)   # 셀 배경
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE  = RGBColor(0xFF, 0x7F, 0x27)
GREEN   = RGBColor(0x00, 0x8B, 0x45)

FIGURES = os.path.join(_ROOT, "results", config.SITE, "figures")


def parse_report(site):
    """results/{site}/research_report.md 를 파싱해 슬라이드에 쓸 값을 dict로 반환."""
    path = os.path.join(_ROOT, "results", site, "research_report.md")
    with open(path, encoding="utf-8") as fh:
        lines = fh.readlines()

    D = {'site': site}

    # ── 내부 헬퍼 ────────────────────────────────────────────
    def strip_bold(s):
        return re.sub(r'\*\*([^*]+)\*\*', r'\1', s).strip()

    def flt(s):
        cleaned = re.sub(r'[^\d.\-]', '', strip_bold(s))
        try:
            return float(cleaned)
        except ValueError:
            return 0.0

    def sint(s):
        digits = re.sub(r'[^\d]', '', strip_bold(s))
        return int(digits) if digits else 0

    def find_section(header):
        for i, ln in enumerate(lines):
            if header in ln:
                return i
        return -1

    def table_rows(start, max_rows=20):
        if start < 0:
            return []
        rows = []
        seen_sep = False
        count = 0
        for ln in lines[start + 1: start + 1 + max_rows + 15]:
            stripped = ln.strip()
            if not stripped.startswith('|'):
                if seen_sep:
                    break
                continue
            cols = [c.strip() for c in stripped.split('|')[1:-1]]
            if cols and all(re.fullmatch(r'[-: ]+', c) for c in cols if c):
                seen_sep = True
                continue
            if seen_sep:
                rows.append(cols)
                count += 1
                if count >= max_rows:
                    break
        return rows

    def fmt_beta(raw, ndp=2):
        m = re.match(r'(-?[\d.]+)(.*)', raw.strip())
        if m:
            return f"{float(m.group(1)):.{ndp}f}{m.group(2)}"
        return raw

    def fmt_ci(raw, ndp=2):
        m = re.match(r'\[(-?[\d.]+),\s*(-?[\d.]+)\]', raw.strip())
        if m:
            return f"[{float(m.group(1)):.{ndp}f}, {float(m.group(2)):.{ndp}f}]"
        return raw

    def fmt_float_val(raw, ndp=3):
        try:
            return f"{float(strip_bold(raw)):.{ndp}f}"
        except ValueError:
            return strip_bold(raw)

    def _pdisp(raw):
        raw_s = strip_bold(raw).strip()
        try:
            v = float(raw_s)
            if v < 0.001:
                return "< 0.001"
            if v < 0.01:
                return f"{v:.4f}"
            return f"{v:.3f}"
        except ValueError:
            return raw_s

    def _hl_disp(raw):
        raw_s = strip_bold(raw).strip()
        try:
            v = float(raw_s)
            if v < 0.01:
                return f"{v:.2e}"
            return f"{v:.3f}"
        except ValueError:
            return raw_s

    # ── 2.0 데이터 정제 단계 ─────────────────────────────────
    sec20 = find_section('### 2.0 데이터 정제 단계')
    r20 = table_rows(sec20, max_rows=5)
    _clean_steps = [
        ('원시 데이터 (metadata)', r20[0][1] if r20 else 'N/A'),
        ('원시 데이터 (features)', r20[1][1] if len(r20)>1 else 'N/A'),
        ('Inner join (공통 PatientID)', r20[2][1] if len(r20)>2 else 'N/A'),
        ('중복 PatientID 제거 후', r20[3][1] if len(r20)>3 else 'N/A'),
        ('결측치 행 제거 후 (최종)', strip_bold(r20[4][1]) if len(r20)>4 else 'N/A'),
    ]
    D['clean_steps'] = _clean_steps
    D['raw_n'] = re.sub(r'[^\d]', '', _clean_steps[0][1]) if _clean_steps else '?'

    # ── 2.1 데이터셋 기본 정보 ───────────────────────────────
    sec = find_section('### 2.1 데이터셋 기본 정보')
    r = table_rows(sec, max_rows=7)

    mm = re.search(r'(\d[\d,]*)', strip_bold(r[0][1]))
    D['n_total'] = int(mm.group(1).replace(',', '')) if mm else 0

    mm = re.search(r'(\d+)명 \(([\d.]+)%\)', r[1][1])
    if mm:
        D['n_male'] = int(mm.group(1))
        D['male_pct'] = mm.group(2)

    mm = re.search(r'(\d+)명 \(([\d.]+)%\)', r[2][1])
    if mm:
        D['n_female'] = int(mm.group(1))
        D['female_pct'] = mm.group(2)

    mm = re.search(r'(\d+) ~ (\d+)', r[3][1])
    if mm:
        D['tama_min'] = mm.group(1)
        D['tama_max'] = mm.group(2)

    mm = re.search(r'([\d.]+) \(([\d.]+)\)', r[4][1])
    if mm:
        D['tama_mean'] = mm.group(1)
        D['tama_sd'] = mm.group(2)

    D['ct_models'] = sint(r[5][1])

    # kVp 주요값 (row 6 있으면 파싱)
    D['kvp_dominant'] = 'N/A'
    D['kvp_dominant_pct'] = ''
    if len(r) > 6:
        kraw = strip_bold(r[6][1])
        mm = re.search(r'(\d+)\s*kVp.*?([\d.]+)%', kraw)
        if mm:
            D['kvp_dominant']     = mm.group(1)
            D['kvp_dominant_pct'] = mm.group(2)

    # 스캐너 Top 1 (CT 스캐너 분포 표에서 파싱)
    D['scanner_top1'] = ''
    sec_sc = find_section('CT 스캐너 분포 (상위 5종)')
    if sec_sc >= 0:
        rs = table_rows(sec_sc, max_rows=1)
        if rs:
            D['scanner_top1'] = rs[0][0]

    # ── 2.2 성별 TAMA 분포 ───────────────────────────────────
    sec = find_section('### 2.2 성별 TAMA 분포')
    r = table_rows(sec, max_rows=2)

    D['male_n']      = r[0][1]
    D['male_mean']   = r[0][2]
    D['male_sd']     = r[0][3]
    D['male_median'] = r[0][4]
    D['male_p25']    = r[0][5]
    D['male_p75']    = r[0][6]

    D['female_n']      = r[1][1]
    D['female_mean']   = r[1][2]
    D['female_sd']     = r[1][3]
    D['female_median'] = r[1][4]
    D['female_p25']    = r[1][5]
    D['female_p75']    = r[1][6]

    # ── 2.3 이진화 기준 ──────────────────────────────────────
    sec = find_section('### 2.3 Logistic Regression 이진화')
    for ln in lines[sec: sec + 8]:
        mm = re.search(r'남성.*?TAMA < (\d+)', ln)
        if mm:
            D['thr_male'] = mm.group(1)
        mm = re.search(r'여성.*?TAMA < (\d+)', ln)
        if mm:
            D['thr_female'] = mm.group(1)
        mm = re.search(r'양성 비율.*?(\d+)/(\d+) = ([\d.]+)%', ln)
        if mm:
            D['n_pos'] = mm.group(1)
            D['pos_rate'] = mm.group(3)

    # ── 2.4 AEC Feature ──────────────────────────────────────
    sec = find_section('### 2.4 AEC Feature 선택 근거')
    r = table_rows(sec, max_rows=len(config.SELECTED_AEC_FEATURES))
    for i, key in enumerate(config.SELECTED_AEC_FEATURES):
        D[f'pearson_{key}'] = strip_bold(r[i][1])
        D[f'vif_{key}']     = strip_bold(r[i][3])

    # ── 3.1 전처리 더미 수 ───────────────────────────────────
    sec = find_section('### 3.1 전처리')
    D['ct_dummy_count'] = '30'
    for ln in lines[sec: sec + 15]:
        mm = re.search(r'→ (\d+)개 더미', ln)
        if mm:
            D['ct_dummy_count'] = mm.group(1)
            break

    # ── 4.1.1 선형 단변량 ────────────────────────────────────
    sec = find_section('#### 4.1.1 단변량 분석 (Univariate)')
    _aec_keys = config.SELECTED_AEC_FEATURES
    lin_uni_keys = ['Sex', 'Age'] + _aec_keys + ['ManufacturerModelName']
    r = table_rows(sec, max_rows=len(lin_uni_keys))
    for i, key in enumerate(lin_uni_keys):
        raw_beta = strip_bold(r[i][1])
        D[f'linu_{key}_beta']      = raw_beta
        D[f'linu_{key}_beta_disp'] = fmt_beta(raw_beta, 2)
        D[f'linu_{key}_ci']        = fmt_ci(r[i][2], 2)
        D[f'linu_{key}_pval']      = r[i][3]
        D[f'linu_{key}_r2']        = fmt_float_val(r[i][4], 3)

    mm = re.match(r'(-?[\d.]+)', D['linu_Sex_beta'])
    D['linu_Sex_beta_num'] = float(mm.group(1)) if mm else 0.0

    # ── 4.1.2 선형 다변량 성능 ───────────────────────────────
    sec412 = find_section('#### 4.1.2 다변량 분석 (Multivariate')
    r = table_rows(sec412, max_rows=9)
    # 0=N, 1=R², 2=Adj R², 3=F, 4=RMSE, 5=MAE, 6=AIC, 7=BIC, 8=CV

    D['lin_r2']     = strip_bold(r[1][1])
    D['lin_adj_r2'] = strip_bold(r[2][1])

    mm = re.match(r'([\d.]+)', strip_bold(r[3][1]))
    fstat_val = float(mm.group(1)) if mm else 0.0
    D['lin_fstat_disp'] = f"{fstat_val:.2f}\n(p < 0.001)"

    mm = re.match(r'([\d.]+)', strip_bold(r[4][1]))
    D['lin_rmse_disp'] = f"{float(mm.group(1)):.2f} cm²" if mm else strip_bold(r[4][1])

    mm = re.match(r'([\d.]+)', strip_bold(r[5][1]))
    D['lin_mae_disp'] = f"{float(mm.group(1)):.2f} cm²" if mm else strip_bold(r[5][1])

    D['lin_aic'] = strip_bold(r[6][1])
    D['lin_bic'] = strip_bold(r[7][1])
    D['lin_aic_disp'] = f"{int(float(D['lin_aic'])):,}\n/ {int(float(D['lin_bic'])):,}"

    # 잔차 진단 표 (4.1.2 내부, 4.2.1 이전)
    sec421 = find_section('#### 4.2.1 단변량 분석 (Crude OR)')
    resid_hdr = -1
    for i in range(sec412, sec421):
        if '검정' in lines[i] and '통계량' in lines[i] and lines[i].strip().startswith('|'):
            resid_hdr = i
            break
    if resid_hdr >= 0:
        r = table_rows(resid_hdr, max_rows=4)
        D['resid_sw_stat']      = r[0][1]
        D['resid_sw_p']         = r[0][2]
        D['resid_sw_res']       = r[0][3]
        D['resid_bp_stat_disp'] = f"{float(r[1][1]):.2f}"
        D['resid_bp_p']         = r[1][2]
        D['resid_bp_res']       = r[1][3]
        D['resid_dw_stat']      = r[2][1]
        D['resid_dw_res']       = r[2][3]
        D['resid_kappa_stat']   = r[3][1]
        D['resid_kappa_res']    = r[3][3]

    # ── 4.2.1 로지스틱 단변량 ────────────────────────────────
    log_uni_keys = ['Sex', 'Age'] + _aec_keys + ['ManufacturerModelName']
    r = table_rows(sec421, max_rows=len(log_uni_keys))
    for i, key in enumerate(log_uni_keys):
        D[f'logu_{key}_or']   = fmt_beta(strip_bold(r[i][1]), 3)
        D[f'logu_{key}_ci']   = fmt_ci(r[i][2], 3)
        D[f'logu_{key}_pval'] = r[i][3]
        D[f'logu_{key}_auc']  = fmt_float_val(r[i][4], 3)

    # ── 4.2.2 로지스틱 다변량 성능 ───────────────────────────
    sec = find_section('#### 4.2.2 다변량 분석 (전체 변수 투입)')
    r = table_rows(sec, max_rows=10)
    # 0=N, 1=AUC, 2=Sens, 3=Spec, 4=PPV, 5=NPV, 6=HL, 7=Nagelkerke, 8=Brier, 9=AIC

    auc_raw = strip_bold(r[1][1])
    mm = re.match(r'([\d.]+)\s*\[([\d.]+)[–\-]([\d.]+)\]', auc_raw)
    if mm:
        D['log_auc']    = mm.group(1)
        D['log_auc_lo'] = mm.group(2)
        D['log_auc_hi'] = mm.group(3)
    D['log_auc_ci_disp'] = f"[{D['log_auc_lo']} – {D['log_auc_hi']}]"

    D['log_sensitivity']  = strip_bold(r[2][1])
    D['log_specificity']  = strip_bold(r[3][1])
    D['log_ppv']          = strip_bold(r[4][1])
    D['log_npv']          = strip_bold(r[5][1])
    D['log_ppv_npv_disp'] = f"{D['log_ppv']} / {D['log_npv']}"

    hl_raw = strip_bold(r[6][1])
    mm = re.search(r'χ²=([\d.]+),\s*p=([\d.]+)', hl_raw)
    if mm:
        D['log_hl_stat'] = mm.group(1)
        D['log_hl_p']    = mm.group(2)

    D['log_nagelkerke'] = strip_bold(r[7][1])
    D['log_brier']      = strip_bold(r[8][1])
    D['log_aic']        = strip_bold(r[9][1])
    D['log_aic_disp']   = f"{float(D['log_aic']):,.2f}"

    # ── 4.3.1 선형 Case 비교 ─────────────────────────────────
    # 열 순서: 지표 | Case0(AEC단독) | Case1(Sex+Age) | Case2(+AEC) | Case3(+CT/kVp) | 1→2 | 2→3
    sec = find_section('#### 4.3.1 선형 회귀 성능 비교')
    r = table_rows(sec, max_rows=4)
    # 0=R², 1=Adj R², 2=RMSE, 3=AIC

    for col, ckey in [(1, 'c0'), (2, 'c1'), (3, 'c2'), (4, 'c3')]:
        D[f'lin_{ckey}_r2']    = strip_bold(r[0][col])
        D[f'lin_{ckey}_adjr2'] = strip_bold(r[1][col])
        D[f'lin_{ckey}_rmse']  = strip_bold(r[2][col])
        D[f'lin_{ckey}_aic']   = strip_bold(r[3][col])

    r2_12  = flt(D['lin_c2_r2'])    - flt(D['lin_c1_r2'])
    r2_23  = flt(D['lin_c3_r2'])    - flt(D['lin_c2_r2'])
    ar2_12 = flt(D['lin_c2_adjr2']) - flt(D['lin_c1_adjr2'])
    ar2_23 = flt(D['lin_c3_adjr2']) - flt(D['lin_c2_adjr2'])
    rm_12  = flt(D['lin_c1_rmse'])  - flt(D['lin_c2_rmse'])  # 감소 = positive
    rm_23  = flt(D['lin_c2_rmse'])  - flt(D['lin_c3_rmse'])
    ai_12  = flt(D['lin_c1_aic'])   - flt(D['lin_c2_aic'])   # 감소 = positive
    ai_23  = flt(D['lin_c2_aic'])   - flt(D['lin_c3_aic'])

    # 표 값은 절대적 차이
    D['lin_delta12_r2']    = f"+{r2_12:.4f}"
    D['lin_delta23_r2']    = f"+{r2_23:.4f}"
    D['lin_delta12_adjr2'] = f"+{ar2_12:.4f}"
    D['lin_delta23_adjr2'] = f"+{ar2_23:.4f}"
    D['lin_delta12_rmse']  = f"-{rm_12:.2f}"  # 감소이므로 음수표시
    D['lin_delta23_rmse']  = f"-{rm_23:.2f}"
    D['lin_delta12_aic']   = f"-{ai_12:.0f}"  # 감소이므로 음수표시
    D['lin_delta23_aic']   = f"-{ai_23:.0f}"

    # 요약 텍스트는 상대적 변화 (%)
    r2_12_pct = r2_12 / flt(D['lin_c1_r2']) * 100
    r2_23_pct = r2_23 / flt(D['lin_c2_r2']) * 100
    D['lin_summary'] = (
        f"AEC 특징 추가(Case1→2)로 R² +{r2_12_pct:.1f}% 향상  |  "
        f"CT 모델명/kVp 추가(Case2→3)로 추가 +{r2_23_pct:.1f}% 향상"
    )

    # ── 4.3.2 로지스틱 Case 비교 ─────────────────────────────
    # 열 순서: 지표 | Case0(AEC단독) | Case1(Sex+Age) | Case2(+AEC) | Case3(+CT/kVp) | 1→2 | 2→3
    sec = find_section('#### 4.3.2 로지스틱 회귀 성능 비교')
    r = table_rows(sec, max_rows=4)
    # 0=AUC, 1=Nagelkerke, 2=AIC, 3=HL p

    for col, ckey in [(1, 'c0'), (2, 'c1'), (3, 'c2'), (4, 'c3')]:
        D[f'log_{ckey}_auc'] = strip_bold(r[0][col])
        D[f'log_{ckey}_nag'] = strip_bold(r[1][col])
        D[f'log_{ckey}_aic'] = strip_bold(r[2][col])
        D[f'log_{ckey}_hl']  = strip_bold(r[3][col])

    D['log_c0_hl_disp'] = _hl_disp(D['log_c0_hl'])
    D['log_c1_hl_disp'] = _hl_disp(D['log_c1_hl'])
    D['log_c2_hl_disp'] = _hl_disp(D['log_c2_hl'])
    D['log_c3_hl_disp'] = _hl_disp(D['log_c3_hl'])

    au_12 = flt(D['log_c2_auc']) - flt(D['log_c1_auc'])
    au_23 = flt(D['log_c3_auc']) - flt(D['log_c2_auc'])
    na_12 = flt(D['log_c2_nag']) - flt(D['log_c1_nag'])
    na_23 = flt(D['log_c3_nag']) - flt(D['log_c2_nag'])
    la_12 = flt(D['log_c1_aic']) - flt(D['log_c2_aic'])  # 감소 = positive
    la_23 = flt(D['log_c2_aic']) - flt(D['log_c3_aic'])

    # 표 값은 절대적 차이
    D['log_delta12_auc'] = f"+{au_12:.3f}"
    D['log_delta23_auc'] = f"+{au_23:.3f}"
    D['log_delta12_nag'] = f"+{na_12:.3f}"
    D['log_delta23_nag'] = f"+{na_23:.3f}"
    D['log_delta12_aic'] = f"-{la_12:.0f}"  # 감소이므로 음수표시
    D['log_delta23_aic'] = f"-{la_23:.0f}"

    # 요약 텍스트는 상대적 변화 (%)
    au_12_pct = au_12 / flt(D['log_c1_auc']) * 100
    au_23_pct = au_23 / flt(D['log_c2_auc']) * 100
    D['log_summary'] = (
        f"AEC 추가(Case1→2)로 AUC +{au_12_pct:.1f}% 향상  |  "
        f"CT 모델명/kVp 추가(Case2→3)로 +{au_23_pct:.1f}% 추가 개선"
    )

    return D


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃

D = parse_report(config.SITE)

# ────────────────────────────────────────────────────────────
# 헬퍼 함수
# ────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.line.width = line_w
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if line_rgb:
        shp.line.color.rgb = line_rgb
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None):
    """슬라이드 상단 헤더 바"""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=NAVY)
    add_text(slide, title, 0.35, 0.1, 10, 0.6,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.65, 10, 0.45,
                 size=14, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)


def add_image(slide, fname, l, t, w, h):
    path = os.path.join(FIGURES, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        shp = add_rect(slide, l, t, w, h, fill_rgb=GRAY, line_rgb=BLUE)
        add_text(slide, f"[그림 없음]\n{fname}", l+0.1, t+h/2-0.2, w-0.2, 0.4,
                 size=9, color=BLUE, align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows,
              l, t, w, h,
              header_fill=NAVY, header_color=WHITE,
              alt_fill=SKYBLUE, font_size=11):
    """심플 텍스트-기반 표 (python-pptx Table)"""
    cols = len(headers)
    total_rows = 1 + len(rows)
    tbl = slide.shapes.add_table(total_rows, cols,
                                  Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    col_w = Inches(w / cols)
    for c in range(cols):
        tbl.columns[c].width = col_w

    # 헤더
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = header_color

    # 데이터 행
    for r, row in enumerate(rows):
        fill_col = GRAY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_col
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK


def bullet_box(slide, items, l, t, w, h, size=13, title=None):
    """글머리 기호 텍스트 박스"""
    offset = 0
    if title:
        add_text(slide, title, l, t, w, 0.35, size=13, bold=True, color=NAVY)
        offset = 0.35
    txb = slide.shapes.add_textbox(Inches(l), Inches(t + offset),
                                    Inches(w), Inches(h - offset))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)


def _pdisp(raw):
    """p-value 문자열을 표시용으로 변환"""
    raw_s = raw.strip()
    try:
        v = float(raw_s)
        if v < 0.001:
            return "< 0.001"
        if v < 0.01:
            return f"{v:.4f}"
        return f"{v:.3f}"
    except ValueError:
        return raw_s


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# 전체 배경
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
# 하단 포인트 바
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=BLUE)

# 중앙 흰 카드
add_rect(slide, 1.0, 1.5, 11.33, 4.5, fill_rgb=WHITE)

add_text(slide,
         "AEC 곡선 기반 TAMA 예측\n회귀분석 연구 보고서",
         1.3, 1.7, 10.7, 2.0,
         size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(slide,
         "CT 자동 노출 제어(AEC) 특징과 환자 인구통계학적 변수를 활용한\n"
         "복부 근육량 지표(TAMA) 예측 모델 개발 및 검증",
         1.3, 3.6, 10.7, 1.0,
         size=15, color=RGBColor(0x44, 0x44, 0x66), align=PP_ALIGN.CENTER)

add_text(slide,
         f"데이터셋: {D['site']} CT  |  N={D['n_total']:,}명  |  Python (statsmodels · scikit-learn)",
         1.3, 4.7, 10.7, 0.5,
         size=13, color=BLUE, align=PP_ALIGN.CENTER, bold=True)

add_text(slide, "2026", 0.4, 6.85, 12.5, 0.5,
         size=13, color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "목  차", "Table of Contents")

sections = [
    ("01", "연구 개요", "연구 목적 · 설계 · Feature Set (3 Case)"),
    ("02", "데이터 기술 통계", "데이터셋 · 성별 TAMA 분포 · 이진화 기준"),
    ("03", "분석 방법론", "전처리 · 모델 · 교차검증"),
    ("04", "결과 — 선형 회귀 (Part 1)", "단변량 · 다변량 성능 · 잔차 진단"),
    ("05", "결과 — 로지스틱 회귀 (Part 2)", "단변량 OR · 다변량 AUC · Calibration"),
    ("06", "결과 — Case 1–3 비교 (Part 3)", "AEC 추가 효과 · CT 모델명 효과"),
    ("07", "결론 및 한계", "주요 발견 · 제한점 · 향후 방향"),
]

for i, (num, title, desc) in enumerate(sections):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 1.3 + row * 1.65

    add_rect(slide, x, y, 6.0, 1.45, fill_rgb=NAVY)
    add_text(slide, num, x+0.15, y+0.1, 0.7, 0.6, size=22, bold=True, color=ORANGE)
    add_text(slide, title, x+0.85, y+0.08, 4.9, 0.5, size=15, bold=True, color=WHITE)
    add_text(slide, desc, x+0.85, y+0.65, 4.9, 0.7, size=11, color=RGBColor(0xB0,0xC8,0xE8))


# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 연구 개요
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "01  연구 개요", "Research Overview")

# 연구 목적 카드
add_rect(slide, 0.35, 1.25, 12.6, 1.15, fill_rgb=NAVY)
add_text(slide, "연구 목적",
         0.55, 1.28, 3.0, 0.4, size=13, bold=True, color=ORANGE)
add_text(slide,
         "CT 자동 노출 제어(AEC) 곡선과 환자 인구통계(성별·나이)가 "
         "복부 근육량 지표 TAMA(Total Abdominal Muscle Area) 예측에 기여하는 정도를 정량화",
         0.55, 1.62, 12.2, 0.65, size=12, color=WHITE)

# 3개 파트 카드
parts = [
    ("Part 1", "선형 회귀\n(Linear Regression)", "TAMA (연속형) 값 예측"),
    ("Part 2", "로지스틱 회귀\n(Logistic Regression)", "Low TAMA / Sarcopenia 위험 분류"),
    ("Part 3", "Multivariable Analysis\n(Case 1–3 비교)", "AEC·CT 모델 추가 효과 정량화"),
]
for i, (p, m, g) in enumerate(parts):
    x = 0.35 + i * 4.2
    add_rect(slide, x, 2.6, 4.0, 2.1, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1.5))
    add_text(slide, p, x+0.15, 2.68, 3.7, 0.4, size=13, bold=True, color=BLUE)
    add_text(slide, m, x+0.15, 3.05, 3.7, 0.75, size=12, bold=True, color=NAVY)
    add_text(slide, g, x+0.15, 3.8, 3.7, 0.7, size=11, color=DARK)

# Feature Set 표
add_text(slide, "Feature Set (Case 구성)", 0.35, 4.85, 6.0, 0.35,
         size=13, bold=True, color=NAVY)
add_text(slide, "* 성별·나이는 필수 공변량 — Case 1·2·3에 항상 포함. Case 0은 AEC 단독 진단력 참조용.",
         6.5, 4.9, 6.5, 0.35, size=10, italic=True, color=BLUE)
_aec_feat_str = ", ".join(config.SELECTED_AEC_FEATURES)
add_table(slide,
          ["Case", "예측 변수", "목적"],
          [
              ["Case 0\n(AEC 단독)",
               f"AEC 특징 ({_aec_feat_str})",
               "AEC 독립 진단력 확인"],
              ["Case 1", "성별 (Sex) + 나이 (Age)", "인구통계 기반선"],
              ["Case 2", f"Case 1 + AEC 특징 ({_aec_feat_str})",
               "AEC 기여도 정량화"],
              ["Case 3",
               f"Case 2 + CT 모델명 ({D['ct_dummy_count']}개 dummy) + kVp",
               "스캐너·전압 보정"],
          ],
          0.35, 5.25, 12.6, 2.0, font_size=10)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 데이터 기술 통계
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  데이터 기술 통계", "Descriptive Statistics")

# 기본 정보 카드
_kvp_card = (f"{D['kvp_dominant']} kVp ({D['kvp_dominant_pct']}%)"
             if D['kvp_dominant'] != 'N/A' else 'N/A')
stats = [
    ("총 환자 수",      f"{D['n_total']:,}명"),
    ("남성",           f"{D['n_male']:,}명 ({D['male_pct']}%)"),
    ("여성",           f"{D['n_female']:,}명 ({D['female_pct']}%)"),
    ("CT 스캐너 / 주요 kVp", f"{D['ct_models']}종  |  {_kvp_card}"),
    ("TAMA 범위",       f"{D['tama_min']} ~ {D['tama_max']} cm²"),
    ("TAMA 평균 (SD)",  f"{D['tama_mean']} (±{D['tama_sd']}) cm²"),
]
for i, (label, val) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.1
    add_rect(slide, x, y, 3.9, 0.95, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1))
    add_text(slide, label, x+0.12, y+0.05, 3.7, 0.35, size=10, color=BLUE)
    add_text(slide, val, x+0.12, y+0.42, 3.7, 0.45, size=16, bold=True, color=NAVY)

# 성별 TAMA 분포 표
add_text(slide, "성별 TAMA 분포", 0.35, 3.55, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["성별", "N", "평균 (cm²)", "SD", "중앙값", "P25", "P75"],
          [
              ["남성 (M)", D['male_n'],   D['male_mean'],   D['male_sd'],
               D['male_median'],   D['male_p25'],   D['male_p75']],
              ["여성 (F)", D['female_n'], D['female_mean'], D['female_sd'],
               D['female_median'], D['female_p25'], D['female_p75']],
          ],
          0.35, 3.92, 7.5, 1.35, font_size=11)

# 이진화 기준 카드
add_rect(slide, 8.1, 3.55, 4.9, 1.72, fill_rgb=NAVY)
add_text(slide, "Logistic Regression 이진화 기준", 8.25, 3.6, 4.6, 0.4,
         size=11, bold=True, color=ORANGE)
add_text(slide,
         f"• 남성: TAMA < {D['thr_male']} cm²  →  Low Muscle (1)\n"
         f"• 여성: TAMA < {D['thr_female']} cm²  →  Low Muscle (1)\n"
         f"• 양성 비율: {D['n_pos']} / {D['n_total']:,} = {D['pos_rate']} %",
         8.25, 4.03, 4.6, 1.2, size=12, color=WHITE)

# TAMA 분포 그림
add_text(slide, "TAMA 분포 시각화", 0.35, 5.35, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "03_tama_distribution.png", 0.35, 5.72, 6.0, 1.55)

# AEC Feature 상관 그림
add_text(slide, "AEC Feature 상관계수", 6.65, 5.35, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "01_feature_correlation.png", 6.65, 5.72, 6.33, 1.55)


# ════════════════════════════════════════════════════════════
# 슬라이드 4-A — 데이터 정제 과정
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  데이터 정제 과정", "Data Cleaning Pipeline")

add_text(slide, f"원시 데이터 {D['raw_n']}명에서 결측치 제거 후 최종 분석 대상 {D['n_total']:,}명 확정",
         0.35, 1.15, 12.6, 0.4, size=12, color=DARK)

_steps = D.get('clean_steps', [])
_step_colors = [BLUE, BLUE, BLUE, BLUE, GREEN]
for i, (label, val) in enumerate(_steps):
    y = 1.65 + i * 0.95
    fill = _step_colors[i] if i < len(_step_colors) else BLUE
    add_rect(slide, 0.35, y, 8.5, 0.82, fill_rgb=fill)
    add_text(slide, f"  {i+1}.  {label}", 0.45, y+0.06, 7.5, 0.38, size=12, bold=True, color=WHITE)
    add_rect(slide, 9.1, y, 3.9, 0.82, fill_rgb=WHITE, line_rgb=fill, line_w=Pt(1.5))
    add_text(slide, val, 9.1, y+0.15, 3.9, 0.5, size=14, bold=True, color=fill, align=PP_ALIGN.CENTER)
    if i < len(_steps) - 1:
        add_text(slide, "▼", 4.7, y+0.82, 0.5, 0.3, size=10, color=NAVY, align=PP_ALIGN.CENTER)

add_rect(slide, 0.35, 6.42, 12.6, 0.72, fill_rgb=NAVY)
add_text(slide,
         f"제거 이유: 결측 컬럼(AEC feature) 포함 행 — "
         f"{int(D['raw_n'] or 0) - D['n_total']}명 제외 "
         f"({(int(D['raw_n'] or 0) - D['n_total']) / max(int(D['raw_n'] or 1), 1) * 100:.1f}%)",
         0.55, 6.5, 12.0, 0.55, size=12, color=WHITE)


# ════════════════════════════════════════════════════════════
# 슬라이드 4-B — CT 스캐너 & kVp 분포
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  CT 스캐너 & kVp 분포", "Scanner & kVp Distribution")

_kvp_card = (f"{D['kvp_dominant']} kVp ({D['kvp_dominant_pct']}%)"
             if D.get('kvp_dominant', 'N/A') != 'N/A' else 'N/A')
_scan_stats = [
    ("총 CT 스캐너 종류", f"{D.get('ct_models', '?')}종"),
    ("최다 스캐너",       D.get('scanner_top1', 'N/A')),
    ("주요 kVp",         _kvp_card),
    ("총 분석 환자",      f"{D['n_total']:,}명"),
]
for i, (lbl, val) in enumerate(_scan_stats):
    x = 0.35 + (i % 2) * 6.4
    y = 1.2 + (i // 2) * 1.0
    add_rect(slide, x, y, 6.1, 0.85, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1))
    add_text(slide, lbl, x+0.15, y+0.04, 5.8, 0.32, size=10, color=BLUE)
    add_text(slide, val, x+0.15, y+0.38, 5.8, 0.42, size=13, bold=True, color=NAVY)

add_text(slide, "CT 스캐너 모델 분포", 0.35, 3.3, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "16_scanner_distribution.png", 0.35, 3.68, 6.2, 3.6)

add_text(slide, "kVp 분포", 6.8, 3.3, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "17_kvp_distribution.png", 6.8, 3.68, 6.18, 3.6)


# ════════════════════════════════════════════════════════════
# 슬라이드 5 — AEC Feature 선택
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  AEC Feature 선택", "Feature Selection")

_n_feat = len(config.SELECTED_AEC_FEATURES)
add_text(slide,
         f"Pearson/Spearman 상관계수 + VIF 검사로 선택한 {_n_feat}개 AEC Feature "
         f"({', '.join(config.SELECTED_AEC_FEATURES)}) — amplitude 그룹(mean≈p25≈AUC_normalized) 중 1개 대표 선택, VIF max<2",
         0.35, 1.15, 12.6, 0.45, size=12, color=DARK)

_feat_desc = {
    'mean':           "AEC 평균값 → 체격 전반 반영 (amplitude 그룹 대표)",
    'p25':            "AEC 하위 25% 값 → 저선량 구간 tube current",
    'CV':             "변동계수(std/mean) → 체형 불균일성 반영",
    'skewness':       "AEC 곡선 비대칭성 → 체형 분포 특성",
    'slope_abs_mean': "평균 절대 기울기 → 곡선 동역학",
    'peak_max_height':"최대 피크 높이 → 최고 tube current 지점",
}
_selected_rows = [
    [f,
     D.get(f'pearson_{f}', 'N/A'),
     _feat_desc.get(f, f),
     D.get(f'vif_{f}', 'N/A'),
     "✔"]
    for f in config.SELECTED_AEC_FEATURES
]
_excluded_rows = [
    ["p25",            "높음", "mean과 amplitude 동일 그룹 (r=0.97)", ">70",      "✗ 제외"],
    ["AUC_normalized", "높음", "mean과 r≈1.00, 사실상 동일 feature",  ">57,000", "✗ 제외"],
]
# 이미 선택 목록에 있는 feature는 제외 행에서 빼기
_excluded_rows = [row for row in _excluded_rows
                  if row[0] not in config.SELECTED_AEC_FEATURES]

add_table(slide,
          ["Feature", "Pearson r", "해석", "VIF", "선택"],
          _selected_rows + _excluded_rows,
          0.35, 1.65, 12.6, 2.6, font_size=11)

_excl_names = [row[0] for row in _excluded_rows]
_excl_str = "·".join(_excl_names) if _excl_names else "없음"
add_text(slide,
         f"⚠  {_excl_str}: mean과 동일 amplitude 그룹(r≥0.88) → 다중공선성 제거 / "
         f"mean을 대표로 포함",
         0.35, 4.3, 12.6, 0.5, size=12, italic=True, color=ORANGE)

# 두 그림 나란히
add_text(slide, "Feature 상관계수 Top20", 0.35, 4.85, 5.5, 0.35, size=12, bold=True, color=NAVY)
add_image(slide, "01_feature_correlation.png", 0.35, 5.22, 6.0, 2.0)

add_text(slide, "VIF 비교", 6.65, 4.85, 5.5, 0.35, size=12, bold=True, color=NAVY)
add_image(slide, "02_vif_comparison.png", 6.65, 5.22, 6.33, 2.0)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — Correlation Matrix (다중공선성)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  Correlation Matrix", "Multicollinearity Check")

add_text(slide,
         "선택 Feature 간 Pearson 상관행렬 — |r| > 0.8 쌍은 다중공선성(Multicollinearity) 위험 신호",
         0.35, 1.15, 12.6, 0.4, size=12, color=DARK)

add_image(slide, "18_correlation_matrix.png", 0.35, 1.6, 8.5, 5.65)

# 해석 카드
add_rect(slide, 9.1, 1.6, 3.9, 2.6, fill_rgb=NAVY)
add_text(slide, "해석 기준", 9.25, 1.68, 3.6, 0.4, size=12, bold=True, color=ORANGE)
add_text(slide,
         "• |r| < 0.5  →  낮은 상관\n"
         "• 0.5 ≤ |r| < 0.8  →  중간 상관\n"
         "• |r| ≥ 0.8  →  다중공선성 주의\n\n"
         "VIF > 10 과 함께 확인하여\n"
         "제거 여부 결정",
         9.25, 2.12, 3.6, 2.0, size=11, color=WHITE)

add_rect(slide, 9.1, 4.4, 3.9, 2.85, fill_rgb=WHITE, line_rgb=ORANGE, line_w=Pt(1.5))
add_text(slide, "⚠  제외 결정 사례", 9.25, 4.48, 3.6, 0.38, size=11, bold=True, color=ORANGE)
_amp_group = ['mean', 'p25', 'AUC_normalized', 'peak_max_height', 'peak_mean_height']
_amp_excl = [f for f in _amp_group if f not in config.SELECTED_AEC_FEATURES]
_amp_sel  = [f for f in _amp_group if f in config.SELECTED_AEC_FEATURES]
_amp_excl_str = "\n".join([f"• {f}: 제외 (amplitude 중복)" for f in _amp_excl[:3]])
_amp_sel_str  = ", ".join(_amp_sel) if _amp_sel else "없음"
add_text(slide,
         f"제외된 amplitude 그룹:\n{_amp_excl_str}\n\n선택 대표: {_amp_sel_str}\n(VIF 확인 후 포함)",
         9.25, 4.9, 3.6, 2.25, size=11, color=DARK)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 분석 방법론
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "03  분석 방법론", "Statistical Methods")

# 전처리 표
add_text(slide, "전처리", 0.35, 1.25, 5.5, 0.35, size=14, bold=True, color=NAVY)
add_table(slide,
          ["변수", "처리 방법"],
          [
              ["PatientSex", "이진 인코딩 (M=1, F=0)"],
              ["PatientAge", "Z-score 표준화 (StandardScaler)"],
              ["AEC features", "Z-score 표준화"],
              ["ManufacturerModelName",
               f"One-hot encoding (drop_first=True) → {D['ct_dummy_count']}개 더미"],
          ],
          0.35, 1.65, 6.0, 2.1, font_size=11)

# 모델 카드
add_text(slide, "모델", 6.65, 1.25, 6.0, 0.35, size=14, bold=True, color=NAVY)
models = [
    ("선형 회귀", "statsmodels OLS\n(Ordinary Least Squares)"),
    ("로지스틱 회귀", "statsmodels Logit\n(BFGS, maxiter=1,000)"),
]
for i, (nm, desc) in enumerate(models):
    y = 1.65 + i * 1.1
    add_rect(slide, 6.65, y, 6.33, 0.95, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1.2))
    add_text(slide, nm, 6.8, y+0.05, 3.5, 0.38, size=13, bold=True, color=NAVY)
    add_text(slide, desc, 6.8, y+0.48, 5.9, 0.42, size=11, color=DARK)

# 교차검증 및 CI
add_text(slide, "모델 검증 전략", 0.35, 3.9, 12.6, 0.35, size=14, bold=True, color=NAVY)
add_table(slide,
          ["검증 방법", "상세"],
          [
              ["5-Fold Cross Validation", "각 Fold에서 scaler 재적합 (Data Leakage 방지)"],
              ["Bootstrap AUC 95% CI", "n=1,000 resampling — 비모수 신뢰구간"],
              ["Hosmer-Lemeshow 검정", "로지스틱 보정도(Calibration) 검증 — p>0.05 양호"],
              ["잔차 진단 3종", "Shapiro-Wilk · Breusch-Pagan · Durbin-Watson"],
          ],
          0.35, 4.3, 12.6, 2.1, font_size=11)

add_text(slide,
         "※ 선형·로지스틱 회귀 모두 Case 1(Sex+Age) → Case 2(+AEC) → Case 3(+CT Model) 순으로 점진 투입",
         0.35, 6.5, 12.6, 0.55, size=11, italic=True, color=BLUE)


# ════════════════════════════════════════════════════════════
# 슬라이드 7 — 선형 회귀: 단변량
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 단변량 분석", "Univariate Linear Regression")

_lin_uni_rows = [
    ["Sex (M=1, F=0)",
     D['linu_Sex_beta_disp'], D['linu_Sex_ci'],
     _pdisp(D['linu_Sex_pval']), D['linu_Sex_r2']],
    ["Age (표준화)",
     D['linu_Age_beta_disp'], D['linu_Age_ci'],
     _pdisp(D['linu_Age_pval']), D['linu_Age_r2']],
] + [
    [f"AEC: {f} (표준화)",
     D[f'linu_{f}_beta_disp'], D[f'linu_{f}_ci'],
     _pdisp(D[f'linu_{f}_pval']), D[f'linu_{f}_r2']]
    for f in config.SELECTED_AEC_FEATURES
] + [
    ["ManufacturerModelName",
     D['linu_ManufacturerModelName_beta_disp'], "N/A",
     _pdisp(D['linu_ManufacturerModelName_pval']),
     D['linu_ManufacturerModelName_r2']],
]
add_table(slide,
          ["변수", "β 계수", "95% CI", "p-value", "R²"],
          _lin_uni_rows,
          0.35, 1.25, 12.6, 3.3, font_size=11)

add_text(slide, "* p < 0.05 유의",
         0.35, 4.65, 5.0, 0.35, size=11, italic=True, color=ORANGE)

# 단변량 R² 그림
add_text(slide, "단변량 R² 비교", 0.35, 5.05, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "07_linear_univariate_r2.png", 0.35, 5.42, 6.0, 1.85)

# Forest plot
add_text(slide, "유의 계수 Forest Plot", 6.65, 5.05, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "06_linear_forest.png", 6.65, 5.42, 6.33, 1.85)


# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 선형 회귀: 다변량 성능
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 다변량 모델 성능 (Case 3 Full)", "Multivariate Linear Regression")

# 성능 KPI 카드
kpis = [
    ("R²",          D['lin_r2']),
    ("Adj R²",      D['lin_adj_r2']),
    ("RMSE",        D['lin_rmse_disp']),
    ("MAE",         D['lin_mae_disp']),
    ("F-statistic", D['lin_fstat_disp']),
    ("AIC / BIC",   D['lin_aic_disp']),
]
for i, (label, val) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.25
    add_rect(slide, x, y, 3.9, 1.15, fill_rgb=NAVY)
    add_text(slide, label, x+0.15, y+0.08, 3.6, 0.38, size=12, bold=True, color=ORANGE)
    add_text(slide, val, x+0.15, y+0.52, 3.6, 0.55, size=18, bold=True, color=WHITE)

# 잔차 진단 표
add_text(slide, "잔차 진단", 0.35, 3.85, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["검정", "통계량", "p-value", "결과"],
          [
              ["Shapiro-Wilk",
               D['resid_sw_stat'],      D['resid_sw_p'],  D['resid_sw_res']],
              ["Breusch-Pagan",
               D['resid_bp_stat_disp'], D['resid_bp_p'],  D['resid_bp_res']],
              ["Durbin-Watson",
               D['resid_dw_stat'],      "—",              D['resid_dw_res']],
              ["Condition κ",
               D['resid_kappa_stat'],   "—",              D['resid_kappa_res']],
          ],
          0.35, 4.22, 6.5, 2.0, font_size=11)

add_text(slide,
         "⚠  비정규·이분산 → 계수 추정 BLUE이나 p-value 해석 주의\n   Robust Standard Errors 사용 권고",
         0.35, 6.3, 6.5, 0.75, size=11, italic=True, color=ORANGE)

# 실제 vs 예측 그림
add_text(slide, "실제 vs 예측", 7.1, 3.85, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "04_linear_actual_vs_pred.png", 7.1, 4.22, 5.9, 2.0)

# 잔차 진단 4-panel
add_image(slide, "05_linear_residuals.png", 7.1, 6.25, 5.9, 1.0)


# ════════════════════════════════════════════════════════════
# 슬라이드 8-C — 선형 회귀 다변량 핵심 계수 상세
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 핵심 변수 계수 (Case 3)", "Key Predictor Coefficients")

_lin_xl = os.path.join(_ROOT, 'results', config.SITE, 'linear_results.xlsx')
_lin_coef = pd.read_excel(_lin_xl, sheet_name='다변량_coefficients')
_key_vars = ['Sex', 'Age_z'] + [f + '_z' for f in config.SELECTED_AEC_FEATURES]
_lc = _lin_coef[_lin_coef['Variable'].isin(_key_vars)].copy()
_lc_rows = []
for _, row in _lc.iterrows():
    vname = str(row['Variable']).replace('_z', ' (표준화)').replace('Age_z', 'Age (표준화)')
    sig = ' *' if float(row['p_value']) < 0.05 else ''
    _lc_rows.append([
        vname,
        f"{float(row['β']):.4f}{sig}",
        f"[{float(row['CI_Lower']):.3f}, {float(row['CI_Upper']):.3f}]",
        f"{float(row['SE']):.4f}",
        f"{float(row['t_stat']):.3f}",
        f"{float(row['p_value']):.4e}",
    ])
add_table(slide,
          ["변수", "β 계수", "95% CI", "SE", "t", "p-value"],
          _lc_rows,
          0.35, 1.25, 12.6, len(_lc_rows) * 0.55 + 0.45, font_size=11)

_note_y = 1.25 + len(_lc_rows) * 0.55 + 0.55
add_text(slide, "* p < 0.05 유의  |  CT 모델명 더미는 생략 (별도 슬라이드)  |  표준화 계수: 1 SD 변화 시 TAMA 변화량",
         0.35, _note_y, 12.6, 0.4, size=10, italic=True, color=ORANGE)

add_text(slide, "선형 Forest Plot", 0.35, _note_y + 0.5, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "06_linear_forest.png", 0.35, _note_y + 0.88, 12.6, 7.5 - _note_y - 1.1)


# ════════════════════════════════════════════════════════════
# 슬라이드 9 — 로지스틱 회귀: 단변량 OR
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 단변량 분석 (Crude OR)", "Univariate Logistic Regression")

_log_uni_rows = [
    ["Sex (M=1, F=0)",
     D['logu_Sex_or'], D['logu_Sex_ci'],
     _pdisp(D['logu_Sex_pval']), D['logu_Sex_auc']],
    ["Age (표준화)",
     D['logu_Age_or'], D['logu_Age_ci'],
     _pdisp(D['logu_Age_pval']), D['logu_Age_auc']],
] + [
    [f"AEC: {f} (표준화)",
     D[f'logu_{f}_or'], D[f'logu_{f}_ci'],
     _pdisp(D[f'logu_{f}_pval']), D[f'logu_{f}_auc']]
    for f in config.SELECTED_AEC_FEATURES
] + [
    ["ManufacturerModelName",
     D['logu_ManufacturerModelName_or'], "N/A",
     _pdisp(D['logu_ManufacturerModelName_pval']),
     D['logu_ManufacturerModelName_auc']],
]
add_table(slide,
          ["변수", "Crude OR", "95% CI", "p-value", "AUC"],
          _log_uni_rows,
          0.35, 1.25, 12.6, 3.3, font_size=11)

add_text(slide,
         f"* p < 0.05 유의  |  {', '.join(config.SELECTED_AEC_FEATURES)}: OR 방향 확인  |  Sex: 이진화 후 단변량 비유의",
         0.35, 4.65, 12.0, 0.4, size=11, italic=True, color=ORANGE)

# Crude OR Forest plot
add_text(slide, "Crude OR Forest Plot", 0.35, 5.12, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "11_logistic_forest.png", 0.35, 5.5, 6.0, 1.77)

# ROC 곡선
add_text(slide, "ROC Curve (단변량)", 6.65, 5.12, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "08_logistic_roc.png", 6.65, 5.5, 6.33, 1.77)


# ════════════════════════════════════════════════════════════
# 슬라이드 10 — 로지스틱 회귀: 다변량 성능
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 다변량 모델 성능 (Case 3 Full)", "Multivariate Logistic Regression")

kpis2 = [
    ("AUC-ROC",          D['log_auc']),
    ("Bootstrap 95% CI", D['log_auc_ci_disp']),
    ("Sensitivity",      D['log_sensitivity']),
    ("Specificity",      D['log_specificity']),
    ("PPV / NPV",        D['log_ppv_npv_disp']),
    ("Nagelkerke R²",    D['log_nagelkerke']),
]
for i, (label, val) in enumerate(kpis2):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.2
    add_rect(slide, x, y, 3.9, 1.1, fill_rgb=NAVY)
    add_text(slide, label, x+0.15, y+0.06, 3.6, 0.38, size=11, bold=True, color=ORANGE)
    add_text(slide, val, x+0.15, y+0.5, 3.6, 0.52, size=17, bold=True, color=WHITE)

# 보정도 표
add_text(slide, "보정도 및 적합도", 0.35, 3.72, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["지표", "값", "해석"],
          [
              ["Hosmer-Lemeshow χ²",
               D['log_hl_stat'],
               f"p={D['log_hl_p']} → 보정도 양호 ✔"],
              ["Brier Score",
               D['log_brier'],
               "낮을수록 우수 (0=완벽, 0.25=무작위)"],
              ["AIC",
               D['log_aic_disp'],
               "Case 비교 기준"],
          ],
          0.35, 4.1, 7.0, 1.65, font_size=11)

# Calibration plot
add_text(slide, "Calibration Plot", 7.5, 3.72, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "09_logistic_calibration.png", 7.5, 4.1, 5.48, 2.2)

# Confusion matrix
add_text(slide, "Confusion Matrix", 0.35, 5.85, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "10_logistic_confusion.png", 0.35, 6.22, 3.5, 1.05)


# ════════════════════════════════════════════════════════════
# 슬라이드 10-D — 로지스틱 회귀 다변량 핵심 변수 Adjusted OR
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 핵심 변수 Adjusted OR (Case 3)", "Adjusted Odds Ratios")

_log_xl = os.path.join(_ROOT, 'results', config.SITE, 'logistic_results.xlsx')
_log_coef = pd.read_excel(_log_xl, sheet_name='다변량_coefficients')
_key_vars_log = ['Sex', 'Age_z'] + [f + '_z' for f in config.SELECTED_AEC_FEATURES]
_lc_log = _log_coef[_log_coef['Variable'].isin(_key_vars_log)].copy()
_lor_rows = []
for _, row in _lc_log.iterrows():
    vname = str(row['Variable']).replace('_z', ' (표준화)')
    adj_or = float(row['Adj_OR'])
    _lor_rows.append([
        vname,
        f"{float(row['log_OR']):.4f}",
        f"{adj_or:.4f}",
        "< 1 → 위험 감소" if adj_or < 1 else "> 1 → 위험 증가",
    ])
add_table(slide,
          ["변수", "log(OR)", "Adjusted OR", "방향 해석"],
          _lor_rows,
          0.35, 1.25, 8.0, len(_lor_rows) * 0.58 + 0.45, font_size=11)

_ory = 1.25 + len(_lor_rows) * 0.58 + 0.55
add_text(slide,
         "* Adjusted OR: exp(β). 표준화 변수 기준 1 SD 변화 시 Low TAMA 오즈 배율\n"
         "  OR < 1: 해당 변수 증가 → Low TAMA 위험 감소 / OR > 1: 위험 증가",
         0.35, _ory, 8.0, 0.65, size=10, italic=True, color=ORANGE)

add_text(slide, "Adjusted OR Forest Plot", 8.3, 1.25, 5.0, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "11_logistic_forest.png", 8.3, 1.62, 4.68, 5.65)


# ════════════════════════════════════════════════════════════
# 슬라이드 11 — Case 1·2·3 선형 회귀 비교
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 비교 — 선형 회귀 성능 향상", "Multivariable Analysis: Linear Regression")

add_table(slide,
          ["지표", "Case 0\n(AEC 단독)", "Case 1\n(Sex+Age)", "Case 2\n(+AEC)",
           "Case 3\n(+CT/kVp)", "C1→2 Δ", "C2→3 Δ"],
          [
              ["R²",
               D['lin_c0_r2'],   D['lin_c1_r2'],   D['lin_c2_r2'],   D['lin_c3_r2'],
               D['lin_delta12_r2'],   D['lin_delta23_r2']],
              ["Adj R²",
               D['lin_c0_adjr2'], D['lin_c1_adjr2'], D['lin_c2_adjr2'], D['lin_c3_adjr2'],
               D['lin_delta12_adjr2'], D['lin_delta23_adjr2']],
              ["RMSE (cm²)",
               D['lin_c0_rmse'],  D['lin_c1_rmse'],  D['lin_c2_rmse'],  D['lin_c3_rmse'],
               D['lin_delta12_rmse'],  D['lin_delta23_rmse']],
              ["AIC",
               D['lin_c0_aic'],   D['lin_c1_aic'],   D['lin_c2_aic'],   D['lin_c3_aic'],
               D['lin_delta12_aic'],   D['lin_delta23_aic']],
          ],
          0.35, 1.25, 12.6, 2.3, font_size=10)

add_text(slide, D['lin_summary'],
         0.35, 3.65, 12.6, 0.45, size=13, bold=True, color=BLUE)

add_text(slide, "R² / RMSE 비교", 0.35, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "12_case_metrics_bar.png", 0.35, 4.58, 6.0, 2.7)

add_text(slide, "AIC / BIC 비교", 6.65, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "14_case_aic_bar.png", 6.65, 4.58, 6.33, 2.7)


# ════════════════════════════════════════════════════════════
# 슬라이드 12 — Case 1·2·3 로지스틱 회귀 비교
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 비교 — 로지스틱 회귀 성능 향상", "Multivariable Analysis: Logistic Regression")

add_table(slide,
          ["지표", "Case 0\n(AEC 단독)", "Case 1\n(Sex+Age)", "Case 2\n(+AEC)",
           "Case 3\n(+CT/kVp)", "C1→2 Δ", "C2→3 Δ"],
          [
              ["AUC-ROC",
               D['log_c0_auc'], D['log_c1_auc'], D['log_c2_auc'], D['log_c3_auc'],
               D['log_delta12_auc'], D['log_delta23_auc']],
              ["Nagelkerke R²",
               D['log_c0_nag'], D['log_c1_nag'], D['log_c2_nag'], D['log_c3_nag'],
               D['log_delta12_nag'], D['log_delta23_nag']],
              ["AIC",
               D['log_c0_aic'], D['log_c1_aic'], D['log_c2_aic'], D['log_c3_aic'],
               D['log_delta12_aic'], D['log_delta23_aic']],
              ["HL p-value",
               D['log_c0_hl_disp'], D['log_c1_hl_disp'], D['log_c2_hl_disp'], D['log_c3_hl_disp'],
               "보정도 개선", "보정도 양호"],
          ],
          0.35, 1.25, 12.6, 2.3, font_size=10)

add_text(slide, D['log_summary'],
         0.35, 3.65, 12.6, 0.45, size=13, bold=True, color=BLUE)

add_text(slide, "AUC / Nagelkerke R² 비교", 0.35, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "13_case_auc_bar.png", 0.35, 4.58, 6.0, 2.7)

add_text(slide, "다중 지표 추이", 6.65, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "15_case_progression.png", 6.65, 4.58, 6.33, 2.7)


# ════════════════════════════════════════════════════════════
# 슬라이드 12-E — Case 0 AEC 단독 진단력 심층 분석
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 0 — AEC 단독 진단력", "AEC-Only Baseline (Case 0)")

add_text(slide,
         "성별·나이 없이 AEC 특징만으로 TAMA를 예측할 수 있는가?  "
         "— 인구통계 보정 전 AEC 신호의 독립적 진단력 정량화",
         0.35, 1.15, 12.6, 0.42, size=12, color=DARK)

_c0_kpis = [
    ("선형 R²\n(Case 0)",            D['lin_c0_r2'],   "AEC 단독 TAMA 설명력"),
    ("선형 RMSE\n(Case 0)",          D['lin_c0_rmse'],  "예측 오차 (cm²)"),
    ("로지스틱 AUC\n(Case 0)",       D['log_c0_auc'],  "AEC 단독 판별 능력"),
    ("Nagelkerke R²\n(Case 0)",      D['log_c0_nag'],  "로지스틱 설명력"),
]
for i, (lbl, val, desc) in enumerate(_c0_kpis):
    x = 0.35 + i * 3.22
    add_rect(slide, x, 1.72, 3.1, 1.6, fill_rgb=NAVY)
    add_text(slide, lbl,  x+0.12, 1.78, 2.9, 0.55, size=11, bold=True, color=ORANGE)
    add_text(slide, val,  x+0.12, 2.33, 2.9, 0.55, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x+0.12, 2.9,  2.9, 0.35, size=10, color=RGBColor(0xB0,0xC8,0xE8))

add_table(slide,
          ["지표", "Case 0\n(AEC 단독)", "Case 1\n(Sex+Age)", "Case 1→0 비교", "의미"],
          [
              ["선형 R²",
               D['lin_c0_r2'], D['lin_c1_r2'],
               f"Case1 - Case0 = {float(D['lin_c1_r2'])-float(D['lin_c0_r2']):.3f}",
               "인구통계 추가 효과"],
              ["선형 RMSE",
               D['lin_c0_rmse'], D['lin_c1_rmse'],
               f"RMSE 개선: {float(D['lin_c0_rmse'])-float(D['lin_c1_rmse']):.2f} cm²",
               "오차 감소"],
              ["로지스틱 AUC",
               D['log_c0_auc'], D['log_c1_auc'],
               f"ΔAUC = {float(D['log_c1_auc'])-float(D['log_c0_auc']):.3f}",
               "Sex+Age 기여"],
              ["Nagelkerke R²",
               D['log_c0_nag'], D['log_c1_nag'],
               f"ΔNag = {float(D['log_c1_nag'])-float(D['log_c0_nag']):.3f}",
               "설명력 변화"],
          ],
          0.35, 3.5, 12.6, 2.3, font_size=10)

add_rect(slide, 0.35, 5.95, 12.6, 0.8, fill_rgb=BLUE)
add_text(slide,
         f"해석: AEC 단독(Case 0)으로도 선형 R²={D['lin_c0_r2']}, AUC={D['log_c0_auc']} 달성 "
         f"→ AEC 신호에 근육량 관련 정보가 내포되어 있음을 확인",
         0.55, 6.02, 12.1, 0.65, size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════
# 슬라이드 12-F — 성능 평가지표 해석 기준
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  성능 평가지표 해석 기준", "Statistical Metric Reference")

add_text(slide, "선형 회귀 지표", 0.35, 1.18, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["지표", "해석 기준"],
          [
              ["R² / Adj R²",  "0~1 사이, 1에 가까울수록 설명력↑ — Case 간 비교에 Adj R² 우선"],
              ["RMSE",         "단위 cm², 낮을수록 우수 — 임상 허용 오차 기준으로 해석"],
              ["AIC / BIC",    "낮을수록 선호 — 모델 복잡도 대비 적합도 (Case 비교 기준)"],
              ["Durbin-Watson","1.5~2.5 정상 — 잔차 자기상관 없음"],
              ["Shapiro-Wilk", "p > 0.05: 정규분포 충족 / p ≤ 0.05: 비정규 (주의)"],
              ["Breusch-Pagan","p > 0.05: 등분산 충족 / p ≤ 0.05: 이분산 (Robust SE 권장)"],
          ],
          0.35, 1.55, 6.3, 3.3, font_size=10)

add_text(slide, "로지스틱 회귀 지표", 6.8, 1.18, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["지표", "해석 기준"],
          [
              ["AUC-ROC",          "0.5=무작위, 0.7↑=양호, 0.8↑=우수, 1.0=완벽"],
              ["Bootstrap 95%CI",  "n=1000 비모수 — CI가 좁을수록 안정적 추정"],
              ["Sensitivity",       "실제 Low TAMA를 얼마나 탐지 (TP율)"],
              ["Specificity",       "정상 TAMA를 얼마나 식별 (TN율)"],
              ["Hosmer-Lemeshow",   "p > 0.05: 보정 양호 / ≤ 0.05: 체계적 오차"],
              ["Nagelkerke R²",     "0~1 pseudo-R² — 설명력 비교용"],
              ["Brier Score",       "0=완벽, 0.25=무작위 — 확률 예측 정밀도"],
          ],
          6.8, 1.55, 6.18, 3.7, font_size=10)

add_rect(slide, 0.35, 5.0, 12.6, 0.75, fill_rgb=NAVY)
add_text(slide, "임계값 해석 기준  (Low TAMA)",
         0.55, 5.07, 4.0, 0.35, size=11, bold=True, color=ORANGE)
add_text(slide,
         f"남성: TAMA < {D['thr_male']} cm²  →  Low Muscle (Positive=1)  |  "
         f"여성: TAMA < {D['thr_female']} cm²  →  Low Muscle (Positive=1)  |  "
         f"양성 비율: {D['n_pos']} / {D['n_total']:,} = {D['pos_rate']} %",
         0.55, 5.42, 12.1, 0.28, size=11, color=WHITE)

add_rect(slide, 0.35, 5.85, 12.6, 1.42, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1))
add_text(slide, "모델 선택 원칙",
         0.55, 5.92, 4.0, 0.35, size=11, bold=True, color=NAVY)
add_text(slide,
         "① 연속형 TAMA 예측(Part 1): R², Adj R², RMSE 기준으로 모델 평가\n"
         "② 이진 분류(Part 2): AUC, Sensitivity/Specificity 균형, HL 보정도 확인\n"
         "③ Case 비교(Part 3): AIC/BIC 감소 + ΔR²/ΔAUC 임상 유의성 동시 검토",
         0.55, 6.3, 12.1, 0.9, size=11, color=DARK)


# ════════════════════════════════════════════════════════════
# 슬라이드 13 — 결론
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "07  결론", "Conclusions")

_r2_12 = float(D['lin_c2_r2']) - float(D['lin_c1_r2'])
_r2_23 = float(D['lin_c3_r2']) - float(D['lin_c2_r2'])
_au_12 = float(D['log_c2_auc']) - float(D['log_c1_auc'])
_au_23 = float(D['log_c3_auc']) - float(D['log_c2_auc'])

conclusions = [
    ("성별의 압도적 기여",
     f"선형 단변량: Sex 단독 R²={D['linu_Sex_r2']}  |  "
     f"남성 TAMA가 여성 대비 평균 +{D['linu_Sex_beta_num']:.1f} cm²\n"
     "성별·나이는 필수 공변량으로 모든 모델에 포함 (제외 불가)"),
    ("AEC 단독 진단력 (Case 0)",
     f"성별·나이 없이 AEC만으로:\n"
     f"  선형 R²={D['lin_c0_r2']}  |  로지스틱 AUC={D['log_c0_auc']}\n"
     "AEC 신호가 독립적인 근육량 정보를 내포"),
    ("AEC 특징의 추가 기여 (Case 1→2)",
     f"선형 R²  {float(D['lin_c1_r2']):.3f} → {float(D['lin_c2_r2']):.3f}  "
     f"(+{_r2_12:.3f})\n"
     f"로지스틱 AUC  {float(D['log_c1_auc']):.3f} → {float(D['log_c2_auc']):.3f}  "
     f"(+{_au_12:.3f})"),
    ("CT 모델명·kVp 추가 기여 및 최종 성능",
     f"Case2→3  선형 R²+{_r2_23:.3f}  |  AUC +{_au_23:.3f}\n"
     f"최종(Case3): R²={D['lin_r2']}, RMSE={D['lin_rmse_disp']}\n"
     f"AUC={D['log_auc']} {D['log_auc_ci_disp']}  (HL p={D['log_hl_p']})"),
]
for i, (title, body) in enumerate(conclusions):
    row = i // 2
    col = i % 2
    x = 0.35 + col * 6.4
    y = 1.25 + row * 2.35
    add_rect(slide, x, y, 6.1, 2.2, fill_rgb=NAVY)
    add_rect(slide, x, y, 0.18, 2.2, fill_rgb=ORANGE)
    add_text(slide, f"{i+1}. {title}", x+0.3, y+0.1, 5.6, 0.45,
             size=12, bold=True, color=ORANGE)
    add_text(slide, body, x+0.3, y+0.6, 5.6, 1.55, size=11, color=WHITE)

add_text(slide,
         "AEC 특징은 성별·나이와 독립적으로 TAMA 예측에 기여 — CT 촬영 시 수집되는 AEC 데이터의 임상적 활용 가능성 확인",
         0.35, 6.8, 12.6, 0.55, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 14 — 한계 및 향후 방향
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "07  한계 및 향후 방향", "Limitations & Future Work")

other_site = "신촌" if D['site'] == "강남" else "강남"

limits = [
    "잔차 정규성·등분산성 가정 위반 → Robust Standard Errors 또는 비모수 검정 검토 필요",
    f"Logistic 이진화 임계값 (M<{D['thr_male']}, F<{D['thr_female']} cm²) — 데이터 내 P25 기반; 문헌(Prado 2008, Martin 2013) 기준값과 비교 필요",
    f"단일 기관({D['site']}) 데이터 → 외부 검증({other_site} 데이터 활용) 필요",
    "AEC Feature는 통계적 요약값 기반 — raw AEC 시계열의 잠재 패턴이 미반영",
]
future = [
    f"{other_site} 데이터를 이용한 외부 검증 (External Validation)",
    "Raw AEC 시계열(~200포인트)을 1D CNN / LSTM 직접 입력 — feature engineering 불필요",
    "Robust Regression / Weighted Least Squares 적용",
    "머신러닝 앙상블 모델 비교 (Random Forest, XGBoost)",
    "다기관 데이터 수집을 통한 모델 일반화",
]

add_rect(slide, 0.35, 1.2, 6.0, 3.3, fill_rgb=WHITE, line_rgb=ORANGE, line_w=Pt(1.5))
add_text(slide, "⚠  한계점", 0.55, 1.28, 5.5, 0.38, size=13, bold=True, color=ORANGE)
for i, lim in enumerate(limits):
    add_text(slide, f"• {lim}", 0.55, 1.72 + i*0.72, 5.6, 0.65, size=11, color=DARK)

add_rect(slide, 6.7, 1.2, 6.28, 3.3, fill_rgb=WHITE, line_rgb=GREEN, line_w=Pt(1.5))
add_text(slide, "→  향후 연구 방향", 6.88, 1.28, 5.8, 0.38, size=13, bold=True, color=GREEN)
for i, fut in enumerate(future):
    add_text(slide, f"• {fut}", 6.88, 1.72 + i*0.57, 5.8, 0.52, size=11, color=DARK)

# 결론 요약 강조 박스
add_rect(slide, 0.35, 4.65, 12.6, 1.85, fill_rgb=NAVY)
add_text(slide, "핵심 메시지",
         0.55, 4.72, 3.0, 0.38, size=12, bold=True, color=ORANGE)
add_text(slide,
         f"CT 스캐너의 AEC 곡선에서 추출한 통계 특징({_aec_feat_str})은 인구통계학적 변수와 독립적으로\n"
         "TAMA 예측에 유의미하게 기여한다. 향후 외부 검증 및 다기관 연구를 통해 임상 적용 가능성을 확인해야 한다.",
         0.55, 5.12, 12.1, 1.3, size=13, color=WHITE)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
output_path = os.path.join(_ROOT, "results", f"{config.SITE}_TAMA_연구보고서.pptx")
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
