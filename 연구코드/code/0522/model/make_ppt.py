"""
SMI 이진 분류 실험 결과 PPT 생성 스크립트.
실행: python make_ppt.py
출력: 연구코드/results/0515/results_summary.pptx
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from copy import deepcopy
import lxml.etree as etree

BASE     = Path(__file__).parent.parent.parent.parent / "results" / "0515"
TEMPLATE = Path(r'C:\Users\jhjun\Downloads\Templete(HEART LAB).pptx')
IMG      = lambda model, case, fname: str(BASE / model / case / fname)

# ── 색상 팔레트 ───────────────────────────────────────────────
C_TITLE_BG  = RGBColor(0x1F, 0x34, 0x6B)   # 진한 네이비 (템플릿 기준)
C_TITLE_FG  = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT    = RGBColor(0x2E, 0x75, 0xB6)
C_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
C_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
C_BEST_BG   = RGBColor(0xFF, 0xF2, 0xCC)
C_BEST_FG   = RGBColor(0x7F, 0x60, 0x00)
C_ROW_ALT   = RGBColor(0xF2, 0xF7, 0xFD)
C_BODY      = RGBColor(0x26, 0x26, 0x26)
C_DIVIDER   = RGBColor(0xBF, 0xBF, 0xBF)
C_SECTION   = RGBColor(0xED, 0x7D, 0x31)
C_GREEN     = RGBColor(0x37, 0x86, 0x10)

W, H = Inches(13.33), Inches(7.5)

# ── 헬퍼 ─────────────────────────────────────────────────────

def _remove_all_slides(prs):
    """템플릿에서 기존 슬라이드를 모두 제거."""
    sldIdLst = prs.slides._sldIdLst
    rIds = [
        s.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        for s in sldIdLst
    ]
    for s in list(sldIdLst):
        sldIdLst.remove(s)
    for rId in rIds:
        if rId:
            prs.part.drop_rel(rId)

def new_prs():
    prs = Presentation(str(TEMPLATE))
    _remove_all_slides(prs)
    return prs

def blank_slide(prs):
    """빈 화면 레이아웃(6) — 제목 슬라이드 전용."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def content_slide(prs):
    """제목 및 내용 레이아웃(1) — 하단 네이비 바 자동 포함."""
    layout = prs.slide_layouts[1]
    return prs.slides.add_slide(layout)

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=14, bold=False, color=C_BODY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def set_slide_title(slide, title, subtitle=""):
    """layout[1] 슬라이드의 제목 placeholder에 제목·부제목 설정."""
    ph = slide.placeholders[0]
    tf = ph.text_frame
    tf.text = title
    if tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.size  = Pt(20)
        run.font.bold  = True
        run.font.color.rgb = RGBColor(0x02, 0x36, 0x70)
    if subtitle:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = subtitle
        r.font.size  = Pt(9)
        r.font.bold  = False
        r.font.color.rgb = RGBColor(0x60, 0x60, 0x60)

def section_label(slide, text, x, y, w=Inches(2), h=Inches(0.35)):
    add_rect(slide, x, y, w, h, fill=C_ACCENT)
    add_text(slide, text, x+Inches(0.1), y, w-Inches(0.1), h,
             size=11, bold=True, color=C_TITLE_FG)

def _set_cell(cell, text, bg=None, fg=C_BODY, size=9, bold=False, align=PP_ALIGN.CENTER):
    tf = cell.text_frame
    tf.word_wrap = False
    for para in tf.paragraphs:
        para.clear()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = fg
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

def results_table(slide, x, y, w, h, rows_data, col_widths,
                  best_row=None, header_rows=1):
    """
    rows_data: list of list of str  (첫 번째가 헤더)
    best_row: 0-based row index (헤더 제외) of best case
    """
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            is_header = (ri < header_rows)
            is_best   = (not is_header) and (ri - header_rows == best_row)
            is_alt    = (not is_header) and (not is_best) and ((ri - header_rows) % 2 == 1)

            if is_header:
                bg, fg, bld = C_HEADER_BG, C_HEADER_FG, True
            elif is_best:
                bg, fg, bld = C_BEST_BG, C_BEST_FG, True
            elif is_alt:
                bg, fg, bld = C_ROW_ALT, C_BODY, False
            else:
                bg, fg, bld = RGBColor(0xFF,0xFF,0xFF), C_BODY, False

            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            _set_cell(cell, val, bg=bg, fg=fg, size=9, bold=bld, align=align)
    return tbl

def add_image(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        add_rect(slide, x, y, w, h, fill=RGBColor(0xEE,0xEE,0xEE), line=C_DIVIDER)
        add_text(slide, f"[{os.path.basename(path)}]", x, y+h//2-Inches(0.2), w, Inches(0.4),
                 size=8, color=C_DIVIDER, align=PP_ALIGN.CENTER)


# ── 슬라이드별 내용 ───────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, fill=C_TITLE_BG)
    add_rect(sl, 0, Inches(3.0), W, Inches(0.05), fill=C_ACCENT)
    add_text(sl, "SMI 근감소증 이진 분류",
             Inches(1.0), Inches(1.6), Inches(11), Inches(0.9),
             size=40, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
    add_text(sl, "실험 결과 요약",
             Inches(1.0), Inches(2.5), Inches(11), Inches(0.7),
             size=28, bold=False, color=RGBColor(0xBD,0xCC,0xE0), align=PP_ALIGN.CENTER)
    add_text(sl, "Model 1 (Clinic)  ·  Model 2 (Clinic+AEC)  ·  Model 2_2 (Unmatched)  ·  Model 3 (Clinic+Scanner+AEC)",
             Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
             size=14, color=RGBColor(0x9D,0xB4,0xD4), align=PP_ALIGN.CENTER)
    add_text(sl, "5-Fold CV  |  LR 또는 Deep Model  |  2 / 4 / 4 / 8 Scaling Cases",
             Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
             size=12, color=RGBColor(0x9D,0xB4,0xD4), align=PP_ALIGN.CENTER)


def slide_overview(prs):
    sl = content_slide(prs)
    set_slide_title(sl, "실험 설계 개요")

    boxes = [
        ("Model 1", "Clinic Only", "Age · Sex · BMI", "LR 또는 ResNet1D", "2 cases"),
        ("Model 2", "Clinic + AEC\n(Matched)", "Age · Sex · BMI\n+ AEC 256", "LR 또는 CrossAttn", "4 cases"),
        ("Model 2_2", "Clinic + AEC\n(Unmatched)", "AEC 행 shuffle\n음성 대조군", "LR 또는 CrossAttn", "4 cases"),
        ("Model 3", "Clinic + Scanner\n+ AEC", "Age · Sex · BMI\n+ kVp · Mfr · AEC", "LR 또는 CrossAttn3", "8 cases"),
    ]
    bw = Inches(2.9)
    for i, (title, subtitle, features, models, cases) in enumerate(boxes):
        bx = Inches(0.3) + i * (bw + Inches(0.2))
        by = Inches(1.1)
        add_rect(sl, bx, by, bw, Inches(4.2), fill=RGBColor(0xF5,0xF8,0xFD), line=C_ACCENT)
        add_rect(sl, bx, by, bw, Inches(0.55), fill=C_ACCENT)
        add_text(sl, title, bx+Inches(0.1), by+Inches(0.05), bw-Inches(0.2), Inches(0.45),
                 size=16, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)
        add_text(sl, subtitle, bx+Inches(0.1), by+Inches(0.65), bw-Inches(0.2), Inches(0.7),
                 size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(sl, f"입력:\n{features}", bx+Inches(0.15), by+Inches(1.4), bw-Inches(0.3), Inches(1.0),
                 size=10, color=C_BODY)
        add_text(sl, f"모델:\n{models}", bx+Inches(0.15), by+Inches(2.5), bw-Inches(0.3), Inches(0.8),
                 size=10, color=C_BODY)
        add_rect(sl, bx+Inches(0.5), by+Inches(3.5), bw-Inches(1.0), Inches(0.4), fill=C_SECTION)
        add_text(sl, cases, bx+Inches(0.5), by+Inches(3.5), bw-Inches(1.0), Inches(0.4),
                 size=11, bold=True, color=C_TITLE_FG, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.2),
             fill=RGBColor(0xF0,0xF4,0xFA), line=C_DIVIDER)
    add_text(sl, "공통 설정",
             Inches(0.5), Inches(5.65), Inches(2), Inches(0.35),
             size=10, bold=True, color=C_ACCENT)
    info = ("데이터: 강남 SMI 데이터  |  Train 1,097 / Test 275  |  Sarcopenia 10.9%  |  "
            "5-Fold Stratified CV  |  Scaling 대상: 입력 feature만 (label 미적용)  |  "
            "Best case 기준: Test 전체 AUC")
    add_text(sl, info,
             Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.7),
             size=9, color=C_BODY)


def slide_dataset(prs):
    sl = content_slide(prs)
    set_slide_title(sl, "데이터셋 분포",
                    "Train 1,097  ·  Test 275  ·  Sarcopenia 약 10.9%")

    section_label(sl, "클래스 분포", Inches(0.3), Inches(1.05))
    rows = [
        ["Split", "Sex", "n", "Normal", "Normal%", "Sarco", "Sarco%"],
        ["Train", "M",   "402", "332", "82.6%", "70", "17.4%"],
        ["Train", "F",   "695", "645", "92.8%", "50",  "7.2%"],
        ["Train", "All", "1097","977", "89.1%","120", "10.9%"],
        ["Test",  "M",   "112",  "95", "84.8%", "17", "15.2%"],
        ["Test",  "F",   "163", "150", "92.0%", "13",  "8.0%"],
        ["Test",  "All", "275", "245", "89.1%", "30", "10.9%"],
    ]
    cw = [Inches(0.7)]*2 + [Inches(0.55)]*5
    results_table(sl, Inches(0.3), Inches(1.45), Inches(4.5), Inches(2.5),
                  rows, cw, best_row=None)

    section_label(sl, "Age / BMI 통계", Inches(0.3), Inches(4.1))
    rows2 = [
        ["",      "Age Mean±Std",        "BMI Mean±Std"],
        ["Train M","59.81 ± 12.51",     "24.22 ± 3.26"],
        ["Train F","55.36 ± 12.15",     "23.09 ± 3.43"],
        ["Train All","56.99 ± 12.47",   "23.51 ± 3.41"],
        ["Test M", "59.05 ± 12.52",     "24.07 ± 3.30"],
        ["Test F", "56.52 ± 12.29",     "22.99 ± 3.19"],
        ["Test All","57.55 ± 12.45",    "23.43 ± 3.28"],
    ]
    cw2 = [Inches(1.2), Inches(1.7), Inches(1.6)]
    results_table(sl, Inches(0.3), Inches(4.5), Inches(4.5), Inches(2.4),
                  rows2, cw2, best_row=None)

    section_label(sl, "분포 시각화 (Train vs Test)", Inches(5.1), Inches(1.05))
    add_image(sl, str(BASE / "model_1" / "data_distribution.png"),
              Inches(5.0), Inches(1.45), Inches(8.0), Inches(5.4))


def _model_results_slide(prs, model_label, subtitle, all_cases_rows, best_idx_lr, best_idx_deep,
                         deep_label, lr_img_path, deep_img_path, cal_img_lr, cal_img_deep):
    """모델별 결과 슬라이드: 좌=테이블, 우=best case 그림."""
    sl = content_slide(prs)
    set_slide_title(sl, model_label, subtitle)

    section_label(sl, "Logistic Regression", Inches(0.3), Inches(1.05), w=Inches(2.5))
    lr_rows, deep_rows, cw = all_cases_rows
    results_table(sl, Inches(0.3), Inches(1.45),
                  sum(cw), Inches(0.45 * (len(lr_rows))),
                  lr_rows, cw, best_row=best_idx_lr)

    top2 = Inches(1.45) + Inches(0.45 * len(lr_rows)) + Inches(0.25)
    section_label(sl, deep_label, Inches(0.3), top2, w=Inches(2.5))
    results_table(sl, Inches(0.3), top2 + Inches(0.4),
                  sum(cw), Inches(0.45 * (len(deep_rows))),
                  deep_rows, cw, best_row=best_idx_deep)

    rx = Inches(6.5)
    section_label(sl, f"LR Best Case: {lr_rows[best_idx_lr+1][0]}",
                  rx, Inches(1.05), w=Inches(3.2))
    add_image(sl, lr_img_path,  rx,            Inches(1.45), Inches(3.2), Inches(2.2))
    add_image(sl, cal_img_lr,   rx,            Inches(3.7),  Inches(3.2), Inches(2.0))

    rx2 = Inches(9.9)
    section_label(sl, f"{deep_label} Best Case: {deep_rows[best_idx_deep+1][0]}",
                  rx2, Inches(1.05), w=Inches(3.2))
    add_image(sl, deep_img_path, rx2,          Inches(1.45), Inches(3.2), Inches(2.2))
    add_image(sl, cal_img_deep,  rx2,          Inches(3.7),  Inches(3.2), Inches(2.0))


def slide_model1(prs):
    lr_rows = [
        ["Case",        "AUC",   "AUPRC", "Brier", "Acc",   "F1"],
        ["no_scale",    "0.7527","0.3338","0.1992","0.7091","0.3333"],
        ["scale_clinic","0.7531","0.3357","0.1992","0.7055","0.3306"],
    ]
    rn_rows = [
        ["Case",        "AUC",   "AUPRC", "Brier", "Acc",   "F1"],
        ["no_scale",    "0.7287","0.3309","0.5107","0.2145","0.2059"],
        ["scale_clinic","0.7656","0.3153","0.1885","0.7382","0.3455"],
    ]
    cw = [Inches(1.4)] + [Inches(0.8)]*5
    _model_results_slide(
        prs,
        "Model 1 — Clinic Only",
        "입력: Age · Sex · BMI  |  LR 또는 ResNet1D  |  2 scaling cases",
        (lr_rows, rn_rows, cw),
        best_idx_lr=1, best_idx_deep=1,
        deep_label="ResNet1D",
        lr_img_path=IMG("model_1","scale_clinic","test_roc_curves.png"),
        deep_img_path=IMG("model_1","scale_clinic","test_roc_curves.png"),
        cal_img_lr=IMG("model_1","scale_clinic","calibration.png"),
        cal_img_deep=IMG("model_1","scale_clinic","calibration.png"),
    )


def slide_model1_figures(prs):
    """M1 scale_clinic 상세 그림 슬라이드."""
    sl = content_slide(prs)
    set_slide_title(sl, "Model 1 — Best Case (scale_clinic) 상세 결과",
                    "CV ROC · 학습 곡선 · Test ROC · Confusion Matrix · Calibration")

    imgs = [
        ("CV ROC Curves",   "cv_roc_curves.png"),
        ("Training Curves", "training_curves.png"),
        ("Test ROC",        "test_roc_curves.png"),
        ("Confusion Matrix","confusion_matrices.png"),
        ("Calibration",     "calibration.png"),
    ]
    iw, ih = Inches(2.5), Inches(2.8)
    gap = Inches(0.1)
    for i, (lbl, fname) in enumerate(imgs):
        ix = Inches(0.25) + i * (iw + gap)
        add_text(sl, lbl, ix, Inches(1.05), iw, Inches(0.3),
                 size=9, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_image(sl, IMG("model_1","scale_clinic", fname), ix, Inches(1.35), iw, ih)

    section_label(sl, "5-Fold CV (scale_clinic)", Inches(0.25), Inches(4.35), w=Inches(3.5))
    cv_rows = [
        ["",    "LR",   "ResNet1D"],
        ["AUC", "0.7989 ± 0.046", "0.8339 ± 0.027"],
        ["AUPRC","0.3600 ± 0.095","0.4013 ± 0.042"],
        ["Brier","0.1833 ± 0.016","0.1773 ± 0.021"],
        ["Acc", "0.7201 ± 0.034","0.7355 ± 0.056"],
        ["F1",  "0.3513 ± 0.038","0.4145 ± 0.048"],
    ]
    cw2 = [Inches(0.8), Inches(1.8), Inches(1.8)]
    results_table(sl, Inches(0.25), Inches(4.75), Inches(4.4), Inches(2.1),
                  cv_rows, cw2, best_row=None)

    add_rect(sl, Inches(4.9), Inches(4.35), Inches(8.1), Inches(2.5),
             fill=RGBColor(0xF5,0xF8,0xFD), line=C_DIVIDER)
    stats = [
        ("AUC: LR vs ResNet1D", "Δ=+0.035,  t p=0.022 *"),
        ("F1:  LR vs ResNet1D", "Δ=+0.063,  t p=0.021 *"),
        ("Test AUC (LR)",       "scale_clinic: 0.7531  (best)"),
        ("Test AUC (ResNet1D)", "scale_clinic: 0.7656  (best)"),
    ]
    add_text(sl, "Paired t-test (n=5 folds) 핵심 결과",
             Inches(5.1), Inches(4.45), Inches(7.7), Inches(0.4),
             size=11, bold=True, color=C_ACCENT)
    for i, (k, v) in enumerate(stats):
        iy = Inches(4.9) + i * Inches(0.48)
        add_text(sl, k, Inches(5.1), iy, Inches(3.5), Inches(0.45), size=10, bold=True)
        add_text(sl, v, Inches(8.7), iy, Inches(4.0), Inches(0.45), size=10, color=C_ACCENT)


def slide_model2(prs):
    lr_rows = [
        ["Case",        "AUC",   "AUPRC", "Brier", "Acc",   "F1"],
        ["no_scale",    "0.6133","0.2317","0.2039","0.6873","0.2321"],
        ["scale_clinic","0.6354","0.2331","0.1885","0.7345","0.2913"],
        ["scale_aec",   "0.7396","0.3390","0.1916","0.7273","0.3590"],
        ["scale_both",  "0.7388","0.3377","0.1919","0.7273","0.3590"],
    ]
    ca_rows = [
        ["Case",        "AUC",   "AUPRC", "Brier", "Acc",   "F1"],
        ["no_scale",    "0.7565","0.3310","0.1395","0.8145","0.3377"],
        ["scale_clinic","0.7571","0.3351","0.1870","0.7309","0.3729"],
        ["scale_aec",   "0.7278","0.2942","0.2056","0.6800","0.3016"],
        ["scale_both",  "0.7762","0.3816","0.1918","0.6909","0.3511"],
    ]
    cw = [Inches(1.4)] + [Inches(0.8)]*5
    _model_results_slide(
        prs,
        "Model 2 — Clinic + AEC (Matched)",
        "입력: Age · Sex · BMI + AEC 256  |  LR 또는 CrossAttn  |  4 scaling cases",
        (lr_rows, ca_rows, cw),
        best_idx_lr=2, best_idx_deep=3,
        deep_label="CrossAttn",
        lr_img_path=IMG("model_2","scale_aec","test_roc_curves.png"),
        deep_img_path=IMG("model_2","scale_both","test_roc_curves.png"),
        cal_img_lr=IMG("model_2","scale_aec","calibration.png"),
        cal_img_deep=IMG("model_2","scale_both","calibration.png"),
    )


def slide_model2_figures(prs):
    sl = content_slide(prs)
    set_slide_title(sl, "Model 2 — Best Case 상세 결과",
                    "CrossAttn best: scale_both (AUC=0.7762)  |  LR best: scale_aec (AUC=0.7396)")

    imgs = [
        ("CV ROC",   IMG("model_2","scale_both","cv_roc_curves.png")),
        ("Training", IMG("model_2","scale_both","training_curves.png")),
        ("Test ROC", IMG("model_2","scale_both","test_roc_curves.png")),
        ("CM",       IMG("model_2","scale_both","confusion_matrices.png")),
        ("Calibration", IMG("model_2","scale_both","calibration.png")),
    ]
    iw, ih = Inches(2.5), Inches(2.8)
    for i, (lbl, path) in enumerate(imgs):
        ix = Inches(0.25) + i * (iw + Inches(0.1))
        add_text(sl, lbl, ix, Inches(1.05), iw, Inches(0.3),
                 size=9, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_image(sl, path, ix, Inches(1.35), iw, ih)

    section_label(sl, "CV: LR vs CrossAttn (scale_both)", Inches(0.25), Inches(4.35), w=Inches(4.0))
    cv_rows = [
        ["",    "LR",   "CrossAttn", "Δ",   "p (t-test)"],
        ["AUC", "0.771","0.819",     "+0.049","0.046 *"],
        ["AUPRC","0.318","0.414",    "+0.096","0.024 *"],
        ["Brier","0.182","0.205",    "+0.024","0.379"],
        ["Acc", "0.740","0.680",     "−0.060","0.274"],
        ["F1",  "0.354","0.369",     "+0.015","0.718"],
    ]
    cw3 = [Inches(0.8), Inches(1.1), Inches(1.1), Inches(0.8), Inches(1.3)]
    results_table(sl, Inches(0.25), Inches(4.75), Inches(5.1), Inches(2.1),
                  cv_rows, cw3, best_row=None)

    add_rect(sl, Inches(5.7), Inches(4.35), Inches(7.3), Inches(2.5),
             fill=RGBColor(0xF5,0xF8,0xFD), line=C_DIVIDER)
    add_text(sl, "Key Finding",
             Inches(5.9), Inches(4.45), Inches(6.9), Inches(0.4),
             size=11, bold=True, color=C_ACCENT)
    notes = [
        "CrossAttn (scale_both)  Test AUC = 0.7762  — 전체 모델 중 최고",
        "no_scale 케이스에서 CrossAttn vs LR AUC 차이 +0.180 (p=0.004 **)",
        "AEC 스케일링(scale_aec) 적용 시 LR AUC가 0.614→0.740으로 상승",
    ]
    for i, n in enumerate(notes):
        add_text(sl, f"• {n}", Inches(5.9), Inches(4.9)+i*Inches(0.52),
                 Inches(7.0), Inches(0.5), size=10, color=C_BODY)


def slide_model2_2(prs):
    sl = content_slide(prs)
    set_slide_title(sl, "Model 2_2 — Clinic + AEC (Unmatched, 음성 대조군)",
                    "AEC 행 shuffle → Clinic-AEC 환자 대응 해제  |  M2 > M2_2 이면 AEC 정합 정보가 유효")

    section_label(sl, "Test AUC: Model 2 vs Model 2_2 비교", Inches(0.3), Inches(1.05), w=Inches(5.0))
    cmp_rows = [
        ["Case",        "M2 LR", "M2_2 LR", "Δ(LR)", "M2 CA", "M2_2 CA","Δ(CA)"],
        ["no_scale",    "0.6133","0.6189",   "−0.006","0.7565","0.7137",  "+0.043"],
        ["scale_clinic","0.6354","0.5974",   "+0.038","0.7571","0.7592",  "−0.002"],
        ["scale_aec",   "0.7396","0.7261",   "+0.014","0.7278","0.6882",  "+0.040"],
        ["scale_both",  "0.7388","0.7259",   "+0.013","0.7762","0.7562",  "+0.020"],
    ]
    cw = [Inches(1.5)] + [Inches(0.9)]*6
    results_table(sl, Inches(0.3), Inches(1.45), Inches(7.2), Inches(2.2),
                  cmp_rows, cw, best_row=None)

    section_label(sl, "M2_2 Best (CrossAttn: scale_clinic)", Inches(0.3), Inches(3.8), w=Inches(3.5))
    for i, (lbl, fname) in enumerate([
        ("Test ROC", "test_roc_curves.png"),
        ("CM", "confusion_matrices.png"),
        ("Calibration", "calibration.png"),
    ]):
        ix = Inches(0.3) + i * Inches(2.55)
        add_text(sl, lbl, ix, Inches(4.2), Inches(2.4), Inches(0.3),
                 size=9, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_image(sl, IMG("model_2_2","scale_clinic",fname), ix, Inches(4.5), Inches(2.4), Inches(2.4))

    add_rect(sl, Inches(7.8), Inches(1.0), Inches(5.2), Inches(5.9),
             fill=RGBColor(0xF5,0xF8,0xFD), line=C_DIVIDER)
    add_text(sl, "결과 해석", Inches(8.0), Inches(1.1), Inches(4.8), Inches(0.4),
             size=12, bold=True, color=C_ACCENT)
    notes = [
        ("M2 CrossAttn ≥ M2_2 CrossAttn",
         "대부분 케이스에서 Matched가 Unmatched보다 높거나 동등."),
        ("scale_clinic 케이스",
         "M2 CA(0.7571) vs M2_2 CA(0.7592) — 거의 차이 없음."),
        ("no_scale 케이스",
         "M2 CA(0.7565) > M2_2 CA(0.7137) — AEC 정합 효과 확인."),
        ("전반적 해석",
         "AEC 스케일링 조건에 따라 정합 효과가 달라짐.\n일관된 우위를 보이지 않아 추가 분석 필요."),
    ]
    for i, (title, body) in enumerate(notes):
        iy = Inches(1.6) + i * Inches(1.25)
        add_rect(sl, Inches(8.0), iy, Inches(4.8), Inches(0.25), fill=C_ACCENT)
        add_text(sl, title, Inches(8.1), iy, Inches(4.6), Inches(0.25),
                 size=9, bold=True, color=C_TITLE_FG)
        add_text(sl, body, Inches(8.1), iy+Inches(0.27), Inches(4.6), Inches(0.85),
                 size=9, color=C_BODY)


def slide_model3(prs):
    lr_rows = [
        ["Case",            "AUC",   "AUPRC","Brier","Acc",  "F1"],
        ["no_scale",        "0.6713","0.1958","0.1956","0.7154","0.2745"],
        ["scale_clinic",    "0.5770","0.1384","0.2157","0.6769","0.1923"],
        ["scale_aec",       "0.7608","0.2516","0.1888","0.7115","0.3119"],
        ["scale_scan",      "0.6236","0.1591","0.2019","0.6846","0.2264"],
        ["scale_clinic_aec","0.7597","0.2572","0.1880","0.7154","0.3148"],
        ["scale_clinic_scan","0.5914","0.1471","0.1951","0.7500","0.2529"],
        ["scale_aec_scan",  "0.7603","0.2597","0.1883","0.7115","0.3119"],
        ["scale_all",       "0.7599","0.2574","0.1884","0.7115","0.3119"],
    ]
    ca3_rows = [
        ["Case",            "AUC",   "AUPRC","Brier","Acc",  "F1"],
        ["no_scale",        "0.7542","0.2410","0.1717","0.7423","0.2947"],
        ["scale_clinic",    "0.7577","0.2326","0.2184","0.7346","0.3894"],
        ["scale_aec",       "0.7606","0.2351","0.2005","0.6846","0.3387"],
        ["scale_scan",      "0.7492","0.2158","0.1785","0.7500","0.3299"],
        ["scale_clinic_aec","0.7503","0.2393","0.1588","0.7885","0.3373"],
        ["scale_clinic_scan","0.7606","0.2505","0.1634","0.7885","0.3678"],
        ["scale_aec_scan",  "0.7040","0.2031","0.1556","0.7885","0.2857"],
        ["scale_all",       "0.7263","0.2262","0.1654","0.7654","0.2651"],
    ]
    cw = [Inches(1.6)] + [Inches(0.75)]*5
    _model_results_slide(
        prs,
        "Model 3 — Clinic + Scanner + AEC",
        "입력: Age · Sex · BMI + kVp + Manufacturer + AEC 256  |  LR 또는 CrossAttn3  |  8 scaling cases",
        (lr_rows, ca3_rows, cw),
        best_idx_lr=2, best_idx_deep=5,
        deep_label="CrossAttn3",
        lr_img_path=IMG("model_3","scale_aec","test_roc_curves.png"),
        deep_img_path=IMG("model_3","scale_clinic_scan","test_roc_curves.png"),
        cal_img_lr=IMG("model_3","scale_aec","calibration.png"),
        cal_img_deep=IMG("model_3","scale_clinic_scan","calibration.png"),
    )


def slide_best_summary(prs):
    sl = content_slide(prs)
    set_slide_title(sl, "Best Cases Summary",
                    "각 모델 / 서브모델별 Test 전체 AUC 기준 최우수 케이스")

    rows = [
        ["Model",   "Sub-model",  "Best Case",       "AUC",   "AUPRC","Brier","Acc",  "F1"],
        ["M1",      "LR",         "scale_clinic",    "0.7531","0.3357","0.1992","0.7055","0.3306"],
        ["M1",      "ResNet1D",   "scale_clinic",    "0.7656","0.3153","0.1885","0.7382","0.3455"],
        ["M2",      "LR",         "scale_aec",       "0.7396","0.3390","0.1916","0.7273","0.3590"],
        ["M2",      "CrossAttn",  "scale_both",      "0.7762","0.3816","0.1918","0.6909","0.3511"],
        ["M2_2",    "LR",         "scale_aec",       "0.7261","0.2982","0.1919","0.7236","0.3559"],
        ["M2_2",    "CrossAttn",  "scale_clinic",    "0.7592","0.3149","0.2525","0.5855","0.2875"],
        ["M3",      "LR",         "scale_aec",       "0.7608","0.2516","0.1888","0.7115","0.3119"],
        ["M3",      "CrossAttn3", "scale_clinic_scan","0.7606","0.2505","0.1634","0.7885","0.3678"],
    ]
    cw = [Inches(0.7), Inches(1.2), Inches(1.9)] + [Inches(0.85)]*5
    results_table(sl, Inches(0.4), Inches(1.1), Inches(11.9), Inches(4.1),
                  rows, cw, best_row=3)

    add_rect(sl, Inches(0.4), Inches(5.4), Inches(11.9), Inches(1.45),
             fill=RGBColor(0xF5,0xF8,0xFD), line=C_DIVIDER)
    add_text(sl, "Test AUC 시각 비교 (Deep Model)",
             Inches(0.6), Inches(5.45), Inches(5), Inches(0.35),
             size=10, bold=True, color=C_ACCENT)
    bar_data = [
        ("M1 ResNet1D",   0.7656, "scale_clinic"),
        ("M2 CrossAttn",  0.7762, "scale_both"),
        ("M2_2 CrossAttn",0.7592, "scale_clinic"),
        ("M3 CrossAttn3", 0.7606, "scale_clinic_scan"),
    ]
    bar_max_w = Inches(9.5)
    auc_min = 0.70
    for i, (lbl, auc, case) in enumerate(bar_data):
        bx = Inches(2.3)
        by = Inches(5.85) + i * Inches(0.26)
        bw = bar_max_w * (auc - auc_min) / (0.85 - auc_min)
        bg = C_SECTION if auc == max(d[1] for d in bar_data) else C_ACCENT
        add_rect(sl, bx, by, bw, Inches(0.20), fill=bg)
        add_text(sl, lbl, Inches(0.5), by, Inches(1.7), Inches(0.22), size=8, color=C_BODY)
        add_text(sl, f"{auc:.4f}  ({case})", bx+bw+Inches(0.05), by, Inches(2), Inches(0.22),
                 size=8, color=C_BODY)


def slide_key_findings(prs):
    sl = content_slide(prs)
    set_slide_title(sl, "Key Findings")

    findings = [
        (C_ACCENT, "Model 1: Scaling 효과",
         "ResNet1D에서 scale_clinic 적용 시 Test AUC 0.7287 → 0.7656 (+0.037). "
         "CV 단계에서 LR(0.799) vs ResNet1D(0.834)는 AUC p=0.022 유의."),
        (C_SECTION, "Model 2: AEC 기여 확인",
         "CrossAttn (scale_both) Test AUC=0.7762 — 전체 케이스 최고. "
         "no_scale에서 LR 대비 CrossAttn AUC +0.180 (p=0.004 **). "
         "단, scale_aec 시 LR AUC가 0.613→0.740 상승하며 격차 축소."),
        (C_GREEN,  "Model 2_2: 음성 대조 결과",
         "Matched vs Unmatched: CrossAttn 기준 M2 ≥ M2_2 경향 있으나 일관성 부족. "
         "no_scale 케이스에서 차이 확인 (0.757 vs 0.714). "
         "AEC 스케일링 조건별로 결과가 상이해 추가 분석 권장."),
        (C_TITLE_BG, "Model 3: Scanner 정보 효과",
         "LR 기준 scale_aec (0.7608) = M3 최고, scale_clinic_scan 조합에서 LR AUC 하락. "
         "CrossAttn3 (scale_clinic_scan) AUC=0.7606 — M3 deep model 최고. "
         "전반적으로 M2 CrossAttn 대비 추가 향상 미미."),
    ]

    bw = Inches(3.1)
    for i, (color, title, body) in enumerate(findings):
        bx = Inches(0.3) + i * (bw + Inches(0.13))
        by = Inches(1.1)
        add_rect(sl, bx, by, bw, Inches(5.7), fill=RGBColor(0xF8,0xF9,0xFA), line=C_DIVIDER)
        add_rect(sl, bx, by, bw, Inches(0.45), fill=color)
        add_text(sl, title, bx+Inches(0.1), by+Inches(0.05), bw-Inches(0.2), Inches(0.38),
                 size=11, bold=True, color=C_TITLE_FG)
        add_text(sl, body, bx+Inches(0.15), by+Inches(0.55), bw-Inches(0.3), Inches(4.9),
                 size=10, color=C_BODY)


# ── 메인 ─────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_title(prs)
    slide_overview(prs)
    slide_dataset(prs)
    slide_model1(prs)
    slide_model1_figures(prs)
    slide_model2(prs)
    slide_model2_figures(prs)
    slide_model2_2(prs)
    slide_model3(prs)
    slide_best_summary(prs)
    slide_key_findings(prs)

    out = BASE / "results_summary.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
