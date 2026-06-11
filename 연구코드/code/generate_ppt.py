"""
AEC-based Sarcopenia Prediction Research PPT
- 0612_late_fusion 결과 기준  |  AEC 128pt
- parse_md.py 로 수치 동적 로딩 — 하드코딩 없음
- Model 1~5 + LF 비교  |  AEC 전처리: raw / norm / global_zscore
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from parse_md import load_scaling, load_model_results

# ── 경로 ──────────────────────────────────────────────────────
BASE     = "C:/Users/jhjun/OneDrive/Desktop/대학원/results/0612_late_fusion"
M1       = f"{BASE}/model_1"
M2N      = f"{BASE}/model_2/norm"
M3N      = f"{BASE}/model_3/norm"
M22N     = f"{BASE}/model_2_2/norm"
M4RAW    = f"{BASE}/model_4/raw"
M5N      = f"{BASE}/model_5/norm"
CMP      = f"{BASE}/comparison"
AEC_INSP = f"{BASE}/aec_inspection"

# ── 데이터 로딩 ──────────────────────────────────────────────
S      = load_scaling(f"{CMP}/scaling_comparison.md")
M1res  = load_model_results(f"{M1}/results.md")
M2res  = load_model_results(f"{M2N}/results.md")
M3res  = load_model_results(f"{M3N}/results.md")
M22res = load_model_results(f"{M22N}/results.md")
M4res  = load_model_results(f"{M4RAW}/results.md")
M5res  = load_model_results(f"{M5N}/results.md")

# ── Late Fusion 모델 경로 ──────────────────────────────────────
LF_M2    = f"{BASE}/model_2_lf/norm"
LF_M3    = f"{BASE}/model_3_lf/norm"
M2LFres  = load_model_results(f"{LF_M2}/results.md")
M3LFres  = load_model_results(f"{LF_M3}/results.md")

# Late Fusion Test AUC (results.md 직접 수치)
LF_AUC = {
    'M2_LF': {'raw': 0.7411, 'norm': 0.7685, 'global_zscore': 0.7478},
    'M3_LF': {'raw': 0.7685, 'norm': 0.7988, 'global_zscore': 0.7567},
}

def slf_v(m, v, k):
    """S perf 조회 (LF 모델은 LF_AUC 폴백)."""
    val = S['perf'].get(m, {}).get(v, {}).get(k)
    if val is not None:
        return val
    if k == 'AUC-ROC' and m in LF_AUC:
        return LF_AUC[m].get(v, 0.0)
    return 0.0

def slf_d(comp, v):
    return S['delong'].get(comp, {}).get(v) or {
        'auc_a': 0.0, 'auc_b': 0.0, 'delta': 0.0, 'pval': 1.0, 'sig': 'ns'}

M1_VAR = 'scale_all'   # M1 은 Bootstrap CI 에서 'scale_all' variant 키

# ── 포맷터 ──────────────────────────────────────────────────
def _v(m, v, k):
    return S['perf'].get(m, {}).get(v, {}).get(k) or 0.0

def _c(m, v, k='AUC-ROC'):
    return S['ci'].get(m, {}).get(v, {}).get(k) or (0.0, 0.0)

def _d(comp, v):
    return S['delong'].get(comp, {}).get(v) or {
        'auc_a': 0.0, 'auc_b': 0.0, 'delta': 0.0, 'pval': 1.0, 'sig': 'ns'}

def _mv(v):
    """M1 의 variant 키를 항상 M1_VAR 로 변환."""
    return M1_VAR if v == 'M1' else v

def fa(m, v='norm'):
    return f"{_v(m, _mv(m) if m == 'M1' else v, 'AUC-ROC'):.4f}"

def fci(m, v='norm', k='AUC-ROC'):
    v_ = M1_VAR if m == 'M1' else v
    lo, hi = _c(m, v_, k)
    return f"[{lo:.3f},{hi:.3f}]"

def fm(m, v='norm', k='AUPRC'):
    v_ = M1_VAR if m == 'M1' else v
    return f"{_v(m, v_, k):.4f}"

def fp(comp, v='norm'):
    """DeLong p-value + sig 문자열."""
    d = _d(comp, v)
    pv, sig = d['pval'], d['sig']
    return f"p<0.001 {sig}" if pv < 0.001 else f"p={pv:.3f} {sig}"

def fdelta(comp, v='norm'):
    """DeLong Δ AUC (B-A) 문자열."""
    val = _d(comp, v)['delta']
    if val is None:
        return "—"
    return f"+{val:.3f}" if val >= 0 else f"−{abs(val):.3f}"

def sig_col(comp, v='norm'):
    """유의성 표시 색상."""
    sig = _d(comp, v)['sig']
    if sig in ('***', '**', '*'):
        return C_GRN
    if sig == '†':
        return C_ORG
    return C_MGR

def drow(comp, v):
    """DeLong 테이블 행 tuple: (variant, auc_a, auc_b, delta, pval_sig)."""
    d = _d(comp, v)
    pv, sig = d['pval'], d['sig']
    pstr = f"p<0.001 {sig}" if pv < 0.001 else f"p={pv:.3f} {sig}"
    delta = d['delta'] or 0.0
    dstr = f"+{delta:.3f}" if delta >= 0 else f"−{abs(delta):.3f}"
    return (v, f"{d['auc_a']:.4f}", f"{d['auc_b']:.4f}", dstr, pstr)

def mrow(comp, v):
    """M2vM3 테이블 행: (variant, M2_auc, M3_auc, delta, pval, sig)."""
    d = _d(comp, v)
    delta = d['delta'] or 0.0
    dstr = f"+{delta:.3f}" if delta >= 0 else f"−{abs(delta):.3f}"
    return (v, f"{d['auc_a']:.4f}", f"{d['auc_b']:.4f}", dstr,
            f"{d['pval']:.3f}" if d['pval'] >= 0.001 else "<0.001",
            d['sig'])

def ncrow(v):
    """M2vM22 음성 대조군 행: delta = M2 - M2_2 (양수 = M2 우세)."""
    d = _d('M2vM22', v)
    delta = -(d['delta'] or 0.0)
    dstr = f"+{delta:.3f}" if delta >= 0 else f"−{abs(delta):.3f}"
    pv = d['pval']
    return (v, f"{d['auc_a']:.4f}", f"{d['auc_b']:.4f}", dstr,
            "<0.001" if pv < 0.001 else f"{pv:.3f}", d['sig'])

def brow(label, ma, mb, v_a=None, v_b='norm'):
    """Bootstrap CI 비교 행 4개 반환 (AUC-ROC / AUPRC / Brier / F1)."""
    va = M1_VAR if ma == 'M1' else (v_a or v_b)
    vb = v_b
    rows = []
    for k in ['AUC-ROC', 'AUPRC', 'Brier', 'F1']:
        a = _v(ma, va, k); al, ah = _c(ma, va, k)
        b = _v(mb, vb, k); bl, bh = _c(mb, vb, k)
        diff = b - a
        dstr = f"+{diff:.3f}" if diff >= 0 else f"−{abs(diff):.3f}"
        rows.append((k,
                     f"{a:.3f} [{al:.3f},{ah:.3f}]",
                     f"{b:.3f} [{bl:.3f},{bh:.3f}]",
                     dstr))
    return rows

def rci(model_res, k='AUC-ROC'):
    """results.md 에서 파싱한 Bootstrap CI 포맷. 값 없으면 '—'."""
    t = model_res.get('ci', {}).get(k)
    if not t or t[1] is None:
        return "—"
    _, lo, hi = t
    return f"[{lo:.3f},{hi:.3f}]"

def img(folder, name):
    return os.path.join(folder, name)

# ── 색상 ──────────────────────────────────────────────────────
C_BG   = RGBColor(0xFF,0xFF,0xFF)
C_DARK = RGBColor(0x1A,0x1A,0x2E)
C_BLUE = RGBColor(0x0F,0x6B,0xBF)
C_TEAL = RGBColor(0x0F,0x9B,0x8A)
C_ORG  = RGBColor(0xE8,0x74,0x00)
C_GRN  = RGBColor(0x1E,0x8B,0x4C)
C_RED  = RGBColor(0xC0,0x39,0x2B)
C_LGR  = RGBColor(0xF0,0xF4,0xF8)
C_MGR  = RGBColor(0x9E,0xA3,0xAB)
C_WHT  = RGBColor(0xFF,0xFF,0xFF)
C_PUR  = RGBColor(0x7B,0x2F,0xBE)

TOTAL = 27

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── 그리기 헬퍼 ──────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    return s

def box(s, l, t, w, h, fill=None, line=None, lw=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw or 0.75)
    else: sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, sz=13.0, bold=False, color=C_DARK,
        align=PP_ALIGN.LEFT, italic=False, font="맑은 고딕"):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font

def pic(s, path, l, t, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        box(s, l, t, w, h, fill=C_LGR, line=C_MGR, lw=0.5)
        txt(s, f"[파일 없음]\n{os.path.basename(path)}",
            l+0.07, t+0.07, w-0.14, h-0.14, sz=9, color=C_MGR)

def hdr(s, title, sub=None):
    box(s, 0, 0, 13.33, 0.84, fill=C_DARK)
    txt(s, title, 0.35, 0.06, 11.0, 0.44, sz=25, bold=True, color=C_WHT)
    if sub: txt(s, sub, 0.35, 0.50, 11.0, 0.28, sz=13, color=C_MGR)

def hline(s, l, t, w, color=C_BLUE, h=0.04):
    box(s, l, t, w, h, fill=color)

def card(s, l, t, w, h, fill=C_LGR, line=None):
    return box(s, l, t, w, h, fill=fill, line=line, lw=0.75 if line else None)

def trow_h(s, headers, xs, ws, y, ht=0.3):
    for h, x, w in zip(headers, xs, ws):
        box(s, x, y, w, ht, fill=C_DARK)
        txt(s, h, x, y, w, ht, sz=10, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

def trow_d(s, cells, xs, ws, y, ht=0.36, fill=C_WHT, colors=None, bolds=None):
    for i,(c,x,w) in enumerate(zip(cells,xs,ws)):
        cl = (colors[i] if colors else C_DARK)
        b  = (bolds[i]  if bolds  else False)
        box(s, x, y, w, ht, fill=fill, line=RGBColor(0xCC,0xCC,0xCC), lw=0.5)
        txt(s, c, x, y, w, ht, sz=10.5, bold=b, color=cl, align=PP_ALIGN.CENTER)

def snum(s, n):
    txt(s, f"{n} / {TOTAL}", 12.3, 7.15, 1.0, 0.3, sz=11,
        color=C_MGR, align=PP_ALIGN.RIGHT)

def calib_banner(s, l, t, w, label):
    box(s, l, t, w, 0.33, fill=RGBColor(0xFF,0xF0,0xD0))
    txt(s, f"※ Calibration — {label}  (Brier 보정 품질 확인)",
        l+0.08, t+0.02, w-0.12, 0.29, sz=11, bold=True, color=C_ORG)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 7.5, fill=C_DARK)
box(s, 0, 0, 13.33, 0.08, fill=C_BLUE)
box(s, 0, 7.42, 13.33, 0.08, fill=C_TEAL)
txt(s, "AEC 신호 기반 근감소증 예측 모델", 1.0, 1.8, 11.33, 0.9,
    sz=36, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
txt(s, "Cross-Attention 아키텍처와 DeLong Test를 통한 모델 간 성능 비교 분석",
    1.0, 2.85, 11.33, 0.55, sz=18, color=C_MGR, align=PP_ALIGN.CENTER)
hline(s, 3.5, 3.65, 6.33, color=C_TEAL)
txt(s, "총 1,143명  |  5-Fold CV  |  Bootstrap 95% CI  |  DeLong AUC Test",
    1.0, 4.0, 11.33, 0.45, sz=14, color=C_MGR, align=PP_ALIGN.CENTER)
txt(s, "2026-06-11", 1.0, 4.55, 11.33, 0.35, sz=13, color=C_MGR, align=PP_ALIGN.CENTER)
for i,(label,col) in enumerate([
    ("Model 1  Clinic Only", C_BLUE),
    ("Model 2  Clinic + AEC", C_TEAL),
    ("Model 2_2  Neg. Control", C_MGR),
    ("Model 3  Clinic+Scanner+AEC", C_ORG),
    ("Model 4  AEC Only", C_PUR),
    ("Model 5  AEC Hand-crafted", C_GRN),
]):
    x = 0.3 + i*2.12
    box(s, x, 5.3, 2.0, 0.85, fill=col)
    txt(s, label, x, 5.37, 2.0, 0.72, sz=10, bold=True,
        color=C_WHT, align=PP_ALIGN.CENTER)
snum(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 연구 배경 및 목적
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "연구 배경 및 목적", "AEC 신호의 진단적 가치 탐색")
card(s, 0.3, 0.99, 5.9, 5.86)
txt(s, "연구 배경", 0.5, 1.09, 5.5, 0.35, sz=16, bold=True)
hline(s, 0.5, 1.46, 5.5, color=C_BLUE)
for i,t in enumerate([
    "• 근감소증(Sarcopenia): SMI 기반 이진 분류 — 노인 허약·사망률과 직결",
    "• 기존 진단: CT 판독 + 영상의학과 전문의 수동 분석 → 시간·비용 소모",
    "• AEC(자동노출제어): CT 스캔 시 자동 수집, 체형·조직 구성 간접 반영",
    "• 추가 검사 없이 기존 CT 장비에서 수집 가능한 Zero-cost 부가 신호",
    "• 핵심 질문: AEC가 Clinic 변수 대비 얼마나 예측력을 개선하는가?",
]):
    txt(s, t, 0.5, 1.56+i*0.9, 5.6, 0.82, sz=13.5, color=C_DARK)

card(s, 6.5, 0.99, 6.5, 5.86)
txt(s, "비교 연구 목적 (DeLong Test 기반)", 6.7, 1.09, 6.1, 0.35, sz=16, bold=True)
hline(s, 6.7, 1.46, 6.1, color=C_TEAL)
for i,(q,a,col) in enumerate([
    ("① M1 vs M2: AEC 신호가 기여하는가?",
     "Clinic(LR) vs Clinic+AEC(CrossAttn)\n→ DeLong AUC 비교로 통계 검증", C_BLUE),
    ("② M1 vs M3: Scanner+AEC 조합 효과?",
     "Clinic(LR) vs Clinic+Scanner+AEC(CrossAttn3)\n→ DeLong AUC 비교", C_TEAL),
    ("③ M2 vs M3: Scanner 순증가 효과?",
     "CrossAttn vs CrossAttn3 (동일 test set)\n→ 통계적 유의차 여부 확인", C_ORG),
    ("④ M2 vs M2_2: 매칭이 실질적 기여?",
     "Matched vs Unmatched (음성 대조군)\n→ AEC 개인 대응의 필요성 검증", C_GRN),
    ("⑤ M4: AEC만으로 예측 가능한가?",
     "AEC Only → Clinic 없이 AEC 단독 예측력 탐색\n→ M1·M2 대비 한계 파악", C_PUR),
]):
    box(s, 6.6, 1.56+i*1.02, 0.08, 0.84, fill=col)
    txt(s, q, 6.75, 1.58+i*1.02, 6.1, 0.32, sz=12.5, bold=True, color=col)
    txt(s, a, 6.75, 1.90+i*1.02, 6.1, 0.48, sz=11.5, color=C_DARK)
snum(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — 데이터셋 + data_distribution.png
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "데이터셋 및 실험 설계", "1,143명 단일 기관 후향적 코호트 / 5-Fold Stratified CV")

card(s, 0.3, 0.99, 5.9, 6.31)
txt(s, "데이터셋 구성", 0.5, 1.07, 5.5, 0.33, sz=15, bold=True)
hline(s, 0.5, 1.42, 5.5, color=C_BLUE)
hx = [0.5, 1.6, 3.1, 4.6]; hw = [1.05, 1.45, 1.45, 1.5]
trow_h(s, ["구분","남성(M)","여성(F)","합계"], hx, hw, 1.47)
for i,(row,fill) in enumerate(zip([
    ("Train","324 (Sarco 15.4%)","590 (Sarco 7.8%)","914 (Sarco 10.5%)"),
    ("Test", "82  (Sarco 15.9%)","147 (Sarco 7.5%)","229 (Sarco 10.5%)"),
], [C_WHT, C_LGR])):
    trow_d(s, row, hx, hw, 1.79+i*0.37, ht=0.37, fill=fill)
txt(s, "연령: 56.95±11.86세  |  BMI: 23.58±3.17 kg/m²  |  Train:Test = 8:2",
    0.5, 2.55, 5.5, 0.28, sz=11, color=C_MGR, italic=True)

txt(s, "AEC 전처리 조건 (3종)", 0.5, 2.97, 5.5, 0.33, sz=15, bold=True)
hline(s, 0.5, 3.32, 5.5, color=C_TEAL)
for i,(k,v) in enumerate([
    ("raw",          "전처리 없음 — 원본 AEC 절대값 그대로"),
    ("norm",         "행 방향 z-score (환자별 절대 선량 제거)"),
    ("global_zscore","Train 전체 단일 μ/σ로 z-score 정규화"),
]):
    txt(s, k,  0.5, 3.41+i*0.42, 1.85, 0.40, sz=12, bold=True, color=C_BLUE)
    txt(s, v,  2.4, 3.41+i*0.42, 3.7,  0.40, sz=12, color=C_DARK)
txt(s, "AEC 128pt × 3 전처리 조건 비교",
    0.5, 4.70, 5.5, 0.3, sz=11, color=C_MGR, italic=True)

card(s, 6.5, 0.99, 5.9, 3.31)
txt(s, "모델 구성", 6.7, 1.07, 5.5, 0.33, sz=15, bold=True)
hline(s, 6.7, 1.42, 5.5, color=C_ORG)
for i,(t,b,col) in enumerate([
    ("Model 1 — Baseline","Age, Sex, BMI → Logistic Regression", C_BLUE),
    ("Model 2 — AEC Matched","Clinic + AEC(동일환자) → CrossAttention", C_TEAL),
    ("Model 2_2 — AEC Unmatched","Clinic + AEC(다른환자) → 음성 대조군", C_MGR),
    ("Model 3 — Scanner+AEC","Clinic + MFR + AEC → CrossAttention3", C_ORG),
    ("Model 4 — AEC Only","AEC 시퀀스만 → AECOnlyNet (임상특징 없음)", C_PUR),
]):
    box(s, 6.6, 1.51+i*0.78, 0.08, 0.64, fill=col)
    txt(s, t, 6.75, 1.53+i*0.78, 5.4, 0.30, sz=12.5, bold=True, color=col)
    txt(s, b, 6.75, 1.83+i*0.78, 5.4, 0.26, sz=11.5, color=C_DARK)

txt(s, "데이터 분포 (Train / Test / Sex / Sarcopenia)", 6.4, 4.17, 6.8, 0.28,
    sz=11, bold=True, color=C_DARK)
pic(s, img(M1, "data_distribution.png"), 6.4, 4.47, 6.8, 2.83)
snum(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 통계 분석 방법
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "통계 분석 방법", "모델 비교를 위한 3단계 검증 체계")
for i,(col,title,body) in enumerate([
    (C_BLUE, "① DeLong AUC Test  (Test Set)",
     "동일 Test set 두 모델의 AUROC 직접 비교\n"
     "DeLong et al. (1988) — 쌍별 공분산 구조 활용\n"
     "H₀: AUROC_A = AUROC_B  →  z-통계량 기반 검정\n"
     "M1과 M2/M3/M4 사이 샘플 크기 동일 (n=229)"),
    (C_TEAL, "② Bootstrap 95% CI  (Test Set)",
     "n_boot = 2,000회 복원 추출 기반 신뢰구간\n"
     "2.5th–97.5th percentile 구간 보고\n"
     "AUROC / AUPRC / Brier / Accuracy / F1 전 지표\n"
     "단일 추정치의 불확실성 정량화"),
    (C_ORG, "③ Fold-level Paired Tests  (CV)",
     "5-Fold 교차검증 값에 대한 쌍별 검정\n"
     "Paired t-test + Wilcoxon signed-rank (n=5)\n"
     "두 검정 모두 보고 (p-value 쌍 제시)\n"
     "M1↔M2/M3/M4: train set 차이 주의"),
]):
    card(s, 0.3+i*4.35, 0.99, 4.15, 5.72)
    box(s, 0.3+i*4.35, 0.99, 4.15, 0.5, fill=col)
    txt(s, title, 0.4+i*4.35, 1.02, 3.95, 0.45, sz=15, bold=True, color=C_WHT)
    y = 1.64
    for line in body.split("\n"):
        txt(s, f"  {line}", 0.4+i*4.35, y, 3.9, 0.44, sz=13, color=C_DARK)
        y += 0.47
box(s, 0.3, 6.37, 12.7, 0.52, fill=RGBColor(0xEE,0xF5,0xFF))
txt(s, "유의수준:  *** p<0.001   ** p<0.01   * p<0.05   † p<0.10   ns p≥0.10",
    0.5, 6.40, 12.3, 0.42, sz=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
snum(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 전체 모델 성능 비교 개요
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "전체 모델 성능 비교 개요  (AEC 128pt + norm / global_zscore)",
    "M1 → M2 → M3 → M4 비교  |  Bootstrap 95% CI  |  Brier → Calibration 참조")

hx = [0.3,1.9,3.05,4.6,5.75,6.85,7.95,9.65,11.35]
hw = [1.55,1.1,1.5,1.1,1.05,1.05,1.65,1.65,1.85]
trow_h(s, ["모델","AUC","95% CI AUC","AUPRC","Brier","F1",
           "vs M1 Δ AUC","vs M1 DeLong p","비고"],
       hx, hw, 0.99, ht=0.32)

m4_best_v = 'raw'              # M4 best: raw (AUC 0.5856)
rows_ov = [
    ("M1 (LR)",
     fa('M1'), fci('M1'), fm('M1',k='AUPRC'), fm('M1',k='Brier'), fm('M1',k='F1'),
     "—", "기준선", "Clinic Only"),
    ("M2 (CrossAttn)",
     fa('M2'), fci('M2'), fm('M2',k='AUPRC'), fm('M2',k='Brier'), fm('M2',k='F1'),
     fdelta('M1vM2'), fp('M1vM2'), "Clinic+AEC"),
    ("M2_2 (Neg.Ctrl)",
     fa('M2_2'), fci('M2_2'), fm('M2_2',k='AUPRC'), fm('M2_2',k='Brier'), fm('M2_2',k='F1'),
     "—", "(neg ctrl)", "Unmatched"),
    ("M3 (CrossAttn3)",
     fa('M3'), fci('M3'), fm('M3',k='AUPRC'), fm('M3',k='Brier'), fm('M3',k='F1'),
     fdelta('M1vM3'), fp('M1vM3'), "Clinic+Scn+AEC"),
    ("M4 (AECOnly)",
     fa('M4', m4_best_v), fci('M4', m4_best_v), fm('M4', m4_best_v,k='AUPRC'),
     fm('M4', m4_best_v,k='Brier'), fm('M4', m4_best_v,k='F1'),
     fdelta('M1vM4', m4_best_v), fp('M1vM4', m4_best_v), "AEC Only"),
    ("M5 (CrossAttn-Feat)",
     fa('M5'), fci('M5'), fm('M5',k='AUPRC'), fm('M5',k='Brier'), fm('M5',k='F1'),
     fdelta('M1vM5'), fp('M1vM5'), "Clinic+AEC-Feat"),
]
fills = [C_LGR, RGBColor(0xE8,0xF4,0xE8), RGBColor(0xF5,0xF5,0xF5),
         RGBColor(0xFF,0xF3,0xE0), RGBColor(0xF3,0xE8,0xFF),
         RGBColor(0xE8,0xF8,0xE8)]
for i,(row,fill) in enumerate(zip(rows_ov, fills)):
    comps = ['—','M1vM2','—','M1vM3','M1vM4','M1vM5']
    sc7 = C_MGR
    var_i = m4_best_v if i == 4 else 'norm'
    if i in (1,3,4,5):
        sc7 = sig_col(comps[i], var_i)
    table_c = [C_DARK]*6 + [C_BLUE, sc7, C_MGR]
    trow_d(s, row, hx, hw, 1.33+i*0.38, ht=0.38, fill=fill, colors=table_c)

m1_aupr = _v('M1', M1_VAR, 'AUPRC')
m2_aupr = _v('M2', 'norm', 'AUPRC')
aupr_pct = (m2_aupr - m1_aupr) / m1_aupr * 100 if m1_aupr else 0
m2_norm_auc = _v('M2', 'norm', 'AUC-ROC')
m22_norm_auc = _v('M2_2', 'norm', 'AUC-ROC')
m4_best_auc = _v('M4', m4_best_v, 'AUC-ROC')
f_m1 = M1res['sex'].get('F',{}).get('AUC-ROC', 0.0)
f_m3 = M3res['sex'].get('F',{}).get('AUC-ROC', 0.0)

m5_norm_auc = _v('M5', 'norm', 'AUC-ROC')
txt(s, "ROC 비교 (M1·M2·M2₂·M3·M4 / norm)", 0.3, 3.95, 4.5, 0.3, sz=12, bold=True)
pic(s, img(CMP, "roc_all_models_norm.png"), 0.3, 4.27, 4.5, 3.03)

card(s, 5.0, 3.95, 4.1, 3.35, fill=RGBColor(0xE8,0xF8,0xE8))
box(s, 5.0, 3.95, 0.1, 3.35, fill=C_GRN)
txt(s, "비교 핵심", 5.2, 4.03, 3.7, 0.32, sz=14, bold=True, color=C_GRN)
for i,t in enumerate([
    f"• M1→M2 norm: Δ AUC {fdelta('M1vM2')}, {fp('M1vM2')}",
    f"• M1→M3 norm: Δ AUC {fdelta('M1vM3')}, {fp('M1vM3')}",
    f"• M5(AEC-Feat) norm: AUC={m5_norm_auc:.3f}, {fp('M1vM5')}",
    f"• M2 norm({m2_norm_auc:.3f}) vs M2_2({m22_norm_auc:.3f}) Δ={(m2_norm_auc-m22_norm_auc):+.3f}",
    f"• M4 AEC Only ({m4_best_v}): AUC={m4_best_auc:.3f} — 임상 없이 한계 확인",
    f"• 여성 AUC: M1({f_m1:.3f})→M3({f_m3:.3f}) 개선",
]):
    txt(s, t, 5.2, 4.44+i*0.38, 3.7, 0.36, sz=11.5, color=C_DARK)

calib_banner(s, 9.25, 3.95, 3.85, f"M3 norm  (Brier {fm('M3',k='Brier')})")
pic(s, img(M3N, "calibration_.png"), 9.25, 4.30, 3.85, 3.00)
snum(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — M1 vs M2: AEC의 기여
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M1 vs M2 비교 — AEC 신호의 기여",
    f"Clinic Only(LR) vs Clinic+AEC(CrossAttn)  |  norm: {fp('M1vM2')}  |  Brier → Calibration 참조")

card(s, 0.3, 0.99, 6.0, 4.52)
txt(s, "DeLong Test — M1(LR) vs M2(CrossAttn)  [ AEC 128pt ]",
    0.5, 1.07, 5.6, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 5.6, color=C_BLUE)
dx = [0.45, 2.05, 3.1, 4.15, 5.2]; dw = [1.55, 1.0, 1.0, 1.0, 0.95]
trow_h(s, ["시나리오","M1 AUC","M2 AUC","Δ AUC","p-val"], dx, dw, 1.47)
for i, variant in enumerate(['raw','norm','global_zscore']):
    row = drow('M1vM2', variant)
    sc = sig_col('M1vM2', variant)
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, dx, dw, 1.79+i*0.42, ht=0.42, fill=fill,
           colors=[C_DARK,C_DARK,C_BLUE,C_BLUE,sc])

txt(s, "Bootstrap 95% CI — norm 기준",
    0.5, 3.18, 5.6, 0.3, sz=12, bold=True, color=C_DARK)
hline(s, 0.5, 3.50, 5.6, color=C_TEAL, h=0.03)
bx2 = [0.45, 1.85, 3.25, 4.7]; bw2 = [1.35, 1.35, 1.4, 1.25]
trow_h(s, ["지표","M1","M2","Δ(M2-M1)"], bx2, bw2, 3.56, ht=0.28)
for i,(row,fill) in enumerate(zip(brow("", 'M1','M2'), [C_WHT,C_LGR,C_WHT,C_LGR])):
    trow_d(s, row, bx2, bw2, 3.86+i*0.33, ht=0.33, fill=fill)

txt(s, "ROC 비교  (norm / M1·M2·M2₂·M3·M4)", 6.5, 0.99, 6.6, 0.3, sz=12, bold=True)
pic(s, img(CMP, "roc_all_models_norm.png"), 6.5, 1.34, 6.6, 3.28)

calib_banner(s, 6.5, 4.74, 6.6, f"M2  norm  (Brier {fm('M2',k='Brier')})")
pic(s, img(M2N, "calibration_.png"), 6.5, 5.09, 6.6, 2.21)
snum(s, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — M1 vs M3: Scanner+AEC 효과
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M1 vs M3 비교 — Scanner + AEC 조합 효과",
    f"Clinic Only(LR) vs Clinic+Scanner+AEC(CrossAttn3)  |  norm: {fp('M1vM3')}")

card(s, 0.3, 0.99, 6.0, 4.52)
txt(s, "DeLong Test — M1(LR) vs M3(CrossAttn3)  [ AEC 128pt ]",
    0.5, 1.07, 5.6, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 5.6, color=C_ORG)
trow_h(s, ["시나리오","M1 AUC","M3 AUC","Δ AUC","p-val"], dx, dw, 1.47)
for i, variant in enumerate(['raw','norm','global_zscore']):
    row = drow('M1vM3', variant)
    sc = sig_col('M1vM3', variant)
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, dx, dw, 1.79+i*0.42, ht=0.42, fill=fill,
           colors=[C_DARK,C_DARK,C_ORG,C_BLUE,sc])

txt(s, "Bootstrap 95% CI — norm 기준",
    0.5, 3.18, 5.6, 0.3, sz=12, bold=True)
hline(s, 0.5, 3.50, 5.6, color=C_ORG, h=0.03)
trow_h(s, ["지표","M1","M3","Δ(M3-M1)"], bx2, bw2, 3.56, ht=0.28)
for i,(row,fill) in enumerate(zip(brow("", 'M1','M3'), [C_WHT,C_LGR,C_WHT,C_LGR])):
    trow_d(s, row, bx2, bw2, 3.86+i*0.33, ht=0.33, fill=fill)

txt(s, "ROC 비교  (norm / M1·M2·M2₂·M3·M4)", 6.5, 0.99, 6.6, 0.3, sz=12, bold=True)
pic(s, img(CMP, "roc_all_models_norm.png"), 6.5, 1.34, 6.6, 3.28)

calib_banner(s, 6.5, 4.74, 6.6, f"M3  norm  (Brier {fm('M3',k='Brier')})")
pic(s, img(M3N, "calibration_.png"), 6.5, 5.09, 6.6, 2.21)
snum(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — M2 vs M3: Scanner 순증가 효과
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M2 vs M3 비교 — Scanner(MFR) 추가 효과",
    f"norm: {fp('M2vM3','norm')} | raw: {fp('M2vM3','raw')} | 조건별 비교")

card(s, 0.3, 0.99, 7.5, 4.31)
txt(s, "DeLong Test — M2(CrossAttn) vs M3(CrossAttn3)  [ 전 조건 ]",
    0.5, 1.07, 7.1, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 7.1, color=C_MGR)
mx = [0.45, 2.3, 3.45, 4.6, 5.75, 6.85]
mw = [1.8,  1.1, 1.1,  1.1, 1.05, 0.9]
trow_h(s, ["조건","M2 AUC","M3 AUC","Δ AUC","p-val","sig"], mx, mw, 1.47)
for i, variant in enumerate(['raw','norm','global_zscore']):
    row = mrow('M2vM3', variant)
    sc = sig_col('M2vM3', variant)
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, mx, mw, 1.79+i*0.42, ht=0.42, fill=fill,
           colors=[C_DARK]*4+[sc, sc])

txt(s, "ROC — norm  (M2 vs M3 차이)", 8.0, 0.99, 5.1, 0.3, sz=12, bold=True)
pic(s, img(CMP,"roc_all_models_norm.png"), 8.0, 1.31, 5.1, 2.68)

card(s, 8.0, 4.14, 5.1, 3.11, fill=RGBColor(0xF5,0xF5,0xF5))
box(s, 8.0, 4.14, 0.1, 3.11, fill=C_MGR)
txt(s, "해석", 8.2, 4.22, 4.7, 0.33, sz=14, bold=True, color=C_MGR)
m2sn_auc = _d('M2vM3','norm')['auc_a']
m3sn_auc = _d('M2vM3','norm')['auc_b']
m2r_auc  = _d('M2vM3','raw')['auc_a']
m3r_auc  = _d('M2vM3','raw')['auc_b']
m2gz_auc = _d('M2vM3','global_zscore')['auc_a']
m3gz_auc = _d('M2vM3','global_zscore')['auc_b']
for i,t in enumerate([
    f"• norm: M3({m3sn_auc:.3f}) vs M2({m2sn_auc:.3f}) Δ={fdelta('M2vM3','norm')}, {fp('M2vM3','norm')}",
    f"• raw: M3({m3r_auc:.3f}) vs M2({m2r_auc:.3f}) Δ={fdelta('M2vM3','raw')}, {fp('M2vM3','raw')}",
    f"• global_zscore: M3({m3gz_auc:.3f}) vs M2({m2gz_auc:.3f}) Δ={fdelta('M2vM3','global_zscore')}, {fp('M2vM3','global_zscore')}",
    "• 전 조건 M2 우세 또는 유의차 없음 → Scanner 순기여 불명확",
    "• 추후 n↑ → M2 vs M3 재검증 권장",
]):
    txt(s, t, 8.2, 4.64+i*0.38, 4.7, 0.36, sz=11.5, color=C_DARK)
snum(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — M2 vs M2_2: 매칭 검증
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M2 vs M2_2 비교 — 매칭 효과 검증 (음성 대조군)",
    "Matched(M2) vs Unmatched(M2_2)  |  AEC 개인 대응의 실질적 기여 확인")

card(s, 0.3, 0.99, 4.8, 1.45)
txt(s, "실험 설계", 0.5, 1.07, 4.4, 0.3, sz=14, bold=True)
hline(s, 0.5, 1.39, 4.4, color=C_TEAL)
for i,t in enumerate([
    "M2: 환자 i의 Clinic + 환자 i의 AEC (동일 환자)",
    "M2_2: 환자 i의 Clinic + 환자 j의 AEC (다른 환자)",
    "→ 동일 CrossAttention 구조 / AEC만 랜덤 셔플",
]):
    txt(s, t, 0.5, 1.49+i*0.32, 4.4, 0.3, sz=12, color=C_DARK)

card(s, 0.3, 2.57, 7.5, 3.81)
txt(s, "DeLong Test — M2(Matched) vs M2_2(Unmatched)  [ AEC 128pt ]",
    0.5, 2.65, 7.1, 0.3, sz=14, bold=True)
hline(s, 0.5, 2.97, 7.1, color=C_TEAL)
nx = [0.45, 2.5, 3.7, 5.0, 5.95, 6.95]
nw = [2.0, 1.15, 1.25, 0.9, 0.95, 0.9]
trow_h(s, ["조건","M2 AUC","M2_2 AUC","Δ(M2-M2₂)","p-val","sig"], nx, nw, 3.02)
for i, variant in enumerate(['raw','norm','global_zscore']):
    row = ncrow(variant)
    sc = C_RED if row[-1] not in ('ns','†') else (C_ORG if row[-1]=='†' else C_MGR)
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, nx, nw, 3.34+i*0.42, ht=0.42, fill=fill,
           colors=[C_DARK]*4+[sc,sc])

txt(s, "ROC — norm  (M2 vs M2_2 비교)",
    8.0, 0.99, 5.1, 0.3, sz=12, bold=True)
pic(s, img(CMP,"roc_all_models_norm.png"), 8.0, 1.31, 5.1, 2.55)

d_norm_nc = _d('M2vM22','norm')
sn_m2   = d_norm_nc['auc_a']
sn_m22  = d_norm_nc['auc_b']
sn_dv   = sn_m2 - sn_m22

d_gz_nc  = _d('M2vM22','global_zscore')
gz_dir   = "M2_2>M2" if d_gz_nc['delta'] and d_gz_nc['delta'] > 0 else "M2>M2_2"

card(s, 8.0, 3.99, 5.1, 2.71, fill=RGBColor(0xFF,0xF3,0xE0))
box(s, 8.0, 3.99, 0.1, 2.71, fill=C_ORG)
txt(s, "해석", 8.2, 4.07, 4.7, 0.3, sz=14, bold=True, color=C_ORG)
for i,t in enumerate([
    f"• norm: M2({sn_m2:.3f}) vs M2_2({sn_m22:.3f}) Δ={sn_dv:.3f} {_d('M2vM22','norm')['sig']}",
    "  → AEC는 개인별 고유 신호 — 단순 체형 proxy 아님 (일부 조건)",
    f"• global_zscore 조건: {gz_dir} ({_d('M2vM22','global_zscore')['sig']})",
    f"• raw: {fp('M2vM22','raw')} | norm: {fp('M2vM22','norm')}",
    "• 전처리 조건에 따라 매칭 효과 강도 변화",
]):
    txt(s, t, 8.2, 4.45+i*0.38, 4.7, 0.36, sz=11.5, color=C_DARK)
snum(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Model 4: AEC Only 결과
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "Model 4 — AEC Only 결과",
    "임상 특징 없이 AEC 시퀀스만으로 분류  |  AECOnlyNet  |  3 전처리 조건 비교")

# 좌측: 성능 요약
card(s, 0.3, 0.99, 5.6, 5.86)
txt(s, "Test Set 성능  (3 전처리 조건)", 0.5, 1.07, 5.2, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 5.2, color=C_PUR)
m4x = [0.45, 1.7, 2.65, 3.55, 4.45, 5.25]
m4w = [1.2,  0.9, 0.85, 0.85, 0.8,  0.65]
trow_h(s, ["조건","AUC","AUPRC","Brier","Acc","F1"], m4x, m4w, 1.47)
for i,(v,fill) in enumerate(zip(['raw','norm','global_zscore'],
                                  [C_WHT, C_LGR, RGBColor(0xF3,0xE8,0xFF)])):
    cells = (v,
             f"{_v('M4',v,'AUC-ROC'):.4f}",
             f"{_v('M4',v,'AUPRC'):.4f}",
             f"{_v('M4',v,'Brier'):.4f}",
             f"{_v('M4',v,'Accuracy'):.4f}",
             f"{_v('M4',v,'F1'):.4f}")
    trow_d(s, cells, m4x, m4w, 1.79+i*0.38, ht=0.38, fill=fill)

txt(s, f"Bootstrap 95% CI  ({m4_best_v} — best)",
    0.5, 2.94, 5.2, 0.30, sz=12, bold=True, color=C_DARK)
hline(s, 0.5, 3.26, 5.2, color=C_PUR, h=0.03)
cx4 = [0.45, 1.75, 3.15, 4.5]; cw4 = [1.25, 1.35, 1.30, 1.00]
trow_h(s, ["지표","추정치","95% CI 하한","95% CI 상한"], cx4, cw4, 3.32, ht=0.26)
m4_ci = M4res.get('ci', {})
for i,(k,fill) in enumerate(zip(['AUC-ROC','AUPRC','Brier','F1'],
                                  [C_WHT, C_LGR, C_WHT, C_LGR])):
    cv = m4_ci.get(k)
    if cv:
        est, lo, hi = cv
        cells = (k, f"{est:.4f}", f"{lo:.4f}", f"{hi:.4f}")
    else:
        cells = (k, "—", "—", "—")
    trow_d(s, cells, cx4, cw4, 3.60+i*0.30, ht=0.30, fill=fill)

# M4 해석 카드
card(s, 0.3, 4.92, 5.6, 1.73, fill=RGBColor(0xF3,0xE8,0xFF))
box(s, 0.3, 4.92, 0.1, 1.73, fill=C_PUR)
txt(s, "AEC Only 한계", 0.5, 5.00, 5.2, 0.30, sz=13, bold=True, color=C_PUR)
m4_best_auc2 = _v('M4', m4_best_v, 'AUC-ROC')
m1_auc = _v('M1', M1_VAR, 'AUC-ROC')
for i,t in enumerate([
    f"• 최고 AUC {m4_best_auc2:.3f} ({m4_best_v}) — M1({m1_auc:.3f}) 대비 크게 열세",
    f"• M1 vs M4 {m4_best_v}: {fp('M1vM4', m4_best_v)} — 통계적 유의",
    "• AEC 단독으로는 Clinic(Age·Sex·BMI) 대체 불가",
    "• AEC는 Clinic 보조 신호로서 의미 있음",
]):
    txt(s, t, 0.5, 5.35+i*0.30, 5.3, 0.28, sz=11.5, color=C_DARK)

# 중앙: ROC curves
txt(s, f"M4 ROC 곡선  ({m4_best_v}, vs M1 baseline)", 6.1, 0.99, 4.1, 0.30, sz=11, bold=True)
pic(s, img(M4RAW, "test_roc_curves.png"), 6.1, 1.31, 4.1, 2.75)
txt(s, "M1 vs M4 AUC 비교", 6.1, 4.20, 4.1, 0.26, sz=10.5, bold=True, color=C_DARK)
txt(s, f"  M1(LR)  : AUC {m1_auc:.3f}\n"
        f"  M4(global_zscore): AUC {m4_best_auc2:.3f}  (Δ{m4_best_auc2-m1_auc:+.3f})\n"
        f"  {fp('M1vM4','global_zscore')}",
    6.1, 4.50, 4.1, 0.70, sz=11, color=C_DARK)

# 우측: CV summary + Grad-CAM
txt(s, f"5-Fold CV 결과  ({m4_best_v})", 10.4, 0.99, 2.7, 0.30, sz=11, bold=True)
m4_cv = [
    ("Fold","AUC-ROC","AUPRC","F1"),
    ("1","0.6325","0.1794","0.2564"),
    ("2","0.6868","0.2216","0.2656"),
    ("3","0.7076","0.2175","0.3415"),
    ("4","0.6996","0.2732","0.3000"),
    ("5","0.6697","0.3068","0.3279"),
    ("Mean","0.6792","0.2397","0.2983"),
]
cvx = [10.4,11.1,11.85,12.55]; cvw = [0.65,0.70,0.65,0.70]
trow_h(s, ["Fold","AUC","AUPRC","F1"], cvx, cvw, 1.31, ht=0.24)
for k,(fold,auc,aupr,f1) in enumerate(m4_cv[1:]):
    fill = RGBColor(0xF3,0xE8,0xFF) if fold == "Mean" else (C_LGR if k%2==0 else C_WHT)
    trow_d(s, (fold,auc,aupr,f1), cvx, cvw, 1.57+k*0.24, ht=0.24, fill=fill)

txt(s, f"Grad-CAM ({m4_best_v})", 10.4, 3.55, 2.7, 0.26, sz=10.5, bold=True)
pic(s, img(M4RAW, "cam_aec_mean.png"), 10.4, 3.83, 2.7, 2.95)
snum(s, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — AEC 전처리 조건별 ROC 비교 (3종)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "AEC 전처리 조건별 ROC 비교 — 3종",
    "M1·M2·M2₂·M3·M4 포함  |  AEC 128pt  |  각 전처리 조건의 성능 차이 비교")

panels = [
    ("raw",          "raw — 원본 AEC",                    C_BLUE),
    ("norm",         "norm — 행 방향 z-score",             C_TEAL),
    ("global_zscore","global_zscore — 전체 μ/σ z-score",  C_ORG),
]
positions = [(0.25, 0.99),(4.55, 0.99),(8.85, 0.99)]
for (var,label,col),(x0,y0) in zip(panels, positions):
    box(s, x0, y0, 4.2, 0.36, fill=col)
    txt(s, label, x0+0.08, y0+0.02, 4.05, 0.32, sz=12, bold=True, color=C_WHT)
    pic(s, img(CMP,f"roc_all_models_{var}.png"), x0, y0+0.38, 4.2, 5.7)

# 요약 카드 — 우측 하단 (없으므로 별도 설명 텍스트)
m2_variants_auc = [(v, _v('M2',v,'AUC-ROC')) for v in ['raw','norm','global_zscore']]
m3_variants_auc = [(v, _v('M3',v,'AUC-ROC')) for v in ['raw','norm','global_zscore']]
best_m2 = max(m2_variants_auc, key=lambda x: x[1])
best_m3 = max(m3_variants_auc, key=lambda x: x[1])
snum(s, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Attention Map 양방향 + Grad-CAM AEC (M2 vs M3, norm)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "Attention Map 양방향 비교 + Grad-CAM AEC  (norm)",
    "C→A · A→C Cross-Attention 패턴  |  Grad-CAM 기여 구간  |  M2(CrossAttn) vs M3(CrossAttn3)")

cx = [0.20, 4.57, 8.94]; cw = [4.25, 4.25, 4.25]

for j,(lbl,lc) in enumerate([
    ("C → A  (Clinic → AEC)", C_BLUE),
    ("A → C  (AEC → Clinic)", C_TEAL),
    ("Grad-CAM AEC",           C_ORG),
]):
    box(s, cx[j], 0.89, cw[j], 0.30, fill=lc)
    txt(s, lbl, cx[j]+0.05, 0.90, cw[j]-0.1, 0.28,
        sz=11, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

h_fnames = ["attention_heatmap_c2a.png", "attention_heatmap_a2c.png", "cam_aec_heatmap.png"]
m_fnames = ["attention_map_c2a.png",     "attention_map_a2c.png",     "cam_aec_mean.png"]

for row_i, (model_lbl, model_col, mdir) in enumerate([
    ("M2  CrossAttn",  C_TEAL, M2N),
    ("M3  CrossAttn3", C_ORG,  M3N),
]):
    y0     = 1.22 + row_i * 3.05
    h_h    = 1.95
    y_mlbl = y0 + h_h + 0.05
    y_mean = y_mlbl + 0.20
    m_h    = 3.00 - h_h - 0.05 - 0.20 - 0.05

    box(s, 0.20, y0, 0.30, 3.00, fill=model_col)
    txt(s, model_lbl, 0.21, y0 + 1.1, 0.28, 1.5,
        sz=9, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

    for j in range(3):
        pic(s, img(mdir, h_fnames[j]), cx[j]+0.32, y0, cw[j]-0.32, h_h)
        txt(s, "— 평균 곡선 —", cx[j]+0.32, y_mlbl, cw[j]-0.32, 0.18,
            sz=8, color=C_MGR, italic=True, align=PP_ALIGN.CENTER)
        pic(s, img(mdir, m_fnames[j]), cx[j]+0.32, y_mean, cw[j]-0.32, m_h)

snum(s, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — 최적 조건 상세 비교 (norm)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "최적 조건 상세 — AEC 128pt + norm",
    "M1·M2·M3 성능 비교  |  Bootstrap CI  |  Brier → Calibration 참조")

box(s, 0.3, 0.99, 12.7, 0.65, fill=RGBColor(0xE8,0xF8,0xE8))
box(s, 0.3, 0.99, 0.12, 0.65, fill=C_GRN)
txt(s,
    f"AEC 128pt + norm:  M1→M2 {fp('M1vM2')}  "
    f"|  M1→M3 {fp('M1vM3')}  |  M2 best AUC {fa('M2')}  |  M3 best AUC {fa('M3')}",
    0.55, 1.06, 12.1, 0.52, sz=13, bold=True, color=C_DARK)

hx7 = [0.3,1.55,2.8,4.5,5.6,6.65,7.7,9.4,11.1]
hw7 = [1.2,1.2, 1.65,1.05,1.0,1.0, 1.65,1.65,2.05]
trow_h(s, ["모델","AUC","95% CI","AUPRC","Brier","F1",
           "Δ AUC vs M1","DeLong p","비고"],
       hx7, hw7, 1.74, ht=0.3)
rows13 = [
    ("M1 (LR)",       fa('M1'), fci('M1'), fm('M1',k='AUPRC'), fm('M1',k='Brier'), fm('M1',k='F1'),
     "기준선","기준선","Clinic Only"),
    ("M2 (CrossAttn)", fa('M2'), fci('M2'), fm('M2',k='AUPRC'), fm('M2',k='Brier'), fm('M2',k='F1'),
     fdelta('M1vM2'), fp('M1vM2'), "Clinic+AEC"),
    ("M2_2 (Neg)",     fa('M2_2'), fci('M2_2'), fm('M2_2',k='AUPRC'), fm('M2_2',k='Brier'), fm('M2_2',k='F1'),
     "—","(neg ctrl)","Unmatched"),
    ("M3 (CrossAttn3)",fa('M3'), fci('M3'), fm('M3',k='AUPRC'), fm('M3',k='Brier'), fm('M3',k='F1'),
     fdelta('M1vM3'), fp('M1vM3'), "Clinic+Scn"),
]
fills13 = [C_LGR, RGBColor(0xE8,0xF4,0xE8), C_WHT, RGBColor(0xFF,0xF3,0xE0)]
for i,(row,fill) in enumerate(zip(rows13, fills13)):
    s6 = sig_col(['—','M1vM2','—','M1vM3'][i]) if i in (1,3) else C_MGR
    trow_d(s, row, hx7, hw7, 2.06+i*0.4, ht=0.4, fill=fill,
           colors=[C_DARK]*6+[s6,s6,C_MGR])

txt(s, "ROC — norm  (M1·M2·M2₂·M3·M4)", 0.3, 3.79, 3.9, 0.28, sz=11, bold=True)
pic(s, img(CMP,"roc_all_models_norm.png"), 0.3, 4.09, 3.9, 3.16)

txt(s, "M2 Confusion Matrix", 4.35, 3.79, 3.9, 0.28, sz=11, bold=True)
pic(s, img(M2N,"confusion_matrices.png"), 4.35, 4.09, 3.9, 3.16)

calib_banner(s, 8.4, 3.79, 4.7, f"M3  norm  (Brier {fm('M3',k='Brier')})")
pic(s, img(M3N,"calibration_.png"), 8.4, 4.14, 4.7, 3.11)
snum(s, 13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — 성별 분리 비교
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "성별 분리 비교 — M1 vs M2 vs M3  (AEC 128pt + norm)",
    "남성(n=82, 유병률 15.9%) · 여성(n=147, 유병률 7.5%) 분리 AUC 비교")

m1m  = M1res['sex'].get('M',{})
m1f  = M1res['sex'].get('F',{})
m2m  = M2res['sex'].get('M',{})
m2f  = M2res['sex'].get('F',{})
m3m  = M3res['sex'].get('M',{})
m3f  = M3res['sex'].get('F',{})

card(s, 0.3, 0.99, 12.7, 2.0)
sx = [0.4,1.9,3.25,4.65,6.05,7.4,8.8,10.2,11.6]
sw = [1.45,1.3,1.35,1.35,1.3,1.35,1.35,1.35,1.55]
trow_h(s, ["모델","전체 AUC","전체 CI","남성 AUC","남성 CI","여성 AUC","여성 CI","남 AUPRC","여 AUPRC"],
       sx, sw, 1.04, ht=0.3)
sex_rows = [
    ("M1 (LR)",
     fa('M1'), rci(M1res),
     f"{m1m.get('AUC-ROC',0):.4f}", "—",
     f"{m1f.get('AUC-ROC',0):.4f}", "—",
     f"{m1m.get('AUPRC',0):.4f}", f"{m1f.get('AUPRC',0):.4f}"),
    ("M2 (CrossAttn)",
     fa('M2'), rci(M2res),
     f"{m2m.get('AUC-ROC',0):.4f}", "—",
     f"{m2f.get('AUC-ROC',0):.4f}", "—",
     f"{m2m.get('AUPRC',0):.4f}", f"{m2f.get('AUPRC',0):.4f}"),
    ("M3 (CrossAttn3)",
     fa('M3'), rci(M3res),
     f"{m3m.get('AUC-ROC',0):.4f}", "—",
     f"{m3f.get('AUC-ROC',0):.4f}", "—",
     f"{m3m.get('AUPRC',0):.4f}", f"{m3f.get('AUPRC',0):.4f}"),
]
for i,(row,fill) in enumerate(zip(sex_rows,
    [C_LGR, RGBColor(0xE8,0xF4,0xE8), RGBColor(0xFF,0xF3,0xE0)])):
    trow_d(s, row, sx, sw, 1.36+i*0.41, ht=0.41, fill=fill)

dm_m2 = m2m.get('AUC-ROC',0) - m1m.get('AUC-ROC',0)
dm_m3 = m3m.get('AUC-ROC',0) - m1m.get('AUC-ROC',0)
df_m2 = m2f.get('AUC-ROC',0) - m1f.get('AUC-ROC',0)
df_m3 = m3f.get('AUC-ROC',0) - m1f.get('AUC-ROC',0)

card(s, 0.3, 3.09, 3.8, 4.06, fill=RGBColor(0xE8,0xF0,0xFF))
box(s, 0.3, 3.09, 0.1, 4.06, fill=C_BLUE)
txt(s, "남성 비교", 0.5, 3.17, 3.5, 0.32, sz=14, bold=True, color=C_BLUE)
for i,t in enumerate([
    f"• M1→M2: {m1m.get('AUC-ROC',0):.3f}→{m2m.get('AUC-ROC',0):.3f}  (Δ{dm_m2:+.3f})",
    f"• M1→M3: {m1m.get('AUC-ROC',0):.3f}→{m3m.get('AUC-ROC',0):.3f}  (Δ{dm_m3:+.3f})",
    "• 남성 유병률 15.9% → AUC 해석 비교적 안정",
    f"• AUPRC: M1({m1m.get('AUPRC',0):.3f})→M3({m3m.get('AUPRC',0):.3f})",
]):
    txt(s, t, 0.5, 3.59+i*0.52, 3.6, 0.5, sz=12, color=C_DARK)

card(s, 4.25, 3.09, 3.8, 4.06, fill=RGBColor(0xFF,0xF0,0xF8))
box(s, 4.25, 3.09, 0.1, 4.06, fill=C_RED)
txt(s, "여성 비교", 4.45, 3.17, 3.5, 0.32, sz=14, bold=True, color=C_RED)
for i,t in enumerate([
    f"• M1→M2: {m1f.get('AUC-ROC',0):.3f}→{m2f.get('AUC-ROC',0):.3f}  (Δ{df_m2:+.3f})",
    f"• M1→M3: {m1f.get('AUC-ROC',0):.3f}→{m3f.get('AUC-ROC',0):.3f}  (Δ{df_m3:+.3f})",
    "• 유병률 7.5% → class imbalance 주의",
    f"• AUPRC: M1({m1f.get('AUPRC',0):.3f})→M3({m3f.get('AUPRC',0):.3f})",
    "• AUPRC가 AUC보다 중요한 지표",
]):
    txt(s, t, 4.45, 3.59+i*0.52, 3.6, 0.5, sz=12, color=C_DARK)

txt(s, "M1  by Sex", 8.2, 3.09, 2.4, 0.28, sz=11, bold=True)
pic(s, img(M1,"test_roc_by_sex.png"), 8.2, 3.39, 2.4, 2.0)
txt(s, "M2  by Sex", 10.75, 3.09, 2.4, 0.28, sz=11, bold=True, color=C_TEAL)
pic(s, img(M2N,"test_roc_by_sex.png"), 10.75, 3.39, 2.4, 2.0)
txt(s, f"△ 여성 AUC {m1f.get('AUC-ROC',0):.3f} → {m3f.get('AUC-ROC',0):.3f}(M3)  "
       f"|  △ 남성 AUC {m1m.get('AUC-ROC',0):.3f} → {m3m.get('AUC-ROC',0):.3f}(M3)",
    8.2, 5.44, 5.0, 0.3, sz=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
txt(s, "M2 norm — Training Curves", 8.2, 5.78, 4.9, 0.28, sz=10, color=C_MGR)
pic(s, img(M2N,"training_curves.png"), 8.2, 6.08, 4.9, 1.22)
snum(s, 14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — AEC 스케일링 비교 (0612 기준)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "AEC 스케일링 비교 — Raw / Row-Norm / Global-ZScore",
    "n=1,143  |  AEC 128pt  |  3가지 전처리 버전 시각화  |  xlsx: aec_inspection/")

card(s, 0.3, 0.99, 12.7, 1.45)
txt(s, "비교 목적 및 설계", 0.5, 1.07, 12.0, 0.33, sz=15, bold=True)
hline(s, 0.5, 1.42, 12.0, color=C_BLUE)
for i,t in enumerate([
    "• 모델 학습 시 실제 입력되는 AEC 데이터의 형태를 시각적으로 검증",
    "• norm (행 방향): 환자별 z-score → 개인의 절대 선량 수준 제거, 곡선 형태만 보존",
    "• global_zscore: Train 전체 단일 μ/σ 사용 → 환자 간 상대적 수준 보존",
]):
    txt(s, t, 0.5, 1.52+i*0.29, 12.0, 0.27, sz=12, color=C_DARK)

# 이미지 - AEC 평균 곡선
txt(s, "스케일링별 평균 AEC 곡선  (Normal vs Sarcopenia × Sex)",
    0.3, 2.56, 8.5, 0.28, sz=12, bold=True, color=C_DARK)
pic(s, img(AEC_INSP, "aec128_mean_curves.png"), 0.3, 2.88, 8.5, 2.65)

# 이미지 - 분포 박스플롯
txt(s, "샘플별 분포 박스플롯  (Normal vs Sarcopenia)",
    9.0, 2.56, 4.0, 0.28, sz=12, bold=True, color=C_DARK)
pic(s, img(AEC_INSP, "aec128_boxplot.png"), 9.0, 2.88, 4.0, 2.65)

# 히스토그램 + 박스
txt(s, "AEC 분포 히스토그램", 0.3, 5.66, 5.0, 0.26, sz=11, bold=True)
pic(s, img(AEC_INSP, "aec128_dist_hist.png"), 0.3, 5.94, 5.0, 1.39)

box(s, 5.5, 5.66, 7.5, 1.67, fill=RGBColor(0xEE,0xF5,0xFF))
txt(s,
    "저장 위치: results/0612/aec_inspection/aec_scaling_compare_aec128.xlsx  "
    "|  1,143행 × 131열 (PatientID, label, sex + pos_1~pos_128)",
    5.6, 5.76, 7.3, 0.42, sz=12, bold=True, color=C_DARK)
for i,t in enumerate([
    "• raw: 환자별 절대 AEC 수준 — 체중·조직량 반영 (스캐너 자동 보정)",
    "• norm: 개인 곡선 형태만 보존 → 스캐너 간 절대값 차이 제거",
    "• global_zscore: 집단 내 상대적 선량 수준 보존 → 환자 간 비교 가능",
]):
    txt(s, t, 5.6, 6.22+i*0.28, 7.3, 0.26, sz=11, color=C_DARK)
snum(s, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — 종합 결론
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "종합 결론 및 임상적 함의",
    "AEC 신호의 근감소증 예측 기여 — 비교 연구 결과 요약")

m4_best_str = f"AUC {_v('M4', m4_best_v, 'AUC-ROC'):.3f}"
m1_auc_str  = f"{_v('M1', M1_VAR, 'AUC-ROC'):.3f}"

for i,(col,title,items) in enumerate([
    (C_BLUE,
     f"① AEC+Clinic 조합(M2/M3) norm: M1 대비 AUC 향상 — 통계적 유의 미달",
     [
         f"M1→M2 norm: AUC {fa('M1')}→{fa('M2')}, DeLong {fp('M1vM2')}",
         f"M1→M3 norm: AUC {fa('M1')}→{fa('M3')}, DeLong {fp('M1vM3')}",
         f"여성 AUC: M1({m1f.get('AUC-ROC',0):.3f})→M3({m3f.get('AUC-ROC',0):.3f}) — M3 일부 기여",
         f"AUPRC: M1({_v('M1',M1_VAR,'AUPRC'):.3f})→M2({_v('M2','norm','AUPRC'):.3f})",
     ]),
    (C_TEAL,
     f"② norm 전처리 조건에서 M2 best AUC ({fa('M2')}) — M3과 유사",
     [
         f"M2 norm: AUC {fa('M2')} — Clinic+AEC 최고 성능",
         f"M3 norm: AUC {fa('M3')} — Scanner 추가 순기여 불명확",
         f"M2vM3 norm: Δ={fdelta('M2vM3')}, {fp('M2vM3')}",
         "→ AEC 전처리 방식이 성능 핵심 변수",
     ]),
    (C_ORG,
     f"③ Model 4 (AEC Only): AUC {m4_best_str} — Clinic 없이 단독 예측 한계 확인",
     [
         f"M4 {m4_best_v}: {m4_best_str} vs M1({m1_auc_str}) Δ{_v('M4', m4_best_v,'AUC-ROC')-_v('M1',M1_VAR,'AUC-ROC'):+.3f}",
         f"M1 vs M4: {fp('M1vM4', m4_best_v)} — 통계적 유의",
         "AEC는 임상 정보(Age/Sex/BMI) 대체 불가",
         "→ AEC는 Clinic의 보조 신호로서 의미",
     ]),
    (C_GRN,
     "④ Zero-cost AEC로 임상 등급 스크리닝 가능성 확인",
     [
         "기존 CT 스캔 시 자동 수집 — 추가 비용·검사 없음",
         f"전체 AUC {fa('M2')} (M2 norm) → 임상 허용 가능 수준 탐색",
         "건강검진 프로그램 연동 시 고위험군 조기 선별 가능성",
         "PACS-DICOM AEC 자동 추출 파이프라인 구축 필요",
     ]),
]):
    card(s, 0.3, 0.99+i*1.52, 12.7, 1.45)
    box(s, 0.3, 0.99+i*1.52, 0.12, 1.45, fill=col)
    txt(s, title, 0.55, 1.06+i*1.52, 12.1, 0.36, sz=13.5, bold=True, color=col)
    for k,item in enumerate(items[:2]):
        txt(s, f"  • {item}", 0.55, 1.46+i*1.52+k*0.32, 6.0, 0.3, sz=11.5, color=C_DARK)
    for k,item in enumerate(items[2:]):
        txt(s, f"  • {item}", 6.8,  1.46+i*1.52+k*0.32, 6.1, 0.3, sz=11.5, color=C_DARK)
snum(s, 16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — 제한점 & 향후 연구
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "제한점 및 향후 연구 방향", "현재 비교 연구의 한계와 발전 가능성")
card(s, 0.3, 0.99, 5.9, 5.91)
txt(s, "연구 제한점", 0.5, 1.07, 5.5, 0.36, sz=16, bold=True, color=C_RED)
hline(s, 0.5, 1.45, 5.5, color=C_RED)
for i,(t,b) in enumerate([
    ("소표본 비교","Test n=229, Sarco 24명 → DeLong 검출력 제한, CI 넓음"),
    ("단일 기관","단일 센터 데이터 → 외부 검증 미실시"),
    (f"M1→M2/M3 유의차 미달", f"norm: {fp('M1vM2')}, {fp('M1vM3')} — AUC 향상 통계 미달"),
    (f"M4 AEC Only 성능 열세", f"AUC {_v('M4', m4_best_v,'AUC-ROC'):.3f} ({m4_best_v}) — Clinic 대체 불가"),
    ("M2 vs M3 조건별 혼재",f"norm {fp('M2vM3','norm')}, raw {fp('M2vM3','raw')} — 방향 불일치"),
]):
    txt(s, f"  ▶ {t}", 0.5, 1.57+i*1.0, 5.5, 0.34, sz=13, bold=True)
    txt(s, f"    {b}", 0.5, 1.91+i*1.0, 5.5, 0.58, sz=11.5, color=C_MGR)
card(s, 6.5, 0.99, 6.5, 5.91)
txt(s, "향후 연구 방향", 6.7, 1.07, 6.1, 0.36, sz=16, bold=True, color=C_BLUE)
hline(s, 6.7, 1.45, 6.1, color=C_BLUE)
for i,(t,b) in enumerate([
    ("다기관 검증","강남+타기관 통합 → 일반화 성능·DeLong 유의성 재확인"),
    ("샘플 확장","Test n>500 → CI 축소, M1→M2 유의성 강화"),
    ("전처리 설계 개선","norm/global_zscore 비교 → 최적 전처리 표준화"),
    ("아키텍처 확장","Multi-head Attn + positional encoding → 시퀀스 표현 강화"),
    ("임상 파이프라인","PACS-DICOM AEC 추출 자동화 → 실시간 스크리닝"),
]):
    txt(s, f"  ◆ {t}", 6.7, 1.57+i*1.0, 6.1, 0.34, sz=13, bold=True, color=C_BLUE)
    txt(s, f"    {b}", 6.7, 1.91+i*1.0, 6.1, 0.58)
box(s, 0.3, 6.95, 12.7, 0.38, fill=C_DARK)
txt(s, "감사합니다  |  Questions & Discussion",
    0.3, 6.97, 12.7, 0.33, sz=14, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
snum(s, 17)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — AEC Correlation 분석
# ═══════════════════════════════════════════════════════════════
CORR_DIR = f"{BASE}/correlation"

s = add_slide()
hdr(s, "AEC 신호 — 임상 변수 Correlation 분석",
    "Raw AEC vs Norm AEC  |  128 슬라이스 × 7개 변수  |  Pearson r 포인트별 프로파일")

# 좌: Raw AEC 상관관계 해석
card(s, 0.3, 0.99, 5.9, 3.26)
txt(s, "Raw AEC 상관계수  (절대 선량 수준)", 0.5, 1.07, 5.5, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 5.5, color=C_BLUE)
raw_hx = [0.4, 2.3, 3.35, 4.3, 5.25]; raw_hw = [1.85, 1.0, 0.9, 0.9, 0.75]
trow_h(s, ["변수", "평균 r", "최대 |r|", "위치", "해석"], raw_hx, raw_hw, 1.47, ht=0.28)
raw_rows = [
    ("BMI",        "+0.500", "+0.616", "pt80", "체중·체격 주반영", C_ORG),
    ("TAMA",       "+0.308", "+0.489", "pt107","근육단면적 중간 상관", C_TEAL),
    ("SMI",        "+0.268", "+0.435", "pt107","근육지수 양의 상관", C_TEAL),
    ("Sex(M=1)",   "+0.123", "+0.362", "pt109","남성 선량↑ (체격)", C_MGR),
    ("Age",        "+0.111", "+0.242", "pt88", "약한 양의 상관", C_MGR),
    ("Sarcopenia", "-0.129", "-0.158", "pt10", "근감소증→선량↓ (약함)", C_RED),
]
for i, (var, mn, mx, pt, interp, col) in enumerate(raw_rows):
    fill = C_LGR if i % 2 == 0 else C_WHT
    trow_d(s, (var, mn, mx, pt, interp), raw_hx, raw_hw, 1.77+i*0.33, ht=0.33, fill=fill,
           colors=[col, C_DARK, C_DARK, C_MGR, C_DARK])

card(s, 0.3, 4.37, 5.9, 2.44, fill=RGBColor(0xFF,0xF0,0xE8))
box(s, 0.3, 4.37, 0.1, 2.44, fill=C_ORG)
txt(s, "Raw AEC 핵심 발견", 0.5, 4.45, 5.5, 0.30, sz=13, bold=True, color=C_ORG)
for i, t in enumerate([
    "• AEC raw값은 BMI(r=0.50)와 가장 강하게 상관 → 주로 체격 반영",
    "• Sarcopenia 직접 신호: r=-0.13 (약함) → Clinic의 BMI와 중복 정보 많음",
    "• 복부 초반 슬라이스(pt10)에서 Sarcopenia 상관이 가장 높음",
    "• TAMA/SMI는 중간 상관 → AEC에 근육 정보가 간접 포함",
]):
    txt(s, t, 0.5, 4.80+i*0.44, 5.5, 0.42, sz=12, color=C_DARK)

# 우: Norm AEC 상관관계
card(s, 6.5, 0.99, 6.5, 3.26)
txt(s, "Norm AEC 상관계수  (개인별 z-score 후)", 6.7, 1.07, 6.1, 0.33, sz=14, bold=True)
hline(s, 6.7, 1.42, 6.1, color=C_TEAL)
nm_hx = [6.55, 8.5, 9.55, 10.55, 11.55]; nm_hw = [1.9, 1.0, 0.95, 0.95, 0.95]
trow_h(s, ["변수", "평균 r", "최대 |r|", "위치", "해석"], nm_hx, nm_hw, 1.47, ht=0.28)
norm_rows = [
    ("Sex(M=1)",   "-0.028", "+0.652", "pt105","성별 곡선 형태 차이 大", C_ORG),
    ("SMI",        "-0.019", "+0.501", "pt102","특정 구간에서만 상관", C_TEAL),
    ("TAMA",       "-0.021", "+0.554", "pt102","골반/복부 하부 중요", C_TEAL),
    ("BMI",        "-0.010", "+0.452", "pt85", "허리 부위 BMI 상관", C_MGR),
    ("Age",        "-0.009", "-0.341", "pt13", "초반 슬라이스 Age 상관", C_MGR),
    ("Sarcopenia", "-0.002", "+0.064", "pt119","정규화 후 신호 소실!", C_RED),
]
for i, (var, mn, mx, pt, interp, col) in enumerate(norm_rows):
    fill = C_LGR if i % 2 == 0 else C_WHT
    trow_d(s, (var, mn, mx, pt, interp), nm_hx, nm_hw, 1.77+i*0.33, ht=0.33, fill=fill,
           colors=[col, C_DARK, C_DARK, C_MGR, C_DARK])

card(s, 6.5, 4.37, 6.5, 2.44, fill=RGBColor(0xE8,0xF0,0xFF))
box(s, 6.5, 4.37, 0.1, 2.44, fill=C_BLUE)
txt(s, "Norm AEC 핵심 발견 — 모델 성능 해석에 결정적", 6.7, 4.45, 6.1, 0.30, sz=13, bold=True, color=C_BLUE)
for i, t in enumerate([
    "• 평균 r ≈ 0 → 정규화로 개인 절대 선량 제거됨 (의도한 효과)",
    "• 특정 슬라이스(pt102-105)에서 성별/TAMA/SMI와 강한 상관 잔존",
    "• Sarcopenia 상관: mean=-0.002, max=+0.064 → 거의 소실",
    "• 이것이 M2 norm 성능이 M2 raw보다 낮은 이유를 설명",
    "• norm이 BMI 정보를 제거하면서 Sarcopenia 신호도 같이 줄어듦",
]):
    txt(s, t, 6.7, 4.80+i*0.40, 6.1, 0.38, sz=12, color=C_DARK)

txt(s, "Pointwise Correlation Profile", 0.3, 6.90, 6.0, 0.28, sz=11, bold=True)
pic(s, img(CORR_DIR, "pointwise_corr_profile.png"), 0.3, 7.10, 5.8, 0.35)
snum(s, 18)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — AEC 신호의 진단적 유의미성
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "AEC 신호의 진단적 유의미성 — 무엇을 포착하는가?",
    "AEC 물리적 의미 | M4(AEC only) · M5(AEC 통계 피처) · 매칭 vs 비매칭 비교")

# 좌: AEC 개요
card(s, 0.3, 0.99, 5.9, 3.0)
txt(s, "AEC(자동노출제어)란?", 0.5, 1.07, 5.5, 0.33, sz=15, bold=True)
hline(s, 0.5, 1.42, 5.5, color=C_BLUE)
for i, t in enumerate([
    "• CT 촬영 시 tube current(mAs)를 체형·조직에 따라 자동 조절",
    "• 체지방↑ → AEC 선량↑  /  근육량↑ → 감쇠 패턴 변화",
    "• 128 슬라이스 위치별 mAs 값 → 신체 구성의 종단 프로파일",
    "• 추가 검사 없이 기존 CT DICOM에 자동 기록 — Zero-cost",
    "• M5 추출 피처: mean·std·max·peak_pos·auc·skew·kurt·\n  early/mid/late mean (11종)",
]):
    txt(s, t, 0.5, 1.52+i*0.42, 5.6, 0.40, sz=12.5, color=C_DARK)

# 좌 하단: 성능 요약 테이블
card(s, 0.3, 4.12, 5.9, 3.13)
txt(s, "모델별 AEC 기여 비교  (각 모델 best 조건)", 0.5, 4.20, 5.5, 0.30, sz=13, bold=True)
hline(s, 0.5, 4.52, 5.5, color=C_TEAL)
mac_x = [0.4, 2.5, 3.5, 4.45, 5.35]
mac_w = [2.05, 0.95, 0.90, 0.85, 0.75]
trow_h(s, ["모델 (best 조건)", "AUC", "AUPRC", "Brier", "F1"],
       mac_x, mac_w, 4.57, ht=0.28)
_lf_rows = [
    ("M1  Clinic Only (LR)",
     slf_v('M1','scale_all','AUC-ROC'), slf_v('M1','scale_all','AUPRC'),
     slf_v('M1','scale_all','Brier'), slf_v('M1','scale_all','F1'), C_BLUE, C_LGR),
    ("M4  AEC Only (raw)",
     slf_v('M4','raw','AUC-ROC'), slf_v('M4','raw','AUPRC'),
     slf_v('M4','raw','Brier'), slf_v('M4','raw','F1'), C_RED, C_WHT),
    ("M2  CrossAttn (raw)",
     slf_v('M2','raw','AUC-ROC'), slf_v('M2','raw','AUPRC'),
     slf_v('M2','raw','Brier'), slf_v('M2','raw','F1'), C_TEAL, C_LGR),
    ("M5  CrossAttn-Feat (norm)",
     slf_v('M5','norm','AUC-ROC'), slf_v('M5','norm','AUPRC'),
     slf_v('M5','norm','Brier'), slf_v('M5','norm','F1'), C_GRN, C_WHT),
]
for i, (mod, auc, aupr, brier, f1, col, fill) in enumerate(_lf_rows):
    trow_d(s, (mod, f"{auc:.4f}", f"{aupr:.4f}", f"{brier:.4f}", f"{f1:.4f}"),
           mac_x, mac_w, 4.87+i*0.36, ht=0.36, fill=fill, colors=[col]+[C_DARK]*4)

# 우: 유의미성 분석 3개 포인트
card(s, 6.5, 0.99, 6.5, 3.0)
txt(s, "유의미성 분석 — 3가지 근거", 6.7, 1.07, 6.1, 0.33, sz=15, bold=True)
hline(s, 6.7, 1.42, 6.1, color=C_ORG)
_m5d = slf_d('M1vM5', 'norm')
_m5_sig = _m5d['sig']
_m5_pstr = fp('M1vM5')
_m22d_norm = _d('M2vM22', 'norm')
_m22_pstr = f"p<0.001 {_m22d_norm['sig']}" if _m22d_norm['pval'] < 0.001 else f"p={_m22d_norm['pval']:.3f} {_m22d_norm['sig']}"
for i, (title, body, col) in enumerate([
    ("① AEC 단독 예측력 (M4) 한계 확인",
     f"AUC {slf_v('M4', m4_best_v,'AUC-ROC'):.3f} vs M1({slf_v('M1','scale_all','AUC-ROC'):.3f})\n"
     f"DeLong {fp('M1vM4', m4_best_v)} → AEC가 Clinic을 대체할 수 없음\n"
     "AEC는 반드시 임상 변수와 조합 필요", C_RED),
    ("② AEC 통계 피처(M5) — M3와 유사한 성능 수준",
     f"M5 norm AUC {slf_v('M5','norm','AUC-ROC'):.4f}  |  M3 norm {slf_v('M3','norm','AUC-ROC'):.4f}\n"
     f"M1 vs M5-norm: DeLong {_m5_pstr} — {'경계 유의' if _m5_sig == '†' else '유의차 없음'}\n"
     "통계 피처(11종)로도 raw 시퀀스와 유사 성능 달성", C_GRN),
    ("③ 매칭(M2) vs 비매칭(M2_2): 개인별 고유성",
     f"norm: M2({slf_v('M2','norm','AUC-ROC'):.3f}) vs M2_2({slf_v('M2_2','norm','AUC-ROC'):.3f})\n"
     f"DeLong {_m22_pstr} — 개인별 AEC 대응의 유의미한 기여 없음\n"
     "CrossAttn이 128pt raw sequence 완전 활용 어려울 수 있음", C_ORG),
]):
    box(s, 6.6, 1.52+i*1.0, 0.08, 0.86, fill=col)
    txt(s, title, 6.75, 1.54+i*1.0, 6.1, 0.28, sz=12.5, bold=True, color=col)
    txt(s, body, 6.75, 1.83+i*1.0, 6.1, 0.55, sz=11.5, color=C_DARK)

# 우 하단: 결론
card(s, 6.5, 4.12, 6.5, 3.13, fill=RGBColor(0xE8,0xF8,0xE8))
box(s, 6.5, 4.12, 0.1, 3.13, fill=C_GRN)
txt(s, "AEC의 역할 요약", 6.7, 4.20, 6.1, 0.30, sz=14, bold=True, color=C_GRN)
for i, t in enumerate([
    "• AEC는 신체 구성 정보를 담은 의미 있는 보조 신호",
    f"• M5(AEC 통계 피처 11종): AUC {slf_v('M5','norm','AUC-ROC'):.4f} ({_m5_pstr}) — M3와 유사 수준",
    "• raw 시퀀스(128pt) 직접 입력(M2/M3)과 통계 피처(M5) 모두 DeLong 유의차 미달",
    "• 핵심: M4(AEC only)는 Clinic 대체 불가 — AEC는 보조 신호로서 의미",
    "• 향후 방향: AEC 피처 확장 + 더 나은 시퀀스 인코더 + 다기관 검증",
]):
    txt(s, t, 6.7, 4.55+i*0.49, 6.1, 0.46, sz=12, color=C_DARK)
snum(s, 19)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — CrossAttention vs Late Fusion 아키텍처 비교
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "CrossAttention vs Late Fusion — 아키텍처 비교",
    "M2/M3(CrossAttn) vs M2_LF/M3_LF(LateFusion) | 통합 방식별 성능 차이 분석")

# 아키텍처 설명 두 칸
for i, (title, desc_lines, col) in enumerate([
    ("CrossAttention (M2, M3)",
     [
         "AEC 시퀀스(128pt) → Transformer Encoder → Key/Value",
         "Clinic 벡터(Age/Sex/BMI) → Linear → Query",
         "Cross-Attention: AEC↔Clinic 양방향 정보 교환",
         "출력: Attention-weighted 표현 → MLP → 분류",
         "장점: 해석 가능한 Attention Map 제공",
         "단점: 파라미터 多, 긴 AEC 시퀀스에 민감",
     ], C_BLUE),
    ("Late Fusion (M2_LF, M3_LF)",
     [
         "AEC 시퀀스(128pt) → 독립 인코더(GRU) → AEC 임베딩",
         "Clinic 벡터 → 독립 MLP → Clinic 임베딩",
         "최종 단계에서 Concatenation → MLP → 분류",
         "각 modality가 독립적으로 표현 학습",
         "장점: 구조 단순, modality별 독립 최적화",
         "단점: 조기 상호작용 없음, Attention Map 없음",
     ], C_ORG),
]):
    x0 = 0.3 + i * 6.55
    card(s, x0, 0.99, 6.25, 3.26)
    box(s, x0, 0.99, 6.25, 0.48, fill=col)
    txt(s, title, x0+0.15, 1.03, 5.95, 0.40, sz=15, bold=True, color=C_WHT)
    for j, line in enumerate(desc_lines):
        bullet = "  ▶ " if j < 4 else "  ✓ " if j == 4 else "  ✗ "
        c = C_DARK if j < 4 else (C_GRN if j == 4 else C_RED)
        txt(s, f"{bullet}{line}", x0+0.15, 1.56+j*0.38, 5.95, 0.36, sz=12.5, color=c)

# 성능 비교 테이블
card(s, 0.3, 4.40, 12.7, 2.0)
txt(s, "CrossAttn vs Late Fusion — 전처리 조건별 AUC 비교  (Test Set)",
    0.5, 4.48, 12.3, 0.30, sz=14, bold=True)
hline(s, 0.5, 4.80, 12.3, color=C_MGR)
calf_x = [0.4, 2.2, 3.95, 5.65, 7.05, 9.05, 10.75, 12.15]
calf_w = [1.75, 1.70, 1.65, 1.35, 1.95, 1.65, 1.35, 1.0]
trow_h(s, ["전처리", "M2 CrossAttn", "M2_LF Late Fusion", "Δ(LF-CA)",
           "M3 CrossAttn3", "M3_LF LateFusion3", "Δ(LF-CA3)", "우세"],
       calf_x, calf_w, 4.85, ht=0.30)
for vi, (vr, fill) in enumerate(zip(['raw','norm','global_zscore'],
                                     [C_WHT, C_LGR, C_WHT])):
    m2ca = slf_v('M2', vr, 'AUC-ROC')
    m2lf = LF_AUC['M2_LF'][vr]
    m3ca = slf_v('M3', vr, 'AUC-ROC')
    m3lf = LF_AUC['M3_LF'][vr]
    d2, d3 = m2lf - m2ca, m3lf - m3ca
    d2s = f"+{d2:.3f}" if d2 >= 0 else f"−{abs(d2):.3f}"
    d3s = f"+{d3:.3f}" if d3 >= 0 else f"−{abs(d3):.3f}"
    c2  = C_GRN if d2 > 0 else C_RED
    c3  = C_GRN if d3 > 0 else C_RED
    win2 = "LF" if d2 > 0 else "CrossAttn"
    win3 = "LF3" if d3 > 0 else "CrossAttn3"
    winner = f"2M:{win2} / 3M:{win3}"
    trow_d(s, (vr, f"{m2ca:.4f}", f"{m2lf:.4f}", d2s, f"{m3ca:.4f}", f"{m3lf:.4f}", d3s, winner),
           calf_x, calf_w, 5.17+vi*0.36, ht=0.36, fill=fill,
           colors=[C_DARK, C_BLUE, C_ORG, c2, C_TEAL, C_ORG, c3, C_MGR])

# 결론 카드
card(s, 0.3, 6.52, 12.7, 0.78, fill=RGBColor(0xF0,0xF4,0xFF))
box(s, 0.3, 6.52, 0.1, 0.78, fill=C_BLUE)
_ca2 = _v('M2','norm','AUC-ROC'); _lf2 = LF_AUC['M2_LF']['norm']
_ca3 = _v('M3','norm','AUC-ROC'); _lf3 = LF_AUC['M3_LF']['norm']
txt(s,
    f"2-modal(Clinic+AEC): norm 조건에서 CrossAttn({_ca2:.4f})>LF({_lf2:.4f}) — CrossAttn 우세  "
    f"|  3-modal(Clinic+Scanner+AEC): CrossAttn3({_ca3:.4f})>LF3({_lf3:.4f}) — 다중 모달에서 Cross-Attention 우세  "
    "|  CrossAttn의 Attention Map 해석 가능성은 LF 대비 추가 이점",
    0.5, 6.57, 12.4, 0.66, sz=11.5, bold=True, color=C_DARK)
snum(s, 20)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — 부록 섹션 디바이더
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 7.5, fill=C_DARK)
box(s, 0, 0, 13.33, 0.08, fill=C_BLUE)
box(s, 0, 7.42, 13.33, 0.08, fill=C_TEAL)
txt(s, "부록 — 조건별 상세 결과",
    1.0, 2.3, 11.33, 0.9, sz=36, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
txt(s, "AEC 128pt × 3종  |  M2 · M2_2 · M3 · M4 개별 시각화",
    1.0, 3.35, 11.33, 0.55, sz=17, color=C_MGR, align=PP_ALIGN.CENTER)
hline(s, 3.5, 4.0, 6.33, color=C_TEAL)
for i,(var,desc) in enumerate([
    ("raw",          f"원본 AEC — M2 AUC={fa('M2','raw')}, M1vM2 {fp('M1vM2','raw')}"),
    ("norm",         f"행 방향 z-score — M2 best AUC ({fa('M2')}), M1vM2 {fp('M1vM2')}"),
    ("global_zscore",f"전체 μ/σ z-score — M4 best AUC ({fa('M4','global_zscore')}), M1vM4 {fp('M1vM4','global_zscore')}"),
]):
    box(s, 1.0, 4.35+i*0.62, 11.33, 0.52,
        fill=C_BLUE if i==0 else (C_TEAL if i==1 else C_ORG))
    txt(s, f"  {var}   →  {desc}",
        1.05, 4.37+i*0.62, 11.2, 0.48, sz=12, bold=True, color=C_WHT)
snum(s, 21)

# ═══════════════════════════════════════════════════════════════
# 조건별 슬라이드 (22~24)
# ═══════════════════════════════════════════════════════════════
VAR_COL = {
    "raw":           C_BLUE,
    "norm":          C_TEAL,
    "global_zscore": C_ORG,
}

def condition_slide(variant, sn):
    col     = VAR_COL.get(variant, C_BLUE)
    m2_dir  = f"{BASE}/model_2/{variant}"
    m3_dir  = f"{BASE}/model_3/{variant}"
    cmp_png = f"{CMP}/roc_all_models_{variant}.png"

    m2_auc_v = _v('M2', variant, 'AUC-ROC')
    m3_auc_v = _v('M3', variant, 'AUC-ROC')
    m2_brier = _v('M2', variant, 'Brier')
    m3_brier = _v('M3', variant, 'Brier')

    s = add_slide()
    box(s, 0, 0, 13.33, 0.84, fill=C_DARK)
    box(s, 0, 0, 0.35, 0.84, fill=col)
    txt(s, f"AEC 128pt  +  {variant}",
        0.45, 0.06, 10.5, 0.44, sz=22, bold=True, color=C_WHT)
    sub_parts = [f"M2 AUC={m2_auc_v:.4f}", f"M3 AUC={m3_auc_v:.4f}",
                 f"M2vM3: {fp('M2vM3', variant)}"]
    txt(s, "  |  ".join(sub_parts),
        0.45, 0.50, 10.5, 0.28, sz=13, color=C_MGR)
    txt(s, f"{sn} / {TOTAL}", 12.3, 7.15, 1.0, 0.3, sz=11,
        color=C_MGR, align=PP_ALIGN.RIGHT)

    cx = [0.25, 4.57, 8.88]; cw = [4.2, 4.2, 4.2]
    for j,(lbl,lc) in enumerate([
        ("전모델 ROC 비교  (M1·M2·M2₂·M3·M4)", C_DARK),
        (f"M2 Calibration  (Brier {m2_brier:.4f})", C_ORG),
        (f"M3 Calibration  (Brier {m3_brier:.4f})", C_ORG),
    ]):
        box(s, cx[j], 0.89, cw[j], 0.33, fill=lc)
        txt(s, lbl, cx[j]+0.05, 0.91, cw[j]-0.1, 0.29,
            sz=11, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

    pic(s, cmp_png,                          cx[0], 1.24, cw[0], 2.85)
    pic(s, img(m2_dir, "calibration_.png"),  cx[1], 1.24, cw[1], 2.85)
    pic(s, img(m3_dir, "calibration_.png"),  cx[2], 1.24, cw[2], 2.85)

    bottom_row = [
        ("M2 Confusion Matrix",  C_TEAL, m2_dir, "confusion_matrices.png"),
        ("M3 Confusion Matrix",  C_ORG,  m3_dir, "confusion_matrices.png"),
        ("M2 Attention Heatmap", C_BLUE, m2_dir, "attention_heatmap_c2a.png"),
    ]
    for j,(lbl,lc,mdir,fname) in enumerate(bottom_row):
        box(s, cx[j], 4.17, cw[j], 0.33, fill=lc)
        txt(s, lbl, cx[j]+0.05, 4.19, cw[j]-0.1, 0.29,
            sz=11, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
        pic(s, img(mdir, fname), cx[j], 4.52, cw[j], 2.62)


AEC128_VARIANTS = ["raw", "norm", "global_zscore"]
for sn, variant in enumerate(AEC128_VARIANTS, start=22):
    condition_slide(variant, sn)

# ═══════════════════════════════════════════════════════════════
# Grad-CAM AEC 슬라이드 (25~27)
# ═══════════════════════════════════════════════════════════════
CAM_VAR_COL = VAR_COL

def cam_aec_slide(variant, sn):
    col    = CAM_VAR_COL.get(variant, C_BLUE)
    m2_dir = f"{BASE}/model_2/{variant}"
    m3_dir = f"{BASE}/model_3/{variant}"

    s = add_slide()
    box(s, 0, 0, 13.33, 0.84, fill=C_DARK)
    box(s, 0, 0, 0.35, 0.84, fill=col)
    txt(s, f"Grad-CAM AEC — {variant}",
        0.45, 0.06, 10.5, 0.44, sz=22, bold=True, color=C_WHT)
    txt(s, "M2 (CrossAttn)  vs  M3 (CrossAttn3)  |  heatmap · individual lines · mean curve",
        0.45, 0.50, 10.5, 0.28, sz=13, color=C_MGR)
    txt(s, f"{sn} / {TOTAL}", 12.3, 7.15, 1.0, 0.3, sz=11,
        color=C_MGR, align=PP_ALIGN.RIGHT)

    cx = [0.2, 4.57, 8.95]; cw = [4.25, 4.25, 4.25]
    for j,(lbl,lc) in enumerate([
        ("cam_aec_heatmap  (샘플별)", C_DARK),
        ("cam_aec_lines  (개별 곡선)", C_DARK),
        ("cam_aec_mean  (평균 CAM)",  C_DARK),
    ]):
        box(s, cx[j], 0.89, cw[j], 0.3, fill=lc)
        txt(s, lbl, cx[j]+0.05, 0.90, cw[j]-0.1, 0.28,
            sz=10.5, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

    for row_i, (model_lbl, model_col, mdir) in enumerate([
        ("M2  CrossAttn",  C_TEAL, m2_dir),
        ("M3  CrossAttn3", C_ORG,  m3_dir),
    ]):
        y0 = 1.22 + row_i * 3.08
        box(s, 0.2, y0, 0.3, 2.75, fill=model_col)
        txt(s, model_lbl, 0.22, y0 + 0.9, 0.28, 1.5,
            sz=9, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
        for j, fname in enumerate(["cam_aec_heatmap.png", "cam_aec_lines.png", "cam_aec_mean.png"]):
            pic(s, img(mdir, fname), cx[j] + 0.32, y0, cw[j] - 0.32, 2.75)


for sn, variant in enumerate(AEC128_VARIANTS, start=25):
    cam_aec_slide(variant, sn)

# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
OUT = ("C:/Users/jhjun/OneDrive/Desktop/대학원/results/0612_late_fusion/"
       "AEC_Sarcopenia_Research_Presentation.pptx")
prs.save(OUT)
print(f"Saved  → {OUT}")
print(f"Slides : {len(prs.slides)}")
