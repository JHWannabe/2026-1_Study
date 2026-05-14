"""
260514_AEC_SMI_Analysis.pptx 보완 스크립트
실행: python make_ppt_v2.py
출력: 연구자료/260515_AEC_SMI_Analysis_v2.pptx (45슬라이드)
"""

import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = r'c:\Users\jhjun\OneDrive\Desktop\2026-1_Study\연구자료\260514_AEC_SMI_Analysis.pptx'
DST = r'c:\Users\jhjun\OneDrive\Desktop\2026-1_Study\연구자료\260515_AEC_SMI_Analysis_v2.pptx'

# ─── 색상 ──────────────────────────────────────────────────────
HDR_COLOR  = RGBColor(0x1F, 0x4E, 0x79)  # 진한 파랑 (헤더 배경)
HDR_FG     = RGBColor(0xFF, 0xFF, 0xFF)  # 흰색 (헤더 글자)
ROW_A      = RGBColor(0xD6, 0xE4, 0xF0)  # 연한 파랑 (짝수 행)
ROW_B      = RGBColor(0xFF, 0xFF, 0xFF)  # 흰색 (홀수 행)
STAR_BG    = RGBColor(0xFF, 0xF2, 0xCC)  # 노랑 (강조 행)
WARN_BG    = RGBColor(0xFF, 0xE0, 0xB2)  # 주황 (경고 행)
GREEN_BG   = RGBColor(0xE8, 0xF5, 0xE9)  # 연두 (유의 행)

# ─── 슬라이드 레이아웃 상수 ─────────────────────────────────────
TITLE_ONLY = 6   # 제목만 있는 레이아웃 인덱스 (보통 5 또는 6)

# ─── 유틸 ──────────────────────────────────────────────────────

def _rgb_to_hex(rgb: RGBColor) -> str:
    return f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'


def _set_cell_bg(cell, color: RGBColor):
    """셀 배경색 설정."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # 기존 solidFill 제거
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', _rgb_to_hex(color))


def _set_cell_text(cell, text: str, bold: bool = False, font_size_pt: int = 10,
                   font_color=None, align=PP_ALIGN.CENTER):
    """셀 텍스트·서식 설정."""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color


def add_table_slide(prs, layout_idx,
                    title: str, slide_num: str,
                    header: list, rows: list,
                    star_row=None,
                    warn_fn=None,
                    green_fn=None,
                    font_size: int = 10,
                    note: str = ''):
    """
    깨끗한 새 슬라이드를 만들고 제목·번호·테이블을 추가한다.
    star_row: 0-based 데이터 행 인덱스 (노란 강조)
    warn_fn : row→bool, True이면 주황 배경
    green_fn: row→bool, True이면 연두 배경
    """
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # 레이아웃이 제공하는 placeholder 정리
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)

    W = prs.slide_width
    H = prs.slide_height

    # ── 제목 텍스트 박스 ────────────────────────────────────────
    txBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.3),
        W - Inches(0.8), Inches(0.45)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '3. Detailed progress'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # ── 슬라이드 번호 ───────────────────────────────────────────
    numBox = slide.shapes.add_textbox(
        W - Inches(0.55), H - Inches(0.45),
        Inches(0.45), Inches(0.35)
    )
    tf2 = numBox.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = slide_num
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # ── 슬라이드 제목 (내용) ────────────────────────────────────
    titleBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.75),
        W - Inches(0.8), Inches(0.5)
    )
    tf3 = titleBox.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = title
    run3.font.size = Pt(14)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x26, 0x4F, 0x78)

    if note:
        noteBox = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.2),
            W - Inches(0.8), Inches(0.3)
        )
        tf_n = noteBox.text_frame
        p_n = tf_n.paragraphs[0]
        run_n = p_n.add_run()
        run_n.text = note
        run_n.font.size = Pt(10)
        run_n.font.italic = True
        run_n.font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    # ── 테이블 ──────────────────────────────────────────────────
    top_offset = Inches(1.45) if note else Inches(1.2)
    tbl_h = H - top_offset - Inches(0.5)
    tbl_w = W - Inches(0.8)

    n_rows = len(rows) + 1   # 헤더 포함
    n_cols = len(header)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                       Inches(0.4), top_offset,
                                       tbl_w, tbl_h)
    tbl = tbl_shape.table

    # 열 너비 균등 배분
    col_w = int(tbl_w / n_cols)
    for col in tbl.columns:
        col.width = col_w

    # 헤더 행
    for ci, val in enumerate(header):
        cell = tbl.cell(0, ci)
        _set_cell_bg(cell, HDR_COLOR)
        _set_cell_text(cell, val, bold=True,
                       font_size_pt=font_size,
                       font_color=HDR_FG)

    # 데이터 행
    for ri, row in enumerate(rows):
        # 배경색 결정
        if star_row is not None and ri == star_row:
            bg = STAR_BG
        elif warn_fn is not None and warn_fn(row):
            bg = WARN_BG
        elif green_fn is not None and green_fn(row):
            bg = GREEN_BG
        else:
            bg = ROW_A if ri % 2 == 0 else ROW_B

        for ci, val in enumerate(row):
            if ci >= n_cols:
                break
            cell = tbl.cell(ri + 1, ci)
            _set_cell_bg(cell, bg)
            _bold = (ci == 0) or (star_row is not None and ri == star_row)
            _set_cell_text(cell, val, bold=_bold, font_size_pt=font_size)

    return slide


def add_text_slide(prs, layout_idx, title: str, slide_num: str, lines: list):
    """텍스트 전용 슬라이드를 만든다."""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)

    W = prs.slide_width
    H = prs.slide_height

    # 제목
    txBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.3),
        W - Inches(0.8), Inches(0.45)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '3. Detailed progress'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = HDR_COLOR

    # 슬라이드 번호
    numBox = slide.shapes.add_textbox(
        W - Inches(0.55), H - Inches(0.45),
        Inches(0.45), Inches(0.35)
    )
    tf2 = numBox.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = slide_num
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # 내용 텍스트 박스
    bodyBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.85),
        W - Inches(1.0), H - Inches(1.3)
    )
    tf3 = bodyBox.text_frame
    tf3.word_wrap = True

    first = True
    for line in lines:
        if first:
            para = tf3.paragraphs[0]
            first = False
        else:
            para = tf3.add_paragraph()

        is_header = any(line.startswith(c) for c in ('①', '②', '③', '④', '⑤', '⑥'))
        is_title  = line == lines[0]

        run = para.add_run()
        run.text = line
        run.font.bold = is_header or is_title
        run.font.size = Pt(13 if is_title else (11 if is_header else 10))
        if is_title:
            run.font.color.rgb = HDR_COLOR
        elif is_header:
            run.font.color.rgb = RGBColor(0x1A, 0x5E, 0x91)

    return slide


# ─── 슬라이드 복사 (원본 유지용 — 기존 슬라이드만 재사용) ─────────
def clone_slide_safe(prs, slide_idx):
    """
    기존 슬라이드를 안전하게 복제한다.
    shape ID를 슬라이드 내에서 고유하게 재부여하여 PowerPoint 오류를 방지한다.
    """
    src = prs.slides[slide_idx]
    layout = src.slide_layout

    new_slide = prs.slides.add_slide(layout)

    # 레이아웃에서 자동 추가된 placeholder 제거
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)

    # 원본 shape 복사 + ID 재부여
    cloned = []
    for child in src.shapes._spTree:
        elem = copy.deepcopy(child)
        cloned.append(elem)

    # 슬라이드 내 최대 ID 파악 후 고유 ID 부여
    next_id = 1000
    for elem in cloned:
        for cNvPr in elem.iter(qn('p:cNvPr')):
            cNvPr.set('id', str(next_id))
            next_id += 1
        sp_tree.append(elem)

    return new_slide


# ─── 데이터 ────────────────────────────────────────────────────

CV_HEADER   = ['조건', 'CV AUC (mean±std)', 'CV AUPRC', 'CV Brier', 'CV F1']
TEST_HEADER = ['조건', 'AUC-ROC', 'AUPRC', 'Brier', 'Accuracy', 'F1', 'Sens.', 'Spec.']

NORM_CV = [
    ['M2 | Clinic',   '0.8047 ± 0.0633', '0.3868', '0.1988', '0.3709'],
    ['M2 | Both',     '0.7984 ± 0.0676', '0.3758', '0.1898', '0.3654'],
    ['M2_2 | Clinic', '0.8082 ± 0.0757', '0.4050', '0.1939', '0.3580'],
    ['M2_2 | Both',   '0.8150 ± 0.0600', '0.3376', '0.2468', '0.3490'],
    ['M3 | Clinic',   '0.7991 ± 0.0743', '0.4093', '0.1865', '0.3717'],
    ['M3 | Both',     '0.7957 ± 0.0723', '0.3899', '0.2210', '0.3471'],
]
NORM_TEST = [
    ['M2 | Clinic',   '0.8290', '0.3217', '0.1741', '0.7373', '0.3853', '0.750', '0.736'],
    ['M2 | Both',     '0.8021', '0.3010', '0.2031', '0.6902', '0.3577', '0.786', '0.678'],
    ['M2_2 | Clinic', '0.8038', '0.3033', '0.2126', '0.6784', '0.3492', '0.786', '0.665'],
    ['M2_2 | Both',   '0.8150', '0.3376', '0.2468', '0.6196', '0.3490', '0.821', '0.601'],
    ['M3 | Clinic',   '0.8059', '0.3199', '0.1656', '0.7686', '0.4040', '0.714', '0.775'],
    ['M3 | Both',     '0.8117', '0.3233', '0.1689', '0.7765', '0.4000', '0.714', '0.800'],
]

EXCL_CV = [
    ['M2 | Clinic',   '0.8122 ± 0.0445', '0.4011', '0.1876', '0.3655'],
    ['M2 | Both',     '0.8098 ± 0.0500', '0.3296', '0.2402', '0.3030'],
    ['M2_2 | Clinic', '0.7850 ± 0.0375', '0.3575', '0.1833', '0.3378'],
    ['M2_2 | Both',   '0.7783 ± 0.0450', '0.3129', '0.2252', '0.3361'],
    ['M3 | Clinic',   '0.7988 ± 0.0503', '0.3897', '0.1821', '0.3867'],
    ['M3 | Both',     '0.7990 ± 0.0480', '0.3125', '0.1678', '0.4086'],
]
EXCL_TEST = [
    ['M2 | Clinic',   '0.8325', '0.4029', '0.1696', '0.7293', '0.3673', '0.750', '0.727'],
    ['M2 | Both',     '0.8098', '0.3296', '0.2402', '0.5983', '0.3030', '0.750', '0.587'],
    ['M2_2 | Clinic', '0.8213', '0.3113', '0.1838', '0.7249', '0.4112', '0.846', '0.709'],
    ['M2_2 | Both',   '0.7783', '0.3129', '0.2252', '0.6550', '0.3361', '0.692', '0.655'],
    ['M3 | Clinic',   '0.8089', '0.3427', '0.1719', '0.7555', '0.4043', '0.792', '0.751'],
    ['M3 | Both',     '0.7990', '0.3125', '0.1678', '0.7598', '0.4086', '0.833', '0.747'],
]

BEST_HEADER = ['모델', '서브모델', 'Best Case', 'AUC', 'AUPRC', 'Brier', 'Acc', 'F1']
BEST_DATA = [
    ['M1',   'LR',         'scale_clinic',         '0.8291', '0.3249', '0.1931', '0.698', '0.364'],
    ['M1',   'ResNet1D',   'scale_clinic',         '0.7955', '0.3471', '0.1848', '0.753', '0.388'],
    ['M2',   'LR',         'len064/scale_both',    '0.8147', '0.3195', '0.1846', '0.741', '0.400'],
    ['M2',   'CrossAttn',  'len064/scale_both',    '0.8403', '0.3367', '0.2253', '0.678', '0.379'],
    ['M2',   'ResNet1D',   'len064/scale_both',    '0.8394', '0.3445', '0.1869', '0.682', '0.352'],
    ['M2_2', 'LR',         'norm/scale_clinic',    '0.8357', '0.3122', '0.1885', '0.718', '0.390'],
    ['M2_2', 'CrossAttn',  'excl/scale_clinic',    '0.8213', '0.3113', '0.1838', '0.725', '0.411'],
    ['M2_2', 'ResNet1D',   'norm/scale_both',      '0.8219', '0.3356', '0.1806', '0.718', '0.390'],
    ['M3',   'LR',         'len064/scale_both',    '0.8164', '0.3305', '0.1856', '0.737', '0.396'],
    ['M3',   'CrossAttn3', 'len128/scale_both [*]','0.8453', '0.3993', '0.2036', '0.706', '0.380'],
    ['M3',   'ResNet1D',   'len128/scale_both',    '0.8302', '0.3600', '0.1421', '0.792', '0.391'],
]

MATCH_HEADER = ['조건 (CrossAttn)', 'M2 Matched', 'M2_2 Unmatched', 'AUC 차이']
MATCH_DATA = [
    ['len064 / scale_clinic', '0.8140', '0.8037', '-0.010'],
    ['len064 / scale_both',   '0.8403', '0.7683', '-0.072'],
    ['len128 / scale_clinic', '0.8287', '0.7989', '-0.030'],
    ['len128 / scale_both',   '0.8310', '0.8151', '-0.016'],
    ['norm / scale_clinic',   '0.8290', '0.8038', '-0.025'],
    ['norm / scale_both',     '0.8021', '0.8150', '+0.013'],
]

STAT_HEADER = ['비교', 'Metric', 'delta Mean', 't p-val', 'W p-val', 'sig.']
STAT_DATA = [
    ['M1: LR vs ResNet1D',              'F1',  '+0.039', '0.062', '0.063', 'dagger'],
    ['M2 LR vs CrossAttn (len064/cli)', 'AUC', '+0.146', '0.015', '0.063', '*'],
    ['M2 LR vs CrossAttn (len064/cli)', 'F1',  '+0.086', '0.037', '0.063', '*'],
    ['M2 LR vs CrossAttn (len256/cli)', 'AUC', '+0.232', '0.011', '0.063', '*'],
    ['M2 LR vs CrossAttn (crop60/cli)', 'AUC', '+0.145', '0.004', '0.063', '**'],
    ['M3 LR vs CrossAttn (len128/cli)', 'AUC', '+0.180', '0.010', '0.063', '*'],
    ['M3 LR vs CrossAttn (len128/cli)', 'F1',  '+0.106', '0.001', '0.063', '**'],
    ['M3 LR vs CrossAttn (len256/cli)', 'AUC', '+0.220', '0.008', '0.063', '**'],
    ['M3 LR vs CrossAttn (crop60/cli)', 'AUC', '+0.154', '0.001', '0.063', '**'],
]

RATIONALE_LINES = [
    'Why Cross Attention? — Concat+ResNet1D 대비 Cross Attention 사용 근거',
    '',
    '① Modal Heterogeneity (모달 이종성)',
    '   임상변수 (Age, Sex, BMI): 3개 스칼라 / AEC: 64~256개 시계열 포인트.',
    '   이질적 두 모달을 단순 concat하면 모델이 동질 피처로 처리 → 구조적 차이 소실.',
    '',
    '② Dimension Imbalance (차원 불균형)',
    '   concat 시 AEC(64~256차원)가 임상(3차원)을 압도 → 임상 변수의 영향력 희석.',
    '   CrossAttn은 각 모달을 독립 토큰 공간(d_model)으로 투영해 균형 있게 학습.',
    '',
    '③ Bidirectional Selective Interaction (양방향 선택적 상호작용)',
    '   Clinic→AEC: "이 환자에서 AEC 커브의 어느 구간이 중요한가?" 를 학습.',
    '   AEC→Clinic: "이 AEC 패턴은 어떤 임상 변수와 연관되는가?" 를 학습.',
    '   단순 concat 구조는 이 방향성을 명시적으로 모델링할 수 없음.',
    '',
    '④ Temporal Structure Preservation (시간 구조 보존)',
    '   ResNet1DEncoder: AEC를 Conv1D로 인코딩 후 Attention에 공급 → 국소 패턴 학습.',
    '   단순 concat+ResNet1D: AEC 포인트를 독립 피처로 처리, 시계열 순서 무시.',
    '',
    '⑤ Interpretability (해석가능성) — 논문 핵심 강점',
    '   Attention weight → AEC 어느 시점이 예측에 기여하는지 Heatmap 시각화 가능.',
    '   임상적 근거(clinical rationale) 제시 → 논문 설명력 및 리뷰어 설득력.',
    '',
    '⑥ 실험적 검증 (M2 Best: len064/scale_both)',
    '   LR (concat)       : Test AUC 0.8147, Sensitivity 78.6%',
    '   ResNet1D (concat) : Test AUC 0.8394, Sensitivity 82.1%',
    '   CrossAttn         : Test AUC 0.8403, Sensitivity 89.3%  <- 최고 민감도',
    '   근감소증 탐지(높은 Sensitivity)라는 임상 목적에서 CrossAttn이 명확히 우위.',
]


# ─── 시각화 슬라이드 ────────────────────────────────────────────

RESULTS_BASE = r'c:\Users\jhjun\OneDrive\Desktop\2026-1_Study\연구코드\results\0515'

# 각 변형 조건에서 scale_both를 우선, 없으면 scale_clinic 사용
AEC_VARIANTS = ['len064', 'len128', 'len256', 'crop60', 'crop80', 'norm', 'excl_extreme']

# 슬라이드에 넣을 이미지 조합 (파일명, 레이블)
IMG_TYPES = [
    ('test_roc_curves.png',    'Test ROC'),
    ('test_roc_by_sex.png',    'ROC by Sex'),
    ('confusion_matrices.png', 'Confusion'),
]


def _img_path(model: str, variant: str, scale: str, fname: str):
    """results 디렉토리에서 이미지 절대 경로를 반환한다. 없으면 None."""
    import os
    if model == 'model_1':
        p = os.path.join(RESULTS_BASE, model, scale, fname)
    else:
        p = os.path.join(RESULTS_BASE, model, variant, scale, fname)
    return p if os.path.exists(p) else None


def _best_scale(model: str, variant: str, fname: str):
    """scale_both 우선, 없으면 scale_clinic 경로 반환."""
    p = _img_path(model, variant, 'scale_both', fname)
    if p:
        return p, 'scale_both'
    p = _img_path(model, variant, 'scale_clinic', fname)
    if p:
        return p, 'scale_clinic'
    return None, None


def add_image_slide(prs, layout_idx, title: str, slide_num: str,
                    img_entries: list):
    """
    img_entries: [(img_path, label), ...] 최대 4개
    이미지를 2열 격자로 배치한다.
    """
    import os
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)

    W = prs.slide_width
    H = prs.slide_height

    # 제목
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = '3. Detailed progress'
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = HDR_COLOR

    # 슬라이드 부제목
    tb2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.58), W - Inches(0.8), Inches(0.38))
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = title
    r2.font.size = Pt(13); r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x26, 0x4F, 0x78)

    # 슬라이드 번호
    nb = slide.shapes.add_textbox(W - Inches(0.55), H - Inches(0.38), Inches(0.45), Inches(0.3))
    p3 = nb.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    r3 = p3.add_run()
    r3.text = slide_num
    r3.font.size = Pt(11)
    r3.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    # 이미지 배치 (최대 4개, 2×2 격자)
    valid = [(p, l) for p, l in img_entries if p and os.path.exists(p)]
    n = len(valid)
    if n == 0:
        return slide

    cols = 2
    rows = (n + cols - 1) // cols
    margin_left = Inches(0.3)
    margin_top  = Inches(0.98)
    avail_w = W - Inches(0.6)
    avail_h = H - margin_top - Inches(0.4)
    cell_w  = avail_w / cols
    cell_h  = avail_h / rows

    for idx, (img_path, label) in enumerate(valid):
        col = idx % cols
        row = idx // cols
        img_l = margin_left + col * cell_w + Inches(0.05)
        img_t = margin_top  + row * cell_h + Inches(0.22)
        img_w = cell_w - Inches(0.1)
        img_h = cell_h - Inches(0.28)
        try:
            slide.shapes.add_picture(img_path, img_l, img_t, img_w, img_h)
        except Exception as e:
            print(f'  이미지 추가 실패: {img_path} — {e}')
            continue
        # 레이블
        lbl = slide.shapes.add_textbox(img_l, margin_top + row * cell_h,
                                        img_w, Inches(0.2))
        p_l = lbl.text_frame.paragraphs[0]
        p_l.alignment = PP_ALIGN.CENTER
        r_l = p_l.add_run()
        r_l.text = label
        r_l.font.size = Pt(9); r_l.font.bold = True
        r_l.font.color.rgb = HDR_COLOR

    return slide


def add_visualization_slides(prs, layout_idx):
    """모든 변형 조건별 이미지 슬라이드를 추가한다."""
    import os

    slide_num = 43   # 기존 마지막 번호 이후부터

    # ── Model 1 (baseline) ──────────────────────────────────────
    entries = []
    for fname, label in IMG_TYPES:
        p = _img_path('model_1', '', 'scale_clinic', fname)
        entries.append((p, f'M1 | {label}'))
    # calibration 추가
    p_cal = _img_path('model_1', '', 'scale_clinic', 'calibration.png')
    entries.append((p_cal, 'M1 | Calibration'))
    add_image_slide(prs, layout_idx,
                    'Model 1 — Clinic Only (Age, Sex, BMI) | scale_clinic',
                    str(slide_num), entries)
    print(f'  슬라이드 {slide_num}: Model 1')
    slide_num += 1

    # ── Model 2 / 2_2 / 3 per variant ────────────────────────────
    for variant in AEC_VARIANTS:
        entries = []
        for model_tag, model_dir in [('M2', 'model_2'), ('M2_2', 'model_2_2'), ('M3', 'model_3')]:
            for fname, ltype in [('test_roc_curves.png', 'ROC'),
                                  ('confusion_matrices.png', 'Conf')]:
                p, scale = _best_scale(model_dir, variant, fname)
                label = f'{model_tag} | {ltype} ({scale or "-"})'
                entries.append((p, label))

        label_map = {
            'len064': 'len=64',  'len128': 'len=128', 'len256': 'len=256',
            'crop60': 'Crop 60%', 'crop80': 'Crop 80%',
            'norm': 'Z-Score Norm', 'excl_extreme': 'Excl. Extreme',
        }
        title = f'AEC Variant: {label_map.get(variant, variant)} — Test ROC & Confusion (M2 / M2_2 / M3, scale_both)'
        add_image_slide(prs, layout_idx, title, str(slide_num), entries)
        print(f'  슬라이드 {slide_num}: {variant}')
        slide_num += 1

    # ── by_sex 비교 슬라이드 (Best cases) ─────────────────────────
    best_cases = [
        ('model_1', '',         'scale_clinic', 'M1 baseline'),
        ('model_2', 'len064',   'scale_both',   'M2 len064'),
        ('model_3', 'len128',   'scale_both',   'M3 len128 [Best]'),
    ]
    entries = []
    for model_dir, variant, scale, label in best_cases:
        p = _img_path(model_dir, variant, scale, 'test_roc_by_sex.png')
        entries.append((p, f'{label} | by Sex'))
    # calibration for best
    p_cal = _img_path('model_3', 'len128', 'scale_both', 'calibration.png')
    entries.append((p_cal, 'M3 len128 | Calibration'))
    add_image_slide(prs, layout_idx,
                    'Best Models — ROC by Sex & Calibration',
                    str(slide_num), entries)
    print(f'  슬라이드 {slide_num}: best ROC by sex + calibration')
    slide_num += 1

    # ── Training curves (best model) ──────────────────────────────
    entries = []
    for model_tag, model_dir, variant, scale in [
        ('M1',  'model_1', '',       'scale_clinic'),
        ('M2',  'model_2', 'len064', 'scale_both'),
        ('M3',  'model_3', 'len128', 'scale_both'),
    ]:
        p = _img_path(model_dir, variant, scale, 'training_curves.png')
        entries.append((p, f'{model_tag} | Training'))
    # data distribution (model_1 루트)
    p_dd = os.path.join(RESULTS_BASE, 'model_1', 'data_distribution.png')
    entries.append((p_dd if os.path.exists(p_dd) else None, 'Data Distribution'))
    add_image_slide(prs, layout_idx,
                    'Training Curves (Best Models) & Data Distribution',
                    str(slide_num), entries)
    print(f'  슬라이드 {slide_num}: training curves + data dist')


# ─── 메인 ──────────────────────────────────────────────────────

def main():
    prs = Presentation(SRC)
    print(f'원본 슬라이드: {len(prs.slides)}개')

    # 레이아웃 인덱스 자동 탐색 (빈 레이아웃 선호)
    blank_layout_idx = 6
    for i, layout in enumerate(prs.slide_layouts):
        if 'blank' in layout.name.lower() or layout.name in ('빈 화면', 'Blank'):
            blank_layout_idx = i
            break

    def add_tbl(title, num, header, rows,
                star=None, warn_fn=None, green_fn=None, note='', fs=10):
        return add_table_slide(prs, blank_layout_idx,
                               title, num, header, rows,
                               star_row=star, warn_fn=warn_fn,
                               green_fn=green_fn, font_size=fs, note=note)

    # ── 원본 슬라이드 수 기록 ─────────────────────────────────────
    n_orig = len(prs.slides)

    # ── 새 슬라이드 추가 ─────────────────────────────────────────

    # Norm CV/Test
    add_tbl('AEC Z-Score Normalization — 5-fold CV (Cross Attention)', '',
            CV_HEADER, NORM_CV)
    add_tbl('AEC Z-Score Normalization — Test Set Performance', '',
            TEST_HEADER, NORM_TEST)

    # 통계 유의성 (Excl_extreme/Best/Matching은 원본에 이미 포함)
    add_tbl('통계 유의성 요약 — 5-fold Paired t-test + Wilcoxon', '',
            STAT_HEADER, STAT_DATA,
            green_fn=lambda r: r[5] in ('**', '***'),
            note='*** p<0.001  ·  ** p<0.01  ·  * p<0.05  ·  dagger p<0.10  |  연두=** 이상',
            fs=9)

    # ── 시각화 이미지 슬라이드 추가 ───────────────────────────────
    add_visualization_slides(prs, blank_layout_idx)

    # ── 슬라이드 순서 재배열 (동적 계산) ─────────────────────────
    # 원본에서 "4. Conclusion" 또는 "Conclusion" 슬라이드 위치를 찾아
    # 새 슬라이드들을 그 앞에 삽입한다.
    def find_conclusion_idx(prs):
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text
                    if '4. Conclusion' in t or ('Conclusion' in t and '최고 성능' in t):
                        return i
        return n_orig - 3   # 못 찾으면 원본 끝에서 3번째

    sldIdLst = prs.slides._sldIdLst
    all_ids  = list(sldIdLst)

    n_new    = len(all_ids) - n_orig          # 새로 추가한 슬라이드 수
    orig_ids = all_ids[:n_orig]               # 원본 슬라이드 ID 목록
    new_ids  = all_ids[n_orig:]               # 새 슬라이드 ID 목록

    # 원본에서 Conclusion 슬라이드 위치 탐색
    conclusion_idx = find_conclusion_idx(prs)
    # conclusion_idx는 전체 prs.slides 기준 → 원본 내 인덱스로 보정
    # (새 슬라이드가 아직 순서 재배열 전이므로 orig_ids 기준 사용)
    orig_before = orig_ids[:conclusion_idx]
    orig_after  = orig_ids[conclusion_idx:]

    new_order = orig_before + new_ids + orig_after

    for sid in all_ids:
        sldIdLst.remove(sid)
    for sid in new_order:
        sldIdLst.append(sid)

    print(f'최종 슬라이드: {len(prs.slides)}개')
    prs.save(DST)
    print(f'저장 완료: {DST}')


if __name__ == '__main__':
    main()
