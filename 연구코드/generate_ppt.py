"""
AEC-based Sarcopenia Prediction Research PPT
- 0529/liver_pubis 결과 기준  |  AEC 128pt 전용
- 256pt / Focal / ok+missing 제거
- AEC 스케일링 비교 슬라이드 추가
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE   = "C:/Users/jhjun/OneDrive/Desktop/2026-1_Study/연구코드/results/0529/liver_pubis"
M1     = f"{BASE}/model_1"
M2N    = f"{BASE}/model_2/aec128/norm"
M3N    = f"{BASE}/model_3/aec128/norm"
M22N   = f"{BASE}/model_2_2/aec128/norm"
CMP128 = f"{BASE}/comparison/aec128"

def img(folder, name):
    return os.path.join(folder, name)

# ── Colors ────────────────────────────────────────────────────
C_BG   = RGBColor(0xFF,0xFF,0xFF)
C_DARK = RGBColor(0x1A,0x1A,0x2E)
C_BLUE = RGBColor(0x0F,0x6B,0xBF)
C_TEAL = RGBColor(0x0F,0x9B,0x8A)
C_ORG  = RGBColor(0xE8,0x74,0x00)
C_GRN  = RGBColor(0x1E,0x8B,0x4C)
C_RED  = RGBColor(0xC0,0x39,0x2B)
C_LGR  = RGBColor(0xF0,0xF4,0xF8)
C_MGR  = RGBColor(0x9E,0xA3,0xAB)
C_WHT  = RGBColor(0xFF,0xFF,0xFF)

TOTAL = 27

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C_BG
    return s

def box(s, l, t, w, h, fill=None, line=None, lw=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw or 0.75)
    else: sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, sz=13.0, bold=False, color=C_DARK,
        align=PP_ALIGN.LEFT, italic=False, font="맑은 고딕"):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font

def pic(s, path, l, t, w, h):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        box(s, l, t, w, h, fill=C_LGR, line=C_MGR, lw=0.5)
        txt(s, f"[파일 없음]\n{os.path.basename(path)}", l+0.07, t+0.07,
            w-0.14, h-0.14, sz=9, color=C_MGR)

def hdr(s, title, sub=None):
    box(s, 0, 0, 13.33, 0.84, fill=C_DARK)
    txt(s, title, 0.35, 0.06, 11.0, 0.44, sz=25, bold=True, color=C_WHT)
    if sub: txt(s, sub, 0.35, 0.50, 11.0, 0.28, sz=13, color=C_MGR)

def hline(s, l, t, w, color=C_BLUE, h=0.04):
    box(s, l, t, w, h, fill=color)

def card(s, l, t, w, h, fill=C_LGR, line=None):
    return box(s, l, t, w, h, fill=fill, line=line, lw=0.75 if line else None)

def trow_h(s, headers, xs, ws, y, ht=0.3):
    for h, x, w in zip(headers, xs, ws):
        box(s, x, y, w, ht, fill=C_DARK)
        txt(s, h, x, y, w, ht, sz=10, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

def trow_d(s, cells, xs, ws, y, ht=0.36, fill=C_WHT, colors=None, bolds=None):
    for i,(c,x,w) in enumerate(zip(cells,xs,ws)):
        cl = (colors[i] if colors else C_DARK)
        b  = (bolds[i]  if bolds  else False)
        box(s, x, y, w, ht, fill=fill, line=RGBColor(0xCC,0xCC,0xCC), lw=0.5)
        txt(s, c, x, y, w, ht, sz=10.5, bold=b, color=cl, align=PP_ALIGN.CENTER)

def snum(s, n):
    txt(s, f"{n} / {TOTAL}", 12.3, 7.15, 1.0, 0.3, sz=11,
        color=C_MGR, align=PP_ALIGN.RIGHT)

def calib_banner(s, l, t, w, label):
    box(s, l, t, w, 0.33, fill=RGBColor(0xFF,0xF0,0xD0))
    txt(s, f"※ Calibration — {label}  (Brier 보정 품질 확인)",
        l+0.08, t+0.02, w-0.12, 0.29, sz=11, bold=True, color=C_ORG)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 7.5, fill=C_DARK)
box(s, 0, 0, 13.33, 0.08, fill=C_BLUE)
box(s, 0, 7.42, 13.33, 0.08, fill=C_TEAL)
txt(s, "AEC 신호 기반 근감소증 예측 모델", 1.0, 1.8, 11.33, 0.9,
    sz=36, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
txt(s, "Cross-Attention 아키텍처와 DeLong Test를 통한 모델 간 성능 비교 분석",
    1.0, 2.85, 11.33, 0.55, sz=18, color=C_MGR, align=PP_ALIGN.CENTER)
hline(s, 3.5, 3.65, 6.33, color=C_TEAL)
txt(s, "총 1,143명  |  5-Fold CV  |  Bootstrap 95% CI  |  DeLong AUC Test",
    1.0, 4.0, 11.33, 0.45, sz=14, color=C_MGR, align=PP_ALIGN.CENTER)
txt(s, "2026-05-28", 1.0, 4.55, 11.33, 0.35, sz=13, color=C_MGR, align=PP_ALIGN.CENTER)
for i,(label,col) in enumerate([
    ("Model 1  Clinic Only", C_BLUE),
    ("Model 2  Clinic + AEC", C_TEAL),
    ("Model 2_2  Neg. Control", C_MGR),
    ("Model 3  Clinic+Scanner+AEC", C_ORG),
]):
    x = 1.1 + i*2.8
    box(s, x, 5.3, 2.6, 0.85, fill=col)
    txt(s, label, x, 5.37, 2.6, 0.72, sz=12, bold=True,
        color=C_WHT, align=PP_ALIGN.CENTER)
snum(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 연구 배경 및 목적
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "연구 배경 및 목적", "AEC 신호의 진단적 가치 탐색")
card(s, 0.3, 0.99, 5.9, 5.86)
txt(s, "연구 배경", 0.5, 1.09, 5.5, 0.35, sz=16, bold=True)
hline(s, 0.5, 1.46, 5.5, color=C_BLUE)
for i,t in enumerate([
    "• 근감소증(Sarcopenia): SMI 기반 이진 분류 — 노인 허약·사망률과 직결",
    "• 기존 진단: CT 판독 + 영상의학과 전문의 수동 분석 → 시간·비용 소모",
    "• AEC(자동노출제어): CT 스캔 시 자동 수집, 체형·조직 구성 간접 반영",
    "• 추가 검사 없이 기존 CT 장비에서 수집 가능한 Zero-cost 부가 신호",
    "• 핵심 질문: AEC가 Clinic 변수 대비 얼마나 예측력을 개선하는가?",
]):
    txt(s, t, 0.5, 1.56+i*0.9, 5.6, 0.82, sz=13.5, color=C_DARK)

card(s, 6.5, 0.99, 6.5, 5.86)
txt(s, "비교 연구 목적 (DeLong Test 기반)", 6.7, 1.09, 6.1, 0.35, sz=16, bold=True)
hline(s, 6.7, 1.46, 6.1, color=C_TEAL)
for i,(q,a,col) in enumerate([
    ("① M1 vs M2: AEC 신호가 기여하는가?",
     "Clinic(LR) vs Clinic+AEC(CrossAttn)\n→ DeLong AUC 비교로 통계 검증", C_BLUE),
    ("② M1 vs M3: Scanner+AEC 조합 효과?",
     "Clinic(LR) vs Clinic+Scanner+AEC(CrossAttn3)\n→ DeLong AUC 비교", C_TEAL),
    ("③ M2 vs M3: Scanner 순증가 효과?",
     "CrossAttn vs CrossAttn3 (동일 test set)\n→ 통계적 유의차 여부 확인", C_ORG),
    ("④ M2 vs M2_2: 매칭이 실질적 기여?",
     "Matched vs Unmatched (음성 대조군)\n→ AEC 개인 대응의 필요성 검증", C_GRN),
]):
    box(s, 6.6, 1.56+i*1.25, 0.08, 1.05, fill=col)
    txt(s, q, 6.75, 1.58+i*1.25, 6.1, 0.35, sz=13, bold=True, color=col)
    txt(s, a, 6.75, 1.93+i*1.25, 6.1, 0.55, sz=12, color=C_DARK)
snum(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — 데이터셋 + data_distribution.png
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "데이터셋 및 실험 설계", "1,143명 단일 기관 후향적 코호트 / 5-Fold Stratified CV")

card(s, 0.3, 0.99, 5.9, 6.31)
txt(s, "데이터셋 구성", 0.5, 1.07, 5.5, 0.33, sz=15, bold=True)
hline(s, 0.5, 1.42, 5.5, color=C_BLUE)
hx = [0.5, 1.6, 3.1, 4.6]; hw = [1.05, 1.45, 1.45, 1.5]
trow_h(s, ["구분","남성(M)","여성(F)","합계"], hx, hw, 1.47)
for i,(row,fill) in enumerate(zip([
    ("Train","324 (Sarco 15.4%)","590 (Sarco 7.8%)","914 (Sarco 10.5%)"),
    ("Test", "82  (Sarco 15.9%)","147 (Sarco 7.5%)","229 (Sarco 10.5%)"),
], [C_WHT, C_LGR])):
    trow_d(s, row, hx, hw, 1.79+i*0.37, ht=0.37, fill=fill)
txt(s, "연령: 56.95±11.86세  |  BMI: 23.58±3.17 kg/m²  |  Train:Test = 8:2",
    0.5, 2.55, 5.5, 0.28, sz=11, color=C_MGR, italic=True)

txt(s, "AEC 전처리 조건 (5종)", 0.5, 2.97, 5.5, 0.33, sz=15, bold=True)
hline(s, 0.5, 3.32, 5.5, color=C_TEAL)
for i,(k,v) in enumerate([
    ("norm",         "행 방향 z-score 정규화 (곡선 형태만 보존)"),
    ("crop80/60",    "중앙 80% / 60% 구간만 사용 (양끝 제거)"),
    ("len128",       "시퀀스 길이 128pt 선형 보간"),
    ("excl_extreme", "scan-length 상하위 5% 극단 샘플 제외"),
]):
    txt(s, k,  0.5, 3.41+i*0.4, 1.7, 0.38, sz=12, bold=True, color=C_BLUE)
    txt(s, v,  2.2, 3.41+i*0.4, 3.9, 0.38, sz=12, color=C_DARK)
txt(s, "AEC 128pt × 5 전처리 조건 = 5 case",
    0.5, 5.05, 5.5, 0.3, sz=11, color=C_MGR, italic=True)

card(s, 6.5, 0.99, 5.9, 3.31)
txt(s, "모델 구성", 6.7, 1.07, 5.5, 0.33, sz=15, bold=True)
hline(s, 6.7, 1.42, 5.5, color=C_ORG)
for i,(t,b,col) in enumerate([
    ("Model 1 — Baseline","Age, Sex, BMI → Logistic Regression", C_BLUE),
    ("Model 2 — AEC Matched","Clinic + AEC(동일환자) → CrossAttention", C_TEAL),
    ("Model 2_2 — AEC Unmatched","Clinic + AEC(다른환자) → 음성 대조군", C_MGR),
    ("Model 3 — Scanner+AEC","Clinic + MFR + AEC → CrossAttention3", C_ORG),
]):
    box(s, 6.6, 1.51+i*1.0, 0.08, 0.82, fill=col)
    txt(s, t, 6.75, 1.53+i*1.0, 5.4, 0.33, sz=13, bold=True, color=col)
    txt(s, b, 6.75, 1.86+i*1.0, 5.4, 0.28, sz=12, color=C_DARK)

txt(s, "데이터 분포 (Train / Test / Sex / Sarcopenia)", 6.4, 4.17, 6.8, 0.28,
    sz=11, bold=True, color=C_DARK)
pic(s, img(BASE, "data_distribution.png"), 6.4, 4.47, 6.8, 2.83)
snum(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 통계 분석 방법
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "통계 분석 방법", "모델 비교를 위한 3단계 검증 체계")
for i,(col,title,body) in enumerate([
    (C_BLUE, "① DeLong AUC Test  (Test Set)",
     "동일 Test set 두 모델의 AUROC 직접 비교\n"
     "DeLong et al. (1988) — 쌍별 공분산 구조 활용\n"
     "H₀: AUROC_A = AUROC_B  →  z-통계량 기반 검정\n"
     "excl_extreme 조건: test set 크기 불일치로 M1 비교 제외"),
    (C_TEAL, "② Bootstrap 95% CI  (Test Set)",
     "n_boot = 2,000회 복원 추출 기반 신뢰구간\n"
     "2.5th–97.5th percentile 구간 보고\n"
     "AUROC / AUPRC / Brier / Accuracy / F1 전 지표\n"
     "단일 추정치의 불확실성 정량화"),
    (C_ORG, "③ Fold-level Paired Tests  (CV)",
     "5-Fold 교차검증 값에 대한 쌍별 검정\n"
     "Paired t-test + Wilcoxon signed-rank (n=5)\n"
     "두 검정 모두 보고 (p-value 쌍 제시)\n"
     "M1↔M2/M3: train set 차이 주의"),
]):
    card(s, 0.3+i*4.35, 0.99, 4.15, 5.72)
    box(s, 0.3+i*4.35, 0.99, 4.15, 0.5, fill=col)
    txt(s, title, 0.4+i*4.35, 1.02, 3.95, 0.45, sz=15, bold=True, color=C_WHT)
    y = 1.64
    for line in body.split("\n"):
        txt(s, f"  {line}", 0.4+i*4.35, y, 3.9, 0.44, sz=13, color=C_DARK)
        y += 0.47
box(s, 0.3, 6.37, 12.7, 0.52, fill=RGBColor(0xEE,0xF5,0xFF))
txt(s, "유의수준:  *** p<0.001   ** p<0.01   * p<0.05   † p<0.10   ns p≥0.10",
    0.5, 6.40, 12.3, 0.42, sz=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
snum(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 전체 모델 성능 비교 개요
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "전체 모델 성능 비교 개요  (AEC 128pt + norm)",
    "M1 → M2 → M3 순차 비교  |  Bootstrap 95% CI  |  Brier → Calibration 참조")

hx = [0.3,1.9,3.05,4.6,5.75,6.85,7.95,9.65,11.35]
hw = [1.55,1.1,1.5,1.1,1.05,1.05,1.65,1.65,1.85]
trow_h(s, ["모델","AUC","95% CI AUC","AUPRC","Brier","F1",
           "vs M1 Δ AUC","vs M1 DeLong p","비고"],
       hx, hw, 0.99, ht=0.32)
rows_ov = [
    ("M1 (LR)","0.8030","[0.720, 0.872]","0.3123","0.1913","0.3725",
     "—","기준선","Clinic Only"),
    ("M2 (CrossAttn)","0.8030","[0.719, 0.873]","0.2866","0.1930","0.3559",
     "+0.000","p=1.000 ns","Clinic+AEC"),
    ("M2_2 (Neg.Ctrl)","0.7955","[0.706, 0.871]","0.3449","0.2033","0.3542",
     "−0.008","(neg ctrl)","Unmatched"),
    ("M3 (CrossAttn3)","0.8207","[0.744, 0.889]","0.3174","0.1651","0.3738",
     "+0.018","p=0.502 ns","Clinic+Scn+AEC"),
]
fills = [C_LGR,RGBColor(0xE8,0xF4,0xE8),RGBColor(0xF5,0xF5,0xF5),RGBColor(0xFF,0xF3,0xE0)]
sigs  = [C_MGR, C_MGR, C_MGR, C_MGR]
for i,(row,fill,sc) in enumerate(zip(rows_ov,fills,sigs)):
    table_c = [C_DARK]*6 + [C_BLUE, sc, C_MGR]
    trow_d(s, row, hx, hw, 1.33+i*0.4, ht=0.4, fill=fill, colors=table_c,
           bolds=[i==1]*9)

txt(s, "ROC 비교 (M1·M2·M2₂·M3 / norm)", 0.3, 3.87, 4.5, 0.3, sz=12, bold=True)
pic(s, img(CMP128, "roc_all_models_norm.png"), 0.3, 4.19, 4.5, 3.11)

card(s, 5.0, 3.87, 4.1, 3.43, fill=RGBColor(0xE8,0xF8,0xE8))
box(s, 5.0, 3.87, 0.1, 3.43, fill=C_GRN)
txt(s, "비교 핵심", 5.2, 3.95, 3.7, 0.32, sz=14, bold=True, color=C_GRN)
for i,t in enumerate([
    "• M1→M2: Δ AUC 0.000, DeLong p=1.000 (ns)",
    "• M1→M3: Δ AUC +0.018, DeLong p=0.502 (ns)",
    "• M2→M3: Δ AUC +0.018, p=0.234 (ns)",
    "• AUPRC: M1(0.312)→M2(0.287) −7.9%",
    "• M2 > M2_2 excl_extreme: Δ=+0.294 ***",
    "• 여성 AUC: M1(0.848)→M3(0.904) 큰 개선",
]):
    txt(s, t, 5.2, 4.36+i*0.45, 3.7, 0.42, sz=12, color=C_DARK)

calib_banner(s, 9.25, 3.87, 3.85, "M3 norm  (Brier 0.1651)")
pic(s, img(M3N, "calibration.png"), 9.25, 4.22, 3.85, 3.08)
snum(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — M1 vs M2: AEC의 기여
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M1 vs M2 비교 — AEC 신호의 기여",
    "Clinic Only(LR) vs Clinic+AEC(CrossAttn)  |  DeLong AUC Test  |  Brier → Calibration 참조")

card(s, 0.3, 0.99, 6.0, 4.52)
txt(s, "DeLong Test — M1(LR) vs M2(CrossAttn)  [ AEC 128pt ]",
    0.5, 1.07, 5.6, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 5.6, color=C_BLUE)
dx = [0.45, 2.05, 3.1, 4.15, 5.2]; dw = [1.55, 1.0, 1.0, 1.0, 0.95]
trow_h(s, ["시나리오","M1 AUC","M2 AUC","Δ AUC","p-val"], dx, dw, 1.47)
dlong_m1m2 = [
    ("norm",    "0.8030","0.8030","+0.000","p=1.000 ns"),
    ("len128",  "0.8030","0.8150","+0.012","p=0.639 ns"),
    ("crop80",  "0.8030","0.8065","+0.004","p=0.884 ns"),
    ("crop60",  "0.8030","0.7927","−0.010","p=0.703 ns"),
]
sig_c_m1m2 = [C_MGR, C_MGR, C_MGR, C_MGR]
for i,(row,sc) in enumerate(zip(dlong_m1m2, sig_c_m1m2)):
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, dx, dw, 1.79+i*0.4, ht=0.4, fill=fill,
           colors=[C_DARK,C_DARK,C_BLUE,C_BLUE,sc])

txt(s, "Bootstrap 95% CI — norm 기준",
    0.5, 3.49, 5.6, 0.3, sz=12, bold=True, color=C_DARK)
hline(s, 0.5, 3.81, 5.6, color=C_TEAL, h=0.03)
bx2 = [0.45, 1.85, 3.25, 4.7]; bw2 = [1.35, 1.35, 1.4, 1.25]
trow_h(s, ["지표","M1","M2","Δ(M2-M1)"], bx2, bw2, 3.87, ht=0.28)
for i,(row,fill) in enumerate(zip([
    ("AUC-ROC","0.803 [0.720,0.872]","0.803 [0.719,0.873]","+0.000"),
    ("AUPRC",  "0.312 [0.181,0.496]","0.287 [0.177,0.474]","−0.026"),
    ("Brier",  "0.191 [0.165,0.218]","0.193 [0.164,0.224]","+0.002"),
    ("F1",     "0.373 [0.250,0.488]","0.356 [0.240,0.468]","−0.017"),
],[C_WHT,C_LGR,C_WHT,C_LGR])):
    trow_d(s, row, bx2, bw2, 4.17+i*0.33, ht=0.33, fill=fill)

txt(s, "ROC 비교  (norm / M1·M2·M2₂·M3)", 6.5, 0.99, 6.6, 0.3, sz=12, bold=True)
pic(s, img(CMP128, "roc_all_models_norm.png"), 6.5, 1.34, 6.6, 3.28)

calib_banner(s, 6.5, 4.74, 6.6, "M2  norm  (Brier 0.1930)")
pic(s, img(M2N, "calibration.png"), 6.5, 5.09, 6.6, 2.21)
snum(s, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — M1 vs M3: Scanner+AEC 효과
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M1 vs M3 비교 — Scanner + AEC 조합 효과",
    "Clinic Only(LR) vs Clinic+Scanner+AEC(CrossAttn3)  |  DeLong AUC Test  |  Brier → Calibration 참조")

card(s, 0.3, 0.99, 6.0, 4.52)
txt(s, "DeLong Test — M1(LR) vs M3(CrossAttn3)  [ AEC 128pt ]",
    0.5, 1.07, 5.6, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 5.6, color=C_ORG)
trow_h(s, ["시나리오","M1 AUC","M3 AUC","Δ AUC","p-val"], dx, dw, 1.47)
dlong_m1m3 = [
    ("norm",   "0.8030","0.8207","+0.018","p=0.502 ns"),
    ("len128", "0.8030","0.8197","+0.017","p=0.528 ns"),
    ("crop80", "0.8030","0.8112","+0.008","p=0.758 ns"),
    ("crop60", "0.8030","0.7791","−0.024","p=0.437 ns"),
]
for i,row in enumerate(dlong_m1m3):
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, dx, dw, 1.79+i*0.4, ht=0.4, fill=fill,
           colors=[C_DARK,C_DARK,C_ORG,C_BLUE,C_MGR])

txt(s, "Bootstrap 95% CI — norm 기준",
    0.5, 3.49, 5.6, 0.3, sz=12, bold=True)
hline(s, 0.5, 3.81, 5.6, color=C_ORG, h=0.03)
trow_h(s, ["지표","M1","M3","Δ(M3-M1)"], bx2, bw2, 3.87, ht=0.28)
for i,(row,fill) in enumerate(zip([
    ("AUC-ROC","0.803 [0.720,0.872]","0.821 [0.744,0.889]","+0.018"),
    ("AUPRC",  "0.312 [0.181,0.496]","0.317 [0.194,0.516]","+0.005"),
    ("Brier",  "0.191 [0.165,0.218]","0.165 [0.138,0.193]","−0.026"),
    ("F1",     "0.373 [0.250,0.488]","0.374 [0.252,0.491]","+0.001"),
],[C_WHT,C_LGR,C_WHT,C_LGR])):
    trow_d(s, row, bx2, bw2, 4.17+i*0.33, ht=0.33, fill=fill)

txt(s, "ROC 비교  (norm / M1·M2·M2₂·M3)", 6.5, 0.99, 6.6, 0.3, sz=12, bold=True)
pic(s, img(CMP128, "roc_all_models_norm.png"), 6.5, 1.34, 6.6, 3.28)

calib_banner(s, 6.5, 4.74, 6.6, "M3  norm  (Brier 0.1651)")
pic(s, img(M3N, "calibration.png"), 6.5, 5.09, 6.6, 2.21)
snum(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — M2 vs M3: Scanner 순증가 효과
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M2 vs M3 비교 — Scanner(MFR) 추가 효과",
    "CrossAttn vs CrossAttn3  |  모든 조건에서 DeLong ns — 통계적 유의차 없음")

card(s, 0.3, 0.99, 7.5, 4.31)
txt(s, "DeLong Test — M2(CrossAttn) vs M3(CrossAttn3)  [ 전 조건 ]",
    0.5, 1.07, 7.1, 0.33, sz=14, bold=True)
hline(s, 0.5, 1.42, 7.1, color=C_MGR)
mx = [0.45, 2.3, 3.45, 4.6, 5.75, 6.85]
mw = [1.8,  1.1, 1.1,  1.1, 1.05, 0.9]
trow_h(s, ["조건","M2 AUC","M3 AUC","Δ AUC","p-val","sig"], mx, mw, 1.47)
m2m3_rows = [
    ("norm",         "0.8030","0.8207","+0.018","0.234","ns"),
    ("crop60",       "0.7927","0.7791","−0.014","0.534","ns"),
    ("crop80",       "0.8065","0.8112","+0.005","0.824","ns"),
    ("len128",       "0.8150","0.8197","+0.005","0.808","ns"),
    ("excl_extreme", "0.7866","0.7846","−0.002","0.922","ns"),
]
for i,row in enumerate(m2m3_rows):
    fill = C_WHT if i%2==0 else C_LGR
    trow_d(s, row, mx, mw, 1.79+i*0.37, ht=0.37, fill=fill,
           colors=[C_DARK]*4+[C_MGR, C_MGR])

txt(s, "ROC — norm  (M2 vs M3 차이)", 8.0, 0.99, 5.1, 0.3, sz=12, bold=True)
pic(s, img(CMP128,"roc_all_models_norm.png"), 8.0, 1.31, 5.1, 2.68)

card(s, 8.0, 4.14, 5.1, 3.11, fill=RGBColor(0xF5,0xF5,0xF5))
box(s, 8.0, 4.14, 0.1, 3.11, fill=C_MGR)
txt(s, "해석", 8.2, 4.22, 4.7, 0.33, sz=14, bold=True, color=C_MGR)
for i,t in enumerate([
    "• 5개 조건 전체 ns — Scanner 통계 유의차 없음",
    "• M3 AUC ≥ M2 AUC (norm/crop80/len128 3개 조건)",
    "  → Scanner 추가 시 소폭 상승 경향 (비유의)",
    "• 예외: crop60/excl_extreme M2 > M3 (ns)",
    "• 방향성 불일치 → 검출력 부족 가능성",
    "• 추후 n↑ → M2 vs M3 재검증 권장",
]):
    txt(s, t, 8.2, 4.64+i*0.38, 4.7, 0.36, sz=11.5, color=C_DARK)
snum(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — M2 vs M2_2: 매칭 검증 (음성 대조군)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "M2 vs M2_2 비교 — 매칭 효과 검증 (음성 대조군)",
    "Matched(M2) vs Unmatched(M2_2)  |  AEC 개인 대응의 실질적 기여 확인")

card(s, 0.3, 0.99, 4.8, 1.45)
txt(s, "실험 설계", 0.5, 1.07, 4.4, 0.3, sz=14, bold=True)
hline(s, 0.5, 1.39, 4.4, color=C_TEAL)
for i,t in enumerate([
    "M2: 환자 i의 Clinic + 환자 i의 AEC (동일 환자)",
    "M2_2: 환자 i의 Clinic + 환자 j의 AEC (다른 환자)",
    "→ 동일 CrossAttention 구조 / AEC만 랜덤 셔플",
]):
    txt(s, t, 0.5, 1.49+i*0.32, 4.4, 0.3, sz=12, color=C_DARK)

card(s, 0.3, 2.57, 7.5, 3.81)
txt(s, "DeLong Test — M2(Matched) vs M2_2(Unmatched)  [ AEC 128pt ]",
    0.5, 2.65, 7.1, 0.3, sz=14, bold=True)
hline(s, 0.5, 2.97, 7.1, color=C_TEAL)
nx = [0.45, 2.5, 3.7, 5.0, 5.95, 6.95]
nw = [2.0, 1.15, 1.25, 0.9, 0.95, 0.9]
trow_h(s, ["조건","M2 AUC","M2_2 AUC","Δ AUC","p-val","sig"], nx, nw, 3.02)
neg_ctrl = [
    ("norm",         "0.8030","0.7955","+0.008","0.733","ns"),
    ("len128",       "0.8150","0.7707","+0.044","0.116","ns"),
    ("crop80",       "0.8065","0.7955","+0.011","0.684","ns"),
    ("crop60",       "0.7927","0.8083","−0.016","0.527","ns"),
    ("excl_extreme", "0.7866","0.4928","+0.294","<0.001","***"),
]
for i,row in enumerate(neg_ctrl):
    fill = C_WHT if i%2==0 else C_LGR
    sc = C_RED if row[-1] != "ns" else C_MGR
    trow_d(s, row, nx, nw, 3.34+i*0.37, ht=0.37, fill=fill,
           colors=[C_DARK]*4+[sc,sc])

txt(s, "ROC — excl_extreme (가장 극적인 차이)",
    8.0, 0.99, 5.1, 0.3, sz=12, bold=True)
pic(s, img(CMP128,"roc_all_models_excl_extreme.png"), 8.0, 1.31, 5.1, 2.55)

card(s, 8.0, 3.99, 5.1, 2.71, fill=RGBColor(0xFF,0xF3,0xE0))
box(s, 8.0, 3.99, 0.1, 2.71, fill=C_ORG)
txt(s, "해석", 8.2, 4.07, 4.7, 0.3, sz=14, bold=True, color=C_ORG)
for i,t in enumerate([
    "• excl_extreme: M2(0.787) vs M2_2(0.493) Δ=+0.294 ***",
    "  → AEC는 개인별 고유 신호 — 단순 체형 proxy 아님",
    "• norm 조건: Δ=+0.008, ns → 일반 조건은 매칭 효과 미확인",
    "• 극단 샘플 제거 후에야 AEC 고유 정보가 두드러짐",
    "  → AEC 신호의 개인 특이성은 극단 체형에서 강조",
]):
    txt(s, t, 8.2, 4.45+i*0.38, 4.7, 0.36, sz=11.5, color=C_DARK)
snum(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — AEC 전처리 조건별 ROC 비교 (5종)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "AEC 전처리 조건별 ROC 비교 — 5종",
    "M1·M2·M2₂·M3 포함  |  AEC 128pt  |  각 전처리 조건의 성능 차이 비교")

panels = [
    ("norm",         "norm — 행 방향 정규화",      C_BLUE),
    ("crop80",       "crop80 — 중앙 80% 구간",     C_TEAL),
    ("crop60",       "crop60 — 중앙 60% 구간",     C_ORG),
    ("len128",       "len128 — 선형 보간 128pt",   C_GRN),
    ("excl_extreme", "excl_extreme — 극단 제외",   C_RED),
]
positions = [(0.25, 0.99),(4.6, 0.99),(8.95, 0.99),(0.25, 4.09),(4.6, 4.09)]
for (var,label,col),(x0,y0) in zip(panels[:3], positions[:3]):
    box(s, x0, y0, 4.15, 0.36, fill=col)
    txt(s, label, x0+0.08, y0+0.02, 4.0, 0.32, sz=12, bold=True, color=C_WHT)
    pic(s, img(CMP128,f"roc_all_models_{var}.png"), x0, y0+0.38, 4.15, 3.48)
for (var,label,col),(x0,y0) in zip(panels[3:], positions[3:]):
    box(s, x0, y0, 4.15, 0.36, fill=col)
    txt(s, label, x0+0.08, y0+0.02, 4.0, 0.32, sz=12, bold=True, color=C_WHT)
    pic(s, img(CMP128,f"roc_all_models_{var}.png"), x0, y0+0.38, 4.15, 2.97)

card(s, 8.95, 4.09, 4.15, 3.35, fill=RGBColor(0xF0,0xF4,0xF8))
txt(s, "조건별 요약", 9.1, 4.17, 3.8, 0.3, sz=13, bold=True)
hline(s, 9.1, 4.49, 3.8, color=C_BLUE)
for i,t in enumerate([
    "• norm: M3 최고(0.821), 전 모델 ns",
    "• len128: M2(0.815) ≈ M3(0.820), 전 모델 ns",
    "• crop80: M2(0.807), 전 모델 ns",
    "• crop60: 전 모델 ns",
    "• excl_extreme: M2 vs M2_2 p<0.001 ***",
    "  (M2_2 AUC 0.493 — 매칭 효과 검증)",
]):
    txt(s, t, 9.1, 4.60+i*0.42, 3.8, 0.4, sz=11, color=C_DARK)
snum(s, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Attention Map 비교 (M2 vs M3)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "Attention Map 비교 — M2 vs M3",
    "AEC 128pt + norm  |  CrossAttn(M2) · CrossAttn3(M3) 어텐션 패턴 비교")

for j,(mdir,label,col) in enumerate([
    (M2N,  "Model 2  CrossAttn  (Clinic → AEC)",  C_TEAL),
    (M3N,  "Model 3  CrossAttn3  (Clinic+Scanner → AEC)", C_ORG),
]):
    x0 = 0.3 + j*6.5
    box(s, x0, 0.99, 6.3, 0.38, fill=col)
    txt(s, label, x0+0.1, 1.01, 6.1, 0.34, sz=14, bold=True, color=C_WHT)
    txt(s, "Attention Heatmap (샘플별)", x0+0.1, 1.42, 3.0, 0.25, sz=10, color=C_MGR)
    pic(s, img(mdir,"attention_heatmap_c2a.png"), x0, 1.69, 6.2, 2.75)
    txt(s, "Attention Map C→A (평균)", x0+0.1, 4.49, 3.0, 0.25, sz=10, color=C_MGR)
    pic(s, img(mdir,"attention_map_c2a.png"), x0, 4.76, 6.2, 2.59)
snum(s, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — 최적 조건 상세 비교 (norm)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "최적 조건 상세 — AEC 128pt + norm",
    "M1·M2·M3 성능 비교  |  Bootstrap CI  |  Brier → Calibration 참조")

box(s, 0.3, 0.99, 12.7, 0.65, fill=RGBColor(0xE8,0xF8,0xE8))
box(s, 0.3, 0.99, 0.12, 0.65, fill=C_GRN)
txt(s,
    "AEC 128pt + norm:  M1→M2 p=1.000 (ns) — 유의차 없음  "
    "|  M3 best overall AUC 0.821  |  excl_extreme M2 vs M2_2 p<0.001 (***)",
    0.55, 1.06, 12.1, 0.52, sz=13, bold=True, color=C_DARK)

hx7 = [0.3,1.55,2.8,4.5,5.6,6.65,7.7,9.4,11.1]
hw7 = [1.2,1.2, 1.65,1.05,1.0,1.0, 1.65,1.65,2.05]
trow_h(s, ["모델","AUC","95% CI","AUPRC","Brier","F1",
           "Δ AUC vs M1","DeLong p","비고"],
       hx7, hw7, 1.74, ht=0.3)
for i,(row,fill) in enumerate(zip([
    ("M1 (LR)","0.8030","[0.720,0.872]","0.3123","0.1913","0.3725",
     "기준선","기준선","Clinic Only"),
    ("M2 (CrossAttn)","0.8030","[0.719,0.873]","0.2866","0.1930","0.3559",
     "+0.000","p=1.000 ns","Clinic+AEC"),
    ("M2_2 (Neg)","0.7955","[0.706,0.871]","0.3449","0.2033","0.3542",
     "−0.008","(neg ctrl)","Unmatched"),
    ("M3 (CrossAttn3)","0.8207","[0.744,0.889]","0.3174","0.1651","0.3738",
     "+0.018","p=0.502 ns","Clinic+Scn"),
],[C_LGR,RGBColor(0xE8,0xF4,0xE8),C_WHT,RGBColor(0xFF,0xF3,0xE0)])):
    s6 = C_GRN if "†" in row[6] or "*" in row[6] else C_MGR
    s7 = C_GRN if "†" in row[7] or "*" in row[7] else C_MGR
    trow_d(s, row, hx7, hw7, 2.06+i*0.4, ht=0.4, fill=fill,
           colors=[C_DARK]*6+[s6,s7,C_MGR], bolds=[i==1]*9)

txt(s, "ROC — norm  (M1·M2·M2₂·M3)", 0.3, 3.79, 3.9, 0.28, sz=11, bold=True)
pic(s, img(CMP128,"roc_all_models_norm.png"), 0.3, 4.09, 3.9, 3.16)

txt(s, "M2 Confusion Matrix", 4.35, 3.79, 3.9, 0.28, sz=11, bold=True)
pic(s, img(M2N,"confusion_matrices.png"), 4.35, 4.09, 3.9, 3.16)

calib_banner(s, 8.4, 3.79, 4.7, "M3  norm  (Brier 0.1651)")
pic(s, img(M3N,"calibration.png"), 8.4, 4.14, 4.7, 3.11)
snum(s, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — 성별 분리 비교 (M1 vs M2 vs M3)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "성별 분리 비교 — M1 vs M2 vs M3  (AEC 128pt + norm)",
    "남성(n=82, 유병률 15.9%) · 여성(n=147, 유병률 7.5%) 분리 AUC 비교")

card(s, 0.3, 0.99, 12.7, 2.0)
sx = [0.4,1.9,3.25,4.65,6.05,7.4,8.8,10.2,11.6]
sw = [1.45,1.3,1.35,1.35,1.3,1.35,1.35,1.35,1.55]
trow_h(s, ["모델","전체 AUC","전체 CI","남성 AUC","남성 CI","여성 AUC","여성 CI","남 AUPRC","여 AUPRC"],
       sx, sw, 1.04, ht=0.3)
for i,(row,fill) in enumerate(zip([
    ("M1 (LR)",        "0.8030","[0.720,0.872]","0.7035","—","0.8483","—","0.3155","0.4748"),
    ("M2 (CrossAttn)", "0.8030","[0.719,0.873]","0.7269","—","0.8650","—","0.3027","0.4323"),
    ("M3 (CrossAttn3)","0.8207","[0.744,0.889]","0.7191","—","0.9037","—","0.3083","0.4549"),
],[C_LGR,RGBColor(0xE8,0xF4,0xE8),RGBColor(0xFF,0xF3,0xE0)])):
    trow_d(s, row, sx, sw, 1.36+i*0.41, ht=0.41, fill=fill, bolds=[i>0]*9)

card(s, 0.3, 3.09, 3.8, 4.06, fill=RGBColor(0xE8,0xF0,0xFF))
box(s, 0.3, 3.09, 0.1, 4.06, fill=C_BLUE)
txt(s, "남성 비교", 0.5, 3.17, 3.5, 0.32, sz=14, bold=True, color=C_BLUE)
for i,t in enumerate([
    "• M1→M2: 0.704→0.727  (Δ+0.023)",
    "• M1→M3: 0.704→0.719  (Δ+0.016)",
    "• M3 < M2 for male — Scanner 역효과?",
    "• 유병률 15.9% → AUC 해석 비교적 안정",
    "• AUPRC: M1(0.316)→M2(0.303) −3.9%",
    "  → M3(0.308)도 M1 대비 소폭 감소",
]):
    txt(s, t, 0.5, 3.59+i*0.52, 3.6, 0.5, sz=12, color=C_DARK)

card(s, 4.25, 3.09, 3.8, 4.06, fill=RGBColor(0xFF,0xF0,0xF8))
box(s, 4.25, 3.09, 0.1, 4.06, fill=C_RED)
txt(s, "여성 비교", 4.45, 3.17, 3.5, 0.32, sz=14, bold=True, color=C_RED)
for i,t in enumerate([
    "• M1→M2: 0.848→0.865  (Δ+0.017)",
    "• M1→M3: 0.848→0.904  (Δ+0.056)",
    "  → 여성에서 M3(Scanner+AEC) 기여 뚜렷",
    "• AUPRC: M1(0.475)→M3(0.455) 소폭 감소",
    "• 유병률 7.5% → class imbalance 주의",
    "  → AUPRC가 AUC보다 중요한 지표",
]):
    txt(s, t, 4.45, 3.59+i*0.52, 3.6, 0.5, sz=12, color=C_DARK)

txt(s, "M1  by Sex", 8.2, 3.09, 2.4, 0.28, sz=11, bold=True)
pic(s, img(M1,"test_roc_by_sex.png"), 8.2, 3.39, 2.4, 2.0)
txt(s, "M2  by Sex", 10.75, 3.09, 2.4, 0.28, sz=11, bold=True, color=C_TEAL)
pic(s, img(M2N,"test_roc_by_sex.png"), 10.75, 3.39, 2.4, 2.0)
txt(s, "△ 여성 AUC 0.848 → 0.904(M3)  |  △ 남성 AUC 0.704 → 0.727(M2)",
    8.2, 5.44, 5.0, 0.3, sz=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
txt(s, "M2 norm — Training Curves", 8.2, 5.78, 4.9, 0.28, sz=10, color=C_MGR)
pic(s, img(M2N,"training_curves.png"), 8.2, 6.08, 4.9, 1.22)
snum(s, 13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — AEC 스케일링 비교
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "AEC 스케일링 비교 — Raw / StandardScaler / Std+RowNorm",
    "n=1,143  |  AEC 128pt  |  3가지 전처리 버전 비교  |  xlsx: aec_inspection/")

card(s, 0.3, 0.99, 12.7, 1.75)
txt(s, "비교 목적 및 설계", 0.5, 1.07, 12.0, 0.33, sz=15, bold=True)
hline(s, 0.5, 1.42, 12.0, color=C_BLUE)
for i,t in enumerate([
    "• 모델 학습 시 실제 입력되는 AEC 데이터의 형태를 시각적으로 검증",
    "• StandardScaler (열 방향): 각 시간 포지션별로 표준화 → 학습 시 scale_aec=True 와 동일",
    "• Std+RowNorm: 열 표준화 후 각 환자별 z-score → aec_variant('norm')과 동일한 행 방향 정규화",
]):
    txt(s, t, 0.5, 1.52+i*0.36, 12.0, 0.34, sz=12.5, color=C_DARK)

sx2 = [0.3, 2.3, 5.1, 7.9, 10.7]; sw2 = [1.95, 2.75, 2.75, 2.75, 2.5]
trow_h(s, ["시트명","원본 특성","처리 방법","결과 형태","주요 목적"], sx2, sw2, 2.84, ht=0.32)
for i,(row,fill) in enumerate(zip([
    ("raw",        "DICOM 직접 추출값\n(mAs 기반 절대 선량)", "없음",
     "환자별 절대적 AEC 수준\n수치 범위 넓음",
     "원본 분포 확인\n스캐너 간 차이 관찰"),
    ("std_scaled", "raw 그대로",
     "열 방향 StandardScaler\n(포지션별 μ=0, σ=1)",
     "각 포지션의 상대적 값\n환자 간 비교 가능",
     "모델 실제 입력값 확인\n모든 AEC 모델 공통 적용"),
    ("std_norm",   "std_scaled 결과",
     "+행 방향 z-score\n(환자별 μ=0, σ=1)",
     "곡선 형태만 보존\n절대 수준 차이 제거",
     "norm variant 입력값 확인\n스캐너 간 절대값 차이 제거"),
],[C_LGR, C_WHT, RGBColor(0xE8,0xF4,0xE8)])):
    trow_d(s, row, sx2, sw2, 3.18+i*0.92, ht=0.9, fill=fill)

box(s, 0.3, 6.00, 12.7, 0.52, fill=RGBColor(0xEE,0xF5,0xFF))
txt(s,
    "저장 위치: results/0529/liver_pubis/aec_inspection/aec_scaling_compare_aec128.xlsx  "
    "|  1,143행 × 131열 (PatientID, label, sex + pos_1~pos_128)",
    0.5, 6.03, 12.3, 0.42, sz=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
snum(s, 14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — 종합 결론
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "종합 결론 및 임상적 함의",
    "AEC 신호의 근감소증 예측 기여 — 비교 연구 결과 요약")
for i,(col,title,items) in enumerate([
    (C_BLUE, "① AEC/Scanner 추가 시 AUC 소폭 향상 — 통계 유의성 미확인 (전 조건 ns)", [
        "AEC 128pt + norm: M1→M2 DeLong p=1.000 (ns) — 유의차 없음",
        "M1→M3 p=0.502 ns — Scanner 추가 유의성 미확인",
        "음성 대조군(M2_2 excl_extreme) Δ=+0.294 *** → 매칭 효과 통계 검증",
        "여성 AUC: M1(0.848)→M3(0.904) — M3(Scanner+AEC)가 여성 예측에 기여",
    ]),
    (C_TEAL, "② norm/len128 전처리 조건이 우수 — M3 norm 전체 최고 AUC (0.821)", [
        "M3 norm: AUC 0.821, Brier 0.165 — Brier 기준 최우수 calibration",
        "M2 len128: AUC 0.815 — M2 중 최고, 시퀀스 전체 활용 유리",
        "excl_extreme 조건: AUC 0.787 but M2 vs M2_2 매칭 효과 가장 극명",
        "→ 매칭 효과 연구에는 excl_extreme이 유용한 음성 대조 설계",
    ]),
    (C_ORG, "③ M2 vs M3: Scanner 추가 효과 미입증 — 5개 조건 전체 ns", [
        "norm 조건: M2(0.803) vs M3(0.821) p=0.234 ns",
        "남성 AUC에서는 M3 < M2 → Scanner 정보가 도움되지 않을 수 있음",
        "여성에서 M3(0.904) > M2(0.865) 큰 차이 → 성별별 효과 차별화 가능성",
        "→ 다기관 데이터 확보 시 M3 효과 강화 예상 (단일 기관 한계)",
    ]),
    (C_GRN, "④ Zero-cost AEC로 임상 등급 스크리닝 가능성 확인", [
        "기존 CT 스캔 시 자동 수집 — 추가 비용·검사 없음",
        "전체 AUC 0.821 / 여성 AUC 0.904(M3) → 임상 허용 가능 수준",
        "건강검진 프로그램 연동 시 고위험군 조기 선별 가능",
        "PACS-DICOM AEC 자동 추출 파이프라인 구축 필요",
    ]),
]):
    card(s, 0.3, 0.99+i*1.52, 12.7, 1.45)
    box(s, 0.3, 0.99+i*1.52, 0.12, 1.45, fill=col)
    txt(s, title, 0.55, 1.06+i*1.52, 12.1, 0.36, sz=14, bold=True, color=col)
    for k,item in enumerate(items[:2]):
        txt(s, f"  • {item}", 0.55, 1.46+i*1.52+k*0.32, 6.0, 0.3, sz=11.5, color=C_DARK)
    for k,item in enumerate(items[2:]):
        txt(s, f"  • {item}", 6.8,  1.46+i*1.52+k*0.32, 6.1, 0.3, sz=11.5, color=C_DARK)
snum(s, 15)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — 제한점 & 향후 연구
# ═══════════════════════════════════════════════════════════════
s = add_slide()
hdr(s, "제한점 및 향후 연구 방향", "현재 비교 연구의 한계와 발전 가능성")
card(s, 0.3, 0.99, 5.9, 5.91)
txt(s, "연구 제한점", 0.5, 1.07, 5.5, 0.36, sz=16, bold=True, color=C_RED)
hline(s, 0.5, 1.45, 5.5, color=C_RED)
for i,(t,b) in enumerate([
    ("소표본 비교","Test n=229, Sarco 24명 → DeLong 검출력 제한, CI 넓음"),
    ("단일 기관","단일 센터 데이터 → 외부 검증 미실시"),
    ("excl_extreme 비교 제한","Test set 크기 변화 → M1 vs M2/M3 DeLong 직접 비교 불가"),
    ("M1→M2 유의차 없음","p=1.000 (ns) — norm 조건에서 M1 대비 개선 미확인"),
    ("M2 vs M3 미유의","Scanner 순증가 효과 미확인 → 검출력 부족 가능성"),
]):
    txt(s, f"  ▶ {t}", 0.5, 1.57+i*1.0, 5.5, 0.34, sz=13, bold=True)
    txt(s, f"    {b}", 0.5, 1.91+i*1.0, 5.5, 0.58, sz=11.5, color=C_MGR)
card(s, 6.5, 0.99, 6.5, 5.91)
txt(s, "향후 연구 방향", 6.7, 1.07, 6.1, 0.36, sz=16, bold=True, color=C_BLUE)
hline(s, 6.7, 1.45, 6.1, color=C_BLUE)
for i,(t,b) in enumerate([
    ("다기관 검증","강남+타기관 통합 → 일반화 성능·DeLong 유의성 재확인"),
    ("샘플 확장","Test n>500 → CI 축소, M1→M2 유의성 강화"),
    ("전처리 설계 개선","동일 test set 유지 설계 → excl_extreme 공정 비교"),
    ("아키텍처 확장","Multi-head Attn + positional encoding → 시퀀스 표현 강화"),
    ("임상 파이프라인","PACS-DICOM AEC 추출 자동화 → 실시간 스크리닝"),
]):
    txt(s, f"  ◆ {t}", 6.7, 1.57+i*1.0, 6.1, 0.34, sz=13, bold=True, color=C_BLUE)
    txt(s, f"    {b}", 6.7, 1.91+i*1.0, 6.1, 0.58)
box(s, 0.3, 6.95, 12.7, 0.38, fill=C_DARK)
txt(s, "감사합니다  |  Questions & Discussion",
    0.3, 6.97, 12.7, 0.33, sz=14, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
snum(s, 16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — 부록 섹션 디바이더
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 7.5, fill=C_DARK)
box(s, 0, 0, 13.33, 0.08, fill=C_BLUE)
box(s, 0, 7.42, 13.33, 0.08, fill=C_TEAL)
txt(s, "부록 — 조건별 상세 결과",
    1.0, 2.3, 11.33, 0.9, sz=36, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
txt(s, "AEC 128pt × 전처리 5종  |  M2 · M2_2 · M3 개별 시각화",
    1.0, 3.35, 11.33, 0.55, sz=17, color=C_MGR, align=PP_ALIGN.CENTER)
hline(s, 3.5, 4.0, 6.33, color=C_TEAL)
for i,(var,desc) in enumerate([
    ("norm",         "행 방향 정규화 — M3 최고 AUC (0.821), Brier 최소 (0.165)"),
    ("crop80",       "중앙 80% 구간 — 양끝 10% 제거"),
    ("crop60",       "중앙 60% 구간 — 양끝 20% 제거"),
    ("len128",       "선형 보간 128pt — baseline 해상도"),
    ("excl_extreme", "상하위 5% 극단 제거 — M2 vs M2_2 매칭 효과 극명"),
]):
    box(s, 1.0, 4.35+i*0.52, 11.33, 0.42, fill=C_BLUE if i==0 else C_TEAL if i<3 else C_ORG)
    txt(s, f"  {var}   →  {desc}",
        1.05, 4.37+i*0.52, 11.2, 0.38, sz=12, bold=True, color=C_WHT)
snum(s, 17)

# ═══════════════════════════════════════════════════════════════
# 조건별 슬라이드 생성 (AEC 128pt × 5 variants)
# ═══════════════════════════════════════════════════════════════
VAR_COL = {
    "norm":         C_BLUE,
    "crop80":       C_TEAL,
    "crop60":       C_TEAL,
    "len128":       C_ORG,
    "excl_extreme": C_GRN,
}

def condition_slide(variant, sn):
    col    = VAR_COL.get(variant, C_BLUE)
    m2_dir = f"{BASE}/model_2/aec128/{variant}"
    m3_dir = f"{BASE}/model_3/aec128/{variant}"
    cmp_png = f"{CMP128}/roc_all_models_{variant}.png"

    s = add_slide()
    box(s, 0, 0, 13.33, 0.84, fill=C_DARK)
    box(s, 0, 0, 0.35, 0.84, fill=col)
    txt(s, f"AEC 128pt  +  {variant}",
        0.45, 0.06, 10.5, 0.44, sz=22, bold=True, color=C_WHT)
    txt(s, "M2 (Matched) · M3 (Clinic+Scanner+AEC)  |  Brier → Calibration 포함",
        0.45, 0.50, 10.5, 0.28, sz=13, color=C_MGR)
    txt(s, f"{sn} / {TOTAL}", 12.3, 7.15, 1.0, 0.3, sz=11,
        color=C_MGR, align=PP_ALIGN.RIGHT)

    cx = [0.25, 4.57, 8.88]; cw = [4.2, 4.2, 4.2]
    for j,(lbl,lc) in enumerate([
        ("전모델 ROC 비교  (M1·M2·M2₂·M3)", C_DARK),
        ("M2 Calibration  ※ Brier 보정 확인", C_ORG),
        ("M3 Calibration  ※ Brier 보정 확인", C_ORG),
    ]):
        box(s, cx[j], 0.89, cw[j], 0.33, fill=lc)
        txt(s, lbl, cx[j]+0.05, 0.91, cw[j]-0.1, 0.29,
            sz=11, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

    pic(s, cmp_png,                         cx[0], 1.24, cw[0], 2.85)
    pic(s, img(m2_dir, "calibration.png"),  cx[1], 1.24, cw[1], 2.85)
    pic(s, img(m3_dir, "calibration.png"),  cx[2], 1.24, cw[2], 2.85)

    if variant == "excl_extreme":
        bottom_row = [
            ("M2 Confusion Matrix",             C_TEAL, m2_dir, "confusion_matrices.png"),
            ("M3 Confusion Matrix",             C_ORG,  m3_dir, "confusion_matrices.png"),
            ("데이터 분포 (excl_extreme 필터 후)", C_RED,  m2_dir, "data_distribution.png"),
        ]
    else:
        bottom_row = [
            ("M2 Confusion Matrix",   C_TEAL, m2_dir, "confusion_matrices.png"),
            ("M3 Confusion Matrix",   C_ORG,  m3_dir, "confusion_matrices.png"),
            ("M2 Attention Heatmap",  C_BLUE, m2_dir, "attention_heatmap_c2a.png"),
        ]
    for j,(lbl,lc,mdir,fname) in enumerate(bottom_row):
        box(s, cx[j], 4.17, cw[j], 0.33, fill=lc)
        txt(s, lbl, cx[j]+0.05, 4.19, cw[j]-0.1, 0.29,
            sz=11, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
        pic(s, img(mdir, fname), cx[j], 4.52, cw[j], 2.62)


AEC128_VARIANTS = ["norm", "crop60", "crop80", "len128", "excl_extreme"]
for sn, variant in enumerate(AEC128_VARIANTS, start=18):
    condition_slide(variant, sn)

# ═══════════════════════════════════════════════════════════════
# Grad-CAM AEC 슬라이드 (variant별, M2 vs M3)
# 슬라이드 23~27
# ═══════════════════════════════════════════════════════════════
CAM_VAR_COL = {
    "norm":         C_BLUE,
    "crop80":       C_TEAL,
    "crop60":       C_TEAL,
    "len128":       C_ORG,
    "excl_extreme": C_GRN,
}

def cam_aec_slide(variant, sn):
    col    = CAM_VAR_COL.get(variant, C_BLUE)
    m2_dir = f"{BASE}/model_2/aec128/{variant}"
    m3_dir = f"{BASE}/model_3/aec128/{variant}"

    s = add_slide()
    box(s, 0, 0, 13.33, 0.84, fill=C_DARK)
    box(s, 0, 0, 0.35, 0.84, fill=col)
    txt(s, f"Grad-CAM AEC — {variant}",
        0.45, 0.06, 10.5, 0.44, sz=22, bold=True, color=C_WHT)
    txt(s, "M2 (CrossAttn)  vs  M3 (CrossAttn3)  |  heatmap · individual lines · mean curve",
        0.45, 0.50, 10.5, 0.28, sz=13, color=C_MGR)
    txt(s, f"{sn} / {TOTAL}", 12.3, 7.15, 1.0, 0.3, sz=11,
        color=C_MGR, align=PP_ALIGN.RIGHT)

    # 열 헤더 (heatmap | lines | mean)
    cx = [0.2, 4.57, 8.95]; cw = [4.25, 4.25, 4.25]
    for j,(lbl,lc) in enumerate([
        ("cam_aec_heatmap  (샘플별)", C_DARK),
        ("cam_aec_lines  (개별 곡선)", C_DARK),
        ("cam_aec_mean  (평균 CAM)",  C_DARK),
    ]):
        box(s, cx[j], 0.89, cw[j], 0.3, fill=lc)
        txt(s, lbl, cx[j]+0.05, 0.90, cw[j]-0.1, 0.28,
            sz=10.5, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)

    # 행 레이블 + 이미지 (M2 위, M3 아래)
    for row_i, (model_lbl, model_col, mdir) in enumerate([
        ("M2  CrossAttn",  C_TEAL, m2_dir),
        ("M3  CrossAttn3", C_ORG,  m3_dir),
    ]):
        y0 = 1.22 + row_i * 3.08
        box(s, 0.2, y0, 0.3, 2.75, fill=model_col)
        txt(s, model_lbl, 0.22, y0 + 0.9, 0.28, 1.5,
            sz=9, bold=True, color=C_WHT, align=PP_ALIGN.CENTER)
        for j, fname in enumerate(["cam_aec_heatmap.png", "cam_aec_lines.png", "cam_aec_mean.png"]):
            pic(s, img(mdir, fname), cx[j] + 0.32, y0, cw[j] - 0.32, 2.75)


for sn, variant in enumerate(AEC128_VARIANTS, start=23):
    cam_aec_slide(variant, sn)

# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
OUT = ("C:/Users/jhjun/OneDrive/Desktop/2026-1_Study/연구코드/"
       "AEC_Sarcopenia_Research_Presentation.pptx")
prs.save(OUT)
print(f"Saved  → {OUT}")
print(f"Slides : {len(prs.slides)}")
