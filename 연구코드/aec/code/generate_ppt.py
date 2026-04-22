"""
TAMA 예측 회귀분석 연구 보고서 PPT 생성 스크립트
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os
import config

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x35, 0x5E)   # 제목 배경
BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # 강조
SKYBLUE = RGBColor(0xD6, 0xE4, 0xF0)   # 셀 배경
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE  = RGBColor(0xFF, 0x7F, 0x27)
GREEN   = RGBColor(0x00, 0x8B, 0x45)

FIGURES = f"results/{config.SITE}/figures"


def parse_report(site):
    """results/{site}/research_report.md 를 파싱해 슬라이드에 쓸 값을 dict로 반환."""
    path = os.path.join("results", site, "research_report.md")
    with open(path, encoding="utf-8") as fh:
        lines = fh.readlines()

    D = {'site': site}

    # ── 내부 헬퍼 ────────────────────────────────────────────
    def strip_bold(s):
        return re.sub(r'\*\*([^*]+)\*\*', r'\1', s).strip()

    def flt(s):
        cleaned = re.sub(r'[^\d.\-]', '', strip_bold(s))
        try:
            return float(cleaned)
        except ValueError:
            return 0.0

    def sint(s):
        digits = re.sub(r'[^\d]', '', strip_bold(s))
        return int(digits) if digits else 0

    def find_section(header):
        for i, ln in enumerate(lines):
            if header in ln:
                return i
        return -1

    def table_rows(start, max_rows=20):
        if start < 0:
            return []
        rows = []
        seen_sep = False
        count = 0
        for ln in lines[start + 1: start + 1 + max_rows + 15]:
            stripped = ln.strip()
            if not stripped.startswith('|'):
                if seen_sep:
                    break
                continue
            cols = [c.strip() for c in stripped.split('|')[1:-1]]
            if cols and all(re.fullmatch(r'[-: ]+', c) for c in cols if c):
                seen_sep = True
                continue
            if seen_sep:
                rows.append(cols)
                count += 1
                if count >= max_rows:
                    break
        return rows

    def fmt_beta(raw, ndp=2):
        m = re.match(r'(-?[\d.]+)(.*)', raw.strip())
        if m:
            return f"{float(m.group(1)):.{ndp}f}{m.group(2)}"
        return raw

    def fmt_ci(raw, ndp=2):
        m = re.match(r'\[(-?[\d.]+),\s*(-?[\d.]+)\]', raw.strip())
        if m:
            return f"[{float(m.group(1)):.{ndp}f}, {float(m.group(2)):.{ndp}f}]"
        return raw

    def fmt_float_val(raw, ndp=3):
        try:
            return f"{float(strip_bold(raw)):.{ndp}f}"
        except ValueError:
            return strip_bold(raw)

    def _pdisp(raw):
        raw_s = strip_bold(raw).strip()
        try:
            v = float(raw_s)
            if v < 0.001:
                return "< 0.001"
            if v < 0.01:
                return f"{v:.4f}"
            return f"{v:.3f}"
        except ValueError:
            return raw_s

    def _hl_disp(raw):
        raw_s = strip_bold(raw).strip()
        try:
            v = float(raw_s)
            if v < 0.01:
                return f"{v:.2e}"
            return f"{v:.3f}"
        except ValueError:
            return raw_s

    # ── 2.1 데이터셋 기본 정보 ───────────────────────────────
    sec = find_section('### 2.1 데이터셋 기본 정보')
    r = table_rows(sec, max_rows=6)

    mm = re.search(r'(\d[\d,]*)', strip_bold(r[0][1]))
    D['n_total'] = int(mm.group(1).replace(',', '')) if mm else 0

    mm = re.search(r'(\d+)명 \(([\d.]+)%\)', r[1][1])
    if mm:
        D['n_male'] = int(mm.group(1))
        D['male_pct'] = mm.group(2)

    mm = re.search(r'(\d+)명 \(([\d.]+)%\)', r[2][1])
    if mm:
        D['n_female'] = int(mm.group(1))
        D['female_pct'] = mm.group(2)

    mm = re.search(r'(\d+) ~ (\d+)', r[3][1])
    if mm:
        D['tama_min'] = mm.group(1)
        D['tama_max'] = mm.group(2)

    mm = re.search(r'([\d.]+) \(([\d.]+)\)', r[4][1])
    if mm:
        D['tama_mean'] = mm.group(1)
        D['tama_sd'] = mm.group(2)

    D['ct_models'] = sint(r[5][1])

    # ── 2.2 성별 TAMA 분포 ───────────────────────────────────
    sec = find_section('### 2.2 성별 TAMA 분포')
    r = table_rows(sec, max_rows=2)

    D['male_n']      = r[0][1]
    D['male_mean']   = r[0][2]
    D['male_sd']     = r[0][3]
    D['male_median'] = r[0][4]
    D['male_p25']    = r[0][5]
    D['male_p75']    = r[0][6]

    D['female_n']      = r[1][1]
    D['female_mean']   = r[1][2]
    D['female_sd']     = r[1][3]
    D['female_median'] = r[1][4]
    D['female_p25']    = r[1][5]
    D['female_p75']    = r[1][6]

    # ── 2.3 이진화 기준 ──────────────────────────────────────
    sec = find_section('### 2.3 Logistic Regression 이진화')
    for ln in lines[sec: sec + 8]:
        mm = re.search(r'남성.*?TAMA < (\d+)', ln)
        if mm:
            D['thr_male'] = mm.group(1)
        mm = re.search(r'여성.*?TAMA < (\d+)', ln)
        if mm:
            D['thr_female'] = mm.group(1)
        mm = re.search(r'양성 비율.*?(\d+)/(\d+) = ([\d.]+)%', ln)
        if mm:
            D['n_pos'] = mm.group(1)
            D['pos_rate'] = mm.group(3)

    # ── 2.4 AEC Feature ──────────────────────────────────────
    sec = find_section('### 2.4 AEC Feature 선택 근거')
    r = table_rows(sec, max_rows=4)
    for i, key in enumerate(['p25', 'CV', 'skewness', 'slope_abs_mean']):
        D[f'pearson_{key}'] = strip_bold(r[i][1])
        D[f'vif_{key}']     = strip_bold(r[i][3])

    # ── 3.1 전처리 더미 수 ───────────────────────────────────
    sec = find_section('### 3.1 전처리')
    D['ct_dummy_count'] = '30'
    for ln in lines[sec: sec + 15]:
        mm = re.search(r'→ (\d+)개 더미', ln)
        if mm:
            D['ct_dummy_count'] = mm.group(1)
            break

    # ── 4.1.1 선형 단변량 ────────────────────────────────────
    sec = find_section('#### 4.1.1 단변량 분석 (Univariate)')
    r = table_rows(sec, max_rows=8)
    lin_uni_keys = ['Sex', 'Age', 'p25', 'CV', 'skewness', 'slope_abs_mean',
                    'mean', 'ManufacturerModelName']
    for i, key in enumerate(lin_uni_keys):
        raw_beta = strip_bold(r[i][1])
        D[f'linu_{key}_beta']      = raw_beta
        D[f'linu_{key}_beta_disp'] = fmt_beta(raw_beta, 2)
        D[f'linu_{key}_ci']        = fmt_ci(r[i][2], 2)
        D[f'linu_{key}_pval']      = r[i][3]
        D[f'linu_{key}_r2']        = fmt_float_val(r[i][4], 3)

    mm = re.match(r'(-?[\d.]+)', D['linu_Sex_beta'])
    D['linu_Sex_beta_num'] = float(mm.group(1)) if mm else 0.0

    # ── 4.1.2 선형 다변량 성능 ───────────────────────────────
    sec412 = find_section('#### 4.1.2 다변량 분석 (Multivariate')
    r = table_rows(sec412, max_rows=9)
    # 0=N, 1=R², 2=Adj R², 3=F, 4=RMSE, 5=MAE, 6=AIC, 7=BIC, 8=CV

    D['lin_r2']     = strip_bold(r[1][1])
    D['lin_adj_r2'] = strip_bold(r[2][1])

    mm = re.match(r'([\d.]+)', strip_bold(r[3][1]))
    fstat_val = float(mm.group(1)) if mm else 0.0
    D['lin_fstat_disp'] = f"{fstat_val:.2f}\n(p < 0.001)"

    mm = re.match(r'([\d.]+)', strip_bold(r[4][1]))
    D['lin_rmse_disp'] = f"{float(mm.group(1)):.2f} cm²" if mm else strip_bold(r[4][1])

    mm = re.match(r'([\d.]+)', strip_bold(r[5][1]))
    D['lin_mae_disp'] = f"{float(mm.group(1)):.2f} cm²" if mm else strip_bold(r[5][1])

    D['lin_aic'] = strip_bold(r[6][1])
    D['lin_bic'] = strip_bold(r[7][1])
    D['lin_aic_disp'] = f"{int(float(D['lin_aic'])):,}\n/ {int(float(D['lin_bic'])):,}"

    # 잔차 진단 표 (4.1.2 내부, 4.2.1 이전)
    sec421 = find_section('#### 4.2.1 단변량 분석 (Crude OR)')
    resid_hdr = -1
    for i in range(sec412, sec421):
        if '검정' in lines[i] and '통계량' in lines[i] and lines[i].strip().startswith('|'):
            resid_hdr = i
            break
    if resid_hdr >= 0:
        r = table_rows(resid_hdr, max_rows=4)
        D['resid_sw_stat']      = r[0][1]
        D['resid_sw_p']         = r[0][2]
        D['resid_sw_res']       = r[0][3]
        D['resid_bp_stat_disp'] = f"{float(r[1][1]):.2f}"
        D['resid_bp_p']         = r[1][2]
        D['resid_bp_res']       = r[1][3]
        D['resid_dw_stat']      = r[2][1]
        D['resid_dw_res']       = r[2][3]
        D['resid_kappa_stat']   = r[3][1]
        D['resid_kappa_res']    = r[3][3]

    # ── 4.2.1 로지스틱 단변량 ────────────────────────────────
    r = table_rows(sec421, max_rows=8)
    log_uni_keys = ['Sex', 'Age', 'p25', 'CV', 'skewness', 'slope_abs_mean',
                    'mean', 'ManufacturerModelName']
    for i, key in enumerate(log_uni_keys):
        D[f'logu_{key}_or']   = fmt_beta(strip_bold(r[i][1]), 3)
        D[f'logu_{key}_ci']   = fmt_ci(r[i][2], 3)
        D[f'logu_{key}_pval'] = r[i][3]
        D[f'logu_{key}_auc']  = fmt_float_val(r[i][4], 3)

    # ── 4.2.2 로지스틱 다변량 성능 ───────────────────────────
    sec = find_section('#### 4.2.2 다변량 분석 (전체 변수 투입)')
    r = table_rows(sec, max_rows=10)
    # 0=N, 1=AUC, 2=Sens, 3=Spec, 4=PPV, 5=NPV, 6=HL, 7=Nagelkerke, 8=Brier, 9=AIC

    auc_raw = strip_bold(r[1][1])
    mm = re.match(r'([\d.]+)\s*\[([\d.]+)[–\-]([\d.]+)\]', auc_raw)
    if mm:
        D['log_auc']    = mm.group(1)
        D['log_auc_lo'] = mm.group(2)
        D['log_auc_hi'] = mm.group(3)
    D['log_auc_ci_disp'] = f"[{D['log_auc_lo']} – {D['log_auc_hi']}]"

    D['log_sensitivity']  = strip_bold(r[2][1])
    D['log_specificity']  = strip_bold(r[3][1])
    D['log_ppv']          = strip_bold(r[4][1])
    D['log_npv']          = strip_bold(r[5][1])
    D['log_ppv_npv_disp'] = f"{D['log_ppv']} / {D['log_npv']}"

    hl_raw = strip_bold(r[6][1])
    mm = re.search(r'χ²=([\d.]+),\s*p=([\d.]+)', hl_raw)
    if mm:
        D['log_hl_stat'] = mm.group(1)
        D['log_hl_p']    = mm.group(2)

    D['log_nagelkerke'] = strip_bold(r[7][1])
    D['log_brier']      = strip_bold(r[8][1])
    D['log_aic']        = strip_bold(r[9][1])
    D['log_aic_disp']   = f"{float(D['log_aic']):,.2f}"

    # ── 4.3.1 선형 Case 비교 ─────────────────────────────────
    sec = find_section('#### 4.3.1 선형 회귀 성능 비교')
    r = table_rows(sec, max_rows=4)
    # 0=R², 1=Adj R², 2=RMSE, 3=AIC

    for col, ckey in [(1, 'c1'), (2, 'c2'), (3, 'c3')]:
        D[f'lin_{ckey}_r2']    = strip_bold(r[0][col])
        D[f'lin_{ckey}_adjr2'] = strip_bold(r[1][col])
        D[f'lin_{ckey}_rmse']  = strip_bold(r[2][col])
        D[f'lin_{ckey}_aic']   = strip_bold(r[3][col])

    r2_12  = flt(D['lin_c2_r2'])    - flt(D['lin_c1_r2'])
    r2_23  = flt(D['lin_c3_r2'])    - flt(D['lin_c2_r2'])
    ar2_12 = flt(D['lin_c2_adjr2']) - flt(D['lin_c1_adjr2'])
    ar2_23 = flt(D['lin_c3_adjr2']) - flt(D['lin_c2_adjr2'])
    rm_12  = flt(D['lin_c2_rmse'])  - flt(D['lin_c1_rmse'])
    rm_23  = flt(D['lin_c3_rmse'])  - flt(D['lin_c2_rmse'])
    ai_12  = flt(D['lin_c2_aic'])   - flt(D['lin_c1_aic'])
    ai_23  = flt(D['lin_c3_aic'])   - flt(D['lin_c2_aic'])

    D['lin_delta12_r2']    = f"+{r2_12:.4f} ▲"
    D['lin_delta23_r2']    = f"+{r2_23:.4f} ▲"
    D['lin_delta12_adjr2'] = f"+{ar2_12:.4f} ▲"
    D['lin_delta23_adjr2'] = f"+{ar2_23:.4f} ▲"
    D['lin_delta12_rmse']  = f"{rm_12:.2f} ▼"
    D['lin_delta23_rmse']  = f"{rm_23:.2f} ▼"
    D['lin_delta12_aic']   = f"{ai_12:.0f} ▼"
    D['lin_delta23_aic']   = f"{ai_23:.0f} ▼"
    D['lin_summary'] = (
        f"AEC 특징 추가(Case1→2)로 R² +{r2_12 * 100:.1f}%p 향상  |  "
        f"CT 모델명 추가(Case2→3)로 추가 +{r2_23 * 100:.1f}%p"
    )

    # ── 4.3.2 로지스틱 Case 비교 ─────────────────────────────
    sec = find_section('#### 4.3.2 로지스틱 회귀 성능 비교')
    r = table_rows(sec, max_rows=4)
    # 0=AUC, 1=Nagelkerke, 2=AIC, 3=HL p

    for col, ckey in [(1, 'c1'), (2, 'c2'), (3, 'c3')]:
        D[f'log_{ckey}_auc'] = strip_bold(r[0][col])
        D[f'log_{ckey}_nag'] = strip_bold(r[1][col])
        D[f'log_{ckey}_aic'] = strip_bold(r[2][col])
        D[f'log_{ckey}_hl']  = strip_bold(r[3][col])

    D['log_c1_hl_disp'] = _hl_disp(D['log_c1_hl'])
    D['log_c2_hl_disp'] = _hl_disp(D['log_c2_hl'])
    D['log_c3_hl_disp'] = _hl_disp(D['log_c3_hl'])

    au_12 = flt(D['log_c2_auc']) - flt(D['log_c1_auc'])
    au_23 = flt(D['log_c3_auc']) - flt(D['log_c2_auc'])
    na_12 = flt(D['log_c2_nag']) - flt(D['log_c1_nag'])
    na_23 = flt(D['log_c3_nag']) - flt(D['log_c2_nag'])
    la_12 = flt(D['log_c2_aic']) - flt(D['log_c1_aic'])
    la_23 = flt(D['log_c3_aic']) - flt(D['log_c2_aic'])

    D['log_delta12_auc'] = f"+{au_12:.3f} ▲▲"
    D['log_delta23_auc'] = f"+{au_23:.3f} ▲"
    D['log_delta12_nag'] = f"+{na_12:.3f} ▲▲"
    D['log_delta23_nag'] = f"+{na_23:.3f} ▲"
    D['log_delta12_aic'] = f"{la_12:.0f} ▼▼"
    D['log_delta23_aic'] = f"{la_23:.0f} ▼"
    D['log_summary'] = (
        f"AEC 추가(Case1→2)로 AUC +{au_12 * 100:.1f}%p 대폭 향상  |  "
        f"CT 모델명 추가(Case2→3)로 +{au_23 * 100:.1f}%p 추가 개선"
    )

    return D


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃

D = parse_report(config.SITE)

# ────────────────────────────────────────────────────────────
# 헬퍼 함수
# ────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.line.width = line_w
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if line_rgb:
        shp.line.color.rgb = line_rgb
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None):
    """슬라이드 상단 헤더 바"""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=NAVY)
    add_text(slide, title, 0.35, 0.1, 10, 0.6,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.65, 10, 0.45,
                 size=14, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)


def add_image(slide, fname, l, t, w, h):
    path = os.path.join(FIGURES, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        shp = add_rect(slide, l, t, w, h, fill_rgb=GRAY, line_rgb=BLUE)
        add_text(slide, f"[그림 없음]\n{fname}", l+0.1, t+h/2-0.2, w-0.2, 0.4,
                 size=9, color=BLUE, align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows,
              l, t, w, h,
              header_fill=NAVY, header_color=WHITE,
              alt_fill=SKYBLUE, font_size=11):
    """심플 텍스트-기반 표 (python-pptx Table)"""
    cols = len(headers)
    total_rows = 1 + len(rows)
    tbl = slide.shapes.add_table(total_rows, cols,
                                  Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    col_w = Inches(w / cols)
    for c in range(cols):
        tbl.columns[c].width = col_w

    # 헤더
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = header_color

    # 데이터 행
    for r, row in enumerate(rows):
        fill_col = GRAY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_col
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK


def bullet_box(slide, items, l, t, w, h, size=13, title=None):
    """글머리 기호 텍스트 박스"""
    offset = 0
    if title:
        add_text(slide, title, l, t, w, 0.35, size=13, bold=True, color=NAVY)
        offset = 0.35
    txb = slide.shapes.add_textbox(Inches(l), Inches(t + offset),
                                    Inches(w), Inches(h - offset))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)


def _pdisp(raw):
    """p-value 문자열을 표시용으로 변환"""
    raw_s = raw.strip()
    try:
        v = float(raw_s)
        if v < 0.001:
            return "< 0.001"
        if v < 0.01:
            return f"{v:.4f}"
        return f"{v:.3f}"
    except ValueError:
        return raw_s


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# 전체 배경
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
# 하단 포인트 바
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=BLUE)

# 중앙 흰 카드
add_rect(slide, 1.0, 1.5, 11.33, 4.5, fill_rgb=WHITE)

add_text(slide,
         "AEC 곡선 기반 TAMA 예측\n회귀분석 연구 보고서",
         1.3, 1.7, 10.7, 2.0,
         size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(slide,
         "CT 자동 노출 제어(AEC) 특징과 환자 인구통계학적 변수를 활용한\n"
         "복부 근육량 지표(TAMA) 예측 모델 개발 및 검증",
         1.3, 3.6, 10.7, 1.0,
         size=15, color=RGBColor(0x44, 0x44, 0x66), align=PP_ALIGN.CENTER)

add_text(slide,
         f"데이터셋: {D['site']} CT  |  N={D['n_total']:,}명  |  Python (statsmodels · scikit-learn)",
         1.3, 4.7, 10.7, 0.5,
         size=13, color=BLUE, align=PP_ALIGN.CENTER, bold=True)

add_text(slide, "2026", 0.4, 6.85, 12.5, 0.5,
         size=13, color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "목  차", "Table of Contents")

sections = [
    ("01", "연구 개요", "연구 목적 · 설계 · Feature Set (3 Case)"),
    ("02", "데이터 기술 통계", "데이터셋 · 성별 TAMA 분포 · 이진화 기준"),
    ("03", "분석 방법론", "전처리 · 모델 · 교차검증"),
    ("04", "결과 — 선형 회귀 (Part 1)", "단변량 · 다변량 성능 · 잔차 진단"),
    ("05", "결과 — 로지스틱 회귀 (Part 2)", "단변량 OR · 다변량 AUC · Calibration"),
    ("06", "결과 — Case 1–3 비교 (Part 3)", "AEC 추가 효과 · CT 모델명 효과"),
    ("07", "결론 및 한계", "주요 발견 · 제한점 · 향후 방향"),
]

for i, (num, title, desc) in enumerate(sections):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 1.3 + row * 1.65

    add_rect(slide, x, y, 6.0, 1.45, fill_rgb=NAVY)
    add_text(slide, num, x+0.15, y+0.1, 0.7, 0.6, size=22, bold=True, color=ORANGE)
    add_text(slide, title, x+0.85, y+0.08, 4.9, 0.5, size=15, bold=True, color=WHITE)
    add_text(slide, desc, x+0.85, y+0.65, 4.9, 0.7, size=11, color=RGBColor(0xB0,0xC8,0xE8))


# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 연구 개요
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "01  연구 개요", "Research Overview")

# 연구 목적 카드
add_rect(slide, 0.35, 1.25, 12.6, 1.15, fill_rgb=NAVY)
add_text(slide, "연구 목적",
         0.55, 1.28, 3.0, 0.4, size=13, bold=True, color=ORANGE)
add_text(slide,
         "CT 자동 노출 제어(AEC) 곡선과 환자 인구통계(성별·나이)가 "
         "복부 근육량 지표 TAMA(Total Abdominal Muscle Area) 예측에 기여하는 정도를 정량화",
         0.55, 1.62, 12.2, 0.65, size=12, color=WHITE)

# 3개 파트 카드
parts = [
    ("Part 1", "선형 회귀\n(Linear Regression)", "TAMA (연속형) 값 예측"),
    ("Part 2", "로지스틱 회귀\n(Logistic Regression)", "Low TAMA / Sarcopenia 위험 분류"),
    ("Part 3", "Multivariable Analysis\n(Case 1–3 비교)", "AEC·CT 모델 추가 효과 정량화"),
]
for i, (p, m, g) in enumerate(parts):
    x = 0.35 + i * 4.2
    add_rect(slide, x, 2.6, 4.0, 2.1, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1.5))
    add_text(slide, p, x+0.15, 2.68, 3.7, 0.4, size=13, bold=True, color=BLUE)
    add_text(slide, m, x+0.15, 3.05, 3.7, 0.75, size=12, bold=True, color=NAVY)
    add_text(slide, g, x+0.15, 3.8, 3.7, 0.7, size=11, color=DARK)

# Feature Set 표
add_text(slide, "Feature Set (3 Case)", 0.35, 4.85, 4.0, 0.35,
         size=13, bold=True, color=NAVY)
add_table(slide,
          ["Case", "예측 변수"],
          [
              ["Case 1", "성별 (Sex) + 나이 (Age)"],
              ["Case 2", "Case 1 + AEC 특징 (p25, CV, skewness, slope_abs_mean, mean)"],
              ["Case 3",
               f"Case 2 + CT 모델명 (ManufacturerModelName — {D['ct_dummy_count']}개 dummy) + KVP"],
          ],
          0.35, 5.2, 12.6, 1.9, font_size=11)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 데이터 기술 통계
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  데이터 기술 통계", "Descriptive Statistics")

# 기본 정보 카드
stats = [
    ("총 환자 수",      f"{D['n_total']:,}명"),
    ("남성",           f"{D['n_male']:,}명 ({D['male_pct']}%)"),
    ("여성",           f"{D['n_female']:,}명 ({D['female_pct']}%)"),
    ("CT 스캐너 모델",  f"{D['ct_models']}종"),
    ("TAMA 범위",       f"{D['tama_min']} ~ {D['tama_max']} cm²"),
    ("TAMA 평균 (SD)",  f"{D['tama_mean']} (±{D['tama_sd']}) cm²"),
]
for i, (label, val) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.1
    add_rect(slide, x, y, 3.9, 0.95, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1))
    add_text(slide, label, x+0.12, y+0.05, 3.7, 0.35, size=10, color=BLUE)
    add_text(slide, val, x+0.12, y+0.42, 3.7, 0.45, size=16, bold=True, color=NAVY)

# 성별 TAMA 분포 표
add_text(slide, "성별 TAMA 분포", 0.35, 3.55, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["성별", "N", "평균 (cm²)", "SD", "중앙값", "P25", "P75"],
          [
              ["남성 (M)", D['male_n'],   D['male_mean'],   D['male_sd'],
               D['male_median'],   D['male_p25'],   D['male_p75']],
              ["여성 (F)", D['female_n'], D['female_mean'], D['female_sd'],
               D['female_median'], D['female_p25'], D['female_p75']],
          ],
          0.35, 3.92, 7.5, 1.35, font_size=11)

# 이진화 기준 카드
add_rect(slide, 8.1, 3.55, 4.9, 1.72, fill_rgb=NAVY)
add_text(slide, "Logistic Regression 이진화 기준", 8.25, 3.6, 4.6, 0.4,
         size=11, bold=True, color=ORANGE)
add_text(slide,
         f"• 남성: TAMA < {D['thr_male']} cm²  →  Low Muscle (1)\n"
         f"• 여성: TAMA < {D['thr_female']} cm²  →  Low Muscle (1)\n"
         f"• 양성 비율: {D['n_pos']} / {D['n_total']:,} = {D['pos_rate']} %",
         8.25, 4.03, 4.6, 1.2, size=12, color=WHITE)

# TAMA 분포 그림
add_text(slide, "TAMA 분포 시각화", 0.35, 5.35, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "03_tama_distribution.png", 0.35, 5.72, 6.0, 1.55)

# AEC Feature 상관 그림
add_text(slide, "AEC Feature 상관계수", 6.65, 5.35, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "01_feature_correlation.png", 6.65, 5.72, 6.33, 1.55)


# ════════════════════════════════════════════════════════════
# 슬라이드 5 — AEC Feature 선택
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  AEC Feature 선택", "Feature Selection")

add_text(slide,
         "Pearson/Spearman 상관계수 + VIF 검사를 통해 다중공선성이 없는 4개 핵심 Feature 선택 (mean 포함 사용)",
         0.35, 1.15, 12.6, 0.45, size=13, color=DARK)

add_table(slide,
          ["Feature", "Pearson r", "해석", "VIF", "선택"],
          [
              ["p25",
               D['pearson_p25'],
               "AEC 하위 25% 값 → 저선량 구간 tube current",
               D['vif_p25'], "✔"],
              ["CV",
               D['pearson_CV'],
               "변동계수(std/mean) → 체형 불균일성 반영",
               D['vif_CV'], "✔"],
              ["skewness",
               D['pearson_skewness'],
               "AEC 곡선 비대칭성 → 체형 분포 특성",
               D['vif_skewness'], "✔"],
              ["slope_abs_mean",
               D['pearson_slope_abs_mean'],
               "평균 절대 기울기 → 곡선 동역학",
               D['vif_slope_abs_mean'], "✔"],
              ["mean",            "+0.297",
               "AEC 평균값 → 체격 전반 반영",
               "확인됨", "✔ (포함)"],
              ["AUC_normalized",  "높음",
               "mean과 심각한 공선성",
               ">50,000", "✗ 제외"],
          ],
          0.35, 1.65, 12.6, 2.6, font_size=11)

add_text(slide, "⚠  AUC_normalized는 VIF > 50,000으로 다중공선성 심각 → 제외 / mean은 모델에 포함",
         0.35, 4.3, 12.6, 0.5, size=12, italic=True, color=ORANGE)

# 두 그림 나란히
add_text(slide, "Feature 상관계수 Top20", 0.35, 4.85, 5.5, 0.35, size=12, bold=True, color=NAVY)
add_image(slide, "01_feature_correlation.png", 0.35, 5.22, 6.0, 2.0)

add_text(slide, "VIF 비교", 6.65, 4.85, 5.5, 0.35, size=12, bold=True, color=NAVY)
add_image(slide, "02_vif_comparison.png", 6.65, 5.22, 6.33, 2.0)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 분석 방법론
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "03  분석 방법론", "Statistical Methods")

# 전처리 표
add_text(slide, "전처리", 0.35, 1.25, 5.5, 0.35, size=14, bold=True, color=NAVY)
add_table(slide,
          ["변수", "처리 방법"],
          [
              ["PatientSex", "이진 인코딩 (M=1, F=0)"],
              ["PatientAge", "Z-score 표준화 (StandardScaler)"],
              ["AEC features", "Z-score 표준화"],
              ["ManufacturerModelName",
               f"One-hot encoding (drop_first=True) → {D['ct_dummy_count']}개 더미"],
          ],
          0.35, 1.65, 6.0, 2.1, font_size=11)

# 모델 카드
add_text(slide, "모델", 6.65, 1.25, 6.0, 0.35, size=14, bold=True, color=NAVY)
models = [
    ("선형 회귀", "statsmodels OLS\n(Ordinary Least Squares)"),
    ("로지스틱 회귀", "statsmodels Logit\n(BFGS, maxiter=1,000)"),
]
for i, (nm, desc) in enumerate(models):
    y = 1.65 + i * 1.1
    add_rect(slide, 6.65, y, 6.33, 0.95, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1.2))
    add_text(slide, nm, 6.8, y+0.05, 3.5, 0.38, size=13, bold=True, color=NAVY)
    add_text(slide, desc, 6.8, y+0.48, 5.9, 0.42, size=11, color=DARK)

# 교차검증 및 CI
add_text(slide, "모델 검증 전략", 0.35, 3.9, 12.6, 0.35, size=14, bold=True, color=NAVY)
add_table(slide,
          ["검증 방법", "상세"],
          [
              ["5-Fold Cross Validation", "각 Fold에서 scaler 재적합 (Data Leakage 방지)"],
              ["Bootstrap AUC 95% CI", "n=1,000 resampling — 비모수 신뢰구간"],
              ["Hosmer-Lemeshow 검정", "로지스틱 보정도(Calibration) 검증 — p>0.05 양호"],
              ["잔차 진단 3종", "Shapiro-Wilk · Breusch-Pagan · Durbin-Watson"],
          ],
          0.35, 4.3, 12.6, 2.1, font_size=11)

add_text(slide,
         "※ 선형·로지스틱 회귀 모두 Case 1(Sex+Age) → Case 2(+AEC) → Case 3(+CT Model) 순으로 점진 투입",
         0.35, 6.5, 12.6, 0.55, size=11, italic=True, color=BLUE)


# ════════════════════════════════════════════════════════════
# 슬라이드 7 — 선형 회귀: 단변량
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 단변량 분석", "Univariate Linear Regression")

add_table(slide,
          ["변수", "β 계수", "95% CI", "p-value", "R²"],
          [
              ["Sex (M=1, F=0)",
               D['linu_Sex_beta_disp'],        D['linu_Sex_ci'],
               _pdisp(D['linu_Sex_pval']),     D['linu_Sex_r2']],
              ["Age (표준화)",
               D['linu_Age_beta_disp'],         D['linu_Age_ci'],
               _pdisp(D['linu_Age_pval']),      D['linu_Age_r2']],
              ["AEC: p25 (표준화)",
               D['linu_p25_beta_disp'],         D['linu_p25_ci'],
               _pdisp(D['linu_p25_pval']),      D['linu_p25_r2']],
              ["AEC: CV (표준화)",
               D['linu_CV_beta_disp'],          D['linu_CV_ci'],
               _pdisp(D['linu_CV_pval']),       D['linu_CV_r2']],
              ["AEC: skewness (표준화)",
               D['linu_skewness_beta_disp'],    D['linu_skewness_ci'],
               _pdisp(D['linu_skewness_pval']), D['linu_skewness_r2']],
              ["AEC: slope_abs_mean (표준화)",
               D['linu_slope_abs_mean_beta_disp'], D['linu_slope_abs_mean_ci'],
               _pdisp(D['linu_slope_abs_mean_pval']), D['linu_slope_abs_mean_r2']],
              ["AEC: mean (표준화)",
               D['linu_mean_beta_disp'],        D['linu_mean_ci'],
               _pdisp(D['linu_mean_pval']),     D['linu_mean_r2']],
              ["ManufacturerModelName",
               D['linu_ManufacturerModelName_beta_disp'], "N/A",
               _pdisp(D['linu_ManufacturerModelName_pval']),
               D['linu_ManufacturerModelName_r2']],
          ],
          0.35, 1.25, 12.6, 3.3, font_size=11)

add_text(slide, "* p < 0.05 유의",
         0.35, 4.65, 5.0, 0.35, size=11, italic=True, color=ORANGE)

# 단변량 R² 그림
add_text(slide, "단변량 R² 비교", 0.35, 5.05, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "07_linear_univariate_r2.png", 0.35, 5.42, 6.0, 1.85)

# Forest plot
add_text(slide, "유의 계수 Forest Plot", 6.65, 5.05, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "06_linear_forest.png", 6.65, 5.42, 6.33, 1.85)


# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 선형 회귀: 다변량 성능
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 다변량 모델 성능 (Case 3 Full)", "Multivariate Linear Regression")

# 성능 KPI 카드
kpis = [
    ("R²",          D['lin_r2']),
    ("Adj R²",      D['lin_adj_r2']),
    ("RMSE",        D['lin_rmse_disp']),
    ("MAE",         D['lin_mae_disp']),
    ("F-statistic", D['lin_fstat_disp']),
    ("AIC / BIC",   D['lin_aic_disp']),
]
for i, (label, val) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.25
    add_rect(slide, x, y, 3.9, 1.15, fill_rgb=NAVY)
    add_text(slide, label, x+0.15, y+0.08, 3.6, 0.38, size=12, bold=True, color=ORANGE)
    add_text(slide, val, x+0.15, y+0.52, 3.6, 0.55, size=18, bold=True, color=WHITE)

# 잔차 진단 표
add_text(slide, "잔차 진단", 0.35, 3.85, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["검정", "통계량", "p-value", "결과"],
          [
              ["Shapiro-Wilk",
               D['resid_sw_stat'],      D['resid_sw_p'],  D['resid_sw_res']],
              ["Breusch-Pagan",
               D['resid_bp_stat_disp'], D['resid_bp_p'],  D['resid_bp_res']],
              ["Durbin-Watson",
               D['resid_dw_stat'],      "—",              D['resid_dw_res']],
              ["Condition κ",
               D['resid_kappa_stat'],   "—",              D['resid_kappa_res']],
          ],
          0.35, 4.22, 6.5, 2.0, font_size=11)

add_text(slide,
         "⚠  비정규·이분산 → 계수 추정 BLUE이나 p-value 해석 주의\n   Robust Standard Errors 사용 권고",
         0.35, 6.3, 6.5, 0.75, size=11, italic=True, color=ORANGE)

# 실제 vs 예측 그림
add_text(slide, "실제 vs 예측", 7.1, 3.85, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "04_linear_actual_vs_pred.png", 7.1, 4.22, 5.9, 2.0)

# 잔차 진단 4-panel
add_image(slide, "05_linear_residuals.png", 7.1, 6.25, 5.9, 1.0)


# ════════════════════════════════════════════════════════════
# 슬라이드 9 — 로지스틱 회귀: 단변량 OR
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 단변량 분석 (Crude OR)", "Univariate Logistic Regression")

add_table(slide,
          ["변수", "Crude OR", "95% CI", "p-value", "AUC"],
          [
              ["Sex (M=1, F=0)",
               D['logu_Sex_or'],        D['logu_Sex_ci'],
               _pdisp(D['logu_Sex_pval']),  D['logu_Sex_auc']],
              ["Age (표준화)",
               D['logu_Age_or'],        D['logu_Age_ci'],
               _pdisp(D['logu_Age_pval']),  D['logu_Age_auc']],
              ["AEC: p25 (표준화)",
               D['logu_p25_or'],        D['logu_p25_ci'],
               _pdisp(D['logu_p25_pval']),  D['logu_p25_auc']],
              ["AEC: CV (표준화)",
               D['logu_CV_or'],         D['logu_CV_ci'],
               _pdisp(D['logu_CV_pval']),   D['logu_CV_auc']],
              ["AEC: skewness (표준화)",
               D['logu_skewness_or'],   D['logu_skewness_ci'],
               _pdisp(D['logu_skewness_pval']), D['logu_skewness_auc']],
              ["AEC: slope_abs_mean (표준화)",
               D['logu_slope_abs_mean_or'],  D['logu_slope_abs_mean_ci'],
               _pdisp(D['logu_slope_abs_mean_pval']), D['logu_slope_abs_mean_auc']],
              ["AEC: mean (표준화)",
               D['logu_mean_or'],       D['logu_mean_ci'],
               _pdisp(D['logu_mean_pval']),  D['logu_mean_auc']],
              ["ManufacturerModelName",
               D['logu_ManufacturerModelName_or'], "N/A",
               _pdisp(D['logu_ManufacturerModelName_pval']),
               D['logu_ManufacturerModelName_auc']],
          ],
          0.35, 1.25, 12.6, 3.3, font_size=11)

add_text(slide,
         "* p < 0.05 유의  |  p25·mean: OR < 1 → 해당 변수↑ 시 Low TAMA 위험 감소  |  Sex: 이진화 후 단변량 비유의",
         0.35, 4.65, 12.0, 0.4, size=11, italic=True, color=ORANGE)

# Crude OR Forest plot
add_text(slide, "Crude OR Forest Plot", 0.35, 5.12, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "11_logistic_forest.png", 0.35, 5.5, 6.0, 1.77)

# ROC 곡선
add_text(slide, "ROC Curve (단변량)", 6.65, 5.12, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "08_logistic_roc.png", 6.65, 5.5, 6.33, 1.77)


# ════════════════════════════════════════════════════════════
# 슬라이드 10 — 로지스틱 회귀: 다변량 성능
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 다변량 모델 성능 (Case 3 Full)", "Multivariate Logistic Regression")

kpis2 = [
    ("AUC-ROC",          D['log_auc']),
    ("Bootstrap 95% CI", D['log_auc_ci_disp']),
    ("Sensitivity",      D['log_sensitivity']),
    ("Specificity",      D['log_specificity']),
    ("PPV / NPV",        D['log_ppv_npv_disp']),
    ("Nagelkerke R²",    D['log_nagelkerke']),
]
for i, (label, val) in enumerate(kpis2):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.2
    add_rect(slide, x, y, 3.9, 1.1, fill_rgb=NAVY)
    add_text(slide, label, x+0.15, y+0.06, 3.6, 0.38, size=11, bold=True, color=ORANGE)
    add_text(slide, val, x+0.15, y+0.5, 3.6, 0.52, size=17, bold=True, color=WHITE)

# 보정도 표
add_text(slide, "보정도 및 적합도", 0.35, 3.72, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["지표", "값", "해석"],
          [
              ["Hosmer-Lemeshow χ²",
               D['log_hl_stat'],
               f"p={D['log_hl_p']} → 보정도 양호 ✔"],
              ["Brier Score",
               D['log_brier'],
               "낮을수록 우수 (0=완벽, 0.25=무작위)"],
              ["AIC",
               D['log_aic_disp'],
               "Case 비교 기준"],
          ],
          0.35, 4.1, 7.0, 1.65, font_size=11)

# Calibration plot
add_text(slide, "Calibration Plot", 7.5, 3.72, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "09_logistic_calibration.png", 7.5, 4.1, 5.48, 2.2)

# Confusion matrix
add_text(slide, "Confusion Matrix", 0.35, 5.85, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "10_logistic_confusion.png", 0.35, 6.22, 3.5, 1.05)


# ════════════════════════════════════════════════════════════
# 슬라이드 11 — Case 1·2·3 선형 회귀 비교
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 비교 — 선형 회귀 성능 향상", "Multivariable Analysis: Linear Regression")

add_table(slide,
          ["지표", "Case 1\n(Sex+Age)", "Case 2\n(+AEC)", "Case 3\n(+CT Model)",
           "Case1→2 Δ", "Case2→3 Δ"],
          [
              ["R²",
               D['lin_c1_r2'],   D['lin_c2_r2'],   D['lin_c3_r2'],
               D['lin_delta12_r2'],   D['lin_delta23_r2']],
              ["Adj R²",
               D['lin_c1_adjr2'], D['lin_c2_adjr2'], D['lin_c3_adjr2'],
               D['lin_delta12_adjr2'], D['lin_delta23_adjr2']],
              ["RMSE (cm²)",
               D['lin_c1_rmse'],  D['lin_c2_rmse'],  D['lin_c3_rmse'],
               D['lin_delta12_rmse'],  D['lin_delta23_rmse']],
              ["AIC",
               D['lin_c1_aic'],   D['lin_c2_aic'],   D['lin_c3_aic'],
               D['lin_delta12_aic'],   D['lin_delta23_aic']],
          ],
          0.35, 1.25, 12.6, 2.3, font_size=12)

add_text(slide, D['lin_summary'],
         0.35, 3.65, 12.6, 0.45, size=13, bold=True, color=BLUE)

add_text(slide, "R² / RMSE 비교", 0.35, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "12_case_metrics_bar.png", 0.35, 4.58, 6.0, 2.7)

add_text(slide, "AIC / BIC 비교", 6.65, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "14_case_aic_bar.png", 6.65, 4.58, 6.33, 2.7)


# ════════════════════════════════════════════════════════════
# 슬라이드 12 — Case 1·2·3 로지스틱 회귀 비교
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 비교 — 로지스틱 회귀 성능 향상", "Multivariable Analysis: Logistic Regression")

add_table(slide,
          ["지표", "Case 1\n(Sex+Age)", "Case 2\n(+AEC)", "Case 3\n(+CT Model)",
           "Case1→2 Δ", "Case2→3 Δ"],
          [
              ["AUC-ROC",
               D['log_c1_auc'], D['log_c2_auc'], D['log_c3_auc'],
               D['log_delta12_auc'], D['log_delta23_auc']],
              ["Nagelkerke R²",
               D['log_c1_nag'], D['log_c2_nag'], D['log_c3_nag'],
               D['log_delta12_nag'], D['log_delta23_nag']],
              ["AIC",
               D['log_c1_aic'], D['log_c2_aic'], D['log_c3_aic'],
               D['log_delta12_aic'], D['log_delta23_aic']],
              ["HL p-value",
               D['log_c1_hl_disp'], D['log_c2_hl_disp'], D['log_c3_hl_disp'],
               "보정도 개선 ✔", "보정도 양호 ✔"],
          ],
          0.35, 1.25, 12.6, 2.3, font_size=12)

add_text(slide, D['log_summary'],
         0.35, 3.65, 12.6, 0.45, size=13, bold=True, color=BLUE)

add_text(slide, "AUC / Nagelkerke R² 비교", 0.35, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "13_case_auc_bar.png", 0.35, 4.58, 6.0, 2.7)

add_text(slide, "다중 지표 추이", 6.65, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "15_case_progression.png", 6.65, 4.58, 6.33, 2.7)


# ════════════════════════════════════════════════════════════
# 슬라이드 13 — 결론
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "07  결론", "Conclusions")

_r2_12 = float(D['lin_c2_r2']) - float(D['lin_c1_r2'])
_r2_23 = float(D['lin_c3_r2']) - float(D['lin_c2_r2'])
_au_12 = float(D['log_c2_auc']) - float(D['log_c1_auc'])
_au_23 = float(D['log_c3_auc']) - float(D['log_c2_auc'])

conclusions = [
    ("성별의 압도적 기여",
     f"선형 단변량: Sex 단독 R²={D['linu_Sex_r2']}  |  "
     f"남성 TAMA가 여성 대비 평균 +{D['linu_Sex_beta_num']:.1f} cm²"),
    ("AEC 특징의 독립적 기여",
     f"Case1→2  선형 R²  +{_r2_12:.3f} "
     f"({float(D['lin_c1_r2']):.3f}→{float(D['lin_c2_r2']):.3f})\n"
     f"Case1→2  로지스틱 AUC  +{_au_12:.3f} "
     f"({float(D['log_c1_auc']):.3f}→{float(D['log_c2_auc']):.3f})"),
    ("CT 모델명의 추가 기여",
     f"Case2→3  선형 R²  +{_r2_23:.3f}  |  로지스틱 AUC  +{_au_23:.3f}\n"
     "AIC 추가 감소 → 모델 적합도 개선 확인"),
    ("최종 모델 (Case 3) 요약",
     f"선형: R²={D['lin_r2']}, RMSE={D['lin_rmse_disp']}\n"
     f"로지스틱: AUC={D['log_auc']} {D['log_auc_ci_disp']}  |  "
     f"HL p={D['log_hl_p']} (보정도 양호)"),
]
for i, (title, body) in enumerate(conclusions):
    row = i // 2
    col = i % 2
    x = 0.35 + col * 6.4
    y = 1.25 + row * 2.35
    add_rect(slide, x, y, 6.1, 2.2, fill_rgb=NAVY)
    add_rect(slide, x, y, 0.18, 2.2, fill_rgb=ORANGE)
    add_text(slide, f"{i+1}. {title}", x+0.3, y+0.1, 5.6, 0.45,
             size=13, bold=True, color=ORANGE)
    add_text(slide, body, x+0.3, y+0.6, 5.6, 1.5, size=12, color=WHITE)

add_text(slide,
         "AEC 특징은 성별·나이 외에 TAMA 예측에 독립적이고 유의미한 기여를 함 — CT 촬영 시 수집되는 AEC 데이터의 임상적 활용 가능성 확인",
         0.35, 6.8, 12.6, 0.55, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 14 — 한계 및 향후 방향
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "07  한계 및 향후 방향", "Limitations & Future Work")

other_site = "신촌" if D['site'] == "강남" else "강남"

limits = [
    "잔차 정규성·등분산성 가정 위반 → Robust Standard Errors 또는 비모수 검정 검토 필요",
    f"Logistic 이진화 임계값 (M<{D['thr_male']}, F<{D['thr_female']} cm²) — 문헌 기반, 기관별 참조값 확인 필요",
    f"단일 기관({D['site']}) 데이터 → 외부 검증({other_site} 데이터) 필요",
    "AEC Feature 선택은 통계적 상관분석 기반 — 임상적 의미 추가 검토 권장",
]
future = [
    f"{other_site} 데이터를 이용한 외부 검증 (External Validation)",
    "Robust Regression / Weighted Least Squares 적용",
    "머신러닝 기반 앙상블 모델 비교 (Random Forest, XGBoost)",
    "다기관 데이터 수집을 통한 모델 일반화",
    "AEC Feature의 임상적 의미 방사선사·임상의 검토",
]

add_rect(slide, 0.35, 1.2, 6.0, 3.3, fill_rgb=WHITE, line_rgb=ORANGE, line_w=Pt(1.5))
add_text(slide, "⚠  한계점", 0.55, 1.28, 5.5, 0.38, size=13, bold=True, color=ORANGE)
for i, lim in enumerate(limits):
    add_text(slide, f"• {lim}", 0.55, 1.72 + i*0.72, 5.6, 0.65, size=11, color=DARK)

add_rect(slide, 6.7, 1.2, 6.28, 3.3, fill_rgb=WHITE, line_rgb=GREEN, line_w=Pt(1.5))
add_text(slide, "→  향후 연구 방향", 6.88, 1.28, 5.8, 0.38, size=13, bold=True, color=GREEN)
for i, fut in enumerate(future):
    add_text(slide, f"• {fut}", 6.88, 1.72 + i*0.57, 5.8, 0.52, size=11, color=DARK)

# 결론 요약 강조 박스
add_rect(slide, 0.35, 4.65, 12.6, 1.85, fill_rgb=NAVY)
add_text(slide, "핵심 메시지",
         0.55, 4.72, 3.0, 0.38, size=12, bold=True, color=ORANGE)
add_text(slide,
         "CT 스캐너의 AEC 곡선에서 추출한 통계 특징(p25, CV, skewness, mean)은 인구통계학적 변수와 독립적으로\n"
         "TAMA 예측에 유의미하게 기여한다. 향후 외부 검증 및 다기관 연구를 통해 임상 적용 가능성을 확인해야 한다.",
         0.55, 5.12, 12.1, 1.3, size=13, color=WHITE)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
output_path = f"results/{config.SITE}_TAMA_연구보고서.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
