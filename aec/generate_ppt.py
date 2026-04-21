"""
TAMA 예측 회귀분석 연구 보고서 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x35, 0x5E)   # 제목 배경
BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # 강조
SKYBLUE = RGBColor(0xD6, 0xE4, 0xF0)   # 셀 배경
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE  = RGBColor(0xFF, 0x7F, 0x27)
GREEN   = RGBColor(0x00, 0x8B, 0x45)

FIGURES = "c:/Users/jhjun/OneDrive/Desktop/aec/results/figures"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ────────────────────────────────────────────────────────────
# 헬퍼 함수
# ────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.line.width = line_w
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if line_rgb:
        shp.line.color.rgb = line_rgb
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None):
    """슬라이드 상단 헤더 바"""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=NAVY)
    add_text(slide, title, 0.35, 0.1, 10, 0.6,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.65, 10, 0.45,
                 size=14, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)
    # 페이지 번호 자리 (우측)
    # 슬라이드 번호는 자동 처리 생략


def add_image(slide, fname, l, t, w, h):
    path = os.path.join(FIGURES, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        shp = add_rect(slide, l, t, w, h, fill_rgb=GRAY, line_rgb=BLUE)
        add_text(slide, f"[그림 없음]\n{fname}", l+0.1, t+h/2-0.2, w-0.2, 0.4,
                 size=9, color=BLUE, align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows,
              l, t, w, h,
              header_fill=NAVY, header_color=WHITE,
              alt_fill=SKYBLUE, font_size=11):
    """심플 텍스트-기반 표 (python-pptx Table)"""
    cols = len(headers)
    total_rows = 1 + len(rows)
    tbl = slide.shapes.add_table(total_rows, cols,
                                  Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    col_w = Inches(w / cols)
    for c in range(cols):
        tbl.columns[c].width = col_w

    # 헤더
    for c, hdr in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = header_color

    # 데이터 행
    for r, row in enumerate(rows):
        fill_col = GRAY if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_col
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK


def bullet_box(slide, items, l, t, w, h, size=13, title=None):
    """글머리 기호 텍스트 박스"""
    offset = 0
    if title:
        add_text(slide, title, l, t, w, 0.35, size=13, bold=True, color=NAVY)
        offset = 0.35
    txb = slide.shapes.add_textbox(Inches(l), Inches(t + offset),
                                    Inches(w), Inches(h - offset))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)


# ════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# 전체 배경
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=NAVY)
# 하단 포인트 바
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=BLUE)

# 중앙 흰 카드
add_rect(slide, 1.0, 1.5, 11.33, 4.5, fill_rgb=WHITE)

add_text(slide,
         "AEC 곡선 기반 TAMA 예측\n회귀분석 연구 보고서",
         1.3, 1.7, 10.7, 2.0,
         size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(slide,
         "CT 자동 노출 제어(AEC) 특징과 환자 인구통계학적 변수를 활용한\n"
         "복부 근육량 지표(TAMA) 예측 모델 개발 및 검증",
         1.3, 3.6, 10.7, 1.0,
         size=15, color=RGBColor(0x44, 0x44, 0x66), align=PP_ALIGN.CENTER)

add_text(slide,
         "데이터셋: 강남 CT  |  N=1,672명  |  Python (statsmodels · scikit-learn)",
         1.3, 4.7, 10.7, 0.5,
         size=13, color=BLUE, align=PP_ALIGN.CENTER, bold=True)

add_text(slide, "2026", 0.4, 6.85, 12.5, 0.5,
         size=13, color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "목  차", "Table of Contents")

sections = [
    ("01", "연구 개요", "연구 목적 · 설계 · Feature Set (3 Case)"),
    ("02", "데이터 기술 통계", "데이터셋 · 성별 TAMA 분포 · 이진화 기준"),
    ("03", "분석 방법론", "전처리 · 모델 · 교차검증"),
    ("04", "결과 — 선형 회귀 (Part 1)", "단변량 · 다변량 성능 · 잔차 진단"),
    ("05", "결과 — 로지스틱 회귀 (Part 2)", "단변량 OR · 다변량 AUC · Calibration"),
    ("06", "결과 — Case 1–3 비교 (Part 3)", "AEC 추가 효과 · CT 모델명 효과"),
    ("07", "결론 및 한계", "주요 발견 · 제한점 · 향후 방향"),
]

for i, (num, title, desc) in enumerate(sections):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 1.3 + row * 1.65

    add_rect(slide, x, y, 6.0, 1.45, fill_rgb=NAVY)
    add_text(slide, num, x+0.15, y+0.1, 0.7, 0.6, size=22, bold=True, color=ORANGE)
    add_text(slide, title, x+0.85, y+0.08, 4.9, 0.5, size=15, bold=True, color=WHITE)
    add_text(slide, desc, x+0.85, y+0.65, 4.9, 0.7, size=11, color=RGBColor(0xB0,0xC8,0xE8))

    if i == 6:  # 마지막 하나 남을 때 중앙 정렬
        pass


# ════════════════════════════════════════════════════════════
# 슬라이드 3 — 연구 개요
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "01  연구 개요", "Research Overview")

# 연구 목적 카드
add_rect(slide, 0.35, 1.25, 12.6, 1.15, fill_rgb=NAVY)
add_text(slide, "연구 목적",
         0.55, 1.28, 3.0, 0.4, size=13, bold=True, color=ORANGE)
add_text(slide,
         "CT 자동 노출 제어(AEC) 곡선과 환자 인구통계(성별·나이)가 "
         "복부 근육량 지표 TAMA(Total Abdominal Muscle Area) 예측에 기여하는 정도를 정량화",
         0.55, 1.62, 12.2, 0.65, size=12, color=WHITE)

# 3개 파트 카드
parts = [
    ("Part 1", "선형 회귀\n(Linear Regression)", "TAMA (연속형) 값 예측"),
    ("Part 2", "로지스틱 회귀\n(Logistic Regression)", "Low TAMA / Sarcopenia 위험 분류"),
    ("Part 3", "Multivariable Analysis\n(Case 1–3 비교)", "AEC·CT 모델 추가 효과 정량화"),
]
for i, (p, m, g) in enumerate(parts):
    x = 0.35 + i * 4.2
    add_rect(slide, x, 2.6, 4.0, 2.1, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1.5))
    add_text(slide, p, x+0.15, 2.68, 3.7, 0.4, size=13, bold=True, color=BLUE)
    add_text(slide, m, x+0.15, 3.05, 3.7, 0.75, size=12, bold=True, color=NAVY)
    add_text(slide, g, x+0.15, 3.8, 3.7, 0.7, size=11, color=DARK)

# Feature Set 표
add_text(slide, "Feature Set (3 Case)", 0.35, 4.85, 4.0, 0.35,
         size=13, bold=True, color=NAVY)
add_table(slide,
          ["Case", "예측 변수"],
          [
              ["Case 1", "성별 (Sex) + 나이 (Age)"],
              ["Case 2", "Case 1 + AEC 특징 (p25, CV, skewness, slope_abs_mean)"],
              ["Case 3", "Case 2 + CT 모델명 (ManufacturerModelName — 30개 dummy)"],
          ],
          0.35, 5.2, 12.6, 1.9, font_size=11)


# ════════════════════════════════════════════════════════════
# 슬라이드 4 — 데이터 기술 통계
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  데이터 기술 통계", "Descriptive Statistics")

# 기본 정보 카드
stats = [
    ("총 환자 수", "1,672명"),
    ("남성", "665명 (39.8%)"),
    ("여성", "1,007명 (60.2%)"),
    ("CT 스캐너 모델", "31종"),
    ("TAMA 범위", "14 ~ 299 cm²"),
    ("TAMA 평균 (SD)", "122.54 (±30.48) cm²"),
]
for i, (label, val) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.1
    add_rect(slide, x, y, 3.9, 0.95, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1))
    add_text(slide, label, x+0.12, y+0.05, 3.7, 0.35, size=10, color=BLUE)
    add_text(slide, val, x+0.12, y+0.42, 3.7, 0.45, size=16, bold=True, color=NAVY)

# 성별 TAMA 분포 표
add_text(slide, "성별 TAMA 분포", 0.35, 3.55, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["성별", "N", "평균 (cm²)", "SD", "중앙값", "P25", "P75"],
          [
              ["남성 (M)", "665", "149.52", "27.46", "150.00", "132.00", "166.00"],
              ["여성 (F)", "1,007", "104.71", "15.70", "103.00", "95.00", "114.00"],
          ],
          0.35, 3.92, 7.5, 1.35, font_size=11)

# 이진화 기준 카드
add_rect(slide, 8.1, 3.55, 4.9, 1.72, fill_rgb=NAVY)
add_text(slide, "Logistic Regression 이진화 기준", 8.25, 3.6, 4.6, 0.4,
         size=11, bold=True, color=ORANGE)
add_text(slide,
         "• 남성: TAMA < 170 cm²  →  Low Muscle (1)\n"
         "• 여성: TAMA < 110 cm²  →  Low Muscle (1)\n"
         "• 양성 비율: 1,185 / 1,672 = 70.9 %",
         8.25, 4.03, 4.6, 1.2, size=12, color=WHITE)

# TAMA 분포 그림
add_text(slide, "TAMA 분포 시각화", 0.35, 5.35, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "03_tama_distribution.png", 0.35, 5.72, 6.0, 1.55)

# AEC Feature 상관 그림
add_text(slide, "AEC Feature 상관계수", 6.65, 5.35, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "01_feature_correlation.png", 6.65, 5.72, 6.33, 1.55)


# ════════════════════════════════════════════════════════════
# 슬라이드 5 — AEC Feature 선택
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "02  AEC Feature 선택", "Feature Selection")

add_text(slide,
         "Pearson/Spearman 상관계수 + VIF 검사를 통해 다중공선성이 없는 4개 Feature 선택",
         0.35, 1.15, 12.6, 0.45, size=13, color=DARK)

add_table(slide,
          ["Feature", "Pearson r", "해석", "VIF", "선택"],
          [
              ["p25", "+0.365", "AEC 하위 25% 값 → 저선량 구간 tube current", "~1.9", "✔"],
              ["CV", "-0.349", "변동계수(std/mean) → 체형 불균일성 반영", "~1.8", "✔"],
              ["skewness", "-0.344", "AEC 곡선 비대칭성 → 체형 분포 특성", "~1.6", "✔"],
              ["slope_abs_mean", "~+0.09", "평균 절대 기울기 → 곡선 동역학", "~1.4", "✔"],
              ["mean", "높음", "AUC_normalized와 심각한 공선성", ">50,000", "✗ 제외"],
              ["AUC_normalized", "높음", "mean과 심각한 공선성", ">50,000", "✗ 제외"],
          ],
          0.35, 1.65, 12.6, 2.6, font_size=11)

add_text(slide, "⚠  mean · AUC_normalized는 VIF > 50,000으로 다중공선성 심각 → 제외",
         0.35, 4.3, 12.6, 0.5, size=12, italic=True, color=ORANGE)

# 두 그림 나란히
add_text(slide, "Feature 상관계수 Top20", 0.35, 4.85, 5.5, 0.35, size=12, bold=True, color=NAVY)
add_image(slide, "01_feature_correlation.png", 0.35, 5.22, 6.0, 2.0)

add_text(slide, "VIF 비교", 6.65, 4.85, 5.5, 0.35, size=12, bold=True, color=NAVY)
add_image(slide, "02_vif_comparison.png", 6.65, 5.22, 6.33, 2.0)


# ════════════════════════════════════════════════════════════
# 슬라이드 6 — 분석 방법론
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "03  분석 방법론", "Statistical Methods")

# 전처리 표
add_text(slide, "전처리", 0.35, 1.25, 5.5, 0.35, size=14, bold=True, color=NAVY)
add_table(slide,
          ["변수", "처리 방법"],
          [
              ["PatientSex", "이진 인코딩 (M=1, F=0)"],
              ["PatientAge", "Z-score 표준화 (StandardScaler)"],
              ["AEC features", "Z-score 표준화"],
              ["ManufacturerModelName", "One-hot encoding (drop_first=True) → 30개 더미"],
          ],
          0.35, 1.65, 6.0, 2.1, font_size=11)

# 모델 카드
add_text(slide, "모델", 6.65, 1.25, 6.0, 0.35, size=14, bold=True, color=NAVY)
models = [
    ("선형 회귀", "statsmodels OLS\n(Ordinary Least Squares)"),
    ("로지스틱 회귀", "statsmodels Logit\n(BFGS, maxiter=1,000)"),
]
for i, (nm, desc) in enumerate(models):
    y = 1.65 + i * 1.1
    add_rect(slide, 6.65, y, 6.33, 0.95, fill_rgb=WHITE, line_rgb=BLUE, line_w=Pt(1.2))
    add_text(slide, nm, 6.8, y+0.05, 3.5, 0.38, size=13, bold=True, color=NAVY)
    add_text(slide, desc, 6.8, y+0.48, 5.9, 0.42, size=11, color=DARK)

# 교차검증 및 CI
add_text(slide, "모델 검증 전략", 0.35, 3.9, 12.6, 0.35, size=14, bold=True, color=NAVY)
add_table(slide,
          ["검증 방법", "상세"],
          [
              ["5-Fold Cross Validation", "각 Fold에서 scaler 재적합 (Data Leakage 방지)"],
              ["Bootstrap AUC 95% CI", "n=1,000 resampling — 비모수 신뢰구간"],
              ["Hosmer-Lemeshow 검정", "로지스틱 보정도(Calibration) 검증 — p>0.05 양호"],
              ["잔차 진단 3종", "Shapiro-Wilk · Breusch-Pagan · Durbin-Watson"],
          ],
          0.35, 4.3, 12.6, 2.1, font_size=11)

add_text(slide,
         "※ 선형·로지스틱 회귀 모두 Case 1(Sex+Age) → Case 2(+AEC) → Case 3(+CT Model) 순으로 점진 투입",
         0.35, 6.5, 12.6, 0.55, size=11, italic=True, color=BLUE)


# ════════════════════════════════════════════════════════════
# 슬라이드 7 — 선형 회귀: 단변량
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 단변량 분석", "Univariate Linear Regression")

add_table(slide,
          ["변수", "β 계수", "95% CI", "p-value", "R²"],
          [
              ["Sex (M=1, F=0)", "44.81 *", "[42.74, 46.89]", "< 0.001", "0.518"],
              ["Age (표준화)", "-3.00 *", "[-4.46, -1.55]", "< 0.001", "0.010"],
              ["AEC: p25 (표준화)", "11.13 *", "[9.77, 12.50]", "< 0.001", "0.134"],
              ["AEC: CV (표준화)", "-10.62 *", "[-12.00, -9.25]", "< 0.001", "0.122"],
              ["AEC: skewness (표준화)", "-10.48 *", "[-11.85, -9.10]", "< 0.001", "0.118"],
              ["AEC: slope_abs_mean (표준화)", "1.25", "[-0.22, 2.71]", "0.095", "0.002"],
              ["ManufacturerModelName", "N/A (범주형)*", "N/A", "9.0e-06", "0.045"],
          ],
          0.35, 1.25, 12.6, 3.0, font_size=11)

add_text(slide, "* p < 0.05 유의",
         0.35, 4.35, 5.0, 0.35, size=11, italic=True, color=ORANGE)

# 단변량 R² 그림
add_text(slide, "단변량 R² 비교", 0.35, 4.75, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "07_linear_univariate_r2.png", 0.35, 5.12, 6.0, 2.15)

# Forest plot
add_text(slide, "유의 계수 Forest Plot", 6.65, 4.75, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "06_linear_forest.png", 6.65, 5.12, 6.33, 2.15)


# ════════════════════════════════════════════════════════════
# 슬라이드 8 — 선형 회귀: 다변량 성능
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "04  선형 회귀 — 다변량 모델 성능 (Case 3 Full)", "Multivariate Linear Regression")

# 성능 KPI 카드
kpis = [
    ("R²", "0.6572"),
    ("Adj R²", "0.6497"),
    ("RMSE", "17.84 cm²"),
    ("MAE", "13.18 cm²"),
    ("F-statistic", "87.08\n(p < 0.001)"),
    ("AIC / BIC", "14,455\n/ 14,656"),
]
for i, (label, val) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.25
    add_rect(slide, x, y, 3.9, 1.15, fill_rgb=NAVY)
    add_text(slide, label, x+0.15, y+0.08, 3.6, 0.38, size=12, bold=True, color=ORANGE)
    add_text(slide, val, x+0.15, y+0.52, 3.6, 0.55, size=18, bold=True, color=WHITE)

# 잔차 진단 표
add_text(slide, "잔차 진단", 0.35, 3.85, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["검정", "통계량", "p-value", "결과"],
          [
              ["Shapiro-Wilk", "0.9697", "< 0.001", "비정규 ⚠"],
              ["Breusch-Pagan", "120.47", "< 0.001", "이분산 ⚠"],
              ["Durbin-Watson", "1.9654", "—", "정상 ✔"],
              ["Condition κ", "145.81", "—", "주의(κ<1,000) ✔"],
          ],
          0.35, 4.22, 6.5, 2.0, font_size=11)

add_text(slide,
         "⚠  비정규·이분산 → 계수 추정 BLUE이나 p-value 해석 주의\n   Robust Standard Errors 사용 권고",
         0.35, 6.3, 6.5, 0.75, size=11, italic=True, color=ORANGE)

# 실제 vs 예측 그림
add_text(slide, "실제 vs 예측", 7.1, 3.85, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "04_linear_actual_vs_pred.png", 7.1, 4.22, 5.9, 2.0)

# 잔차 진단 4-panel
add_image(slide, "05_linear_residuals.png", 7.1, 6.25, 5.9, 1.0)


# ════════════════════════════════════════════════════════════
# 슬라이드 9 — 로지스틱 회귀: 단변량 OR
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 단변량 분석 (Crude OR)", "Univariate Logistic Regression")

add_table(slide,
          ["변수", "Crude OR", "95% CI", "p-value", "AUC"],
          [
              ["Sex (M=1, F=0)", "1.844 *", "[1.472, 2.311]", "< 0.001", "0.571"],
              ["Age (표준화)", "1.449 *", "[1.301, 1.615]", "< 0.001", "0.612"],
              ["AEC: p25 (표준화)", "0.562 *", "[0.502, 0.630]", "< 0.001", "0.653"],
              ["AEC: CV (표준화)", "1.305 *", "[1.168, 1.459]", "< 0.001", "0.568"],
              ["AEC: skewness (표준화)", "1.191 *", "[1.074, 1.321]", "9.7e-04", "0.545"],
              ["AEC: slope_abs_mean (표준화)", "1.172 *", "[1.050, 1.308]", "0.0047", "0.548"],
              ["ManufacturerModelName", "LR χ²=57.11 *", "N/A", "0.002", "0.581"],
          ],
          0.35, 1.25, 12.6, 3.0, font_size=11)

add_text(slide, "* p < 0.05 유의  |  AEC p25: OR < 1 → p25↑ 시 Low TAMA 위험 감소",
         0.35, 4.35, 12.0, 0.4, size=11, italic=True, color=ORANGE)

# Crude OR Forest plot
add_text(slide, "Crude OR Forest Plot", 0.35, 4.82, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "11_logistic_forest.png", 0.35, 5.2, 6.0, 2.07)

# ROC 곡선
add_text(slide, "ROC Curve (단변량)", 6.65, 4.82, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "08_logistic_roc.png", 6.65, 5.2, 6.33, 2.07)


# ════════════════════════════════════════════════════════════
# 슬라이드 10 — 로지스틱 회귀: 다변량 성능
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "05  로지스틱 회귀 — 다변량 모델 성능 (Case 3 Full)", "Multivariate Logistic Regression")

kpis2 = [
    ("AUC-ROC", "0.7794"),
    ("Bootstrap 95% CI", "[0.756 – 0.802]"),
    ("Sensitivity", "0.767"),
    ("Specificity", "0.667"),
    ("PPV / NPV", "0.849 / 0.541"),
    ("Nagelkerke R²", "0.2823"),
]
for i, (label, val) in enumerate(kpis2):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.2
    y = 1.25 + row * 1.2
    add_rect(slide, x, y, 3.9, 1.1, fill_rgb=NAVY)
    add_text(slide, label, x+0.15, y+0.06, 3.6, 0.38, size=11, bold=True, color=ORANGE)
    add_text(slide, val, x+0.15, y+0.5, 3.6, 0.52, size=17, bold=True, color=WHITE)

# 보정도 표
add_text(slide, "보정도 및 적합도", 0.35, 3.72, 6.0, 0.35, size=13, bold=True, color=NAVY)
add_table(slide,
          ["지표", "값", "해석"],
          [
              ["Hosmer-Lemeshow χ²", "9.114", "p=0.333 → 보정도 양호 ✔"],
              ["Brier Score", "0.1629", "낮을수록 우수 (0=완벽, 0.25=무작위)"],
              ["AIC", "1,722.81", "Case 비교 기준"],
          ],
          0.35, 4.1, 7.0, 1.65, font_size=11)

# Calibration plot
add_text(slide, "Calibration Plot", 7.5, 3.72, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "09_logistic_calibration.png", 7.5, 4.1, 5.48, 2.2)

# Confusion matrix
add_text(slide, "Confusion Matrix", 0.35, 5.85, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "10_logistic_confusion.png", 0.35, 6.22, 3.5, 1.05)


# ════════════════════════════════════════════════════════════
# 슬라이드 11 — Case 1·2·3 선형 회귀 비교
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 비교 — 선형 회귀 성능 향상", "Multivariable Analysis: Linear Regression")

add_table(slide,
          ["지표", "Case 1\n(Sex+Age)", "Case 2\n(+AEC)", "Case 3\n(+CT Model)",
           "Case1→2 Δ", "Case2→3 Δ"],
          [
              ["R²", "0.5506", "0.6353", "0.6572", "+0.0847 ▲", "+0.0219 ▲"],
              ["Adj R²", "0.5501", "0.6339", "0.6497", "+0.0838 ▲", "+0.0158 ▲"],
              ["RMSE (cm²)", "20.43", "18.41", "17.84", "-2.02 ▼", "-0.56 ▼"],
              ["AIC", "14,840", "14,499", "14,455", "-341 ▼", "-44 ▼"],
          ],
          0.35, 1.25, 12.6, 2.3, font_size=12)

add_text(slide,
         "AEC 특징 추가(Case1→2)로 R² +8.5%p 향상  |  CT 모델명 추가(Case2→3)로 추가 +2.2%p",
         0.35, 3.65, 12.6, 0.45, size=13, bold=True, color=BLUE)

add_text(slide, "R² / RMSE 비교", 0.35, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "12_case_metrics_bar.png", 0.35, 4.58, 6.0, 2.7)

add_text(slide, "AIC / BIC 비교", 6.65, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "14_case_aic_bar.png", 6.65, 4.58, 6.33, 2.7)


# ════════════════════════════════════════════════════════════
# 슬라이드 12 — Case 1·2·3 로지스틱 회귀 비교
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "06  Case 비교 — 로지스틱 회귀 성능 향상", "Multivariable Analysis: Logistic Regression")

add_table(slide,
          ["지표", "Case 1\n(Sex+Age)", "Case 2\n(+AEC)", "Case 3\n(+CT Model)",
           "Case1→2 Δ", "Case2→3 Δ"],
          [
              ["AUC-ROC", "0.6283", "0.7518", "0.7794", "+0.1235 ▲▲", "+0.0276 ▲"],
              ["Nagelkerke R²", "0.0581", "0.2312", "0.2823", "+0.1731 ▲▲", "+0.0511 ▲"],
              ["AIC", "1,954", "1,736", "1,723", "-218 ▼▼", "-13 ▼"],
              ["HL p-value", "≈ 0.000", "0.0894", "0.3328", "보정도 개선 ✔", "보정도 양호 ✔"],
          ],
          0.35, 1.25, 12.6, 2.3, font_size=12)

add_text(slide,
         "AEC 추가(Case1→2)로 AUC +12.4%p 대폭 향상  |  CT 모델명 추가(Case2→3)로 +2.8%p 추가 개선",
         0.35, 3.65, 12.6, 0.45, size=13, bold=True, color=BLUE)

add_text(slide, "AUC / Nagelkerke R² 비교", 0.35, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "13_case_auc_bar.png", 0.35, 4.58, 6.0, 2.7)

add_text(slide, "다중 지표 추이", 6.65, 4.2, 5.5, 0.35, size=13, bold=True, color=NAVY)
add_image(slide, "15_case_progression.png", 6.65, 4.58, 6.33, 2.7)


# ════════════════════════════════════════════════════════════
# 슬라이드 13 — 결론
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "07  결론", "Conclusions")

conclusions = [
    ("성별의 압도적 기여",
     "선형 단변량: Sex 단독 R²=0.518  |  남성 TAMA가 여성 대비 평균 +44.8 cm²"),
    ("AEC 특징의 독립적 기여",
     "Case1→2  선형 R²  +0.085 (0.551→0.635)\n"
     "Case1→2  로지스틱 AUC  +0.124 (0.628→0.752)"),
    ("CT 모델명의 추가 기여",
     "Case2→3  선형 R²  +0.022  |  로지스틱 AUC  +0.028\n"
     "AIC 추가 감소 → 모델 적합도 개선 확인"),
    ("최종 모델 (Case 3) 요약",
     "선형: R²=0.6572, RMSE=17.84 cm²\n"
     "로지스틱: AUC=0.7794 [0.756–0.802]  |  HL p=0.333 (보정도 양호)"),
]
for i, (title, body) in enumerate(conclusions):
    row = i // 2
    col = i % 2
    x = 0.35 + col * 6.4
    y = 1.25 + row * 2.35
    add_rect(slide, x, y, 6.1, 2.2, fill_rgb=NAVY)
    add_rect(slide, x, y, 0.18, 2.2, fill_rgb=ORANGE)
    add_text(slide, f"{i+1}. {title}", x+0.3, y+0.1, 5.6, 0.45,
             size=13, bold=True, color=ORANGE)
    add_text(slide, body, x+0.3, y+0.6, 5.6, 1.5, size=12, color=WHITE)

add_text(slide,
         "AEC 특징은 성별·나이 외에 TAMA 예측에 독립적이고 유의미한 기여를 함 — CT 촬영 시 수집되는 AEC 데이터의 임상적 활용 가능성 확인",
         0.35, 6.8, 12.6, 0.55, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 14 — 한계 및 향후 방향
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=GRAY)
slide_header(slide, "07  한계 및 향후 방향", "Limitations & Future Work")

limits = [
    "잔차 정규성·등분산성 가정 위반 → Robust Standard Errors 또는 비모수 검정 검토 필요",
    "Logistic 이진화 임계값 (M<170, F<110 cm²) — 문헌 기반, 기관별 참조값 확인 필요",
    "단일 기관(강남) 데이터 → 외부 검증(신촌 데이터) 필요",
    "AEC Feature 선택은 통계적 상관분석 기반 — 임상적 의미 추가 검토 권장",
]
future = [
    "신촌 데이터를 이용한 외부 검증 (External Validation)",
    "Robust Regression / Weighted Least Squares 적용",
    "머신러닝 기반 앙상블 모델 비교 (Random Forest, XGBoost)",
    "다기관 데이터 수집을 통한 모델 일반화",
    "AEC Feature의 임상적 의미 방사선사·임상의 검토",
]

add_rect(slide, 0.35, 1.2, 6.0, 3.3, fill_rgb=WHITE, line_rgb=ORANGE, line_w=Pt(1.5))
add_text(slide, "⚠  한계점", 0.55, 1.28, 5.5, 0.38, size=13, bold=True, color=ORANGE)
for i, lim in enumerate(limits):
    add_text(slide, f"• {lim}", 0.55, 1.72 + i*0.72, 5.6, 0.65, size=11, color=DARK)

add_rect(slide, 6.7, 1.2, 6.28, 3.3, fill_rgb=WHITE, line_rgb=GREEN, line_w=Pt(1.5))
add_text(slide, "→  향후 연구 방향", 6.88, 1.28, 5.8, 0.38, size=13, bold=True, color=GREEN)
for i, fut in enumerate(future):
    add_text(slide, f"• {fut}", 6.88, 1.72 + i*0.57, 5.8, 0.52, size=11, color=DARK)

# 결론 요약 강조 박스
add_rect(slide, 0.35, 4.65, 12.6, 1.85, fill_rgb=NAVY)
add_text(slide, "핵심 메시지",
         0.55, 4.72, 3.0, 0.38, size=12, bold=True, color=ORANGE)
add_text(slide,
         "CT 스캐너의 AEC 곡선에서 추출한 통계 특징(p25, CV, skewness)은 인구통계학적 변수와 독립적으로\n"
         "TAMA 예측에 유의미하게 기여한다. 향후 외부 검증 및 다기관 연구를 통해 임상 적용 가능성을 확인해야 한다.",
         0.55, 5.12, 12.1, 1.3, size=13, color=WHITE)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
output_path = "c:/Users/jhjun/OneDrive/Desktop/aec/results/TAMA_연구보고서.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
